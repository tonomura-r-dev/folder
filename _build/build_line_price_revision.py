# -*- coding: utf-8 -*-
"""LINE料金改定に伴う運用最適化・特別プランのご案内（1枚完結）

構成：
  ① 課題とリスク（10月〜） 追加メッセージ単価の改定。一斉配信のままだとコスト増
  ② 弊社の対策提案        セグメント配信への切替／各社様の配信規模に応じた特別調整プラン
  ③ スケジュール          10月＝実データ計測 → 10月下旬＝個別シミュレーション
                          → 11月〜＝特別プラン適用開始

「10月に一律で安くできない理由」を角を立てずに説明し、11月への期待で着地させる。
営業用トークスクリプトはスピーカーノートに入れている。

使い方:
    python _build/build_line_price_revision.py
        → _templates/DYM_LINEOA_FMT.pptx から 1枚もの を生成

    python _build/build_line_price_revision.py <既存.pptx> [出力.pptx] [ページ番号]
        → 既存デッキの指定ページ（既定9）を、この内容で作り直す
"""
import shutil
import sys
from pathlib import Path
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
FMT = str(ROOT / "_templates" / "DYM_LINEOA_FMT.pptx")
STANDALONE_OUT = str(ROOT / "20260826_LINE料金改定に伴う運用最適化・特別プランのご案内.pptx")

TNAVY = "002060"   # タイトル
NAVY = "1F285A"    # 着地・スケジュール
BLUE = "1565C0"    # 対策
RED = "C00000"     # 課題・リスク
GREY = "8C8C8C"    # 動かせない外部要因
INK = "333333"
MUT = "808080"
WHITE = "FFFFFF"
ICE = "C8D0E8"     # 紺地の上の小さい文字
PALE = "F4F7FF"
FADE = "F2F2F2"
BORDER = "D9D9D9"


# ---------- helpers（build_vivical_9p.py と同じ経路） ----------
def set_font(run, size, bold=None, color=INK, name="メイリオ"):
    f = run.font
    f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", name)
    f.color.rgb = RGBColor.from_string(color)


ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
ANCH = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}


def put_text(tf, paras, anchor="t", ml=0.06, mr=0.06, mt=0.03, mb=0.03, wrap=True):
    tf.word_wrap = wrap
    tf.margin_left = Inches(ml)
    tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt)
    tf.margin_bottom = Inches(mb)
    tf.vertical_anchor = ANCH[anchor]
    first = True
    for p in paras:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = ALIGN[p.get("align", "l")]
        if p.get("sa") is not None:
            para.space_after = Pt(p["sa"])
        if p.get("ls") is not None:
            para.line_spacing = p["ls"]
        for (t, sz, b, c) in p["runs"]:
            r = para.add_run()
            r.text = t
            set_font(r, sz, b, c)
    return tf


def add_text(slide, x, y, w, h, paras, anchor="t", **kw):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    put_text(box.text_frame, paras, anchor=anchor, **kw)
    return box


def add_box(slide, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor.from_string(line)
        sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    sp.text_frame.word_wrap = True
    return sp


def clear_slide(slide):
    spTree = slide.shapes._spTree
    for el in list(spTree):
        if el.tag.split("}")[-1] in ("sp", "cxnSp", "pic", "graphicFrame", "grpSp"):
            spTree.remove(el)


# ---------- 中身 ----------
TITLE = "LINE公式アカウント 料金改定に伴う運用最適化・特別プランのご案内"

LEAD = [
    {"runs": [("2026年10月の料金改定により、従来の配信スタイルのままでは配信コストが増大します。", 13.5, True, INK)],
     "align": "c"},
    {"runs": [("10月に実データを計測のうえ、", 13.5, True, INK),
              ("11月から", 13.5, True, RED),
              ("各社様に最適化した特別プランを適用いたします。", 13.5, True, INK)], "align": "c"},
]

# ① 背景と課題（10月〜）：ラベル＋本文の2行
BACKGROUND = [
    ("LINEヤフーの料金体系改定", NAVY,
     "2026年10月より追加メッセージ単価体系が改定。"),
    ("生じるリスク", RED,
     "従来の配信スタイルのまま運用を継続した場合、メッセージ配信コストが大幅に増大する可能性。"),
]

# ② 弊社の対応方針
MEASURES = [
    ("01", "配信の最適化",
     "一斉配信から「セグメント配信」へ移行し、不要コストを圧縮します。"),
    ("02", "特別プランの再設計",
     "新料金下でも費用対効果を最大化できるよう、弊社独自の「調整プラン（特別還元）」をご提供します。"),
]

# ③ スケジュール
# (期間, 小見出し, 本文, 注記, ヘッダー色, カード色, 枠色, 本文色, 注記色, ヘッダー文字色)
SCHEDULE = [
    ("10月", "配信データ計測・検証期",
     "新料金下での実際の配信ボリュームとコストの適正推移を測定・分析します。",
     "原則、通常プラン適用",
     GREY, FADE, BORDER, INK, MUT, WHITE),
    ("10月下旬", "個別シミュレーション提示",
     "10月の配信実績および媒体インセンティブを試算し、貴社に最もメリットが出る特別プランを個別に作成します。",
     "", BLUE, PALE, BLUE, INK, BLUE, WHITE),
    ("11月〜", "特別プラン適用開始",
     "各社様に最適化した特別プランを正式に適用します。",
     "", WHITE, NAVY, NAVY, WHITE, ICE, NAVY),
]

TALK_SCRIPT = """【営業トークスクリプト】LINE料金改定に伴う運用最適化・特別プランのご案内

■ 導入
2026年10月から、LINEヤフー社の料金体系が改定されます。追加メッセージの単価が変わり、
大量配信に対する割引が縮小します。御社のように配信ボリュームが大きいアカウントほど、
影響を受けやすい改定です。

■ 課題とリスク
対策をしないと、これまでと同じ配信をしているだけでコストが上がります。
特に全員に一律で送る「一斉配信」は、反応の薄い層にも同じ単価がかかるので、
改定後はムダが大きく出てしまいます。

■ 弊社の対策
ご提案は2つです。
ひとつめは、一斉配信からセグメント配信への切り替えです。
反応が見込める層に絞って配信することで、通数そのものを減らしながら成果は維持します。
ふたつめが、御社の配信規模に合わせた弊社の特別調整プランです。

■ なぜ11月からなのか（ここが本題）
特別プランの適用は11月からとさせてください。理由は、10月が
「新しい料金体系での実データを取る期間」だからです。
新料金で実際にどれくらいの通数・費用になるかは、動かしてみないと正確に出ません。
10月は原則、通常プランで運用いただき、実データを取らせてください。
そのうえで10月下旬に、実績をもとに媒体インセンティブを算出し、
御社にどれだけ還元できるかを試算します。
その試算に基づいて、11月から御社専用の特別プランを適用します。

■ 着地
10月に一律で値引きをするのではなく、10月の実績を見たうえで、
御社にとって一番効果の大きい形で11月から適用する、という進め方です。
根拠のない一律値引きよりも、実データに基づいた最適化のほうが、
結果的に御社のメリットは大きくなります。

■ 想定QA
Q. 10月分は高いままですか？
A. 10月は原則、通常プランでの適用となります。ただしその1ヶ月の実績が、
   11月以降の調整の根拠になります。無駄な1ヶ月にはしません。

Q. どれくらい下がりますか？
A. 現時点では確約できません。10月の配信実績と媒体インセンティブ次第です。
   10月下旬に、御社専用のシミュレーションをお出しします。

Q. 10月から何かできることはありますか？
A. あります。セグメント配信への切り替えは10月から着手できます。
   ここを進めておくほど、11月の調整余地も大きくなります。
"""

# ---------- レイアウト定数（スライド 10.83 × 7.5 inch） ----------
X0, TOTAL_W = 0.50, 9.83
GAP = 0.30

SEC1_LABEL_Y, SEC1_Y, SEC1_H = 1.58, 1.82, 0.84
SEC2_LABEL_Y, SEC2_Y, SEC2_H = 2.80, 3.04, 1.30
SEC3_LABEL_Y, SEC3_Y, SEC3_H = 4.48, 4.72, 2.08

MEASURE_W = (TOTAL_W - GAP) / 2
STEP_W = (TOTAL_W - GAP * 2) / 3


def section_label(s, y, no, text, color):
    add_text(s, X0, y, 6.0, 0.26,
             [{"runs": [(no + "  ", 12, True, color), (text, 12, True, color)]}],
             anchor="m", ml=0.0, mr=0.0)


def draw_body(s):
    # ===== ① 背景と課題 =====
    section_label(s, SEC1_LABEL_Y, "①", "背景と課題（10月〜）", RED)
    add_box(s, X0, SEC1_Y, TOTAL_W, SEC1_H, fill=FADE, line=BORDER)
    for i, (label, lcol, body) in enumerate(BACKGROUND):
        y = SEC1_Y + 0.10 + i * 0.33
        add_text(s, X0 + 0.30, y, TOTAL_W - 0.60, 0.32,
                 [{"runs": [("■ ", 11, True, lcol),
                            (label, 12, True, lcol),
                            ("：", 12, True, lcol),
                            (body, 12, None, INK)]}],
                 anchor="m", ml=0.0, mr=0.0)

    # ===== ② 弊社の対策提案 =====
    section_label(s, SEC2_LABEL_Y, "②", "弊社の対応方針", BLUE)
    for i, (no, head, body) in enumerate(MEASURES):
        x = X0 + i * (MEASURE_W + GAP)
        add_box(s, x, SEC2_Y, MEASURE_W, SEC2_H, fill=PALE, line=BLUE)
        bd = add_box(s, x + 0.22, SEC2_Y + 0.18, 0.42, 0.42, fill=BLUE)
        put_text(bd.text_frame, [{"runs": [(no, 12, True, WHITE)], "align": "c"}],
                 anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
        add_text(s, x + 0.76, SEC2_Y + 0.16, MEASURE_W - 0.98, 0.46,
                 [{"runs": [(head, 14, True, NAVY)]}], anchor="m", ml=0.0)
        add_text(s, x + 0.76, SEC2_Y + 0.66, MEASURE_W - 0.98, 0.52,
                 [{"runs": [(body, 11, None, INK)], "ls": 1.25}], anchor="t", ml=0.0)

    # ===== ③ スケジュール =====
    section_label(s, SEC3_LABEL_Y, "③", "スケジュール（なぜ11月からの適用なのか）", NAVY)
    for i, (term, head, body, note, hcol, fill, edge, bcol, ncol, htxt) in enumerate(SCHEDULE):
        x = X0 + i * (STEP_W + GAP)
        add_box(s, x, SEC3_Y, STEP_W, SEC3_H, fill=fill, line=edge,
                lw=1.5 if i == 2 else 1.0)

        hd = add_box(s, x + 0.10, SEC3_Y + 0.14, STEP_W - 0.20, 0.42, fill=hcol)
        put_text(hd.text_frame, [{"runs": [(term, 13, True, htxt)], "align": "c"}],
                 anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)

        add_text(s, x + 0.12, SEC3_Y + 0.62, STEP_W - 0.24, 0.34,
                 [{"runs": [(head, 13, True, bcol)], "align": "c"}], anchor="m", ml=0.0, mr=0.0)

        add_text(s, x + 0.14, SEC3_Y + 1.00, STEP_W - 0.28, 0.66,
                 [{"runs": [(body, 10.5, None, bcol)], "ls": 1.28}], anchor="t", ml=0.0, mr=0.0)

        if note:
            add_text(s, x + 0.14, SEC3_Y + 1.70, STEP_W - 0.28, 0.28,
                     [{"runs": [("※ " + note, 9.5, None, ncol)]}], anchor="m", ml=0.0, mr=0.0)

        if i < 2:
            ar = add_box(s, x + STEP_W + 0.04, SEC3_Y + SEC3_H / 2 - 0.15,
                         0.22, 0.30, fill=BORDER, shape=MSO_SHAPE.RIGHT_ARROW)
            ar.line.fill.background()


def set_notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def build_standalone(out=STANDALONE_OUT):
    """DYM汎用FMTをコピー → 1枚に削減 → 中身を作り直す。"""
    shutil.copyfile(FMT, out)
    prs = Presentation(out)
    slides = list(prs.slides)

    div_el = None
    for sh in slides[1].shapes:
        if sh._element.tag.endswith('}cxnSp') and abs(sh.top - 1390675) < 60000:
            div_el = deepcopy(sh._element)
            break
    assert div_el is not None, "divider not found"

    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    keep = ids[1]
    for sldId in ids:
        if sldId is keep:
            continue
        prs.part.drop_rel(sldId.rId)
        sldIdLst.remove(sldId)
    assert len(list(prs.slides)) == 1

    s = list(prs.slides)[0]
    clear_slide(s)
    s.shapes._spTree.append(deepcopy(div_el))
    add_text(s, 0.60, 0.13, 8.60, 0.40,
             [{"runs": [(TITLE, 16, True, TNAVY)]}], anchor="m", ml=0.0, mr=0.0)
    add_text(s, 0.42, 0.60, 10.02, 0.84, LEAD, anchor="m", ml=0.0, mr=0.0)
    draw_body(s)
    set_notes(s, TALK_SCRIPT)
    prs.save(out)
    return out


def build_into_deck(src, out, page=9):
    """既存デッキの指定ページを、この内容で作り直す（他ページは触らない）。"""
    shutil.copyfile(src, out)
    prs = Presentation(out)
    s = prs.slides[page - 1]
    clear_slide(s)
    add_text(s, 0.47, 0.10, 8.60, 0.40,
             [{"runs": [(TITLE, 16, None, TNAVY)]}], anchor="m", ml=0.1, mr=0.0)
    add_text(s, 0.16, 0.60, 10.58, 0.84, LEAD, anchor="m", ml=0.0, mr=0.0)
    draw_body(s)
    set_notes(s, TALK_SCRIPT)
    prs.save(out)
    return out


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("saved:", build_standalone())
    else:
        src = Path(sys.argv[1])
        out = Path(sys.argv[2]) if len(sys.argv) > 2 else src.with_name(src.stem + "_S9修正.pptx")
        page = int(sys.argv[3]) if len(sys.argv) > 3 else 9
        print("saved:", build_into_deck(str(src), str(out), page))
