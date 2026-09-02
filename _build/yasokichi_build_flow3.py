# -*- coding: utf-8 -*-
"""八十吉 パフォーマー指名予約 LINEでの流れ（v3・A案／3工程）"""
import sys, os
SKILL = r'C:\Users\tonomura-r\.claude\skills\dym-format'
sys.path.insert(0, os.path.join(SKILL, 'scripts'))

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from helpers import copy_template, safe_save, delete_slide, duplicate_slide, \
    replace_text_in_slide, COLORS

OUT = sys.argv[1]
prs = copy_template(os.path.join(SKILL, 'assets', 'template.pptx'),
                    os.path.join(os.path.dirname(__file__), '_flow3.pptx'))

BLANK = 3
N_ORIG = len(prs.slides._sldIdLst)
duplicate_slide(prs, BLANK)
for idx in range(N_ORIG - 1, -1, -1):
    delete_slide(prs, idx)
S = list(prs.slides)

NAVY = COLORS['navy_700']
ACCENT = COLORS['accent']
INK = COLORS['ink']
MUTED = COLORS['ink_muted']
SW = 10.8333


def set_title(slide, text):
    replace_text_in_slide(slide, {'［スライドタイトルを入力］': text})
    sh = list(slide.shapes)[0]
    for p in sh.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(28)
            r.font.bold = True
            r.font.name = 'メイリオ'
            r.font.color.rgb = RGBColor.from_string(NAVY)


def arrow(slide, cx, cy, w=0.36):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(cx - w / 2),
                               Inches(cy - 0.12), Inches(w), Inches(0.24))
    a.fill.solid()
    a.fill.fore_color.rgb = RGBColor.from_string(COLORS['navy_300'])
    a.line.fill.background()
    a.shadow.inherit = False


def band(slide, x, y, w, h, fill, lines, size, color, bold=True):
    b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    b.fill.solid()
    b.fill.fore_color.rgb = RGBColor.from_string(fill)
    b.line.fill.background()
    b.shadow.inherit = False
    tf = b.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Pt(14)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(3)
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold
        r.font.name = 'メイリオ'
        r.font.color.rgb = RGBColor.from_string(color)
    return b


def box(slide, x, y, w, h, num, head, who, detail):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor.from_string('FDF0E6')
    sh.line.color.rgb = RGBColor.from_string(ACCENT)
    sh.line.width = Pt(1.75)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Pt(8)
    tf.margin_top = Pt(14)
    rows = [(num, 17, True, ACCENT, 5),
            (head, 19, True, INK, 7),
            (who, 14, True, NAVY, 9),
            (detail, 13, False, MUTED, 0)]
    first = True
    for t, sz, bd, col, sp in rows:
        for j, ln in enumerate(t.split('\n')):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(sp if j == len(t.split('\n')) - 1 else 0)
            r = p.add_run(); r.text = ln
            r.font.size = Pt(sz); r.font.bold = bd
            r.font.name = 'メイリオ'
            r.font.color.rgb = RGBColor.from_string(col)
    return sh


# =====================================================================
s = S[0]
set_title(s, 'LINEでのご予約の流れ')

STEPS = [
    ('①', 'ご予約\nリクエスト', 'お客様',
     '日時・人数・ご指名を\nトークで送る'),
    ('②', 'ご承諾の\nお返事', 'パフォーマー',
     'ご指名に気づいて\n受けられるかを返す'),
    ('③', '確定の\nご連絡', 'パフォーマー',
     'お客様へ直接\nお伝えする'),
]

BW, BH, GAP = 2.85, 2.55, 0.55
total = BW * 3 + GAP * 2
x0 = (SW - total) / 2
y0 = 1.42

for i, (num, head, who, detail) in enumerate(STEPS):
    x = x0 + i * (BW + GAP)
    box(s, x, y0, BW, BH, num, head, who, detail)
    if i < 2:
        arrow(s, x + BW + GAP / 2, y0 + BH / 2)

y = y0 + BH + 0.34
band(s, x0, y, total, 0.78, 'EEF1F7',
     ['店舗様の作業は、お席の確保だけです。ご予約台帳での操作は、これまで通りです。'],
     15, MUTED, bold=False)

y += 0.78 + 0.28
band(s, x0, y, total, 1.32, 'FDF0E6',
     ['ご指名の取り次ぎは要りません。パフォーマーの方に通知が届きます。',
      'チャットは通数に含まれないため、何往復しても月額0円のままです。'],
     18, ACCENT)

safe_save(prs, OUT)
print('SAVED:', OUT, len(prs.slides._sldIdLst), 'slides')
