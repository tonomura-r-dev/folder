# -*- coding: utf-8 -*-
"""EC業界（総合EC・業界汎用）LINEOA施策提案（36枚）

_templates/DYM_LINEOA_FMT.pptx（業界汎用FMT・36枚）をコピーし、
全36枚を _drafts/EC業界_スライド作成プロンプト.md の指定内容で構築する。

  python _build/build_ec.py
  python _build/qa_render.py EC業界_LINEOA施策提案.pptx

【この資料の背骨】
新規獲得の単価は上がり続けている。もう「集める」では勝てない。
Webマーケでできるのは、いま取りこぼしている購入と、一度きりで終わっている顧客を、
追加の広告費ゼロで拾い直すこと。勝負はCPOではなく、LTV ÷ CPO。

【未取得データ（★破線の差込枠にした。数値は捏造していない）】
- S08 広告実績CVR（DYM社内実績・CV地点別）
- S10 Googleトレンド（福袋/ブラックフライデー/セール/母の日）
- S11 LINEヤフー前後検索（通販/定期便/解約/口コミ）
- S12 競合10社の友だち数・配信内容
- S35 LINEヤフー公式のEC導入事例（lycbiz.com）

【落とし穴メモ（_build/README.md・deck-apply スキルより）】
- put_text() は必ず reset_tf() を通す（既存テキストへの追記事故を防ぐ）
- スライドの新規追加はしない。ベース36枚を clear_slide() して作り直す
- 一括置換をしない。数値は行・列を特定して書く
"""
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
SRC = str(ROOT / "_templates" / "DYM_LINEOA_FMT.pptx")
OUT = str(ROOT / "EC業界_LINEOA施策提案.pptx")

# ---- 配色（DYM 3色：ネイビー・炭黒・オレンジ。LINEグリーンはUI部分のみ）----
TNAVY = "002060"   # タイトル
NAVY = "1F285A"    # 打ち手・カード見出し
ORANGE = "ED7D31"  # 強調
RED = "C00000"     # 課題・警告
INK = "333333"     # 本文
MUT = "7F7F7F"
WHITE = "FFFFFF"
PALE = "F4F7FF"    # 淡ネイビー
PORANGE = "FCE4D6"  # 淡オレンジ
GREY = "F2F2F2"
BORDER = "D9D9D9"
GREEN = "06C755"   # LINE UI 専用
BEZEL = "2B2B2B"
SCREEN = "EFF2F7"
PRED = "FDF2F2"

# ---- FMTのレイアウト座標（cm）----
SW, SH = 27.52, 19.05
TITLE_XY = (1.52, 0.38, 24.4, 0.94)
LEAD_XY = (1.20, 1.80, 25.1, 1.90)
DIV_Y = 3.86
CX0, CW = 1.20, 25.12
CY0, CY1 = 4.30, 17.00
FOOT_Y = 17.35

shutil.copyfile(SRC, OUT)
prs = Presentation(OUT)
slides = list(prs.slides)
assert len(slides) == 36, len(slides)
assert (prs.slide_width, prs.slide_height) == (9906000, 6858000)


# ================= helpers（build_chumon_jutaku.py と共通の作法）=================
def set_font(run, size, bold=None, color=INK, name="メイリオ"):
    f = run.font
    f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", name)
    f.color.rgb = RGBColor.from_string(color)


ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
ANCH = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}


def reset_tf(tf):
    """既存の段落・runを全消去（追記事故の防止）"""
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)


def put_text(tf, paras, anchor="t", ml=0.14, mr=0.14, mt=0.06, mb=0.06, wrap=True):
    reset_tf(tf)
    tf.word_wrap = wrap
    tf.margin_left = Cm(ml)
    tf.margin_right = Cm(mr)
    tf.margin_top = Cm(mt)
    tf.margin_bottom = Cm(mb)
    tf.vertical_anchor = ANCH[anchor]
    first = True
    for p in paras:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = ALIGN[p.get("align", "l")]
        if p.get("sa") is not None:
            para.space_after = Pt(p["sa"])
        if p.get("ls") is not None:
            para.line_spacing = p["ls"]
        for item in p["runs"]:
            t, sz, b, c = item
            r = para.add_run()
            r.text = t
            set_font(r, sz, b, c)
    return tf


def T(slide, x, y, w, h, paras, anchor="t", **kw):
    box_ = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    put_text(box_.text_frame, paras, anchor=anchor, **kw)
    return box_


def one(text, sz, b=None, c=INK, align="l", sa=None, ls=None):
    d = {"runs": [(text, sz, b, c)], "align": align}
    if sa is not None:
        d["sa"] = sa
    if ls is not None:
        d["ls"] = ls
    return d


def multi(runs, align="l", sa=None, ls=None):
    """1段落に複数run（色・太さを混在させたいとき）"""
    d = {"runs": runs, "align": align}
    if sa is not None:
        d["sa"] = sa
    if ls is not None:
        d["ls"] = ls
    return d


def box(slide, x, y, w, h, fill=None, line=None, lw=1.0,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06, dash=None):
    sp = slide.shapes.add_shape(shape, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor.from_string(line)
        sp.line.width = Pt(lw)
        if dash:
            sp.line.dash_style = dash
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    reset_tf(sp.text_frame)
    return sp


def card(slide, x, y, w, h, head, body, hcol=NAVY, fill=PALE,
         hsz=11, bsz=9, line=None, anchor="t", ls=1.18):
    sp = box(slide, x, y, w, h, fill=fill, line=line)
    paras = [one(head, hsz, True, hcol, sa=3)]
    for b in (body if isinstance(body, list) else [body]):
        if b:
            paras.append(one(b, bsz, None, INK, ls=ls, sa=1))
    put_text(sp.text_frame, paras, anchor=anchor, ml=0.22, mr=0.18, mt=0.14, mb=0.10)
    return sp


def clear_slide(slide):
    spTree = slide.shapes._spTree
    for el in list(spTree):
        if el.tag.split("}")[-1] in ("sp", "cxnSp", "pic", "graphicFrame", "grpSp"):
            spTree.remove(el)


def frame(slide, title, lead):
    """FMT準拠：タイトル16pt紺（y=0.38）＋リード12pt（y=1.80・2行以内）＋区切り線"""
    clear_slide(slide)
    T(slide, *TITLE_XY, [one(title, 16, True, TNAVY)], anchor="m", ml=0, mr=0)
    T(slide, *LEAD_XY, [one(l, 12, None, INK, ls=1.28) for l in lead],
      anchor="m", ml=0, mr=0)
    ln = slide.shapes.add_connector(1, Cm(0), Cm(DIV_Y), Cm(SW), Cm(DIV_Y))
    ln.line.color.rgb = RGBColor.from_string(BORDER)
    ln.line.width = Pt(1.0)


def foot(slide, text):
    T(slide, CX0, FOOT_Y, CW, 0.9, [one(text, 7.5, None, MUT, ls=1.15)], ml=0, mr=0)


def badge(slide, x, y, w, h, text, fill=PORANGE, col=ORANGE, sz=8.5, dash=None):
    sp = box(slide, x, y, w, h, fill=fill, line=col, lw=1.0, dash=dash)
    put_text(sp.text_frame, [one(text, sz, True, col, align="c")],
             anchor="m", ml=0.06, mr=0.06, mt=0, mb=0)
    return sp


def band(slide, y, text, fill=NAVY, col=WHITE, sz=11.5, h=0.86, x=CX0, w=CW):
    sp = box(slide, x, y, w, h, fill=fill, radius=0.10)
    put_text(sp.text_frame, [one(text, sz, True, col, align="c")],
             anchor="m", ml=0.2, mr=0.2, mt=0, mb=0)
    return sp


def placeholder(slide, x, y, w, h, label, note):
    """★未取得データの差込枠（破線）。捏造しない。"""
    sp = box(slide, x, y, w, h, fill=WHITE, line=MUT, lw=1.25,
              dash=MSO_LINE_DASH_STYLE.DASH)
    put_text(sp.text_frame,
              [one(label, 11, True, MUT, align="c", sa=4),
               one(note, 8.5, None, MUT, align="c", ls=1.25)],
              anchor="m", ml=0.3, mr=0.3, mt=0.1, mb=0.1)
    return sp


def phone(slide, x, y, w, h, header, lines):
    """LINEトーク画面のモック。lines = [(kind, text)]
    kind: 'in'（BOT吹き出し）/ 'chip'（選択肢）/ 'btn'（CTA）/ 'note'（注記）"""
    box(slide, x, y, w, h, fill=BEZEL, radius=0.10)
    box(slide, x + 0.14, y + 0.42, w - 0.28, h - 0.56, fill=SCREEN, radius=0.02,
        shape=MSO_SHAPE.RECTANGLE)
    hd = box(slide, x + 0.14, y + 0.42, w - 0.28, 0.52, fill=GREEN,
             shape=MSO_SHAPE.RECTANGLE)
    put_text(hd.text_frame, [one(header, 7, True, WHITE, align="c")],
             anchor="m", ml=0.05, mr=0.05, mt=0, mb=0)
    cy = y + 1.05
    iw = w - 0.62
    for kind, text in lines:
        nlines = text.count("\n") + 1
        bh = 0.30 + 0.30 * nlines
        if kind == "in":
            sp = box(slide, x + 0.31, cy, iw, bh, fill=WHITE, line=BORDER, lw=0.75)
            put_text(sp.text_frame, [one(l, 6.5, None, INK, ls=1.18)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.14, mr=0.10, mt=0.05, mb=0.05)
        elif kind == "chip":
            sp = box(slide, x + 0.31, cy, iw, bh, fill=WHITE, line=GREEN, lw=0.9)
            put_text(sp.text_frame, [one(l, 6.5, None, "0B7A3B", align="c", ls=1.15)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.08, mr=0.08, mt=0.03, mb=0.03)
        elif kind == "btn":
            sp = box(slide, x + 0.31, cy, iw, bh, fill=GREEN)
            put_text(sp.text_frame, [one(l, 6.8, True, WHITE, align="c", ls=1.15)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.08, mr=0.08, mt=0.03, mb=0.03)
        else:  # note
            sp = slide.shapes.add_textbox(Cm(x + 0.31), Cm(cy), Cm(iw), Cm(bh))
            put_text(sp.text_frame, [one(l, 6.2, None, MUT, ls=1.15)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.05, mr=0.05, mt=0, mb=0)
        cy += bh + 0.12
    return cy


def richmenu(slide, x, y, w, h, tabs):
    """リッチメニュー3タブ×6ボタンのモック。
    tabs = [(タブ名, [ボタン文言6個]), ...]。最初のタブをアクティブ表示。"""
    tab_h = 0.55
    box(slide, x, y, w, h, fill=BEZEL, radius=0.06)
    tw = (w - 0.12) / len(tabs)
    for i, (name, _) in enumerate(tabs):
        active = (i == 0)
        tb = box(slide, x + 0.06 + i * tw, y + 0.06, tw - 0.03, tab_h,
                  fill=GREEN if active else "3F3F3F",
                  shape=MSO_SHAPE.RECTANGLE)
        put_text(tb.text_frame, [one(name, 8, True, WHITE, align="c")],
                 anchor="m", ml=0, mr=0, mt=0, mb=0)
    _, btns = tabs[0]
    gx, gy = x + 0.06, y + 0.06 + tab_h + 0.06
    gw, gh = w - 0.12, h - tab_h - 0.18
    cols, rows = 3, 2
    bw, bh = (gw - 0.06 * (cols - 1)) / cols, (gh - 0.06 * (rows - 1)) / rows
    for i, label in enumerate(btns):
        r, c = divmod(i, cols)
        bb = box(slide, gx + c * (bw + 0.06), gy + r * (bh + 0.06), bw, bh,
                  fill=SCREEN, line=BORDER, lw=0.75, radius=0.04)
        put_text(bb.text_frame, [one(l, 7.2, True, NAVY, align="c", ls=1.1)
                                 for l in label.split("\n")],
                 anchor="m", ml=0.05, mr=0.05, mt=0, mb=0)


def simple_table(slide, x, y, w, h, headers, rows, col_w=None,
                  hsz=9, bsz=8.5, header_fill=NAVY, zebra=PALE,
                  align=None, row_h=None):
    """罫線テーブル（ネイティブpptxテーブル）。headers=[str], rows=[[str]]"""
    n_r, n_c = len(rows) + 1, len(headers)
    gf = slide.shapes.add_table(n_r, n_c, Cm(x), Cm(y), Cm(w), Cm(h))
    tbl = gf.table
    if col_w:
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Cm(cw)
    if row_h:
        for r in tbl.rows:
            r.height = Cm(row_h)
    for j, htext in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor.from_string(header_fill)
        c.margin_left = c.margin_right = Cm(0.12)
        c.margin_top = c.margin_bottom = Cm(0.05)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        put_text(c.text_frame, [one(htext, hsz, True, WHITE,
                                    align=(align[j] if align else "c"))],
                 anchor="m", ml=0.12, mr=0.12, mt=0.02, mb=0.02)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor.from_string(
                zebra if (zebra and i % 2 == 0) else WHITE)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            put_text(c.text_frame, [one(str(val), bsz, None, INK,
                                        align=(align[j] if align else "l"))],
                     anchor="m", ml=0.12, mr=0.12, mt=0.02, mb=0.02)
    return gf


def find_shape(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle in sh.text_frame.text:
            return sh
    return None


# ============================================================
# S01 表紙
# ============================================================
s = slides[0]
sh = find_shape(s, "業界：")
if sh:
    put_text(sh.text_frame, [one("業界：EC（総合EC・業界汎用）", 12, True, TNAVY)], anchor="m")
sh = find_shape(s, "最上位パートナーの知見")
if sh:
    put_text(sh.text_frame,
             [one("広告費ゼロで、カゴ落ちと1回きりの顧客を拾い直し、LTV÷CPOを引き上げる", 13, True, INK)],
             anchor="m")

# ============================================================
# S02 資料アジェンダ
# ============================================================
s = slides[1]
frame(s, "資料アジェンダ", ["この資料が言っていること：新規獲得の単価は上がり続けている。",
                          "もう「集める」では勝てない。勝負はCPOではなく、LTV ÷ CPO。"])
AGENDA = [
    ("1", "なぜLINEか", "S03-06"), ("2", "実態調査", "S07-12"),
    ("3", "全体設計", "S13-15"), ("4", "構築", "S16-20"),
    ("5", "配信設計", "S21-26"), ("6", "歩留まり・工数", "S27-28"),
    ("7", "成果と体制", "S29-32"), ("8", "締め", "S33-36"),
]
gx, gy, gw, gh = CX0, CY0 + 0.3, CW, 11.6
cw = gw / 4
ch = gh / 2
for i, (n, t, r) in enumerate(AGENDA):
    row, col = divmod(i, 4)
    x = gx + col * cw
    y = gy + row * ch
    box(s, x + 0.12, y + 0.12, cw - 0.24, ch - 0.24, fill=PALE)
    T(s, x + 0.12, y + 0.22, cw - 0.24, 0.9,
      [one(n, 22, True, ORANGE, align="c")], anchor="t")
    T(s, x + 0.12, y + 0.95, cw - 0.24, 1.3,
      [one(t, 12.5, True, NAVY, align="c")], anchor="t")
    T(s, x + 0.12, y + ch - 0.55, cw - 0.24, 0.5,
      [one(r, 8.5, None, MUT, align="c")], anchor="t")
foot(s, "全8章・36枚｜主語は「LINE公式アカウントの運用」。広告運用の提案書ではない")

# ============================================================
# S03 市場データ（ニーズ全体図）
# ============================================================
s = slides[2]
frame(s, "市場データ｜ニーズ全体図",
      ["国内BtoC-EC市場は拡大を続けているが、物販系のEC化率はまだ1割弱。",
       "「市場は伸びている。伸びていないのは、取りこぼしを回収する設計」。"])
nums = [
    ("BtoC-EC市場全体", "26.1兆円", "前年比 +5.1%（前年24.8兆円）"),
    ("物販系分野", "15兆2,194億円", "前年比 +3.70%"),
    ("物販系のEC化率", "9.78%", "前年比 +0.40pt"),
]
cw3 = (CW - 0.6) / 3
for i, (label, val, sub) in enumerate(nums):
    x = CX0 + i * (cw3 + 0.3)
    box(s, x, CY0 + 0.2, cw3, 4.6, fill=PALE)
    T(s, x, CY0 + 0.5, cw3, 0.7, [one(label, 12, True, NAVY, align="c")], anchor="m")
    T(s, x, CY0 + 1.3, cw3, 1.7, [one(val, 32, True, ORANGE, align="c")], anchor="m")
    T(s, x, CY0 + 3.2, cw3, 1.3, [one(sub, 10, None, INK, align="c", ls=1.2)], anchor="m")
box(s, CX0, CY0 + 5.2, CW, 1.7, fill="FFF2CC")
T(s, CX0 + 0.3, CY0 + 5.2, CW - 0.6, 1.7,
  [one("EC化率9.78%＝まだ9割はオフライン。市場は伸びている。伸びていないのは回収の設計。", 15, True, TNAVY, ls=1.3)],
  anchor="m")
foot(s, "業界水準｜出典：経済産業省「令和6年度 電子商取引に関する市場調査」2025年8月26日公表。"
        "毎年8月末公表のため、提案時は最新年度（令和7年度）を確認すること")

# ============================================================
# S04 ECのWebマーケ構造（3つの壁）
# ============================================================
s = slides[3]
frame(s, "ECのWebマーケ構造｜3つの壁",
      ["広告費を払って連れてきた人が、カゴに入れたまま多くが消え、",
       "買ってくれた人の多くも2回目を買わずに消えている。一番お金をかけた直後に、一番大きく落ちている。"])
walls = [
    ("① 獲得単価の高騰", "検索広告CPCは10年で168%。\n枠の奪い合いで下がらない", RED),
    ("② 取りこぼし", "カゴ落ち率は多くのサイトで\n60〜70%前後", ORANGE),
    ("③ 使い捨て", "初回購入者の6〜7割が\n2回目を買わずに消える", NAVY),
]
cw3 = (CW - 0.6) / 3
for i, (h, b, c) in enumerate(walls):
    x = CX0 + i * (cw3 + 0.3)
    wb = box(s, x, CY0 + 0.2, cw3, 2.6, fill=WHITE, line=c, lw=1.5)
    put_text(wb.text_frame,
             [one(h, 13, True, c, align="c", sa=6)] +
             [one(l, 10, None, INK, align="c", ls=1.25) for l in b.split("\n")],
             anchor="m")
box(s, CX0, CY0 + 3.2, CW, 3.6, fill=PRED, line=RED, lw=1.0)
T(s, CX0 + 0.4, CY0 + 3.4, CW - 0.8, 1.0,
  [one("カゴ落ち率：多くのサイトで60〜70%前後", 20, True, RED, align="c")], anchor="m")
T(s, CX0 + 0.4, CY0 + 4.5, CW - 0.8, 2.1,
  [one("グローバル平均 約70.22%（Baymard Institute）／国内調査 約63.3%（イー・エージェンシー）", 11, None, INK, align="c", ls=1.3),
   one("広告費をかけて連れてきた人の6〜7割が、購入直前で離脱している", 13, True, TNAVY, align="c", sa=4)],
  anchor="m")
foot(s, "業界水準｜出典：Baymard Institute（グローバル）／イー・エージェンシー（国内）。"
        "検索広告CPCはキーワードマーケティング調べ")

# ============================================================
# S05 EC購入者の心理
# ============================================================
s = slides[4]
frame(s, "EC購入者の心理",
      ["広告費を払って連れてきた人が、カゴに入れたまま約7割消え、買ってくれた人も",
       "6〜7割が2回目を買わずに消えている。しかも全員、連絡先を残さないまま消えている。"])
psy = [
    ("「まだ決めきれない」", "カゴ落ち理由の約43%が\n「ただ見ていた／買う準備が\nできていない」"),
    ("「いくらになるか\n分からないのが怖い」", "カゴ落ち理由1位（53%）が\n送料・税・手数料など\n追加コストの不透明さ"),
    ("「メールは\n登録したくない」", "国内ECのメルマガ開封率は\n20%前後。登録しても\n読まれていない"),
]
cw3 = (CW - 0.6) / 3
for i, (h, b) in enumerate(psy):
    x = CX0 + i * (cw3 + 0.3)
    card(s, x, CY0 + 0.2, cw3, 4.3, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=13, bsz=10.5, ls=1.35)
box(s, CX0, CY0 + 4.9, CW, 2.0, fill=NAVY)
T(s, CX0 + 0.4, CY0 + 4.9, CW - 0.8, 2.0,
  [one("カゴ落ちの過半は「商品の問題」ではなく「情報と接触機会の問題」", 16, True, WHITE, align="c", ls=1.3)],
  anchor="m")
foot(s, "業界水準｜出典：Baymard Institute（カゴ落ち理由）／国内ECメルマガ開封率一般値。"
        "倍率で語らず絶対値のみ使用（LINE当日中開封率と混ぜない）")

# ============================================================
# S06 できること／できないこと ＋ スコープ宣言
# ============================================================
s = slides[5]
frame(s, "Webマーケで「できること／できないこと」",
      ["ここが本提案の誠実さであり、説得力の源。触るのは「カゴに入れた後」と",
       "「1回目を買った後」の2箇所だけ。広告費も流入もLPも触らない。"])
box(s, CX0, CY0 + 0.2, 8.0, 6.5, fill=GREY)
T(s, CX0 + 0.3, CY0 + 0.4, 7.4, 0.6, [one("できないこと", 13, True, MUT)], anchor="m")
for i, t in enumerate(["商品力・価格競争力", "送料・配送スピード・返品条件（物流）",
                       "モール手数料率（言い値）", "在庫・欠品"]):
    T(s, CX0 + 0.3, CY0 + 1.1 + i * 1.3, 7.4, 1.2,
      [one("✕ " + t, 11, None, INK, ls=1.25)], anchor="m")
bx2 = CX0 + 8.3
box(s, bx2, CY0 + 0.2, CW - 8.3, 6.5, fill=PALE)
T(s, bx2 + 0.3, CY0 + 0.4, CW - 8.3 - 0.6, 0.6, [one("できること（＝本提案の範囲）", 13, True, NAVY)], anchor="m")
can = [
    ("カゴ落ちを拾う", "CPO"), ("メールで届かない人に届ける（当日中に約8割開封）", "全段"),
    ("初回購入の背中を押す", "CPO"), ("F2転換を自動化する", "F2転換率・LTV"),
    ("定期の解約を止める／休眠を起こす", "LTV"), ("1stパーティデータを持つ", "翌期のCPO"),
]
for i, (t, kpi) in enumerate(can):
    y = CY0 + 1.05 + i * 0.88
    T(s, bx2 + 0.3, y, CW - 8.3 - 3.3, 0.85,
      [one(f"{i+1}. {t}", 10, None, INK, ls=1.15)], anchor="m")
    badge(s, bx2 + CW - 8.3 - 2.8, y + 0.12, 2.5, 0.55, kpi, fill=WHITE, col=ORANGE, sz=8)
box(s, CX0, CY0 + 6.9, CW, 1.3, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 6.9, CW - 0.6, 1.3,
  [one("触るのは「カゴに入れた後」と「1回目を買った後」の2箇所だけ。広告費も流入もLPも触らない。", 12.5, True, WHITE, align="c", ls=1.25)],
  anchor="m")
foot(s, "意見｜DYMの提案スコープ宣言。「もっと広告を回しましょう」ではなく、左記6項目を自動化してLTV/CPOを上げる提案")

# ============================================================
# S07 CPC推移
# ============================================================
s = slides[6]
frame(s, "市場環境｜検索広告CPCの推移",
      ["広告予算の半分以上がネット広告に集中し、枠の奪い合いで単価は下がらない。",
       "新規獲得の単価は構造的に上がり続ける。"])
nums = [
    ("総広告費", "8兆623億円", "前年比 +5.1%（4年連続で過去最高）"),
    ("インターネット広告費", "4兆459億円", "前年比 +10.8%（初の4兆円超）"),
    ("総広告費に占める構成比", "50.2%", "初の過半数"),
]
cw3 = (CW - 0.6) / 3
for i, (label, val, sub) in enumerate(nums):
    x = CX0 + i * (cw3 + 0.3)
    box(s, x, CY0 + 0.2, cw3, 3.6, fill=PALE if i != 2 else PORANGE)
    T(s, x, CY0 + 0.4, cw3, 0.6, [one(label, 11.5, True, NAVY, align="c")], anchor="m")
    T(s, x, CY0 + 1.05, cw3, 1.4,
      [one(val, 27, True, (ORANGE if i == 2 else NAVY), align="c")], anchor="m")
    T(s, x, CY0 + 2.5, cw3, 1.2, [one(sub, 9.5, None, INK, align="c", ls=1.2)], anchor="m")
box(s, CX0, CY0 + 4.0, CW, 2.9, fill=NAVY)
T(s, CX0 + 0.4, CY0 + 4.15, CW - 0.8, 0.6,
  [one("検索広告の平均クリック単価は10年前比 168%（約1.7倍）", 13, True, WHITE)], anchor="m")
T(s, CX0 + 0.4, CY0 + 4.85, CW - 0.8, 1.9,
  [one("広告予算の半分以上（構成比50.2%）がネット広告に集中している＝枠の奪い合い。", 12, None, "E9EDF7", ls=1.35, sa=3),
   one("だから単価は下がらない。回収側（LTV÷CPO）に回すのがこの資料の背骨。", 12, True, WHITE, ls=1.35)],
  anchor="m")
foot(s, "業界水準｜出典：電通「2025年 日本の広告費」2026年3月5日発表／検索広告CPCはキーワードマーケティング調べ。"
        "海外CPA/CPMデータは使用していない")

# ============================================================
# S08 広告実績CVR（★未取得・差込枠）
# ============================================================
s = slides[7]
frame(s, "市場環境｜広告実績CVR（CV地点記載）",
      ["DYMのEC支援実績から、CV地点別のCPC・CVR・CPOを提示する。",
       "カート投入・初回購入・定期引上のどこで刈れているかを切り分けて見る。"])
simple_table(s, CX0, CY0 + 0.2, CW, 2.0,
             ["CV地点", "CPC", "CVR", "CPO"],
             [["① カート投入", "－", "－", "－"],
              ["② 初回購入", "－", "－", "－"],
              ["③ 定期引上", "－", "－", "－"]],
             col_w=[CW * 0.28, CW * 0.24, CW * 0.24, CW * 0.24], row_h=0.5)
placeholder(s, CX0, CY0 + 3.0, CW, 3.9, "★差込枠｜DYM社内のEC広告実績（未取得）",
            "CV地点別（カート投入／初回購入／定期引上）のCPC・CVR・CPOを記入。\n"
            "CV地点の明記は必須（スキル §6.5）。上司確認①として要確認リストへ")
foot(s, "★未取得｜出典：DYM社内実績（要取得・上司確認①）。取得まで数値は記入しない")

# ============================================================
# S09 広告審査面の懸念（薬機法）
# ============================================================
s = slides[8]
frame(s, "市場環境｜広告審査面の懸念",
      ["「LINEは友だちだけが見るから広告ではない」は誤り。",
       "LINE配信の大半は薬機法上の「広告」に該当すると考えるべき。"])
simple_table(s, CX0, CY0 + 0.2, CW, 2.1,
             ["要件", "中身", "LINE配信での該当性"],
             [["① 誘引性", "顧客の購入意欲を昂進させる意図が明確", "商品案内・購入促進の内容なら該当"],
              ["② 特定性", "特定の商品名が明らかにされている", "送信元が事業者の時点で自動的に該当"],
              ["③ 認知性", "一般人が認知できる状態", "友だちに配信された時点で該当"]],
             col_w=[CW * 0.16, CW * 0.44, CW * 0.40])
box(s, CX0, CY0 + 2.5, CW, 1.5, fill=PRED, line=RED, lw=1.0)
T(s, CX0 + 0.3, CY0 + 2.5, CW - 0.6, 1.5,
  [one("送信元が事業者である時点で特定性は自動的に満たされ、商品案内や購入促進の内容なら誘引性も満たす。", 11.5, True, RED, ls=1.3)],
  anchor="m")
T(s, CX0, CY0 + 4.2, CW, 1.5,
  [one("併せて掛かる法律：景品表示法（優良誤認・有利誤認／No.1表示／ビフォーアフター／打消し表示）、健康増進法（誇大表示）。", 10.5, None, INK, ls=1.35, sa=3),
   one("健康食品は薬機法の直接の規制対象ではないが、医薬品的な効能を標ぼうすると規制対象になる。ECモールの商品ページのレビュー・口コミにも規制が及ぶ場合がある。", 10.5, None, INK, ls=1.35)],
  anchor="t")
T(s, CX0, CY0 + 6.0, CW, 0.9,
  [one("だから配信文面の設計は、最初から表現規制を織り込む必要がある。＝LINEOA運用設計の制約条件。", 12.5, True, TNAVY, ls=1.3)],
  anchor="m")
foot(s, "業界水準｜出典：薬機法の広告3要件（薬監発第148号）／京都府「広告の３要件」")

# ============================================================
# S10 シーズナリティ（年間商戦カレンダー＋★Googleトレンド差込枠）
# ============================================================
s = slides[9]
frame(s, "ニーズ調査｜シーズナリティ",
      ["ECの特徴は「11〜12月に山が2つ立て続けに来る」こと。8月は仕込み月。",
       "S25（年間企画投稿カレンダー）の設計図になる。"])
CAL = [
    ("1月", "初売り・福袋", False), ("2月", "バレンタイン", False),
    ("3月", "楽天SS／新生活", True), ("4月", "新生活後半", False),
    ("5月", "母の日／お買い物マラソン", True), ("6月", "楽天SS／父の日／お中元", True),
    ("7月", "Amazonプライムデー", True), ("8月", "夏枯れ（秋商戦の仕込み）", False),
    ("9月", "楽天SS", False), ("10月", "ハロウィン", False),
    ("11月", "ブラックフライデー", True), ("12月", "楽天大感謝祭／年末", True),
]
cols = 6
cw6 = CW / cols
ch2 = 1.55
for i, (m, ev, peak) in enumerate(CAL):
    r, c = divmod(i, cols)
    x = CX0 + c * cw6
    y = CY0 + 0.2 + r * (ch2 + 0.12)
    box(s, x + 0.05, y, cw6 - 0.1, ch2, fill=(PORANGE if peak else PALE))
    T(s, x + 0.15, y + 0.05, cw6 - 0.3, 0.5,
      [one(m, 11, True, (ORANGE if peak else NAVY))], anchor="t")
    T(s, x + 0.15, y + 0.5, cw6 - 0.3, 1.0,
      [one(ev, 8.5, None, INK, ls=1.15)], anchor="t")
placeholder(s, CX0, CY0 + 3.7, CW, 3.2, "★差込枠｜Googleトレンド（5年＋1年）",
            "KW＝福袋／ブラックフライデー／セール／母の日（すべて一語）。\n"
            "「セール」は他KWを潰すため必ず単独グラフにする")
foot(s, "業界水準｜年間商戦カレンダーは公開の商戦・セール時期。★Googleトレンドは未取得（要取得）")

# ============================================================
# S11 前後検索（★未取得・差込枠）
# ============================================================
s = slides[10]
frame(s, "ニーズ調査｜前後検索",
      ["LINEヤフーの前後検索データで、購入前後の不安の中身を確認する。",
       "本命は「解約」。定期の解約不安がどこで生まれているかを見る。"])
bands3 = [("検索“前”", "口コミ／効果／比較", MUT), ("起点", "通販／定期便", NAVY),
          ("検索“後”", "解約／返品", RED)]
cw3b = (CW - 0.6) / 3
for i, (h, sub, c) in enumerate(bands3):
    x = CX0 + i * (cw3b + 0.3)
    bb = box(s, x, CY0 + 0.2, cw3b, 1.5, fill=c)
    put_text(bb.text_frame, [one(h, 13, True, WHITE, align="c")], anchor="m")
    T(s, x, CY0 + 1.8, cw3b, 0.7, [one(sub, 10, None, INK, align="c")], anchor="m")
placeholder(s, CX0, CY0 + 2.8, CW, 4.1, "★差込枠｜LINEヤフー前後検索データ（未取得）",
            "KW＝通販／定期便／解約／口コミ（すべて一語）。★本命は「解約」。\n"
            "購入後にどんな不安・情報を検索しているかを反映する")
foot(s, "★未取得｜出典：LINEヤフー媒体資料（要取得）")

# ============================================================
# S12 他社分析（10社・友だち数は★差込枠）
# ============================================================
s = slides[11]
frame(s, "他社分析｜競合10社のLINE運用ステータス",
      ["知名度優先で選定した10社。友だち数・配信内容は実測が必要（page.line.meから取得日つきで取得）。",
       "ID連携しているか（購入導線がLINEログインを求めるか）を必ず記録する。"])
COMP = [
    ("オルビス", "診断ドリブン型", "肌診断→出し分け"), ("ファンケル", "会員・定期型", "次回発送前フォロー"),
    ("北の快適工房", "単品リピート通販型", "F2オファーの時期・文面"), ("BASE FOOD", "D2Cサブスク型", "継続率を上げる配信"),
    ("Oisix", "食品定期型", "週次の締切リマインド"), ("ニトリ", "クーポン会員型", "会員証・店舗連動"),
    ("無印良品", "ミニアプリ・会員証型", "LINEミニアプリの使い方"), ("ZOZOTOWN", "配信メディア型", "配信頻度とセグメント"),
    ("ユニクロ", "クーポン型（大手）", "配信頻度・クーポン"), ("SHIRO", "中堅コスメ・運用特徴型", "世界観重視の配信設計"),
]
rows = [[c[0], c[1], c[2], "★未取得", "★未取得"] for c in COMP]
simple_table(s, CX0, CY0 + 0.2, CW, 6.1,
             ["アカウント", "想定の型", "見るポイント", "友だち数", "ID連携"],
             rows, col_w=[CW * 0.16, CW * 0.24, CW * 0.32, CW * 0.14, CW * 0.14],
             bsz=8, row_h=0.5)
foot(s, "選定は知名度優先（スキル§3で事前承認済み）。友だち数・ID連携は page.line.me から取得日つきで実測（★未取得）。"
        "「おすすめ」欄の数値は拾わない")

# ============================================================
# S13 カスタマージャーニー フェーズ8分類 × CV3段
# ============================================================
s = slides[12]
frame(s, "全体設計｜カスタマージャーニー × CV3段",
      ["ECの主戦場は②初回購入ではなく③F2。山（LINEが効く度合い）は",
       "初回購入ではなくF2購入・定期化のフェーズにある。他業界資料と重心が違う。"])
PH = ["認知\n（広告接触）", "サイト来訪", "比較検討\n（カゴに入れる）", "カゴ落ち\n（対策ゾーン）",
      "① 初回購入", "受け取り・\n使用開始", "再検討\n（消費サイクル）", "② F2購入・\n定期化"]
HEI = [1.0, 1.4, 1.8, 2.2, 2.6, 2.2, 2.8, 3.6]
n = len(PH)
cwp = CW / n
base_y = CY0 + 5.6
for i, (label, h) in enumerate(zip(PH, HEI)):
    x = CX0 + i * cwp
    bar = box(s, x + 0.08, base_y - h, cwp - 0.16, h,
              fill=(ORANGE if i == n - 1 else NAVY), radius=0.08)
    T(s, x + 0.02, base_y + 0.08, cwp - 0.04, 1.1,
      [one(l, 8, None, INK, align="c", ls=1.1) for l in label.split("\n")], anchor="t")
badge(s, CX0 + 1 * cwp - 0.55, CY0 + 0.15, 4.2, 0.65, "①軽CV：友だち追加", fill=WHITE, col=NAVY, sz=8.5)
badge(s, CX0 + 4 * cwp - 1.0, CY0 + 0.15, 4.2, 0.65, "②主CV：初回購入", fill=WHITE, col=NAVY, sz=8.5)
badge(s, CX0 + 7 * cwp - 1.1, CY0 + 0.15, 4.2, 0.65, "③最終CV：F2・定期", fill=PORANGE, col=ORANGE, sz=8.5)
box(s, CX0 + 5.6 * cwp, CY0 + 7.2, 2.4 * cwp, 1.1, fill=PORANGE, line=ORANGE, lw=1.0)
T(s, CX0 + 5.6 * cwp, CY0 + 7.2, 2.4 * cwp, 1.1,
  [one("★ここが本資料の主戦場", 10.5, True, RED, align="c", ls=1.2)], anchor="m")
foot(s, "意見｜CV3段はDYM設計。①LINE友だち追加＋初回クーポン／②初回購入（赤字でよい）／"
        "③F2→定期（ここで初めて黒字）")

# ============================================================
# S14 施策全体像 ＋ 対策領域マップ
# ============================================================
s = slides[13]
frame(s, "全体設計｜施策全体像・対策領域マップ",
      ["触るのは2領域だけ。①カゴ落ち対策（取りこぼし回収）と②F2転換・LTV化（使い捨て解消）。",
       "広告費も流入もLPも変えない。"])
areas = [
    ("① カゴ落ち対策", "取りこぼしの回収（CPO）", [
        "サイトポップアップ／購入完了画面でLINE友だち化",
        "あいさつMSGで送料条件を先出し",
        "カゴ落ちリマインド（3通・ID連携が前提）",
    ], NAVY, PALE),
    ("② F2転換・LTV化", "使い捨ての解消（F2転換率・LTV）", [
        "購入者への使用フォロー配信",
        "消費サイクル連動のF2オファー",
        "定期の解約防止・休眠掘り起こし",
    ], ORANGE, PORANGE),
]
cw2 = (CW - 0.4) / 2
for i, (h, sub, items, c, fill) in enumerate(areas):
    x = CX0 + i * (cw2 + 0.4)
    box(s, x, CY0 + 0.2, cw2, 6.6, fill=fill)
    T(s, x + 0.3, CY0 + 0.4, cw2 - 0.6, 0.7, [one(h, 15, True, c)], anchor="m")
    T(s, x + 0.3, CY0 + 1.15, cw2 - 0.6, 0.7, [one(sub, 10.5, None, MUT)], anchor="m")
    for j, it in enumerate(items):
        T(s, x + 0.3, CY0 + 2.1 + j * 1.5, cw2 - 0.6, 1.4,
          [one(f"・{it}", 11, None, INK, ls=1.3)], anchor="t")
foot(s, "意見｜DYM提案の対策領域マップ。S16-20（構築）・S21-26（配信設計）が①②それぞれに対応する")

# ============================================================
# S15 施策展開図（初期・月次）
# ============================================================
s = slides[14]
frame(s, "全体設計｜施策展開図（初期・月次）",
      ["初期構築は1ヶ月、以降は月次運用。やることは決まっていて増え続けない。"])
band(s, CY0 + 0.2, "初期（初動・構築）", fill=NAVY, h=0.75)
init_items = ["あいさつメッセージ設計（送料条件先出し）", "診断・アンケート設計（タグ設計）",
              "リッチメニュー3タブ18ボタン", "LINE ID連携の基盤構築", "カゴ落ちリマインド3通の設計"]
for i, it in enumerate(init_items):
    card(s, CX0 + (i % 3) * (CW / 3), CY0 + 1.2 + (i // 3) * 1.9, CW / 3 - 0.2, 1.7,
         f"{i+1}", it, hcol=NAVY, fill=PALE, hsz=13, bsz=10, ls=1.25)
band(s, CY0 + 5.0, "月次（定例運用）", fill=ORANGE, col=WHITE, h=0.75)
month_items = ["年間企画投稿（商戦カレンダー連動）", "F2オファー配信（消費サイクル連動）",
               "通知メッセージ運用（発送・再入荷）", "定例レポート・改善提案"]
for i, it in enumerate(month_items):
    card(s, CX0 + (i % 4) * (CW / 4), CY0 + 6.0, CW / 4 - 0.2, 1.4,
         f"{i+1}", it, hcol=ORANGE, fill=PORANGE, hsz=12, bsz=9, ls=1.2)
foot(s, "意見｜DYM標準の展開図をEC施策に置換")

# ============================================================
# S16 友だち追加動線
# ============================================================
s = slides[15]
frame(s, "構築｜友だち追加動線",
      ["5本の動線で友だちを増やす。購入完了画面と同梱物QRが最速で貯まる。"])
routes = [
    ("① サイトポップアップ", "離脱意図を検知して\nLINE追加バナーを表示"),
    ("② 購入完了画面", "「発送状況をLINEで\nお届けします」で誘導"),
    ("③ 同梱物QR", "梱包物に同梱。\n受け取り後の最速接点"),
    ("④ LINE広告（CPF）", "友だち追加動線として\nタイムラインに配信"),
    ("⑤ 通知メッセージ経由", "非友だちにも届く通知から\nLINE本体へ誘導"),
]
cw5 = (CW - 0.4 * 4) / 5
for i, (h, b) in enumerate(routes):
    x = CX0 + i * (cw5 + 0.4)
    card(s, x, CY0 + 0.3, cw5, 5.0, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=10.5, bsz=9.5, ls=1.3)
box(s, CX0, CY0 + 5.6, CW, 1.3, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 5.6, CW - 0.6, 1.3,
  [one("③購入完了画面・④同梱物QRが最速で貯まる。ここを外すと以降が全部乗らない。", 12, True, WHITE, align="c", ls=1.25)],
  anchor="m")
foot(s, "意見｜DYM提案の友だち追加動線設計")

# ============================================================
# S17 あいさつメッセージ
# ============================================================
s = slides[16]
frame(s, "構築｜あいさつメッセージ",
      ["初回クーポン＋「メルマガのような配信はしません」宣言＋送料条件の先出し。",
       "カゴ落ち理由1位（追加コストの不透明さ）にここで先回りする。"])
phone(s, CX0 + 0.3, CY0 + 0.2, 8.6, 6.6, "公式アカウント",
      [("in", "友だち追加ありがとうございます！\n初回限定10%OFFクーポンです🎁"),
       ("in", "送料：全国一律550円\n11,000円以上のご購入で送料無料"),
       ("note", "※メルマガのような一斉配信はしません。\nお役立ち情報のみ、月数回お届けします"),
       ("btn", "クーポンを使って探す")])
box(s, CX0 + 9.6, CY0 + 0.2, CW - 9.6, 6.6, fill=PALE)
T(s, CX0 + 9.9, CY0 + 0.4, CW - 9.9 - 0.3, 0.6, [one("設計の要点", 12.5, True, NAVY)], anchor="m")
pts = [
    ("初回クーポンで即メリット", "登録動機を明確化"),
    ("「メルマガのような配信はしません」宣言", "メルマガ開封率20%前後への\n忌避感を先に解消"),
    ("送料条件を最初から先出し", "カゴ落ち理由1位（53%・追加\nコストの不透明さ）に先回り"),
]
for i, (h, b) in enumerate(pts):
    T(s, CX0 + 9.9, CY0 + 1.2 + i * 1.8, CW - 9.9 - 0.3, 1.7,
      [one(f"・{h}", 11, True, NAVY, ls=1.25, sa=2)] +
      [one(l, 9.5, None, INK, ls=1.2) for l in b.split("\n")], anchor="t")
foot(s, "意見｜配信文面は初稿。カゴ落ち理由の出典：Baymard Institute")

# ============================================================
# S18 診断・アンケート
# ============================================================
s = slides[17]
frame(s, "構築｜診断・アンケート",
      ["肌・サイズ・好み・悩みを4問で聞き、タグ化する。以降の配信を全部自動で出し分けられる。"])
Q = [
    ("Q1 肌質・タイプは？", "乾燥／脂性／混合／敏感／わからない", "商品レコメンドの軸"),
    ("Q2 サイズ・容量の目安は？", "少量から試したい／いつも使う量が\n決まっている／大容量派", "F2オファーのタイミング設計に活用"),
    ("Q3 好みのテイスト・悩みは？", "選択式・複数選択可", "企画配信・レコメンドのタグ"),
    ("Q4 購入のきっかけは？", "広告／口コミ／SNS／ギフト", "流入チャネル別の育成シナリオ分岐"),
]
cw2 = (CW - 0.4) / 2
for i, (h, opt, use) in enumerate(Q):
    r, c = divmod(i, 2)
    x = CX0 + c * (cw2 + 0.4)
    y = CY0 + 0.2 + r * 3.3
    card(s, x, y, cw2, 3.0, h, [opt, "→ " + use], hcol=NAVY, fill=PALE, hsz=11.5, bsz=9.5, ls=1.3)
foot(s, "意見｜診断結果→タグ化→以降の配信を自動で出し分け")

# ============================================================
# S19 リッチメニュー
# ============================================================
s = slides[18]
frame(s, "構築｜リッチメニュー",
      ["3タブ×6ボタン＝18ボタンすべてに役割がある。左上は「今すぐ探す」に相当する導線。"])
richmenu(s, CX0 + 1.5, CY0 + 0.3, 12.0, 6.5, [
    ("買う", ["今すぐ\n探す", "新着\n商品", "セール\n情報", "あなたへの\nおすすめ", "クーポン", "再入荷\n通知登録"]),
    ("わたしの", ["注文\n履歴", "配送\n状況確認", "お気に入り", "定期便\n管理", "ポイント", "会員情報"]),
    ("サポート", ["よくある\n質問", "返品・\n交換", "お問い合わせ", "送料・\n支払い", "サイズ\nガイド", "有人\nチャット"]),
])
T(s, CX0 + 14.2, CY0 + 0.5, CW - 14.2, 6.0,
  [one("設計の要点", 12.5, True, NAVY, sa=6),
   one("・買う＝新規購入の導線を集約", 10.5, None, INK, ls=1.3, sa=4),
   one("・わたしの＝配送状況・定期便管理で\n　問い合わせを自己解決", 10.5, None, INK, ls=1.3, sa=4),
   one("・サポート＝返品・サイズ等の\n　不安をKW自動応答で先回り", 10.5, None, INK, ls=1.3, sa=4),
   one("・タブ切替で3倍の情報量を\n　1画面に収める", 10.5, None, INK, ls=1.3)],
  anchor="t")
foot(s, "意見｜DYM提案のリッチメニュー設計")

# ============================================================
# S20 LINE ID連携（★心臓部）
# ============================================================
s = slides[19]
frame(s, "構築｜LINE ID連携（★心臓部）",
      ["ID連携が無いと、S23・S24・S26の施策は全部絵に描いた餅になる。",
       "ECの顧客データとLINE IDを紐づけ、LINE IDをマスタIDにしたCRM基盤にする。"])
simple_table(s, CX0, CY0 + 0.2, CW, 3.0,
             ["施策", "ID連携なし", "ID連携あり"],
             [["カゴ落ちリマインド", "✕ 打てない", "○ カート放棄をトリガーに自動配信"],
              ["購入商品に応じたF2オファー", "✕ 全員に同じ配信", "○ 買った商品に合わせて出し分け"],
              ["定期の次回発送前フォロー", "✕ 打てない", "○ 発送サイクルに合わせて自動"],
              ["購入者を配信から除外", "✕ クーポンを無駄打ち", "○ 除外できる"],
              ["再入荷通知", "✕", "○ 希望者にだけ届く"]],
             col_w=[CW * 0.32, CW * 0.28, CW * 0.40], row_h=0.5)
box(s, CX0, CY0 + 4.2, CW, 2.9, fill=PALE)
T(s, CX0 + 0.3, CY0 + 4.35, CW - 0.6, 0.5, [one("依存関係", 11.5, True, NAVY)], anchor="m")
deps = ["S23 F2オファー", "S24 カゴ落ちリマインド", "S26 通知メッセージ"]
for i, d in enumerate(deps):
    x = CX0 + 0.5 + i * 6.5
    badge(s, x, CY0 + 5.0, 5.6, 0.6, d, fill=WHITE, col=NAVY, sz=9.5)
    ln = s.shapes.add_connector(1, Cm(x + 2.8), Cm(CY0 + 5.6), Cm(CX0 + CW / 2), Cm(CY0 + 6.4))
    ln.line.color.rgb = RGBColor.from_string(ORANGE)
    ln.line.width = Pt(1.25)
badge(s, CX0 + CW / 2 - 3.0, CY0 + 6.45, 6.0, 0.6, "S20 LINE ID連携", fill=NAVY, col=WHITE, sz=10)
foot(s, "意見｜仕組み＝LINEログイン／ミニアプリ経由でECの顧客データとLINE IDを紐づけ")

# ============================================================
# S21 シナリオ2本の設計表
# ============================================================
s = slides[20]
frame(s, "配信設計｜シナリオ2本の設計表（★重点）",
      ["ECのステップ配信は「14日1本」ではない。購入前後で時間軸が別物なので、2本持つ。"])
band(s, CY0 + 0.1, "シナリオ①｜未購入者（友だち追加〜7日）ECの検討は短い。7日で刈る", fill=NAVY, h=0.65, sz=10.5)
sc1 = [
    ("① 初動", "Day0", "あいさつ＋初回クーポン（期限を切る）＋診断への誘導"),
    ("② 不安つぶし", "Day2", "口コミ・使い方・送料と返品条件（カゴ落ち理由1位に先回り）"),
    ("③ 比較対策", "Day4", "他社との違い・FAQ（「解約」「縛り」に先回り／S11と連動）"),
    ("④ クロージング", "Day6・7", "クーポン期限のリマインド。ここが刈り取りの山"),
]
cw4 = CW / 4
for i, (h, d, b) in enumerate(sc1):
    x = CX0 + i * cw4
    card(s, x, CY0 + 0.9, cw4 - 0.15, 2.5, f"{h}｜{d}", b, hcol=NAVY, fill=PALE, hsz=10, bsz=8.7, ls=1.25)
band(s, CY0 + 3.6, "シナリオ②｜購入者（F2狙い・消費サイクル連動）日数固定にしない", fill=ORANGE, col=WHITE, h=0.65, sz=10.5)
sc2 = [
    ("① サンクス", "購入直後", "発送予定・使い方。ここでレビューを求めない"),
    ("② 使用フォロー", "使い始めの頃", "「こう使うと効果が出やすい」＝離脱防止"),
    ("③ F2オファー", "無くなる少し前", "「そろそろ切れる頃です」＋再購入導線。最大の山"),
    ("④ 引き上げ", "F2後", "定期・まとめ買いへの引き上げ／レビュー依頼"),
]
for i, (h, d, b) in enumerate(sc2):
    x = CX0 + i * cw4
    card(s, x, CY0 + 4.4, cw4 - 0.15, 2.5, f"{h}｜{d}", b, hcol=ORANGE, fill=PORANGE, hsz=10, bsz=8.7, ls=1.25)
foot(s, "意見｜③F2オファーのタイミングは商材で変わる。「商品の内容量 ÷ 1日使用量」から逆算する（日数は書かない）")

# ============================================================
# S22 実文面①｜未購入者
# ============================================================
s = slides[21]
frame(s, "配信設計｜実文面①｜未購入者",
      ["Day0のクーポンとDay4の不安つぶし、Day7の期限リマインド。実際に届く文面はこれ。"])
phone(s, CX0 + 0.3, CY0 + 0.2, 8.0, 6.7, "Day0",
      [("in", "友だち追加ありがとうございます！\n初回10%OFFクーポンをお届けします🎁"),
       ("chip", "クーポンを使う"),
       ("note", "通知プレビュー：「友だち追加ありがとうございます！初回1")])
phone(s, CX0 + 8.7, CY0 + 0.2, 8.0, 6.7, "Day4",
      [("in", "「送料はいくら？」というお声を\nよくいただきます。全国一律550円、\n11,000円以上で無料です"),
       ("in", "返品は到着後8日以内、\n未開封なら送料弊社負担です"),
       ("chip", "商品を見てみる")])
phone(s, CX0 + 17.1, CY0 + 0.2, 8.0, 6.7, "Day7",
      [("in", "クーポンは本日23:59まで。\nカートに残っている商品があります"),
       ("btn", "クーポンを使って購入する"),
       ("note", "通知プレビュー：「【本日まで】10%OFFクーポンの")])
foot(s, "意見｜配信文面は初稿。絵文字・改行込みで実際の配信に近い形")

# ============================================================
# S23 実文面②｜購入者
# ============================================================
s = slides[22]
frame(s, "配信設計｜実文面②｜購入者",
      ["サンクス・使用フォロー・F2オファー。ここでレビューを急かさない。"])
phone(s, CX0 + 0.3, CY0 + 0.2, 8.0, 6.7, "購入直後",
      [("in", "ご購入ありがとうございます！\n発送予定は2〜3営業日以内です"),
       ("in", "使い方のコツはこちら\nをご覧ください"),
       ("chip", "使い方ガイドを見る")])
phone(s, CX0 + 8.7, CY0 + 0.2, 8.0, 6.7, "使い始めの頃",
      [("in", "お使いいただけていますか？\nこう使うと効果を実感しやすい\nという声が多いポイントです"),
       ("note", "※ここではレビューを求めない")])
phone(s, CX0 + 17.1, CY0 + 0.2, 8.0, 6.7, "無くなる少し前",
      [("in", "そろそろ無くなる頃かと\nお送りしました"),
       ("btn", "再購入する"),
       ("chip", "定期便に切り替える")])
foot(s, "意見｜F2オファーのタイミングは「内容量÷1日使用量」から逆算（S21参照）")

# ============================================================
# S24 カゴ落ちリマインドの設計
# ============================================================
s = slides[23]
frame(s, "配信設計｜カゴ落ちリマインドの設計",
      ["「何分後・何回・何を出すか」まで具体で書く。S20（ID連携）が前提。"])
remind = [
    ("1通目", "60分後", "カート内容の再掲", "催促しない。「カートに商品が残っています」のみ"),
    ("2通目", "24時間後", "不安つぶし", "送料条件・返品保証など、カゴ落ち理由1位に対応"),
    ("3通目", "48時間後", "期限提示", "「クーポンは本日まで」等、最後の一押し"),
]
cw3 = (CW - 0.6) / 3
for i, (h, t, sub, b) in enumerate(remind):
    x = CX0 + i * (cw3 + 0.3)
    tb = box(s, x, CY0 + 0.2, cw3, 0.9, fill=NAVY)
    put_text(tb.text_frame, [one(f"{h}｜{t}", 12.5, True, WHITE, align="c")], anchor="m")
    card(s, x, CY0 + 1.2, cw3, 3.3, sub, b, hcol=NAVY, fill=PALE, hsz=12, bsz=10, ls=1.35)
box(s, CX0, CY0 + 4.8, CW, 1.2, fill=PORANGE)
T(s, CX0 + 0.3, CY0 + 4.8, CW - 0.6, 1.2,
  [one("★S20（LINE ID連携）が前提。カート放棄をトリガーに自動配信する", 12, True, RED, align="c", ls=1.25)],
  anchor="m")
foot(s, "意見｜通数・タイミングは初稿。配信間隔は実運用データで調整")

# ============================================================
# S25 年間の企画投稿カレンダー
# ============================================================
s = slides[24]
frame(s, "配信設計｜年間の企画投稿カレンダー",
      ["S10の季節性を「◯月：企画（参照データ）」形式に落とす。"])
CAL25 = [
    ("1月", "初売り・福袋特集", "S10商戦データ"), ("2月", "バレンタインギフト特集", "S10商戦データ"),
    ("3月", "楽天SS・新生活準備", "S10商戦データ"), ("4月", "新生活後半フォロー", "S10商戦データ"),
    ("5月", "母の日特集", "S10商戦データ"), ("6月", "父の日・お中元", "S10商戦データ"),
    ("7月", "プライムデー対抗企画", "S10商戦データ"), ("8月", "秋商戦の仕込み案内", "S10「8月は仕込み月」"),
    ("9月", "秋の新商品予告", "S10商戦データ"), ("10月", "ハロウィン特集", "S10商戦データ"),
    ("11月", "ブラックフライデー特集", "S10「年間最大の山」"), ("12月", "年末セール・福袋予告", "S10商戦データ"),
]
cols = 6
cw6 = CW / cols
ch2 = 1.55
for i, (m, ev, ref) in enumerate(CAL25):
    r, c = divmod(i, cols)
    x = CX0 + c * cw6
    y = CY0 + 0.2 + r * (ch2 + 0.12)
    box(s, x + 0.05, y, cw6 - 0.1, ch2, fill=PALE)
    T(s, x + 0.15, y + 0.05, cw6 - 0.3, 0.4, [one(m, 10.5, True, NAVY)], anchor="t")
    T(s, x + 0.15, y + 0.45, cw6 - 0.3, 0.65, [one(ev, 8, None, INK, ls=1.15)], anchor="t")
    T(s, x + 0.15, y + 1.1, cw6 - 0.3, 0.4, [one(ref, 6.8, None, MUT)], anchor="t")
foot(s, "意見｜繁忙期は捌く、閑散期（8月）は仕込む。参照データはS10のシーズナリティ")

# ============================================================
# S26 通知メッセージ
# ============================================================
s = slides[25]
frame(s, "配信設計｜通知メッセージ",
      ["友だちでなくても届く接点。発送・配送状況・再入荷・定期の次回発送前。"])
notif = [
    ("発送通知", "注文確定・発送完了を\n即時通知"), ("配送状況", "配送状況の更新を\nプッシュ通知"),
    ("再入荷通知", "希望者にだけ\n再入荷を通知"), ("定期の次回発送前", "発送サイクルに合わせて\n事前案内（S20が前提）"),
]
cw4 = (CW - 0.3 * 3) / 4
for i, (h, b) in enumerate(notif):
    x = CX0 + i * (cw4 + 0.3)
    card(s, x, CY0 + 0.3, cw4, 3.2, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=12, bsz=10, ls=1.3)
box(s, CX0, CY0 + 3.9, CW, 1.8, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 3.9, CW - 0.6, 1.8,
  [one("★友だちでなくても届く接点がある", 13, True, WHITE, ls=1.3, sa=4),
   one("ブロックされていても、購入者への配送・発送情報は届く。ID連携（S20）が前提。", 10.5, None, "E9EDF7", ls=1.3)],
  anchor="m")
foot(s, "意見｜通知メッセージはAPI連携（Messaging API）が前提。別途費用（S30参照）")

# ============================================================
# S27 F2転換率の改善（★本資料の山）
# ============================================================
s = slides[26]
frame(s, "歩留まり・工数｜F2転換率の改善（★本資料の山）",
      ["初回購入者の6〜7割は2回目を買わない。ここを動かすのがEC最大のレバー。"])
box(s, CX0, CY0 + 0.2, CW, 2.6, fill=NAVY)
T(s, CX0 + 0.4, CY0 + 0.35, CW - 0.8, 0.6, [one("F2転換率", 13, True, "C9D3EC")], anchor="m")
T(s, CX0 + 0.4, CY0 + 0.9, CW - 0.8, 1.7,
  [one("30〜40%", 34, True, WHITE)], anchor="m")
T(s, CX0 + 8.0, CY0 + 0.9, CW - 8.4, 1.7,
  [one("＝ 初回購入者の6〜7割は2回目を買わない", 16, True, ORANGE, ls=1.3)], anchor="m")
levers = [
    ("消費サイクル連動の配信", "「そろそろ無くなる頃」を\n内容量÷1日使用量で逆算配信"),
    ("定期引上", "初回お試し→定期購入への\n引き上げ導線（目安：約20%）"),
    ("休眠掘り起こし", "F2化しなかった層への\n再アプローチ・特典訴求"),
]
cw3 = (CW - 0.6) / 3
for i, (h, b) in enumerate(levers):
    x = CX0 + i * (cw3 + 0.3)
    card(s, x, CY0 + 3.1, cw3, 3.2, h, b.split("\n"), hcol=ORANGE, fill=PORANGE, hsz=12, bsz=10, ls=1.3)
foot(s, "業界水準｜F2転換率（化粧品）25〜35%／（化粧品・健康食品などの単品通販）30〜40%／定期引上率（初回→定期）約20%")

# ============================================================
# S28 工数削減 ＋ 改善モデル
# ============================================================
s = slides[27]
frame(s, "歩留まり・工数｜工数削減＋改善モデル",
      ["問い合わせ・注文/配送照会をセルフ化し、各段の打ち手を1対1で対応させる。"])
T(s, CX0, CY0 + 0.2, CW, 0.5, [one("工数削減", 12.5, True, NAVY)], anchor="m")
auto = [("KW自動応答", "「送料は？」「返品は？」を\n24時間自動応答"),
        ("配送状況セルフ照会", "リッチメニューから\n注文履歴・配送状況を自己解決"),
        ("よくある質問の自動化", "サイズ・支払い方法等を\nメニュー選択式で即回答")]
cw3 = (CW - 0.6) / 3
for i, (h, b) in enumerate(auto):
    x = CX0 + i * (cw3 + 0.3)
    card(s, x, CY0 + 0.8, cw3, 2.1, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=10.5, bsz=9, ls=1.25)
T(s, CX0, CY0 + 3.2, CW, 0.5, [one("改善モデル｜各段に打ち手が1対1で対応", 12.5, True, NAVY)], anchor="m")
model = ["カゴ落ち\n対策", "→", "初回購入\nCVR", "→", "F2転換率\n改善", "→", "定期・LTV\n最大化"]
mw = CW / len(model)
for i, m in enumerate(model):
    x = CX0 + i * mw
    if m == "→":
        T(s, x, CY0 + 4.0, mw, 1.4, [one(m, 18, True, ORANGE, align="c")], anchor="m")
    else:
        bb = box(s, x + 0.1, CY0 + 3.9, mw - 0.2, 1.6, fill=PALE)
        put_text(bb.text_frame, [one(l, 9.5, True, NAVY, align="c", ls=1.15) for l in m.split("\n")], anchor="m")
foot(s, "意見｜DYM提案の改善モデル。工数削減はスタッフ対応ゼロを目指すのではなく、定型対応を減らして個別対応に集中する設計")

# ============================================================
# S29 効果測定の設計
# ============================================================
s = slides[28]
frame(s, "成果と体制｜効果測定の設計",
      ["主KPI＝CPO＋F2転換率。締めはLTV÷CPO。数値は入れず構造だけ図示する。"])
simple_table(s, CX0, CY0 + 0.2, CW, 3.6,
             ["", "Before（広告のみ）", "After（1年目）", "After（2年目）"],
             [["注文数", "広告費÷CPO", "① カゴ落ち回収で増加", "同左＋友だちリスト流入"],
              ["LTV", "≒初回売上", "② F2転換の自動化で上昇", "定期・休眠復活で更に上昇"],
              ["翌期の広告費", "毎期フルで必要", "同左", "③ 友だちリスト分だけ圧縮可"],
              ["LTV ÷ CPO", "伸びない", "改善方向", "★期を追うごとに上がる"]],
             col_w=[CW * 0.16, CW * 0.28, CW * 0.28, CW * 0.28], bsz=8.5, row_h=0.75)
box(s, CX0, CY0 + 4.0, CW, 2.2, fill=PALE)
T(s, CX0 + 0.3, CY0 + 4.15, CW - 0.6, 2.0,
  [one("★肝：広告経由の顧客は1回買って消える。", 11.5, True, RED, ls=1.3, sa=2),
   one("友だち化すると資産として残るので、LTV÷CPOは期を追うごとに上がる。", 11.5, None, INK, ls=1.3, sa=4),
   one("締め：御社の実績をいただければSIMを作ります。", 12, True, TNAVY, ls=1.3)],
  anchor="t")
foot(s, "意見｜数値は入れない（業界汎用のためSIMは作らない）。主KPI＝CPO＋F2転換率")

# ============================================================
# S30 費用プラン
# ============================================================
s = slides[29]
frame(s, "成果と体制｜費用プラン",
      ["無償付帯と別途費用を分けて示す。追加で払うものを隠さない（6ヶ月〜・税抜）。"])
plans = [
    ("① コンサル基本", "初期10万〜", "月20万（3投稿）\n月30万（5投稿）\n月50万（9投稿）"),
    ("② 初動設計＋運用", "初期10万〜", "月5万〜"),
    ("③ 運用代行・効率改善", "初期10万〜", "月0万〜\n（アカウント費のみ）"),
    ("④ 成果報酬型", "初期5万〜", "月＝単価×成果数\n＋固定費"),
]
cw4 = (CW - 0.3 * 3) / 4
for i, (h, init, m) in enumerate(plans):
    x = CX0 + i * (cw4 + 0.3)
    box(s, x, CY0 + 0.2, cw4, 3.3, fill=PALE)
    T(s, x + 0.2, CY0 + 0.35, cw4 - 0.4, 0.9, [one(h, 11, True, NAVY, ls=1.2)], anchor="t")
    T(s, x + 0.2, CY0 + 1.3, cw4 - 0.4, 0.6, [one(init, 13, True, ORANGE)], anchor="t")
    T(s, x + 0.2, CY0 + 1.95, cw4 - 0.4, 1.4, [one(l, 9, None, INK, ls=1.25) for l in m.split("\n")], anchor="t")
T(s, CX0, CY0 + 3.8, CW, 0.5, [one("無償付帯", 11.5, True, NAVY)], anchor="m")
T(s, CX0, CY0 + 4.3, CW, 0.9,
  [one("アカウント開設／プロフィール／リッチメニュー／あいさつ／初期アンケート／KW自動応答／ステップ配信／セグメント配信／タグ管理／GAレポート連携／クリエイティブ／定例会", 9.5, None, INK, ls=1.3)],
  anchor="t")
T(s, CX0, CY0 + 5.3, CW, 0.5, [one("別途費用", 11.5, True, RED)], anchor="m")
T(s, CX0, CY0 + 5.8, CW, 0.9,
  [one("離脱防止（Sitelead）初期10万＋月5万／API連携／通知メッセージ 等。配信ツールはLstep／Hachidoriのいずれかを要件で選定", 9.5, None, INK, ls=1.3)],
  anchor="t")
foot(s, "意見｜DYM標準プラン（6ヶ月〜・税抜）")

# ============================================================
# S31 運用スケジュール
# ============================================================
s = slides[30]
frame(s, "成果と体制｜運用スケジュール",
      ["1ヶ月目に構築、2ヶ月目に稼働、3ヶ月目から改善。6ヶ月で年間配信の型ができる。"])
sched = [
    ("1ヶ月目", "構築", "あいさつMSG・診断・リッチメニュー・ID連携基盤"),
    ("2ヶ月目", "稼働", "友だち追加動線を稼働。カゴ落ちリマインド開始"),
    ("3〜4ヶ月目", "改善①", "F2オファー配信を開始。効果測定・改善"),
    ("5〜6ヶ月目", "改善②", "年間企画投稿カレンダーの型を確立"),
]
cw4 = CW / 4
for i, (h, tag, b) in enumerate(sched):
    x = CX0 + i * cw4
    tb = box(s, x + 0.1, CY0 + 0.3, cw4 - 0.2, 0.8, fill=NAVY)
    put_text(tb.text_frame, [one(h, 11.5, True, WHITE, align="c")], anchor="m")
    T(s, x + 0.1, CY0 + 1.25, cw4 - 0.2, 0.6, [one(tag, 12, True, ORANGE, align="c")], anchor="m")
    card(s, x + 0.1, CY0 + 1.95, cw4 - 0.2, 2.7, "", b, hcol=NAVY, fill=PALE, hsz=1, bsz=9.5, ls=1.3, anchor="m")
    if i < 3:
        T(s, x + cw4 - 0.15, CY0 + 1.9, 0.4, 1.0, [one("▶", 14, True, MUT, align="c")], anchor="m")
foot(s, "意見｜6ヶ月で年間配信の型を確立。以降は月次運用フェーズへ")

# ============================================================
# S32 サポート体制
# ============================================================
s = slides[31]
frame(s, "成果と体制｜サポート体制",
      ["定例会・レポート・クリエイティブ制作まで含む。運用は丸ごと持つ。"])
supp = [
    ("定例会", "月次で振り返り・\n翌月の企画をすり合わせ"),
    ("レポーティング", "GA連携で友だち数・\n配信効果を数値化"),
    ("クリエイティブ制作", "配信文面・リッチメニュー\n画像の制作を内包"),
    ("運用代行", "日々の配信設計・\n自動応答の調整"),
]
cw4 = (CW - 0.3 * 3) / 4
for i, (h, b) in enumerate(supp):
    x = CX0 + i * (cw4 + 0.3)
    card(s, x, CY0 + 0.3, cw4, 3.0, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=11.5, bsz=9.5, ls=1.3)
foot(s, "意見｜DYM標準のサポート体制")

# ============================================================
# S33 飛び道具
# ============================================================
s = slides[32]
frame(s, "締め｜飛び道具",
      ["他社がやっていない打ち手が4つある。"])
tobi = [
    ("A", "カゴ落ち理由をあえて公開", "送料内訳・在庫状況を先に見せる誠実路線。\n「隠さない」ことがカゴ落ち理由1位（追加コスト\nの不透明さ）への最大の対策になる", False),
    ("B", "同梱物QR→開封後アンケート", "受け取り直後の温度が最も高い瞬間に\nミニアプリ会員証化のフックを差し込む", False),
    ("C", "UGC投稿を集めるリッチメニュー導線", "「使用写真を送ると次回クーポン」で\nレビュー・SNS投稿を友だち経由で集約する", False),
    ("D", "「いつでも解約できます」を初回配信で明示", "定期縛りへの不安に先回りする。\n★前後検索で「解約」不安の実在を確認してから展開", True),
]
cw2 = (CW - 0.4) / 2
for i, (tag, h, b, hyp) in enumerate(tobi):
    r, c = divmod(i, 2)
    x = CX0 + c * (cw2 + 0.4)
    y = CY0 + 0.2 + r * 3.4
    card(s, x, y, cw2, 3.1, f"{tag}　{h}", b.split("\n"), hcol=NAVY, fill=PALE, hsz=11.5, bsz=9, ls=1.3)
    if hyp:
        badge(s, x + cw2 - 2.3, y + 0.15, 2.0, 0.5, "仮説", fill=PORANGE, col=ORANGE, sz=8)
foot(s, "意見｜D案は仮説バッジ付き（S11前後検索データの取得後に実在確認）")

# ============================================================
# S34 第2の提案軸｜LINEミニアプリ・LINEギフト
# ============================================================
s = slides[33]
frame(s, "締め｜第2の提案軸｜LINEミニアプリ・LINEギフト",
      ["会員証・注文履歴の一元化と、ギフト需要の取り込み。ID連携（S20）の上に乗る拡張機能。"])
axis = [
    ("LINEミニアプリ", ["会員証・ポイントをLINE内に一元化", "注文履歴・お気に入りをアプリレスで確認",
                    "ID連携基盤（S20）をそのまま活用できる"]),
    ("LINEギフト", ["「送りたい」を逃さない。ギフト需要の取り込み", "受け取り側の新規友だち化にもつながる",
                 "母の日・クリスマス等の季節商戦との相性が良い"]),
]
cw2 = (CW - 0.4) / 2
for i, (h, items) in enumerate(axis):
    x = CX0 + i * (cw2 + 0.4)
    box(s, x, CY0 + 0.2, cw2, 5.6, fill=PALE)
    T(s, x + 0.3, CY0 + 0.4, cw2 - 0.6, 0.7, [one(h, 14, True, NAVY)], anchor="m")
    for j, it in enumerate(items):
        T(s, x + 0.3, CY0 + 1.3 + j * 1.3, cw2 - 0.6, 1.2,
          [one(f"・{it}", 10.5, None, INK, ls=1.3)], anchor="t")
foot(s, "意見｜DYM提案の第2軸。本提案（S16-26）のID連携基盤の上にそのまま構築できる")

# ============================================================
# S35 LINEOA実績（LY公式・EC専用ページは★差込枠）
# ============================================================
s = slides[34]
frame(s, "締め｜LINEOA実績",
      ["LINEヤフー公式データと、EC専用の導入事例を並べる。ECは事例が多い業界。"])
nums = [("LINE国内MAU", "1億人突破", "2025年12月末時点"),
        ("メッセージ開封率（当日中）", "約8割", "受信直後 約2割／3〜6時間で約5割")]
cw2 = (CW - 0.4) / 2
for i, (label, val, sub) in enumerate(nums):
    x = CX0 + i * (cw2 + 0.4)
    box(s, x, CY0 + 0.2, cw2, 2.4, fill=PALE)
    T(s, x + 0.3, CY0 + 0.35, cw2 - 0.6, 0.6, [one(label, 11, True, NAVY)], anchor="m")
    T(s, x + 0.3, CY0 + 0.95, cw2 - 0.6, 1.0, [one(val, 24, True, ORANGE)], anchor="m")
    T(s, x + 0.3, CY0 + 1.8, cw2 - 0.6, 0.5, [one(sub, 9, None, MUT)], anchor="m")
placeholder(s, CX0, CY0 + 2.9, CW, 3.9, "★差込枠｜LINEヤフー公式 EC導入事例（未取得）",
            "取得先：lycbiz.com/jp/service/line-official-account/case-study/ec/（EC専用ページあり）。\n"
            "掲載前に各事例ページの原典で数値を再確認すること")
foot(s, "業界水準｜LINE国内MAU・開封率はLINEヤフー公式値（倍率表現は使用していない）。EC事例は★未取得")


# ============================================================
# 業種残骸の最終チェック（正規表現スキャン）
# ============================================================
import re

RESIDUE_PATTERNS = [
    r"○○業界", r"〇〇業界", r"クロスセル", r"アップセル", r"転職", r"求人",
    r"利用状況タグ", r"誕生日.{0,3}記念日タグ", r"人材\(転職\)業界",
]
residue_hits = []
for i, sl in enumerate(slides, 1):
    for sh in sl.shapes:
        if sh.has_text_frame:
            txt = sh.text_frame.text
            for pat in RESIDUE_PATTERNS:
                if re.search(pat, txt):
                    residue_hits.append((i, pat, txt[:40]))
if residue_hits:
    print("★業種残骸が残っています：")
    for i, pat, txt in residue_hits:
        print(f"  slide {i}: /{pat}/ -> {txt!r}")
else:
    print("業種残骸チェック：残存0")

prs.save(OUT)
print("saved:", OUT)
