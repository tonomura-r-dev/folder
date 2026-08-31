# -*- coding: utf-8 -*-
"""小林住宅株式会社 御中｜LINE構築提案（ラフ版・11枚）

_templates/DYM_LINEOA_FMT.pptx（業界汎用FMT・36枚）をコピーし、
先頭11枚にP1-P11の内容を構築したうえで、残り25枚を削除する。

  python _build/build_kobayashi_jutaku_rough.py
  python _build/qa_render.py "20260901_小林住宅株式会社 御中_LINE構築提案.pptx"

【位置づけ】
データ収集（前後検索・Googleトレンド）とSIM作成の**前**のラフ版。
数値が未確定の箇所は★を付けて空欄のまま残す。埋まっていない数字を推測で埋めない。

【★のまま残した箇所（捏造していない）】
- P4 市場のシーズナリティ：Googleトレンド実数（GWピーク100→直近51のみ既知）
- P5 前後検索で見る検討の流れ：前後検索データ未取得
- P6 競合のLINE運用状況：同価格帯メーカー3社の実査未取得（社名も伏せる＝実データが
  取れるまで実名は出さない、が鉄則）
- P10 想定シミュレーション：SIM未作成（構造のみ）
- サイトUU・資料請求件数/CPA・資料請求→来場転換率・LINE公式アカウント実在確認も未取得

【落とし穴メモ（_build/README.md・deck-apply スキルより）】
- put_text() は必ず reset_tf() を通す（既存テキストへの追記事故を防ぐ）
- 新規スライド追加はしない（パート名重複で壊れる）。11枚に絞るのはFMTの後半25枚を
  削除する形で行う（追加ではなく削除）
- 一括置換をしない。数値は行・列を特定して書く
"""
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
SRC = str(ROOT / "_templates" / "DYM_LINEOA_FMT.pptx")
OUT = str(ROOT / "20260901_小林住宅株式会社 御中_LINE構築提案.pptx")
N_PAGES = 11

# ---- 配色（DYM 3色：ネイビー・炭黒・オレンジ。LINEグリーンはUI部分のみ）----
TNAVY = "002060"   # タイトル
NAVY = "1F285A"    # 打ち手・カード見出し
ORANGE = "ED7D31"  # 強調
RED = "C00000"     # 課題・警告
INK = "333333"     # 本文
MUT = "7F7F7F"
WHITE = "FFFFFF"
PALE = "F4F7FF"    # 淡ネイビー
PORANGE = "FCE4D6"  # 淡オレンジ
GREY = "F2F2F2"
BORDER = "D9D9D9"
GREEN = "06C755"   # LINE UI 専用
BEZEL = "2B2B2B"
SCREEN = "EFF2F7"
PRED = "FDF2F2"

# ---- FMTのレイアウト座標（cm）----
SW, SH = 27.52, 19.05
TITLE_XY = (1.52, 0.38, 24.4, 0.94)
LEAD_XY = (1.20, 1.80, 25.1, 1.90)
DIV_Y = 3.86
CX0, CW = 1.20, 25.12
CY0, CY1 = 4.30, 17.00
FOOT_Y = 17.35

shutil.copyfile(SRC, OUT)
prs = Presentation(OUT)
slides = list(prs.slides)
assert len(slides) == 36, len(slides)
assert (prs.slide_width, prs.slide_height) == (9906000, 6858000)


# ================= helpers（build_chumon_jutaku.py と共通の作法）=================
def set_font(run, size, bold=None, color=INK, name="メイリオ"):
    f = run.font
    f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", name)
    f.color.rgb = RGBColor.from_string(color)


ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
ANCH = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}


def reset_tf(tf):
    """既存の段落・runを全消去（追記事故の防止）"""
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)


def put_text(tf, paras, anchor="t", ml=0.14, mr=0.14, mt=0.06, mb=0.06, wrap=True):
    reset_tf(tf)
    tf.word_wrap = wrap
    tf.margin_left = Cm(ml)
    tf.margin_right = Cm(mr)
    tf.margin_top = Cm(mt)
    tf.margin_bottom = Cm(mb)
    tf.vertical_anchor = ANCH[anchor]
    first = True
    for p in paras:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = ALIGN[p.get("align", "l")]
        if p.get("sa") is not None:
            para.space_after = Pt(p["sa"])
        if p.get("ls") is not None:
            para.line_spacing = p["ls"]
        for item in p["runs"]:
            t, sz, b, c = item
            r = para.add_run()
            r.text = t
            set_font(r, sz, b, c)
    return tf


def T(slide, x, y, w, h, paras, anchor="t", **kw):
    box_ = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    put_text(box_.text_frame, paras, anchor=anchor, **kw)
    return box_


def one(text, sz, b=None, c=INK, align="l", sa=None, ls=None):
    d = {"runs": [(text, sz, b, c)], "align": align}
    if sa is not None:
        d["sa"] = sa
    if ls is not None:
        d["ls"] = ls
    return d


def multi(runs, align="l", sa=None, ls=None):
    """1段落に複数run（色・太さを混在させたいとき）"""
    d = {"runs": runs, "align": align}
    if sa is not None:
        d["sa"] = sa
    if ls is not None:
        d["ls"] = ls
    return d


def box(slide, x, y, w, h, fill=None, line=None, lw=1.0,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06, dash=None):
    sp = slide.shapes.add_shape(shape, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor.from_string(line)
        sp.line.width = Pt(lw)
        if dash:
            sp.line.dash_style = dash
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    reset_tf(sp.text_frame)
    return sp


def card(slide, x, y, w, h, head, body, hcol=NAVY, fill=PALE,
         hsz=11, bsz=9, line=None, anchor="t", ls=1.18):
    sp = box(slide, x, y, w, h, fill=fill, line=line)
    paras = [one(head, hsz, True, hcol, sa=3)]
    for b in (body if isinstance(body, list) else [body]):
        if b:
            paras.append(one(b, bsz, None, INK, ls=ls, sa=1))
    put_text(sp.text_frame, paras, anchor=anchor, ml=0.22, mr=0.18, mt=0.14, mb=0.10)
    return sp


def clear_slide(slide):
    spTree = slide.shapes._spTree
    for el in list(spTree):
        if el.tag.split("}")[-1] in ("sp", "cxnSp", "pic", "graphicFrame", "grpSp"):
            spTree.remove(el)


def frame(slide, title, lead):
    """FMT準拠：タイトル16pt紺（y=0.38）＋リード12pt（y=1.80・2行以内）＋区切り線"""
    clear_slide(slide)
    T(slide, *TITLE_XY, [one(title, 16, True, TNAVY)], anchor="m", ml=0, mr=0)
    T(slide, *LEAD_XY, [one(l, 12, None, INK, ls=1.28) for l in lead],
      anchor="m", ml=0, mr=0)
    ln = slide.shapes.add_connector(1, Cm(0), Cm(DIV_Y), Cm(SW), Cm(DIV_Y))
    ln.line.color.rgb = RGBColor.from_string(BORDER)
    ln.line.width = Pt(1.0)


def foot(slide, text):
    T(slide, CX0, FOOT_Y, CW, 0.9, [one(text, 7.5, None, MUT, ls=1.15)], ml=0, mr=0)


def badge(slide, x, y, w, h, text, fill=PORANGE, col=ORANGE, sz=8.5, dash=None):
    sp = box(slide, x, y, w, h, fill=fill, line=col, lw=1.0, dash=dash)
    put_text(sp.text_frame, [one(text, sz, True, col, align="c")],
             anchor="m", ml=0.06, mr=0.06, mt=0, mb=0)
    return sp


def band(slide, y, text, fill=NAVY, col=WHITE, sz=11.5, h=0.86, x=CX0, w=CW):
    sp = box(slide, x, y, w, h, fill=fill, radius=0.10)
    put_text(sp.text_frame, [one(text, sz, True, col, align="c")],
             anchor="m", ml=0.2, mr=0.2, mt=0, mb=0)
    return sp


def placeholder(slide, x, y, w, h, label, note):
    """★未取得データの差込枠（破線）。捏造しない。"""
    sp = box(slide, x, y, w, h, fill=WHITE, line=MUT, lw=1.25,
              dash=MSO_LINE_DASH_STYLE.DASH)
    put_text(sp.text_frame,
              [one(label, 11, True, MUT, align="c", sa=4),
               one(note, 8.5, None, MUT, align="c", ls=1.25)],
              anchor="m", ml=0.3, mr=0.3, mt=0.1, mb=0.1)
    return sp


def phone(slide, x, y, w, h, header, lines):
    """LINEトーク画面のモック。lines = [(kind, text)]
    kind: 'in'（BOT吹き出し）/ 'chip'（選択肢）/ 'btn'（CTA）/ 'note'（注記）"""
    box(slide, x, y, w, h, fill=BEZEL, radius=0.10)
    box(slide, x + 0.14, y + 0.42, w - 0.28, h - 0.56, fill=SCREEN, radius=0.02,
        shape=MSO_SHAPE.RECTANGLE)
    hd = box(slide, x + 0.14, y + 0.42, w - 0.28, 0.52, fill=GREEN,
             shape=MSO_SHAPE.RECTANGLE)
    put_text(hd.text_frame, [one(header, 7, True, WHITE, align="c")],
             anchor="m", ml=0.05, mr=0.05, mt=0, mb=0)
    cy = y + 1.05
    iw = w - 0.62
    for kind, text in lines:
        nlines = text.count("\n") + 1
        bh = 0.30 + 0.30 * nlines
        if kind == "in":
            sp = box(slide, x + 0.31, cy, iw, bh, fill=WHITE, line=BORDER, lw=0.75)
            put_text(sp.text_frame, [one(l, 6.5, None, INK, ls=1.18)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.14, mr=0.10, mt=0.05, mb=0.05)
        elif kind == "chip":
            sp = box(slide, x + 0.31, cy, iw, bh, fill=WHITE, line=GREEN, lw=0.9)
            put_text(sp.text_frame, [one(l, 6.5, None, "0B7A3B", align="c", ls=1.15)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.08, mr=0.08, mt=0.03, mb=0.03)
        elif kind == "btn":
            sp = box(slide, x + 0.31, cy, iw, bh, fill=GREEN)
            put_text(sp.text_frame, [one(l, 6.8, True, WHITE, align="c", ls=1.15)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.08, mr=0.08, mt=0.03, mb=0.03)
        else:  # note
            sp = slide.shapes.add_textbox(Cm(x + 0.31), Cm(cy), Cm(iw), Cm(bh))
            put_text(sp.text_frame, [one(l, 6.2, None, MUT, ls=1.15)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.05, mr=0.05, mt=0, mb=0)
        cy += bh + 0.12
    return cy


def richmenu(slide, x, y, w, h, tabs):
    """リッチメニュー3タブ×6ボタンのモック。
    tabs = [(タブ名, [ボタン文言6個]), ...]。最初のタブをアクティブ表示。"""
    tab_h = 0.55
    box(slide, x, y, w, h, fill=BEZEL, radius=0.06)
    tw = (w - 0.12) / len(tabs)
    for i, (name, _) in enumerate(tabs):
        active = (i == 0)
        tb = box(slide, x + 0.06 + i * tw, y + 0.06, tw - 0.03, tab_h,
                  fill=GREEN if active else "3F3F3F",
                  shape=MSO_SHAPE.RECTANGLE)
        put_text(tb.text_frame, [one(name, 8, True, WHITE, align="c")],
                 anchor="m", ml=0, mr=0, mt=0, mb=0)
    _, btns = tabs[0]
    gx, gy = x + 0.06, y + 0.06 + tab_h + 0.06
    gw, gh = w - 0.12, h - tab_h - 0.18
    cols, rows = 3, 2
    bw, bh = (gw - 0.06 * (cols - 1)) / cols, (gh - 0.06 * (rows - 1)) / rows
    for i, label in enumerate(btns):
        r, c = divmod(i, cols)
        bb = box(slide, gx + c * (bw + 0.06), gy + r * (bh + 0.06), bw, bh,
                  fill=SCREEN, line=BORDER, lw=0.75, radius=0.04)
        put_text(bb.text_frame, [one(l, 7.2, True, NAVY, align="c", ls=1.1)
                                 for l in label.split("\n")],
                 anchor="m", ml=0.05, mr=0.05, mt=0, mb=0)


def simple_table(slide, x, y, w, h, headers, rows, col_w=None,
                  hsz=9, bsz=8.5, header_fill=NAVY, zebra=PALE,
                  align=None, row_h=None):
    """罫線テーブル（ネイティブpptxテーブル）。headers=[str], rows=[[str]]"""
    n_r, n_c = len(rows) + 1, len(headers)
    gf = slide.shapes.add_table(n_r, n_c, Cm(x), Cm(y), Cm(w), Cm(h))
    tbl = gf.table
    if col_w:
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Cm(cw)
    if row_h:
        for r in tbl.rows:
            r.height = Cm(row_h)
    for j, htext in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor.from_string(header_fill)
        c.margin_left = c.margin_right = Cm(0.12)
        c.margin_top = c.margin_bottom = Cm(0.05)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        put_text(c.text_frame, [one(htext, hsz, True, WHITE,
                                    align=(align[j] if align else "c"))],
                 anchor="m", ml=0.12, mr=0.12, mt=0.02, mb=0.02)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor.from_string(
                zebra if (zebra and i % 2 == 0) else WHITE)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            put_text(c.text_frame, [one(str(val), bsz, None, INK,
                                        align=(align[j] if align else "l"))],
                     anchor="m", ml=0.12, mr=0.12, mt=0.02, mb=0.02)
    return gf


def find_shape(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle in sh.text_frame.text:
            return sh
    return None


def delete_slide(prs, index):
    """指定indexのスライドをプレゼンテーションから削除する（標準的なpython-pptxレシピ）。"""
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    rId = slides_list[index].rId
    prs.part.drop_rel(rId)
    xml_slides.remove(slides_list[index])


def trim_to(prs, slides_, n):
    """先頭n枚を残し、残りを末尾から順に削除する（末尾から消すとindexがズレない）。"""
    for i in range(len(slides_) - 1, n - 1, -1):
        delete_slide(prs, i)


# ============================================================
# P01 表紙
# ============================================================
s = slides[0]
sh = find_shape(s, "業界：")
if sh:
    put_text(sh.text_frame, [one("小林住宅株式会社 御中", 12, True, TNAVY)], anchor="m")
sh = find_shape(s, "最上位パートナーの知見")
if sh:
    put_text(sh.text_frame,
             [one("資料請求は取れているのに、営業が追いきれず来場に繋がっていない。LINEで追い切る", 13, True, INK)],
             anchor="m")
badge(s, CX0, 0.4, 5.2, 0.65, "ラフ版（データ収集前）", fill=PORANGE, col=ORANGE, sz=10)
T(s, CX0, 17.6, CW, 0.8,
  [one("代理店：株式会社プロモ／営業担当：髙橋駿太様", 8.5, None, MUT)], ml=0, mr=0)

# ============================================================
# P02 現状｜いま起きていること
# ============================================================
s = slides[1]
frame(s, "現状｜いま起きていること",
      ["資料請求は取れているが、来場に繋がっていない。",
       "結果として、来場予約1件のCPAが高騰している。"])
nums = [
    ("来場予約CPA", "15万円", "管理画面ベース。実反響では\nもう少し良い数値"),
    ("来場目標に対する実績", "70 / 100件", "達成率70%（全展示場合計）"),
    ("1施設あたりの来場数", "10 → 14件", "月4件の上積みが必要（7施設）"),
]
cw3 = (CW - 0.6) / 3
for i, (label, val, sub) in enumerate(nums):
    x = CX0 + i * (cw3 + 0.3)
    box(s, x, CY0 + 0.2, cw3, 3.6, fill=PALE)
    T(s, x, CY0 + 0.4, cw3, 0.6, [one(label, 11.5, True, NAVY, align="c")], anchor="m")
    T(s, x, CY0 + 1.1, cw3, 1.3, [one(val, 26, True, ORANGE, align="c")], anchor="m")
    T(s, x, CY0 + 2.5, cw3, 1.2, [one(l, 9.5, None, INK, align="c", ls=1.25) for l in sub.split("\n")], anchor="m")
box(s, CX0, CY0 + 4.0, CW, 2.3, fill=NAVY)
T(s, CX0 + 0.4, CY0 + 4.15, CW - 0.8, 2.0,
  [one("資料請求は取れている。問題は、そこから来場までの間。", 13, True, WHITE, ls=1.3, sa=4),
   one("展示場7施設（京都1・大阪5・兵庫1）、広告予算は月100〜150万円（Meta比重やや高め）。"
       "媒体はGoogle・Yahoo・Meta、CV地点は来場予約／カタログ請求／オンライン相談／間取りプラン請求／お問い合わせの5つ。",
       10.5, None, "E9EDF7", ls=1.3)],
  anchor="m")
foot(s, "業界水準｜殿村さんヒアリングによる確定値（2026年8月時点・代理店提供の管理画面数値）")

# ============================================================
# P03 この商材で、いま何が起きているか
# ============================================================
s = slides[2]
frame(s, "この商材で、いま何が起きているか",
      ["追客ができていないのではなく、追える手段がない。"])
cards = [
    ("① 顧客が動く時間に、営業は架電できない",
     ["30〜40代・共働きの子育て世帯", "平日日中は仕事と育児 → 電話に出られない", "検討するのは夜21時以降と週末"]),
    ("② 検討8〜12か月に対して、接点は資料請求の1回",
     ["資料請求 → 8〜12か月ほぼ無接触 → 来場 or 他社決定", "その間に3〜5社を比較"]),
    ("③ 高性能住宅は、一度では伝わらない",
     ["断熱等級6・7／C値／全館空調は理解に時間がかかる", "「外断熱供給実績No.1」も、伝える機会が1回では届かない"]),
    ("④ だから来場予約を直接買うことになる",
     ["資料請求からの引き上げができない → 来場予約を広告で買う", "→ CPA15万円"]),
]
cw2 = (CW - 0.4) / 2
for i, (h, lines) in enumerate(cards):
    r, c = divmod(i, 2)
    x = CX0 + c * (cw2 + 0.4)
    y = CY0 + 0.2 + r * 3.1
    card(s, x, y, cw2, 2.85, h, lines, hcol=NAVY, fill=PALE, hsz=11.5, bsz=9.5, ls=1.3)
box(s, CX0, CY0 + 6.6, CW, 1.5, fill=RED)
T(s, CX0 + 0.4, CY0 + 6.6, CW - 0.8, 1.5,
  [one("CPAが高いのは結果です。原因は、8か月間の接点がないことです。", 14, True, WHITE, align="c", ls=1.3)],
  anchor="m")
foot(s, "意見｜商談ヒアリングに基づく整理")

# ============================================================
# P04 市場のシーズナリティ（★Googleトレンド待ち）
# ============================================================
s = slides[3]
frame(s, "市場のシーズナリティ",
      ["ゴールデンウィークをピークに、検索関心は落ち込んでいる。",
       "★詳細はGoogleトレンド取得後に反映する。"])
box(s, CX0, CY0 + 0.2, CW, 2.4, fill=PALE)
T(s, CX0 + 0.4, CY0 + 0.4, CW - 0.8, 0.6, [one("既知の1点（Googleトレンド・相対指標）", 11.5, True, NAVY)], anchor="m")
gw_x = CX0 + 2.0
box(s, gw_x, CY0 + 1.1, 4.5, 1.2, fill=WHITE, line=ORANGE, lw=1.25)
T(s, gw_x, CY0 + 1.2, 4.5, 1.0, [one("ゴールデンウィーク", 9.5, None, INK, align="c"), one("100", 22, True, ORANGE, align="c")], anchor="m")
T(s, gw_x + 4.7, CY0 + 1.55, 1.0, 0.5, [one("→", 20, True, MUT, align="c")], anchor="m")
box(s, gw_x + 5.9, CY0 + 1.1, 4.5, 1.2, fill=WHITE, line=MUT, lw=1.25)
T(s, gw_x + 5.9, CY0 + 1.2, 4.5, 1.0, [one("直近", 9.5, None, INK, align="c"), one("51", 22, True, INK, align="c")], anchor="m")
placeholder(s, CX0, CY0 + 2.8, CW, 4.3, "★差込枠｜Googleトレンド実データ（未取得）",
            "対象KWの5年＋1年推移を取得後、折れ線グラフを差し込む。\n"
            "スケールの異なるKWは混ぜず別グラフにする（他案件での既知の事故）")
foot(s, "業界水準｜GWピーク・直近値は殿村さん提供の既知情報。それ以外はGoogleトレンド取得後に反映（未取得）")

# ============================================================
# P05 前後検索で見る検討の流れ（★前後検索データ待ち）
# ============================================================
s = slides[4]
frame(s, "前後検索で見る検討の流れ",
      ["★前後検索データ取得後、検討前後で検索KWがどう変わるかを反映する。"])
bands3 = [("検索“前”15〜7日", MUT), ("起点0日", NAVY), ("検索“後”7〜15日", RED)]
cw3b = (CW - 0.6) / 3
for i, (h, c) in enumerate(bands3):
    x = CX0 + i * (cw3b + 0.3)
    bb = box(s, x, CY0 + 0.2, cw3b, 1.1, fill=c)
    put_text(bb.text_frame, [one(h, 11.5, True, WHITE, align="c")], anchor="m")
    placeholder(s, x, CY0 + 1.4, cw3b, 3.0, "★未取得", "検索KWを取得後に記入")
foot(s, "★未取得｜出典：LINEヤフー媒体資料（要取得）。起点KWは施策URL・商材名から要検討")

# ============================================================
# P06 競合のLINE運用状況（★実査待ち）
# ============================================================
s = slides[5]
frame(s, "競合のLINE運用状況",
      ["同価格帯メーカー3社を実査する。競合の実名は、実データが取れるまで出さない。"])
simple_table(s, CX0, CY0 + 0.2, CW, 2.5,
             ["アカウント", "友だち数", "運用の型", "あいさつMSG／リッチメニュー"],
             [["★同価格帯メーカー1（社名未定）", "★未取得", "★未取得", "★未取得"],
              ["★同価格帯メーカー2（社名未定）", "★未取得", "★未取得", "★未取得"],
              ["★同価格帯メーカー3（社名未定）", "★未取得", "★未取得", "★未取得"]],
             col_w=[CW * 0.34, CW * 0.18, CW * 0.22, CW * 0.26], row_h=0.7)
placeholder(s, CX0, CY0 + 3.3, CW, 3.6, "★差込枠｜競合3社の実査（未取得）",
            "page.line.meから取得日つきで実測。断熱等級・C値を訴求する同価格帯の\n"
            "住宅メーカーから選定（社名は実査時に確定・上長確認要）")
foot(s, "★未取得｜対象3社は実査時に選定。実名掲載は社内確認のうえ判断")

# ============================================================
# P07 いま、どこで落ちているか
# ============================================================
s = slides[6]
frame(s, "いま、どこで落ちているか",
      ["落ちているのは「資料請求→来場」の間。7展示場すべてで、月4件の上積みが必要。"])
box(s, CX0, CY0 + 0.2, CW, 1.7, fill=RED)
T(s, CX0 + 0.4, CY0 + 0.3, CW - 0.8, 1.5,
  [multi([("資料請求", 16, True, WHITE), ("　→　★ここで落ちている　→　", 13, True, "FCE4D6"), ("来場", 16, True, WHITE)],
         align="c")],
  anchor="m")
T(s, CX0, CY0 + 2.1, CW, 0.6, [one("7展示場の内訳（1施設あたり平均）", 11.5, True, NAVY)], anchor="m")
venues = ["京都・四条（総合住宅展示場）", "ABCハウジングウェルビーみのお 第1", "ABCハウジングウェルビーみのお 第2",
          "中百舌鳥住宅公園", "ショールーム（本社）", "箕面森町モデルハウス（宿泊体験型）", "ABCハウジング伊丹・昆陽の里住宅公園"]
cw7 = CW / 7
for i, v in enumerate(venues):
    x = CX0 + i * cw7
    box(s, x + 0.05, CY0 + 2.8, cw7 - 0.1, 2.6, fill=PALE)
    T(s, x + 0.12, CY0 + 2.9, cw7 - 0.24, 1.9, [one(v, 8, None, INK, align="c", ls=1.2)], anchor="t")
    T(s, x + 0.12, CY0 + 4.6, cw7 - 0.24, 0.7, [one("10→14件", 9.5, True, ORANGE, align="c")], anchor="m")
foot(s, "業界水準｜殿村さん提供の確定値。展示場ごとの内訳は「1施設あたり平均」を7施設に適用したモデル値")


# ============================================================
# P08 LINE導入後の導線
# ============================================================
s = slides[7]
frame(s, "LINE導入後の導線",
      ["LINEは、その日のうちに約8割が開封します。",
       "夜でも週末でも、相手の都合のいい時間に届きます。（出典：LINEヤフー for Business）"])
tactics = [
    ("資料請求者を自動追客", "資料請求の直後にLINE友だち化。\n営業が追いきれていない間を\nLINEが自動で埋める"),
    ("エリアで展示場を出し分け", "希望エリアに応じて\n7展示場の中から最寄りを提案\n（京都1／大阪5／兵庫1）"),
    ("宿泊体験を来場動機に", "箕面森町モデルハウスの\n宿泊体験型を来場のきっかけに\n（他展示場にはない特徴）"),
]
cw3 = (CW - 0.6) / 3
for i, (h, b) in enumerate(tactics):
    x = CX0 + i * (cw3 + 0.3)
    card(s, x, CY0 + 0.2, cw3, 3.4, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=12.5, bsz=10, ls=1.3)
box(s, CX0, CY0 + 3.9, CW, 1.8, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 3.9, CW - 0.8, 1.8,
  [one("触るのは「資料請求のあと」から「来場」までの区間だけ。広告の設計は変えない。", 12.5, True, WHITE, align="c", ls=1.3)],
  anchor="m")
foot(s, "業界水準｜LINE開封率「当日中に約8割」はLINEヤフー公式値。他の施策は意見（DYM提案）")

# ============================================================
# P09 施策一覧（施策 × 解決する課題 × 想定KPI）
# ============================================================
s = slides[8]
frame(s, "施策一覧",
      ["施策・解決する課題・想定KPIの対応表。すべて「資料請求→来場」の区間に効く施策。"])
simple_table(s, CX0, CY0 + 0.2, CW, 6.0,
             ["施策", "解決する課題", "想定KPI"],
             [["資料請求者への自動追客（LINE誘導）", "資料請求を営業サイドが最後まで追えていない", "資料請求→来場 転換率"],
              ["あいさつMSG＋来場予約への誘導", "検討8〜12か月に対し接点が資料請求の1回のみ", "友だち化率"],
              ["エリア別リッチメニュー（7展示場の出し分け）", "展示場が7施設に分散し情報が届きにくい", "展示場別・来場数"],
              ["高性能住宅の理解を深める継続配信\n（断熱等級・C値・全館空調の解説）", "高性能住宅は一度の接触では伝わらない", "来場予約 CVR"],
              ["宿泊体験（箕面森町）への送客", "来場の動機づけが弱い層への一手が無い", "宿泊体験 申込数"]],
             col_w=[CW * 0.34, CW * 0.38, CW * 0.28], bsz=8.5, row_h=1.0)
foot(s, "意見｜P3の診断（①〜④）に対応する施策案。想定KPIは方向性であり、目標値はSIM作成後に確定")

# ============================================================
# P10 想定シミュレーション（★SIM待ち）
# ============================================================
s = slides[9]
frame(s, "想定シミュレーション",
      ["★SIM未作成。ここでは測る構造だけを示す。数値は御社実績をいただき次第、反映する。"])
simple_table(s, CX0, CY0 + 0.2, CW, 2.6,
             ["", "Before（広告のみ）", "After（LINE導入後）"],
             [["資料請求数", "現状", "同左（広告費は据え置き）"],
              ["資料請求→来場 転換率", "★未取得", "★改善想定（感度分析で提示）"],
              ["来場数", "70件", "★向上想定"],
              ["来場CPA", "15万円", "★低下想定"]],
             col_w=[CW * 0.26, CW * 0.37, CW * 0.37], bsz=8.5, row_h=0.6)
placeholder(s, CX0, CY0 + 3.1, CW, 3.7, "★差込枠｜感度分析（未取得）",
            "資料請求数・転換率の実績値をいただき次第、転換率を◯pt改善した場合の\n"
            "来場数・CPAの変化を複数シナリオで試算する（lineoa-simスキル相当の作業）")
foot(s, "意見｜数値は御社実績（資料請求数・現状CPA・転換率）をいただき次第、確定")

# ============================================================
# P11 費用プラン ＋ まず90日
# ============================================================
s = slides[10]
frame(s, "費用プラン ＋ まず90日",
      ["無償付帯と別途費用を分けて示す（6ヶ月〜・税抜）。導入初月から90日の進め方。"])
plans = [
    ("① コンサル基本", "初期10万〜", "月20万（3投稿）\n月30万（5投稿）\n月50万（9投稿）"),
    ("② 初動設計＋運用", "初期10万〜", "月5万〜"),
    ("③ 運用代行・効率改善", "初期10万〜", "月0万〜\n（アカウント費のみ）"),
]
cw3 = (CW - 0.6) / 3
for i, (h, init, m) in enumerate(plans):
    x = CX0 + i * (cw3 + 0.3)
    box(s, x, CY0 + 0.2, cw3, 2.4, fill=PALE)
    T(s, x + 0.2, CY0 + 0.3, cw3 - 0.4, 0.7, [one(h, 10.5, True, NAVY, ls=1.2)], anchor="t")
    T(s, x + 0.2, CY0 + 1.0, cw3 - 0.4, 0.5, [one(init, 12.5, True, ORANGE)], anchor="t")
    T(s, x + 0.2, CY0 + 1.5, cw3 - 0.4, 0.9, [one(l, 8.5, None, INK, ls=1.2) for l in m.split("\n")], anchor="t")
band(s, CY0 + 2.8, "まず90日", fill=NAVY, h=0.65, sz=11)
days = [
    ("Day1〜30", "構築", "あいさつMSG・友だち追加動線・\nエリア別リッチメニュー設計"),
    ("Day31〜60", "稼働", "資料請求者への自動追客を開始。\n7展示場への出し分け稼働"),
    ("Day61〜90", "計測・改善", "資料請求→来場の転換率を計測。\n次フェーズ（データ収集・SIM）へ"),
]
cw3d = (CW - 0.6) / 3
for i, (h, tag, b) in enumerate(days):
    x = CX0 + i * (cw3d + 0.3)
    tb = box(s, x, CY0 + 3.6, cw3d, 0.7, fill=ORANGE)
    put_text(tb.text_frame, [one(f"{h}｜{tag}", 10.5, True, WHITE, align="c")], anchor="m")
    card(s, x, CY0 + 4.4, cw3d, 1.9, "", [b], hcol=NAVY, fill=PORANGE, hsz=1, bsz=9, ls=1.3, anchor="m")
foot(s, "意見｜DYM標準プラン（6ヶ月〜・税抜）。無償付帯：アカウント開設／プロフィール／リッチメニュー／あいさつ／初期アンケート／KW自動応答／定例会")


# ============================================================
# 業種残骸の最終チェック（正規表現スキャン・保持する11枚のみ）
# ============================================================
import re

RESIDUE_PATTERNS = [
    r"○○業界", r"〇〇業界", r"クロスセル", r"アップセル", r"転職", r"求人",
    r"利用状況タグ", r"誕生日.{0,3}記念日タグ", r"人材\(転職\)業界",
]
residue_hits = []
for i, sl in enumerate(slides[:N_PAGES], 1):
    for sh in sl.shapes:
        if sh.has_text_frame:
            txt = sh.text_frame.text
            for pat in RESIDUE_PATTERNS:
                if re.search(pat, txt):
                    residue_hits.append((i, pat, txt[:40]))
if residue_hits:
    print("★業種残骸が残っています：")
    for i, pat, txt in residue_hits:
        print(f"  slide {i}: /{pat}/ -> {txt!r}")
else:
    print("業種残骸チェック：残存0")

# ============================================================
# 11枚に絞り込み（末尾25枚を削除）
# ============================================================
trim_to(prs, slides, N_PAGES)
assert len(prs.slides) == N_PAGES, len(prs.slides)

prs.save(OUT)
print("saved:", OUT)
print(f"スライド数: {len(Presentation(OUT).slides)}")
