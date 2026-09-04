# -*- coding: utf-8 -*-
"""単品リピート通販（単品通販・D2C）業界汎用 LINEOA施策提案（36枚）

_templates/DYM_LINEOA_FMT.pptx（業界汎用FMT・36枚）をコピーし、
全36枚を _drafts/単品通販業界_スライド作成プロンプト.md の指定内容で構築する。

  python3 _build/build_tanpin_tsuhan.py
  python3 _build/qa_render.py 単品通販業界_LINEOA施策提案.pptx _qa

【この資料の背骨】
単品通販のLTVを決めているのは、新規獲得ではなく解約率。
月次解約率が3%なら顧客は33ヶ月、7%なら14ヶ月しか続かない。
そして法律は「解約させない」を禁じた。引き止めるのではなく、解約の手前で受け止める。
それができれば、限界CPOが上がって、広告が楽になる。

【★絶対に守る1点】
「解約を止める」「引き止める」とは書かない。特商法で契約解除の妨害には罰則がある。
S24 に「それでも解約なら、止めない。すぐ解約導線を出す」を必ず明記する。
→ スクリプト末尾の GUARD チェックで自動検査している。

【未取得データ（★破線の差込枠にした。数値は捏造していない）】
- S08 広告実績CVR（DYM社内実績・CV地点別＝初回購入／定期引上／F2）
- S10 Googleトレンド（サプリ／青汁／定期便／解約）
- S11 LINEヤフー前後検索（サプリ／定期便／解約／口コミ）
- S12 競合8社（本流）の友だち数・解約導線
- S35 LINEヤフー公式のD2C・単品通販導入事例（lycbiz.com）

【落とし穴メモ（_build/README.md より）】
- put_text() は必ず reset_tf() を通す（既存テキストへの追記事故を防ぐ）
- スライドの新規追加はしない。ベース36枚を clear_slide() して作り直す
- 一括置換をしない。数値は行・列を特定して書く
"""
import re
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
SRC = str(ROOT / "_templates" / "DYM_LINEOA_FMT.pptx")
OUT = str(ROOT / "単品通販業界_LINEOA施策提案.pptx")

# ---- 配色（DYM 3色：ネイビー・炭黒・オレンジ。LINEグリーンはUI部分のみ）----
TNAVY = "002060"
NAVY = "1F285A"
ORANGE = "ED7D31"
RED = "C00000"
INK = "333333"
MUT = "7F7F7F"
WHITE = "FFFFFF"
PALE = "F4F7FF"
PORANGE = "FCE4D6"
GREY = "F2F2F2"
BORDER = "D9D9D9"
GREEN = "06C755"
BEZEL = "2B2B2B"
SCREEN = "EFF2F7"
PRED = "FDF2F2"

# ---- FMTのレイアウト座標（cm）----
SW, SH = 27.52, 19.05
TITLE_XY = (1.52, 0.38, 24.4, 0.94)
LEAD_XY = (1.20, 1.80, 25.1, 1.90)
DIV_Y = 3.86
CX0, CW = 1.20, 25.12
CY0, CY1 = 4.30, 17.00
FOOT_Y = 17.35

shutil.copyfile(SRC, OUT)
prs = Presentation(OUT)
slides = list(prs.slides)
assert len(slides) == 36, len(slides)
assert (prs.slide_width, prs.slide_height) == (9906000, 6858000)


# ================= helpers（build_ec.py / build_chumon_jutaku.py と共通）=================
def set_font(run, size, bold=None, color=INK, name="メイリオ"):
    f = run.font
    f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", name)
    f.color.rgb = RGBColor.from_string(color)


ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
ANCH = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}


def reset_tf(tf):
    """既存の段落・runを全消去（追記事故の防止）"""
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)


def put_text(tf, paras, anchor="t", ml=0.14, mr=0.14, mt=0.06, mb=0.06, wrap=True):
    reset_tf(tf)
    tf.word_wrap = wrap
    tf.margin_left = Cm(ml)
    tf.margin_right = Cm(mr)
    tf.margin_top = Cm(mt)
    tf.margin_bottom = Cm(mb)
    tf.vertical_anchor = ANCH[anchor]
    first = True
    for p in paras:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = ALIGN[p.get("align", "l")]
        if p.get("sa") is not None:
            para.space_after = Pt(p["sa"])
        if p.get("ls") is not None:
            para.line_spacing = p["ls"]
        for item in p["runs"]:
            t, sz, b, c = item
            r = para.add_run()
            r.text = t
            set_font(r, sz, b, c)
    return tf


def T(slide, x, y, w, h, paras, anchor="t", **kw):
    box_ = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    put_text(box_.text_frame, paras, anchor=anchor, **kw)
    return box_


def one(text, sz, b=None, c=INK, align="l", sa=None, ls=None):
    d = {"runs": [(text, sz, b, c)], "align": align}
    if sa is not None:
        d["sa"] = sa
    if ls is not None:
        d["ls"] = ls
    return d


def multi(runs, align="l", sa=None, ls=None):
    d = {"runs": runs, "align": align}
    if sa is not None:
        d["sa"] = sa
    if ls is not None:
        d["ls"] = ls
    return d


def box(slide, x, y, w, h, fill=None, line=None, lw=1.0,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06, dash=None):
    sp = slide.shapes.add_shape(shape, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor.from_string(line)
        sp.line.width = Pt(lw)
        if dash:
            sp.line.dash_style = dash
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    reset_tf(sp.text_frame)
    return sp


def card(slide, x, y, w, h, head, body, hcol=NAVY, fill=PALE,
         hsz=11, bsz=9, line=None, anchor="t", ls=1.18):
    sp = box(slide, x, y, w, h, fill=fill, line=line)
    paras = [one(head, hsz, True, hcol, sa=3)]
    for b in (body if isinstance(body, list) else [body]):
        if b:
            paras.append(one(b, bsz, None, INK, ls=ls, sa=1))
    put_text(sp.text_frame, paras, anchor=anchor, ml=0.22, mr=0.18, mt=0.14, mb=0.10)
    return sp


def clear_slide(slide):
    spTree = slide.shapes._spTree
    for el in list(spTree):
        if el.tag.split("}")[-1] in ("sp", "cxnSp", "pic", "graphicFrame", "grpSp"):
            spTree.remove(el)


def frame(slide, title, lead):
    """FMT準拠：タイトル16pt紺（y=0.38）＋リード12pt（y=1.80・2行以内）＋区切り線"""
    clear_slide(slide)
    T(slide, *TITLE_XY, [one(title, 16, True, TNAVY)], anchor="m", ml=0, mr=0)
    T(slide, *LEAD_XY, [one(l, 12, None, INK, ls=1.28) for l in lead],
      anchor="m", ml=0, mr=0)
    ln = slide.shapes.add_connector(1, Cm(0), Cm(DIV_Y), Cm(SW), Cm(DIV_Y))
    ln.line.color.rgb = RGBColor.from_string(BORDER)
    ln.line.width = Pt(1.0)


def foot(slide, text):
    T(slide, CX0, FOOT_Y, CW, 0.9, [one(text, 7.5, None, MUT, ls=1.15)], ml=0, mr=0)


def badge(slide, x, y, w, h, text, fill=PORANGE, col=ORANGE, sz=8.5, dash=None):
    sp = box(slide, x, y, w, h, fill=fill, line=col, lw=1.0, dash=dash)
    put_text(sp.text_frame, [one(text, sz, True, col, align="c")],
             anchor="m", ml=0.06, mr=0.06, mt=0, mb=0)
    return sp


def band(slide, y, text, fill=NAVY, col=WHITE, sz=11.5, h=0.86, x=CX0, w=CW):
    sp = box(slide, x, y, w, h, fill=fill, radius=0.10)
    put_text(sp.text_frame, [one(text, sz, True, col, align="c")],
             anchor="m", ml=0.2, mr=0.2, mt=0, mb=0)
    return sp


def placeholder(slide, x, y, w, h, label, note):
    """★未取得データの差込枠（破線）。捏造しない。"""
    sp = box(slide, x, y, w, h, fill=WHITE, line=MUT, lw=1.25,
             dash=MSO_LINE_DASH_STYLE.DASH)
    put_text(sp.text_frame,
             [one(label, 11, True, MUT, align="c", sa=4),
              one(note, 8.5, None, MUT, align="c", ls=1.25)],
             anchor="m", ml=0.3, mr=0.3, mt=0.1, mb=0.1)
    return sp


def stat(slide, x, y, w, label, val, sub, h=2.4, vsz=24, fill=PALE, vcol=ORANGE):
    """数値カード（結論の数値は主役サイズで置く）"""
    box(slide, x, y, w, h, fill=fill)
    T(slide, x + 0.3, y + 0.15, w - 0.6, 0.6, [one(label, 10.5, True, NAVY)], anchor="m")
    T(slide, x + 0.3, y + 0.75, w - 0.6, 1.0, [one(val, vsz, True, vcol)], anchor="m")
    T(slide, x + 0.3, y + h - 0.65, w - 0.6, 0.55, [one(sub, 8.5, None, MUT, ls=1.15)],
      anchor="m")


def arrow_down(slide, cx, y, h=0.42, col=MUT):
    sp = box(slide, cx - 0.22, y, 0.44, h, fill=col, shape=MSO_SHAPE.DOWN_ARROW)
    return sp


def phone(slide, x, y, w, h, header, lines):
    """LINEトーク画面のモック。lines = [(kind, text)]
    kind: 'in'（BOT吹き出し）/ 'chip'（選択肢）/ 'btn'（CTA）/ 'note'（注記）"""
    box(slide, x, y, w, h, fill=BEZEL, radius=0.10)
    box(slide, x + 0.14, y + 0.42, w - 0.28, h - 0.56, fill=SCREEN, radius=0.02,
        shape=MSO_SHAPE.RECTANGLE)
    hd = box(slide, x + 0.14, y + 0.42, w - 0.28, 0.52, fill=GREEN,
             shape=MSO_SHAPE.RECTANGLE)
    put_text(hd.text_frame, [one(header, 7, True, WHITE, align="c")],
             anchor="m", ml=0.05, mr=0.05, mt=0, mb=0)
    cy = y + 1.05
    iw = w - 0.62
    for kind, text in lines:
        nlines = text.count("\n") + 1
        bh = 0.30 + 0.30 * nlines
        if kind == "in":
            sp = box(slide, x + 0.31, cy, iw, bh, fill=WHITE, line=BORDER, lw=0.75)
            put_text(sp.text_frame, [one(l, 6.5, None, INK, ls=1.18)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.14, mr=0.10, mt=0.05, mb=0.05)
        elif kind == "chip":
            sp = box(slide, x + 0.31, cy, iw, bh, fill=WHITE, line=GREEN, lw=0.9)
            put_text(sp.text_frame, [one(l, 6.5, None, "0B7A3B", align="c", ls=1.15)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.08, mr=0.08, mt=0.03, mb=0.03)
        elif kind == "btn":
            sp = box(slide, x + 0.31, cy, iw, bh, fill=GREEN)
            put_text(sp.text_frame, [one(l, 6.8, True, WHITE, align="c", ls=1.15)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.08, mr=0.08, mt=0.03, mb=0.03)
        else:  # note
            sp = slide.shapes.add_textbox(Cm(x + 0.31), Cm(cy), Cm(iw), Cm(bh))
            put_text(sp.text_frame, [one(l, 6.2, None, MUT, ls=1.15)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.05, mr=0.05, mt=0, mb=0)
        cy += bh + 0.12
    return cy


def richmenu(slide, x, y, w, h, tabs, star_idx=None):
    """リッチメニュー3タブ×6ボタン。star_idx のボタンをオレンジで強調。"""
    tab_h = 0.55
    box(slide, x, y, w, h, fill=BEZEL, radius=0.06)
    tw = (w - 0.12) / len(tabs)
    for i, (name, _) in enumerate(tabs):
        active = (i == 0)
        tb = box(slide, x + 0.06 + i * tw, y + 0.06, tw - 0.03, tab_h,
                 fill=GREEN if active else "3F3F3F", shape=MSO_SHAPE.RECTANGLE)
        put_text(tb.text_frame, [one(name, 8, True, WHITE, align="c")],
                 anchor="m", ml=0, mr=0, mt=0, mb=0)
    _, btns = tabs[0]
    gx, gy = x + 0.06, y + 0.06 + tab_h + 0.06
    gw, gh = w - 0.12, h - tab_h - 0.18
    cols, rows = 3, 2
    bw, bh = (gw - 0.06 * (cols - 1)) / cols, (gh - 0.06 * (rows - 1)) / rows
    for i, label in enumerate(btns):
        r, c = divmod(i, cols)
        hot = (i == star_idx)
        bb = box(slide, gx + c * (bw + 0.06), gy + r * (bh + 0.06), bw, bh,
                 fill=PORANGE if hot else SCREEN, line=ORANGE if hot else BORDER,
                 lw=1.5 if hot else 0.75, radius=0.04)
        put_text(bb.text_frame, [one(l, 7.2, True, ORANGE if hot else NAVY,
                                     align="c", ls=1.1)
                                 for l in label.split("\n")],
                 anchor="m", ml=0.05, mr=0.05, mt=0, mb=0)


def simple_table(slide, x, y, w, h, headers, rows, col_w=None,
                 hsz=9, bsz=8.5, header_fill=NAVY, zebra=PALE,
                 align=None, row_h=None):
    n_r, n_c = len(rows) + 1, len(headers)
    gf = slide.shapes.add_table(n_r, n_c, Cm(x), Cm(y), Cm(w), Cm(h))
    tbl = gf.table
    if col_w:
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Cm(cw)
    if row_h:
        for r in tbl.rows:
            r.height = Cm(row_h)
    for j, htext in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor.from_string(header_fill)
        c.margin_left = c.margin_right = Cm(0.12)
        c.margin_top = c.margin_bottom = Cm(0.05)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        put_text(c.text_frame, [one(htext, hsz, True, WHITE,
                                    align=(align[j] if align else "c"))],
                 anchor="m", ml=0.12, mr=0.12, mt=0.02, mb=0.02)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor.from_string(
                zebra if (zebra and i % 2 == 0) else WHITE)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            put_text(c.text_frame, [one(str(val), bsz, None, INK,
                                        align=(align[j] if align else "l"))],
                     anchor="m", ml=0.12, mr=0.12, mt=0.02, mb=0.02)
    return gf


def find_shape(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle in sh.text_frame.text:
            return sh
    return None


# ============================================================
# S01 表紙
# ============================================================
s = slides[0]
sh = find_shape(s, "業界：")
if sh:
    put_text(sh.text_frame,
             [one("業界：単品リピート通販（単品通販・D2C／業界汎用）", 12, True, TNAVY)],
             anchor="m")
sh = find_shape(s, "最上位パートナーの知見")
if sh:
    put_text(sh.text_frame,
             [one("解約を止めるのではなく、解約の手前で受け止める。LTVを伸ばし、限界CPOを上げる",
                  13, True, INK)],
             anchor="m")

# ============================================================
# S02 資料アジェンダ
# ============================================================
s = slides[1]
frame(s, "資料アジェンダ",
      ["この資料が言っていること：単品通販のLTVを決めているのは、新規獲得ではなく解約率。",
       "引き止めるのではなく、解約の手前で受け止める導線をつくる。"])
AGENDA = [
    ("1", "なぜLINEか", "S03-06"), ("2", "実態調査", "S07-12"),
    ("3", "全体設計", "S13-15"), ("4", "構築", "S16-20"),
    ("5", "配信設計", "S21-26"), ("6", "歩留まり・工数", "S27-28"),
    ("7", "成果と体制", "S29-32"), ("8", "締め", "S33-36"),
]
gx, gy, gw, gh = CX0, CY0 + 0.3, CW, 11.6
cw = gw / 4
ch = gh / 2
for i, (n, t, r) in enumerate(AGENDA):
    row, col = divmod(i, 4)
    x = gx + col * cw
    y = gy + row * ch
    box(s, x + 0.12, y + 0.12, cw - 0.24, ch - 0.24, fill=PALE)
    T(s, x + 0.12, y + 0.22, cw - 0.24, 0.9,
      [one(n, 22, True, ORANGE, align="c")], anchor="t")
    T(s, x + 0.12, y + 0.95, cw - 0.24, 1.3,
      [one(t, 12.5, True, NAVY, align="c")], anchor="t")
    T(s, x + 0.12, y + ch - 0.55, cw - 0.24, 0.5,
      [one(r, 8.5, None, MUT, align="c")], anchor="t")
foot(s, "全8章・36枚｜主語は「LINE公式アカウントの運用」。広告運用の提案書ではない")

# ============================================================
# S03 市場データ（ニーズ全体図）
# ============================================================
s = slides[2]
frame(s, "市場データ｜ニーズ全体図",
      ["通販市場そのものは伸びている。ただし上位は総合ブランドで固まっている。",
       "単品通販の本流（1商品×定期）は、その下で戦っている。"])
band(s, CY0 + 0.1, "化粧品通販 91社合計　5,921億円超／実質成長率 4.1%", fill=NAVY)
cw3 = (CW - 0.6) / 3
tops = [
    ("健康食品通販 1位", "818億円", "サントリーウエルネス（11年連続1位・前年比+5%）"),
    ("健康食品通販 2位", "493億円", "DHC（前年比+3%）"),
    ("化粧品通販 1位", "365億円", "オルビス"),
]
for i, (label, val, sub) in enumerate(tops):
    stat(s, CX0 + i * (cw3 + 0.3), CY0 + 1.2, cw3, label, val, sub, h=2.6, vsz=22)
box(s, CX0, CY0 + 4.1, CW, 2.5, fill=PORANGE)
T(s, CX0 + 0.4, CY0 + 4.3, CW - 0.8, 2.1,
  [one("上位に並ぶのは、実店舗も商品数も持つ「総合ブランド」。", 12, True, ORANGE, sa=4),
   one("単品通販の本流は 1商品 × 定期購入 のモデル。売上規模ではなく、"
       "「1人の顧客をどれだけ長く続けてもらえるか」で勝負している。", 11, None, INK, ls=1.30)],
  anchor="m")
box(s, CX0, CY0 + 7.0, CW, 4.4, fill=PALE)
T(s, CX0 + 0.4, CY0 + 7.2, CW - 0.8, 4.0,
  [one("だから、この資料が見るのは市場規模ではない", 12.5, True, NAVY, sa=6),
   one("・初回を買った人の、どれだけが2回目を買っているか（F2転換率）", 11, None, INK, ls=1.35, sa=2),
   one("・定期を始めた人が、どれだけ続いているか（月次解約率）", 11, None, INK, ls=1.35, sa=2),
   one("この2つが、単品通販のLTVをほぼ決めている。", 11.5, True, ORANGE, ls=1.35)],
  anchor="m")
foot(s, "業界水準｜日本ネット経済新聞／日本流通産業新聞「2025年版 通販売上高ランキング」"
        "（化粧品通販は2024年4月〜2025年3月決算・91社）。提案時に最新年版を確認すること")

# ============================================================
# S04 単品通販の構造（3つの壁）
# ============================================================
s = slides[3]
frame(s, "単品通販の構造｜3つの壁",
      ["新規獲得の単価は下げられない。触れるのは、買った後の2つの壁だけ。"])
walls = [
    ("① CPOの高騰", "触れない",
     ["広告依存・記事LP・アフィリエイトの相場", "検索広告CPCは10年で168%",
      "ネット広告費は構成比50.2%で初の過半数"], MUT, GREY),
    ("② 初回で消える", "触れる",
     ["初回購入後の離脱率は50〜60%超（美容・健康食品）",
      "F2転換率は30〜40%＝6〜7割が1回で終わる",
      "効果実感の前に脱落している"], ORANGE, PORANGE),
    ("③ 解約の受け皿が電話しかない", "触れる",
     ["定期購入の相談は年74,146件（2022年度）",
      "主な苦情は「電話がつながらず解約できない」",
      "★塞ぐと違法（特商法・解約妨害に罰則）"], ORANGE, PORANGE),
]
bw = (CW - 0.6) / 3
for i, (head, tag, items, col, fill) in enumerate(walls):
    x = CX0 + i * (bw + 0.3)
    box(s, x, CY0 + 0.2, bw, 7.4, fill=fill, line=col if col == ORANGE else None,
        lw=1.5)
    T(s, x + 0.3, CY0 + 0.4, bw - 0.6, 1.4, [one(head, 12.5, True, col, ls=1.2)],
      anchor="m")
    badge(s, x + 0.3, CY0 + 1.9, 1.9, 0.6, tag,
          fill=WHITE, col=col, sz=9)
    for j, it in enumerate(items):
        T(s, x + 0.3, CY0 + 2.8 + j * 1.5, bw - 0.6, 1.4,
          [one("・" + it, 10, None, INK, ls=1.32)], anchor="t")
box(s, CX0, CY0 + 8.2, CW, 3.2, fill=NAVY)
T(s, CX0 + 0.5, CY0 + 8.4, CW - 1.0, 2.8,
  [one("②③を潰すと、広告費も新規獲得数も据え置きのまま LTV が上がる。", 14, True, WHITE, sa=6),
   one("そしてLTVが上がると、限界CPO（1件の獲得にかけられる広告費の上限）が上がる。"
       "＝ 同じ広告枠でも勝てる範囲が広がる。", 11.5, None, WHITE, ls=1.32)],
  anchor="m")
foot(s, "業界水準｜離脱率・F2転換率＝単品通販のCRM各社／相談件数＝国民生活センター／"
        "CPC・広告費＝キーワードマーケティング・電通「2025年 日本の広告費」")

# ============================================================
# S05 顧客心理
# ============================================================
s = slides[4]
frame(s, "顧客心理｜なぜ1回で終わるのか",
      ["商品が悪いから辞めるのではない。効果を実感する前に、不安が先に来る。"])
psy = [
    ("「効いているのか分からない」",
     "ヘアケア・スキンケアは効果実感に時間がかかるため短期離脱が多い。"
     "食品・日用品は消費サイクルが明確なので継続率が比較的高い。",
     "→ 使い始めの伴走が効く（S22）"),
    ("「定期に縛られたくない」",
     "解約条件が分かりにくいこと自体が不安の源。"
     "だから購入後に「解約」を検索する人が出る。",
     "→ 先に選択肢を出せば解約にならない（S24）"),
    ("「言い出しにくい・繋がらない」",
     "定期購入の相談は年74,146件。主な苦情は「電話がつながらず解約できない」。"
     "電話しか窓口が無いこと自体が不満を生んでいる。",
     "→ 24時間受けられる窓口を持つ（S24）"),
]
for i, (head, body, arrow) in enumerate(psy):
    y = CY0 + 0.2 + i * 2.9
    box(s, CX0, y, CW * 0.62, 2.6, fill=PALE)
    T(s, CX0 + 0.35, y + 0.2, CW * 0.62 - 0.7, 2.2,
      [one(head, 12.5, True, NAVY, sa=5),
       one(body, 10, None, INK, ls=1.32)], anchor="m")
    box(s, CX0 + CW * 0.62 + 0.3, y, CW * 0.38 - 0.3, 2.6, fill=PORANGE)
    T(s, CX0 + CW * 0.62 + 0.6, y + 0.2, CW * 0.38 - 0.9, 2.2,
      [one(arrow, 11, True, ORANGE, ls=1.3)], anchor="m")
box(s, CX0, CY0 + 9.1, CW, 2.3, fill=GREY)
T(s, CX0 + 0.5, CY0 + 9.3, CW - 1.0, 1.9,
  [multi([("そして、追いかける手段も機能していない。", 12, True, INK),
          ("　メルマガ開封率は20%前後（業界水準）に対し、", 11, None, INK),
          ("LINEは当日中に約8割が開封（LINEヤフー公式）。", 12, True, ORANGE)], ls=1.35)],
  anchor="m")
foot(s, "業界水準｜継続率の商材差＝定期通販CRM各社／相談件数＝国民生活センター／"
        "メルマガ開封率＝国内EC平均。LY公式｜開封率は倍率ではなく実数で記載")

# ============================================================
# S06 できること／できないこと ＋ スコープ宣言
# ============================================================
s = slides[5]
frame(s, "Webマーケでできること／できないこと",
      ["先に線を引く。ここが本提案の誠実さであり、説得力の源。"])
LW = CW * 0.36          # できないこと（左）
RW = CW * 0.64 - 0.3    # できること（右）
RX = CX0 + LW + 0.3
box(s, CX0, CY0 + 0.2, LW, 6.6, fill=GREY)
T(s, CX0 + 0.35, CY0 + 0.4, LW - 0.7, 0.7,
  [one("できないこと", 13, True, MUT)], anchor="m")
cants = [
    "商品力・効果実感そのもの",
    "価格・原価・送料",
    "広告の単価（記事LP・アフィリの相場）",
    "★解約を止めること（法律で禁止）",
]
for j, it in enumerate(cants):
    hot = it.startswith("★")
    T(s, CX0 + 0.35, CY0 + 1.4 + j * 1.28, LW - 0.7, 1.18,
      [one("・" + it, 10.5, True if hot else None, RED if hot else INK, ls=1.28)],
      anchor="m")
box(s, RX, CY0 + 0.2, RW, 6.6, fill=PALE)
T(s, RX + 0.35, CY0 + 0.4, RW - 0.7, 0.7,
  [one("できること（＝本提案の範囲）", 13, True, NAVY)], anchor="m")
# 見出し列と説明列の幅を実測で確保する（説明が折り返して次の行に食い込むのを防ぐ）
HW = RW * 0.44
cans = [
    ("① 初回購入者を友だち化", "同梱物QR。商品が必ず届く"),
    ("② 効果実感まで伴走する", "使い始めのフォロー。短期離脱を潰す"),
    ("③ F2オファーを自動で出す", "消費サイクルに合わせる（S18で逆算）"),
    ("④ ★解約の受け皿をつくる", "アンケート→スキップ／周期変更を先に"),
    ("⑤ 次回発送前フォロー", "届く前に「スキップ」を選べる"),
    ("⑥ 休眠の掘り起こし", "解約後も再開提案ができる"),
]
for j, (h, b) in enumerate(cans):
    y = CY0 + 1.3 + j * 0.88
    hot = "★" in h
    T(s, RX + 0.35, y, HW, 0.82,
      [one(h, 10.5, True, ORANGE if hot else NAVY)], anchor="m")
    T(s, RX + 0.35 + HW, y, RW - 0.7 - HW, 0.82,
      [one(b, 9.5, None, INK, ls=1.2)], anchor="m")
box(s, CX0, CY0 + 7.1, CW, 4.3, fill=NAVY)
T(s, CX0 + 0.6, CY0 + 7.25, CW - 1.2, 4.0,
  [one("スコープ宣言", 11, True, ORANGE, sa=6),
   one("触るのは「買った後」だけ。広告もLPも商品も触りません。", 15, True, WHITE, ls=1.3, sa=4),
   one("そして、解約は止めません。止められないし、止めてはいけない。", 15, True, WHITE, ls=1.3, sa=4),
   one("やるのは、解約の手前に選択肢を置くことです。", 15, True, WHITE, ls=1.3)],
  anchor="m")
foot(s, "法令｜特定商取引法は契約解除の妨害行為に罰則を定めている（2022年6月1日施行の改正）")

# ============================================================
# S07 CPC推移
# ============================================================
s = slides[6]
frame(s, "実態調査①｜広告単価は、もう下がらない",
      ["新規獲得の単価は相場。だから触らない。触るのは回収側（LTV）だけ。"])
cw3 = (CW - 0.6) / 3
nums = [
    ("インターネット広告費", "4兆459億円", "前年比 +10.8%（初の4兆円超）"),
    ("総広告費に占める構成比", "50.2%", "初の過半数＝枠の奪い合い"),
    ("検索広告の平均CPC", "168%", "10年前との比較（約1.7倍）"),
]
for i, (label, val, sub) in enumerate(nums):
    stat(s, CX0 + i * (cw3 + 0.3), CY0 + 0.2, cw3, label, val, sub, h=2.8, vsz=26)
box(s, CX0, CY0 + 3.4, CW, 3.0, fill=PORANGE)
T(s, CX0 + 0.5, CY0 + 3.6, CW - 1.0, 2.6,
  [one("一番効くのは「構成比50.2%」。", 14, True, ORANGE, sa=5),
   one("広告予算の半分以上がネットに集中している＝枠の奪い合いが起きている。"
       "だから単価は下がらない。競合が減らない限り、これは構造として続く。",
       11.5, None, INK, ls=1.32)],
  anchor="m")
box(s, CX0, CY0 + 6.9, CW, 4.5, fill=NAVY)
T(s, CX0 + 0.6, CY0 + 7.1, CW - 1.2, 4.1,
  [one("下げられないなら、上限を上げるしかない。", 15, True, WHITE, sa=6),
   one("限界CPO ＝ 1件の新規顧客獲得にかけられる広告費の上限額（損益分岐点）。",
       11.5, None, WHITE, ls=1.3, sa=3),
   one("LTVが上がる → 限界CPOが上がる → 同じ広告枠でも勝てる範囲が広がる。",
       13, True, ORANGE, ls=1.3)],
  anchor="m")
foot(s, "業界水準｜電通「2025年 日本の広告費」（2026年3月5日発表）／"
        "検索広告CPCはキーワードマーケティングの分析。海外のCPM/CPAデータは本資料では使用していない")

# ============================================================
# S08 広告実績CVR（★差込枠）
# ============================================================
s = slides[7]
frame(s, "実態調査②｜弊社の広告実績（CV地点別）",
      ["CV地点を明記しないと、CVRは比較できない。単品通販は3地点で見る。"])
simple_table(s, CX0, CY0 + 0.2, CW, 3.4,
             ["CV地点", "何を測るか", "この資料での役割"],
             [["初回購入", "トライアル／お試しの申込", "獲得側。CPOの分母"],
              ["定期引上", "初回お試し → 定期購入への引き上げ", "業界水準は約20%"],
              ["F2（2回目購入）", "単発購入者の再購入", "業界水準は30〜40%"]],
             col_w=[5.0, 11.0, 9.12], align=["c", "l", "l"], row_h=0.82)
placeholder(s, CX0, CY0 + 4.1, CW, 4.2, "★差込枠｜弊社の単品通販 広告実績（未提供）",
            "CV地点別（初回購入／定期引上／F2）の CPC・CVR・CPO を社内から取得して差し込む。\n"
            "★CV地点の明記は必須。地点が違う数値を並べると比較にならない")
box(s, CX0, CY0 + 8.8, CW, 2.6, fill=GREY)
T(s, CX0 + 0.5, CY0 + 9.0, CW - 1.0, 2.2,
  [one("なお、本提案が触るのは初回購入より後ろ。", 12, True, INK, sa=4),
   one("この実績は「いま獲得にいくらかかっているか」を示すためのもので、"
       "獲得単価そのものを改善する提案ではない。", 10.5, None, INK, ls=1.3)],
  anchor="m")
foot(s, "★未取得｜提案前に社内実績を差し込むこと。数値は捏造していない")

# ============================================================
# S09 法規制と広告審査（★この業界で一番厚い）
# ============================================================
s = slides[8]
frame(s, "実態調査③｜「LINEは友だちだけが見るから広告ではない」は誤り",
      ["単品通販は、法律がビジネスモデルそのものを規制している業界。",
       "配信文面の設計は、最初から表現規制を織り込む必要がある。"])
steps = [("① 初回赤字 → 定期で回収", "「定期に引き上げること」が事業の生命線"),
         ("② 表示が過激化", "業界全体が「初回を安く見せる」表示に寄った"),
         ("③ 消費者トラブルが激増", "定期購入の相談が年7万件規模に"),
         ("④ 名指しの法改正", "改正の目的そのものが「詐欺的な定期購入商法対策」")]
sw4 = (CW - 3 * 0.55) / 4
for i, (h, b) in enumerate(steps):
    x = CX0 + i * (sw4 + 0.55)
    box(s, x, CY0 + 0.2, sw4, 2.3, fill=PALE if i < 3 else PORANGE,
        line=ORANGE if i == 3 else None, lw=1.5)
    T(s, x + 0.22, CY0 + 0.35, sw4 - 0.44, 2.0,
      [one(h, 10.5, True, ORANGE if i == 3 else NAVY, ls=1.2, sa=4),
       one(b, 8.5, None, INK, ls=1.25)], anchor="m")
    if i < 3:
        sp = box(s, x + sw4 + 0.06, CY0 + 1.15, 0.42, 0.42, fill=MUT,
                 shape=MSO_SHAPE.RIGHT_ARROW)
simple_table(s, CX0, CY0 + 2.9, CW, 3.5,
             ["法律", "何が規制されるか", "他のECは"],
             [["薬機法", "効果効能が書けない（商材が化粧品・健康食品）", "かからない業種も多い"],
              ["景品表示法", "効果の見せ方・No.1表示・打消し表示（優良誤認／有利誤認）", "かかる"],
              ["特商法（定期購入規制）",
               "最終確認画面の表示義務／★契約解除の妨害に罰則", "かからない（都度購入のため）"]],
             col_w=[5.6, 12.5, 7.02], align=["l", "l", "l"], row_h=0.85)
box(s, CX0, CY0 + 6.9, CW * 0.48, 4.5, fill=PALE)
T(s, CX0 + 0.35, CY0 + 7.1, CW * 0.48 - 0.7, 4.1,
  [one("薬機法上の「広告」3要件", 12, True, NAVY, sa=5),
   one("① 誘引性：購入意欲を昂進させる意図が明確", 10, None, INK, ls=1.3, sa=2),
   one("② 特定性：特定の商品名が明らかにされている", 10, None, INK, ls=1.3, sa=2),
   one("③ 認知性：一般人が認知できる状態", 10, None, INK, ls=1.3, sa=5),
   one("3つすべてを満たすと「広告」に該当する", 10, True, INK, ls=1.3)],
  anchor="m")
box(s, CX0 + CW * 0.48 + 0.3, CY0 + 6.9, CW * 0.52 - 0.3, 4.5, fill=PORANGE,
    line=ORANGE, lw=1.5)
T(s, CX0 + CW * 0.48 + 0.65, CY0 + 7.1, CW * 0.52 - 0.95, 4.1,
  [one("だからLINE配信の大半は「広告」に該当する", 12.5, True, ORANGE, ls=1.25, sa=5),
   one("送信元が事業者である時点で特定性は自動的に満たされ、"
       "商品案内や購入促進の内容なら誘引性も満たす。"
       "メールマガジン・DMも同じく広告媒体として扱われる。", 10.5, None, INK, ls=1.32)],
  anchor="m")
foot(s, "法令｜特定商取引法改正（2021年成立・2022年6月1日施行）／薬機法の広告3要件（薬監発第148号）。"
        "★審査に落ちた実例は社内から取得して差し込むこと")

# ============================================================
# S10 シーズナリティ（★差込枠）
# ============================================================
s = slides[9]
frame(s, "実態調査④｜検索の季節性",
      ["年間の山がどこにあるかで、企画投稿カレンダー（S25）の骨格が決まる。"])
placeholder(s, CX0, CY0 + 0.2, CW, 6.4, "★差込枠｜Googleトレンド（未取得）",
            "対象KW（必ず一語）：サプリ／青汁／定期便／解約　※5年＋1年の2本を取得する\n"
            "★Googleトレンドは相対指標。スケールの違うKWを1枚に混ぜないこと（小さい方が0に潰れる）\n"
            "「サプリ」と「解約」はスケールが違う可能性が高いため、必ず別グラフにする")
box(s, CX0, CY0 + 7.1, CW, 4.3, fill=PALE)
T(s, CX0 + 0.5, CY0 + 7.3, CW - 1.0, 3.9,
  [one("この枠で確かめたい仮説", 12, True, NAVY, sa=5),
   one("・「解約」への関心が伸びているなら、S24（解約の受け皿）の必要性がデータで裏づけられる",
       10.5, None, INK, ls=1.32, sa=3),
   one("・「定期便」の推移は、定期モデルそのものへの受容度を示す", 10.5, None, INK, ls=1.32, sa=3),
   one("・季節の山が出れば、S25の企画投稿を「◯月：企画（参照データ）」の形で置ける",
       10.5, None, INK, ls=1.32)],
  anchor="m")
foot(s, "★未取得｜取得手順は本文に記載。生成には _build/make_images.py（trend）を使う")

# ============================================================
# S11 前後検索（★差込枠）
# ============================================================
s = slides[10]
frame(s, "実態調査⑤｜前後検索（購入の前と後で、何を調べているか）",
      ["単品通販の特徴は、不安が「買う前」より「買った後」に出ること。"])
cw2 = (CW - 0.4) / 2
box(s, CX0, CY0 + 0.2, cw2, 2.6, fill=PALE)
T(s, CX0 + 0.35, CY0 + 0.4, cw2 - 0.7, 2.2,
  [one("購入の前", 12, True, NAVY, sa=4),
   one("効果／口コミ／成分／比較／最安", 11, None, INK, ls=1.3)], anchor="m")
box(s, CX0 + cw2 + 0.4, CY0 + 0.2, cw2, 2.6, fill=PORANGE, line=ORANGE, lw=1.5)
T(s, CX0 + cw2 + 0.75, CY0 + 0.4, cw2 - 0.7, 2.2,
  [one("購入の後　★ここが本命", 12, True, ORANGE, sa=4),
   one("解約／返品／休止／縛り", 11, True, INK, ls=1.3)], anchor="m")
placeholder(s, CX0, CY0 + 3.3, CW, 4.4, "★差込枠｜LINEヤフー 前後検索（Journey・未取得）",
            "対象KW（必ず一語）：サプリ／定期便／解約／口コミ\n"
            "★「解約」が本命。解約の前後に何を調べているかが分かれば、\n"
            "S24で先回りすべき内容（スキップ・周期変更・休止のどれを先に出すか）が決まる")
box(s, CX0, CY0 + 8.2, CW, 3.2, fill=NAVY)
T(s, CX0 + 0.6, CY0 + 8.4, CW - 1.2, 2.8,
  [one("「解約」を調べている人は、まだ解約していない。", 14, True, WHITE, sa=5),
   one("調べている時点で接触できれば、スキップや周期変更という選択肢を出せる。"
       "電話しか窓口が無いと、この瞬間に触れない。", 11.5, None, WHITE, ls=1.32)],
  anchor="m")
foot(s, "★未取得｜KWは必ず一語で取得する（二語にすると検索Volが落ちて使えない）")

# ============================================================
# S12 他社分析
# ============================================================
s = slides[11]
frame(s, "実態調査⑥｜競合のLINE運用",
      ["「大手」と「単品通販の本流」は別物。本流の運用を見る。"])
simple_table(s, CX0, CY0 + 0.2, CW, 2.9,
             ["運用の型", "特徴", "解約導線"],
             [["診断ドリブン型", "診断で属性を取得し、検討段階別に出し分ける", "有無を実測で確認"],
              ["特典・会員型", "初回特典・継続特典で友だち化。配信は販促中心", "有無を実測で確認"],
              ["カタログ・配信型", "新商品や特集を定期配信。セグメントは浅い", "有無を実測で確認"],
              ["CRM統合型", "購買・定期ステータスと連携し1to1で出し分ける", "★LINEで完結する型"]],
             col_w=[5.6, 13.0, 6.52], align=["l", "l", "l"], row_h=0.62)
placeholder(s, CX0, CY0 + 3.4, CW * 0.6, 4.4, "★差込枠｜競合8社の友だち数・配信内容（未取得）",
            "北の達人／新日本製薬／再春館製薬所／やずや／\n"
            "世田谷自然食品／キューサイ／健康家族／ステラ漢方\n"
            "page.line.me（公式）から取得し、取得日を必ず併記する")
box(s, CX0 + CW * 0.6 + 0.3, CY0 + 3.4, CW * 0.4 - 0.3, 4.4, fill=PORANGE,
    line=ORANGE, lw=1.5)
T(s, CX0 + CW * 0.6 + 0.65, CY0 + 3.6, CW * 0.4 - 0.95, 4.0,
  [one("★この業界固有の記録項目", 11.5, True, ORANGE, sa=5),
   one("・リッチメニューに「解約」「休止」があるか", 10, None, INK, ls=1.3, sa=3),
   one("・解約導線がLINEで完結するか（電話に逃がしていないか）", 10, None, INK, ls=1.3, sa=3),
   one("・次回発送前に通知が来るか／スキップを選べるか", 10, None, INK, ls=1.3)],
  anchor="m")
box(s, CX0, CY0 + 8.3, CW, 3.1, fill=GREY)
T(s, CX0 + 0.5, CY0 + 8.5, CW - 1.0, 2.7,
  [one("参考｜サントリーウエルネス・DHC・ファンケル・オルビスは総合ブランド。", 11.5, True, INK, sa=4),
   one("売上規模の比較には使えるが、1商品×定期の運用の型としては参考にならない。"
       "混ぜて分析しないこと。", 10.5, None, INK, ls=1.3)],
  anchor="m")
foot(s, "★未取得｜数値は page.line.me（公式）から。第三者サイトの数値は掲載しない。"
        "実名掲載の可否は社内確認が必要")

# ============================================================
# S13 カスタマージャーニー フェーズ8分類 × CV3段
# ============================================================
s = slides[12]
frame(s, "全体設計①｜カスタマージャーニー（8フェーズ）× CV3段",
      ["単品通販は、友だち追加（軽CV）の起点が「購入後」にある。ここが他業界と違う。"])
phases = ["認知", "興味", "比較", "初回購入", "使用開始", "F2", "定期化", "継続／解約検討"]
pw = (CW - 7 * 0.14) / 8
for i, p in enumerate(phases):
    x = CX0 + i * (pw + 0.14)
    hot = i >= 4
    box(s, x, CY0 + 0.2, pw, 1.0, fill=PORANGE if hot else GREY,
        line=ORANGE if hot else None, lw=1.2)
    T(s, x, CY0 + 0.2, pw, 1.0, [one(p, 9.5, True, ORANGE if hot else MUT, align="c")],
      anchor="m", ml=0.04, mr=0.04)
T(s, CX0, CY0 + 1.35, CW, 0.6,
  [one("← 触らない（広告・LP・商品）　｜　★本提案が触る区間 →", 10, True, INK, align="c")],
  anchor="m")
cv = [
    ("① 軽CV", "LINE友だち追加", "同梱物QR／購入完了画面／マイページ。"
     "★購入後に取るのが基本。商品が必ず届くので取りこぼしが少ない", PORANGE, ORANGE),
    ("② 主CV", "F2（2回目購入）／定期引上", "初回離脱50〜60%を削る。"
     "消費サイクルに合わせてオファーを出す", PORANGE, ORANGE),
    ("③ 最終CV／LTV", "定期の継続", "★本提案の主戦場。解約の手前で受け止める。"
     "月次解約率がLTVをほぼ決める", PORANGE, ORANGE),
]
for i, (tag, head, body, fill, col) in enumerate(cv):
    y = CY0 + 2.3 + i * 3.0
    box(s, CX0, y, CW, 2.7, fill=fill, line=col, lw=1.3)
    T(s, CX0 + 0.35, y + 0.2, 3.4, 2.3, [one(tag, 12, True, col)], anchor="m")
    T(s, CX0 + 3.9, y + 0.2, 7.0, 2.3, [one(head, 13, True, NAVY, ls=1.2)], anchor="m")
    T(s, CX0 + 11.2, y + 0.2, CW - 11.6, 2.3, [one(body, 10.5, None, INK, ls=1.3)],
      anchor="m")
foot(s, "意見｜重CVが遠い業種ではないが、軽CVを「買う前」に取れない。"
        "ナーチャリングの起点が購入後にあるという、他業界に無い形")

# ============================================================
# S14 施策全体像 ＋ 対策領域マップ
# ============================================================
s = slides[13]
frame(s, "全体設計②｜施策全体像と対策領域",
      ["どの施策が、誰を拾い、どのKPIに効くのか。1枚で対応させる。"])
simple_table(s, CX0, CY0 + 0.2, CW, 6.4,
             ["施策", "誰を拾うか", "効くKPI", "枚"],
             [["同梱物QR・購入完了画面での友だち化", "初回購入者（未友だち）", "全段の前提", "S16"],
              ["あいさつ＋初期アンケート", "友だち追加直後", "タグ設計の起点", "S17-18"],
              ["LINE ID連携", "全員（定期ステータスの取得）", "★以降すべての前提", "S20"],
              ["使用フォロー配信", "使い始めで不安になっている人", "F2転換率", "S22"],
              ["F2オファー（消費サイクル連動）", "そろそろ無くなる人", "★F2転換率", "S22"],
              ["次回発送前フォロー（スキップ提示）", "定期継続中の人", "★月次解約率", "S23"],
              ["解約シナリオ", "解約を考えている人", "★月次解約率", "S24"],
              ["休眠掘り起こし", "解約後・離脱後の人", "LTV", "S27"]],
             col_w=[9.2, 7.4, 5.6, 2.92], align=["l", "l", "l", "c"], row_h=0.68)
box(s, CX0, CY0 + 7.1, CW, 4.3, fill=NAVY)
T(s, CX0 + 0.6, CY0 + 7.3, CW - 1.2, 3.9,
  [one("対策領域は2つだけ", 12, True, ORANGE, sa=6),
   one("① 買った直後〜F2まで（効果実感の伴走とタイミング）　→ F2転換率",
       13, True, WHITE, ls=1.3, sa=4),
   one("② 定期継続中〜解約検討（選択肢を先に出す）　→ 月次解約率",
       13, True, WHITE, ls=1.3, sa=4),
   one("広告・LP・商品には一切触れない。", 11.5, None, WHITE, ls=1.3)],
  anchor="m")
foot(s, "意見｜離脱防止＝未CV者を軽いCVへ、サンクス系＝CV済み者を重いCVへ、の対応で設計している")

# ============================================================
# S15 施策展開図（初期・月次）
# ============================================================
s = slides[14]
frame(s, "全体設計③｜施策展開図（初期構築と月次運用）",
      ["初期に「基盤」を作り、月次で「回す」。基盤が無いと月次が打てない。"])
box(s, CX0, CY0 + 0.2, CW * 0.48, 5.6, fill=PALE)
T(s, CX0 + 0.35, CY0 + 0.4, CW * 0.48 - 0.7, 0.7,
  [one("初期構築（1〜2ヶ月目）", 13, True, NAVY)], anchor="m")
init = ["アカウント開設・プロフィール整備",
        "★LINE ID連携（購買・定期ステータス）",
        "同梱物QR・購入完了画面の導線設置",
        "あいさつメッセージ／初期アンケート",
        "リッチメニュー3タブ（解約・休止の枠を含む）",
        "シナリオ3本の設計と実装"]
for j, it in enumerate(init):
    hot = it.startswith("★")
    T(s, CX0 + 0.35, CY0 + 1.3 + j * 0.75, CW * 0.48 - 0.7, 0.7,
      [one("・" + it, 10.5, True if hot else None, ORANGE if hot else INK, ls=1.25)],
      anchor="m")
box(s, CX0 + CW * 0.48 + 0.3, CY0 + 0.2, CW * 0.52 - 0.3, 5.6, fill=PALE)
T(s, CX0 + CW * 0.48 + 0.65, CY0 + 0.4, CW * 0.52 - 0.95, 0.7,
  [one("月次運用（3ヶ月目〜）", 13, True, NAVY)], anchor="m")
mon = ["企画投稿（季節・成分・使い方）",
       "F2オファーのタイミング調整（消費サイクル）",
       "解約アンケートの回答分析 → オファー改善",
       "セグメント配信（定期回数・購入商品別）",
       "通知メッセージ（発送・次回発送前）",
       "定例会でF2転換率・月次解約率をレビュー"]
for j, it in enumerate(mon):
    T(s, CX0 + CW * 0.48 + 0.65, CY0 + 1.3 + j * 0.75, CW * 0.52 - 0.95, 0.7,
      [one("・" + it, 10.5, None, INK, ls=1.25)], anchor="m")
box(s, CX0, CY0 + 6.5, CW, 4.9, fill=PORANGE)
T(s, CX0 + 0.6, CY0 + 6.7, CW - 1.2, 4.5,
  [one("★順番を間違えない", 12.5, True, ORANGE, sa=6),
   one("LINE ID連携（S20）が最初。これが無いと「定期何回目か」「次回発送日はいつか」が"
       "分からないので、F2オファーもスキップ提案も解約受付も全部打てない。",
       12, None, INK, ls=1.32, sa=4),
   one("連携なしで始めると、全員に同じ配信を送るだけのアカウントになる。",
       12, True, INK, ls=1.32)],
  anchor="m")
foot(s, "意見｜6ヶ月契約を前提とした標準的な展開。連携の可否で初期の所要期間が変わる")

# ============================================================
# S16 友だち追加動線
# ============================================================
s = slides[15]
frame(s, "構築①｜友だち追加動線（★同梱物QRが主役）",
      ["単品通販は「商品が必ず届く」。この一点が、他業種に無い最強の動線になる。"])
routes = [
    ("★ 同梱物QR", "主役",
     "初回配送に必ず同梱される。開封率が高く、購入者を確実に捕まえられる。"
     "「解約もここでできます」と書くと追加率が上がる", True),
    ("購入完了画面", "併用",
     "決済直後にその場で追加してもらう。離脱が少ない", False),
    ("マイページ", "併用",
     "定期の変更・確認をしに来た人を拾う。解約導線の手前に置く", False),
    ("通知メッセージ経由", "拡張",
     "電話番号ベースで届くので、未友だちの既存顧客にも到達できる（S34）", False),
    ("LINE広告 CPF", "拡張",
     "友だち追加課金。新規の母数を増やしたいときのみ", False),
]
for i, (head, tag, body, hot) in enumerate(routes):
    y = CY0 + 0.2 + i * 2.25
    box(s, CX0, y, CW, 2.05, fill=PORANGE if hot else PALE,
        line=ORANGE if hot else None, lw=1.5)
    T(s, CX0 + 0.35, y + 0.15, 5.4, 1.75,
      [one(head, 12.5, True, ORANGE if hot else NAVY, ls=1.2)], anchor="m")
    badge(s, CX0 + 5.9, y + 0.6, 1.7, 0.7, tag, fill=WHITE,
          col=ORANGE if hot else NAVY, sz=9)
    T(s, CX0 + 8.0, y + 0.15, CW - 8.4, 1.75, [one(body, 10.5, None, INK, ls=1.3)],
      anchor="m")
foot(s, "意見｜LINE広告CPFは友だち追加動線として置いている。広告運用の提案ではない")

# ============================================================
# S17 あいさつメッセージ
# ============================================================
s = slides[16]
frame(s, "構築②｜あいさつメッセージ",
      ["最初の1通で「ここは売り込みの場ではない」と分からせる。",
       "★「解約もここでできます」を、こちらから先に言う。"])
phone(s, CX0 + 0.6, CY0 + 0.2, 7.4, 10.8, "●●（商品名）公式", [
    ("in", "ご購入ありがとうございます。\n担当の●●です。"),
    ("in", "✅ お電話はいたしません\n✅ しつこい配信もしません"),
    ("in", "【このアカウントでできること】\n・発送状況の確認\n・お届け日の変更・スキップ\n"
            "・解約のお手続き\n・使い方のご相談"),
    ("in", "初回のお届けは\n●月●日ごろの予定です。"),
    ("btn", "使い方を見る"),
    ("chip", "お届け日を変更する"),
    ("note", "※ 友だち追加から数秒以内に自動送信"),
])
box(s, CX0 + 8.6, CY0 + 0.2, CW - 8.0, 5.0, fill=PORANGE, line=ORANGE, lw=1.5)
T(s, CX0 + 8.95, CY0 + 0.4, CW - 8.7, 4.6,
  [one("★「解約のお手続き」を最初に書く理由", 12.5, True, ORANGE, sa=5),
   one("解約導線を隠すのは、特商法上のリスクであるだけでなく、"
       "実際に「電話がつながらない」という苦情の原因になっている。", 10.5, None, INK, ls=1.32, sa=4),
   one("先に出しておくと、解約したくなった人が必ずLINEに来る。"
       "＝ スキップや周期変更を提示できる場所に来てくれる。", 10.5, True, INK, ls=1.32)],
  anchor="m")
box(s, CX0 + 8.6, CY0 + 5.6, CW - 8.0, 5.4, fill=PALE)
T(s, CX0 + 8.95, CY0 + 5.8, CW - 8.7, 5.0,
  [one("書かないこと", 12, True, NAVY, sa=5),
   one("・「いつでも解約可能」などの強調表示（誤認表示として禁止されている）",
       10.5, None, INK, ls=1.32, sa=3),
   one("・効果効能を断定する表現（薬機法）", 10.5, None, INK, ls=1.32, sa=3),
   one("・初回価格だけを目立たせる書き方（2回目以降の条件を併記する）",
       10.5, None, INK, ls=1.32)],
  anchor="m")
foot(s, "法令｜「お試し」「いつでも解約可能」等の強調表示で誤認させることは特商法で禁止されている")

# ============================================================
# S18 初期アンケート → タグ設計
# ============================================================
s = slides[17]
frame(s, "構築③｜初期アンケートとタグ設計",
      ["★使用頻度を最初に聞く。これがF2オファーのタイミングを逆算する材料になる。"])
simple_table(s, CX0, CY0 + 0.2, CW, 4.2,
             ["設問", "選択肢の例", "取れるタグ", "何に使うか"],
             [["お悩みは", "乾燥／ハリ／くすみ／疲れ", "悩みタグ", "配信内容の出し分け"],
              ["★1日の使用量は", "1回／2回／不定期", "使用頻度タグ", "★F2オファーの日数を逆算"],
              ["同種の商品の使用歴", "初めて／他社から乗り換え", "経験タグ", "使い方の説明の深さ"],
              ["いつ使うことが多いか", "朝／夜／両方", "生活タグ", "配信時間の最適化"]],
             col_w=[6.0, 7.4, 4.6, 7.12], align=["l", "l", "l", "l"], row_h=0.82)
box(s, CX0, CY0 + 4.9, CW, 3.2, fill=PORANGE, line=ORANGE, lw=1.5)
T(s, CX0 + 0.6, CY0 + 5.1, CW - 1.2, 2.8,
  [one("★F2オファーのタイミングは、日数固定にしない", 13, True, ORANGE, sa=5),
   one("商品の内容量 ÷ 1日の使用量 ＝ 無くなる日　を逆算し、その少し前に出す。"
       "同じ商品でも使用頻度が違えば、無くなる日は変わる。", 11.5, None, INK, ls=1.32)],
  anchor="m")
box(s, CX0, CY0 + 8.6, CW, 2.8, fill=PALE)
T(s, CX0 + 0.6, CY0 + 8.8, CW - 1.2, 2.4,
  [one("設問は4つまで。多いと最後まで答えてもらえない。", 11.5, True, NAVY, sa=4),
   one("残りの属性は、購買データ側（S20のID連携）から自動で取る。"
       "聞かなくて済むことは聞かない。", 10.5, None, INK, ls=1.3)],
  anchor="m")
foot(s, "意見｜業界汎用のため具体的な日数は記載していない。商材ごとに逆算して設定する")

# ============================================================
# S19 リッチメニュー
# ============================================================
s = slides[18]
frame(s, "構築④｜リッチメニュー（3タブ × 6枠＝18ボタン）",
      ["★「解約・休止」を必ず1枠に置く。隠さないことが、結果的に解約を減らす。"])
richmenu(s, CX0 + 0.8, CY0 + 0.3, 14.4, 6.6, [
    ("お届け", ["お届け日を\n変更する", "今回だけ\nスキップ", "★解約・休止\nのご相談",
                "配送状況を\n見る", "お届け間隔を\n変える", "数量を\n変える"]),
    ("使い方", ["正しい\n使い方", "よくある\nご質問", "成分について",
                "効果を感じない\nときは", "保管方法", "相談する"]),
    ("お客様情報", ["ご登録情報", "お支払い方法", "ご購入履歴",
                    "クーポン", "お友だち紹介", "お問い合わせ"]),
], star_idx=2)
box(s, CX0 + 15.8, CY0 + 0.3, CW - 15.4, 6.6, fill=PORANGE, line=ORANGE, lw=1.5)
T(s, CX0 + 16.15, CY0 + 0.5, CW - 15.8, 6.2,
  [one("★なぜ解約を1軍に置くのか", 12.5, True, ORANGE, sa=5),
   one("探させると、人は電話に行く。電話が繋がらないと苦情になる。",
       10.5, None, INK, ls=1.32, sa=4),
   one("LINEに置けば、タップした瞬間に「今回だけスキップ」「お届け間隔を変える」を"
       "先に提示できる。", 10.5, None, INK, ls=1.32, sa=4),
   one("解約を隠す設計より、解約を受ける設計のほうが、結果として残る。",
       11, True, INK, ls=1.32)],
  anchor="m")
box(s, CX0, CY0 + 7.5, CW, 3.9, fill=PALE)
T(s, CX0 + 0.6, CY0 + 7.7, CW - 1.2, 3.5,
  [one("タブの並び順にも意味がある", 12, True, NAVY, sa=5),
   one("1番目「お届け」＝ 定期中の人が一番よく来る用事（変更・スキップ・解約）を最短で。",
       10.5, None, INK, ls=1.3, sa=3),
   one("2番目「使い方」＝ 効果を感じない人の受け皿。短期離脱の防波堤。",
       10.5, None, INK, ls=1.3, sa=3),
   one("3番目「お客様情報」＝ 頻度は低いが必要なもの。", 10.5, None, INK, ls=1.3)],
  anchor="m")
foot(s, "意見｜18ボタンは商材により調整する。ただし「解約・休止」の枠は必ず1軍タブに残すこと")

# ============================================================
# S20 LINE ID連携（★前提）
# ============================================================
s = slides[19]
frame(s, "構築⑤｜LINE ID連携（★以降すべての前提）",
      ["「定期何回目か」「次回発送日はいつか」が分からないと、"
       "この後の施策は1つも打てない。"])
simple_table(s, CX0, CY0 + 0.2, CW, 3.8,
             ["やりたいこと", "ID連携なし", "ID連携あり"],
             [["購入商品に応じたF2オファー", "× 全員に同じ配信しか送れない", "○ 買った商品で出し分け"],
              ["定期の次回発送前フォロー", "× 打てない", "○ 発送サイクルに合わせて自動"],
              ["★解約受付・スキップ受付", "× 打てない", "○ 定期ステータスを見て出し分け"],
              ["購入者を配信から除外", "× クーポンを無駄打ち", "○ 除外できる"]],
             col_w=[9.6, 8.0, 7.52], align=["l", "l", "l"], row_h=0.76)
dep = [("S22 F2オファー", 0), ("S23 スキップ提示", 1), ("S24 解約シナリオ", 2),
       ("S26 通知メッセージ", 3)]
dw = (CW - 3 * 0.3) / 4
for i, (label, _) in enumerate(dep):
    x = CX0 + i * (dw + 0.3)
    box(s, x, CY0 + 4.5, dw, 1.2, fill=PALE)
    T(s, x, CY0 + 4.5, dw, 1.2, [one(label, 10.5, True, NAVY, align="c")],
      anchor="m", ml=0.1, mr=0.1)
    arrow_down(s, x + dw / 2, CY0 + 5.85, 0.5)
box(s, CX0, CY0 + 6.6, CW, 1.6, fill=ORANGE)
T(s, CX0, CY0 + 6.6, CW, 1.6,
  [one("S20　LINE ID連携（購買データ・定期ステータス × LINE）", 14, True, WHITE, align="c")],
  anchor="m")
box(s, CX0, CY0 + 8.7, CW, 2.7, fill=NAVY)
T(s, CX0 + 0.6, CY0 + 8.9, CW - 1.2, 2.3,
  [one("連携しないまま始めると、全員に同じ配信を送るだけのアカウントになる。",
       13.5, True, WHITE, sa=4),
   one("LINE IDをマスタIDにすれば、EC・店舗・SNSを横断したCRM基盤としても使える。",
       11, None, WHITE, ls=1.3)],
  anchor="m")
foot(s, "業界水準｜LINEログイン／LINEミニアプリ経由で顧客IDとLINE IDを紐づける。"
        "カート/定期システム側の連携可否は初期に確認が必要")

# ============================================================
# S21 シナリオ3本の設計表
# ============================================================
s = slides[20]
frame(s, "配信設計①｜シナリオは3本（購入後だけで完結する）",
      ["単品通販はLPで直接買う。だから「買う前」のシナリオを持たない。"])
simple_table(s, CX0, CY0 + 0.2, CW, 11.2,
             ["シナリオ", "タイミング", "方向性", "根拠"],
             [["① 初回購入後\n（F2狙い）", "購入直後",
               "サンクス。発送予定と正しい使い方。レビューは求めない",
               "初回離脱が50〜60%超で最も落ちる区間のため"],
              ["", "使い始めの頃",
               "★使用フォロー。「効果実感までの目安」を伝える",
               "効果実感に時間がかかる商材ほど短期離脱が多い"],
              ["", "★無くなる少し前",
               "★F2オファー。「そろそろ切れる頃です」＋再購入導線",
               "内容量÷1日使用量で逆算（S18の使用頻度）"],
              ["", "F2の後", "定期・まとめ買いへの引き上げ", "定期引上率は業界水準で約20%"],
              ["② 定期継続中\n（解約予防）", "★次回発送の数日前",
               "★「今回はスキップできます」を先に出す",
               "余っていることが解約理由の上位にあるため"],
              ["", "効果実感の停滞期", "使い方の見直し・成分の解説",
               "「効果を感じない」を解約前に拾うため"],
              ["", "節目（3回目・6回目）", "継続特典・プラン変更の提案",
               "継続の意思を確認できる自然なタイミング"],
              ["③ 解約\n（S24で詳述）", "解約意思を検知した瞬間",
               "★アンケート→理由別オファー→それでも解約なら止めない",
               "★特商法。解約妨害には罰則がある"]],
             col_w=[4.4, 5.0, 8.6, 7.12], align=["l", "l", "l", "l"], row_h=1.3,
             bsz=8)
foot(s, "★形式は「◯日後：方向性（〜というデータがあるため）」。"
        "具体的な日数は商材で変わるため、業界汎用版では逆算の型のみ示している")

# ============================================================
# S22 実文面①（サンクス・使用フォロー・F2オファー）
# ============================================================
s = slides[21]
frame(s, "配信設計②｜実文面（初回購入後 → F2）",
      ["★F2オファーは「売り込み」ではなく「切れる頃のお知らせ」として出す。"])
phone(s, CX0 + 0.4, CY0 + 0.2, 7.2, 10.8, "●●公式｜使用フォロー", [
    ("in", "使い始めて1週間ですね。\nいかがですか？"),
    ("in", "【この時期によくあるご質問】\n"
            "Q. まだ変化を感じません\n"
            "A. 肌の生まれ変わりには\n　 時間がかかります。\n　 まずは1本お使いください。"),
    ("in", "使い方のコツを\n動画にまとめました。"),
    ("chip", "使い方の動画を見る"),
    ("note", "※ 効果を断定する表現は使わない（薬機法）"),
])
phone(s, CX0 + 8.6, CY0 + 0.2, 7.2, 10.8, "●●公式｜F2オファー", [
    ("in", "そろそろ残り少なく\nなる頃でしょうか。"),
    ("in", "お使いのペースだと\n●月●日ごろに\n無くなる計算です。"),
    ("in", "切らさずお使いいただくなら、\n今のうちのご注文が安心です。"),
    ("btn", "同じものを注文する"),
    ("chip", "定期にして送料無料にする"),
    ("note", "※ 通知冒頭15字：そろそろ残り少なく"),
])
box(s, CX0 + 16.8, CY0 + 0.2, CW - 16.4, 10.8, fill=PORANGE, line=ORANGE, lw=1.5)
T(s, CX0 + 17.15, CY0 + 0.4, CW - 16.7, 10.4,
  [one("★この2通の作法", 12.5, True, ORANGE, sa=6),
   one("使用フォローでは売らない。", 11, True, INK, ls=1.3, sa=3),
   one("「効果を感じない」を先に拾うことが、いちばんの離脱防止になる。",
       10, None, INK, ls=1.3, sa=6),
   one("F2オファーは日付で語る。", 11, True, INK, ls=1.3, sa=3),
   one("「買ってください」ではなく「●日に無くなります」。"
       "事実を伝えるだけで、必要な人は自分で動く。", 10, None, INK, ls=1.3, sa=6),
   one("薬機法の線引き", 11, True, INK, ls=1.3, sa=3),
   one("効果効能を断定しない。体験談も、広告に該当するため同じ規制がかかる。",
       10, None, INK, ls=1.3)],
  anchor="m")
foot(s, "法令｜LINE配信は薬機法上の「広告」に該当する（誘引性・特定性・認知性）。"
        "文面は必ず薬機法チェックを通すこと")

# ============================================================
# S23 実文面②（次回発送前・節目）
# ============================================================
s = slides[22]
frame(s, "配信設計③｜実文面（定期継続中・解約予防）",
      ["★次回発送の数日前に「今回はスキップできます」を、こちらから先に出す。"])
phone(s, CX0 + 0.6, CY0 + 0.2, 7.4, 10.8, "●●公式｜次回発送前", [
    ("in", "次回のお届けは\n●月●日の予定です。"),
    ("in", "まだ残っていませんか？\n余っているときは\nスキップできます。"),
    ("chip", "今回はスキップする"),
    ("chip", "お届け間隔を変える"),
    ("chip", "このままお届け"),
    ("note", "※ 発送3日前に自動送信"),
])
box(s, CX0 + 8.6, CY0 + 0.2, CW - 8.0, 5.2, fill=PORANGE, line=ORANGE, lw=1.5)
T(s, CX0 + 8.95, CY0 + 0.4, CW - 8.7, 4.8,
  [one("★これは引き止めではない", 13, True, ORANGE, sa=5),
   one("「余っているのに届く」は、解約理由の上位。"
       "届く前にスキップを選べるようにするのは、顧客にとっての利便であって、"
       "解約の妨害ではない。", 11, None, INK, ls=1.32, sa=4),
   one("解約を申し出た人を引き止めるのとは、行為として全く別のもの。",
       11, True, INK, ls=1.32)],
  anchor="m")
box(s, CX0 + 8.6, CY0 + 5.8, CW - 8.0, 5.2, fill=PALE)
T(s, CX0 + 8.95, CY0 + 6.0, CW - 8.7, 4.8,
  [one("節目の配信（3回目・6回目）", 13, True, NAVY, sa=5),
   one("・3回目：ここまでの使用量と、次のステップの提案", 10.5, None, INK, ls=1.3, sa=3),
   one("・6回目：継続特典／まとめ買いへのプラン変更", 10.5, None, INK, ls=1.3, sa=3),
   one("・いずれも「やめる」も同じ画面から選べる状態にしておく",
       10.5, True, INK, ls=1.3)],
  anchor="m")
foot(s, "意見｜スキップ・間隔変更・数量変更は、解約の「代わり」になる選択肢。"
        "これを出せるかどうかで月次解約率が変わる")

# ============================================================
# S24 解約シナリオ（★目玉）
# ============================================================
s = slides[23]
frame(s, "配信設計④｜解約シナリオ（★この提案の核心）",
      ["解約を止めることはできない。できるのは、解約の手前に選択肢を置くことだけ。"])
flow = [
    ("解約意思の検知", "リッチメニュー「解約・休止のご相談」をタップ 等", PALE, NAVY),
    ("① 解約アンケート", "理由を1タップで。入力の手間を増やさない", PALE, NAVY),
    ("② 理由に応じたオファー",
     "余っている→スキップ／周期変更　　高い→プラン変更／数量変更\n"
     "効果を感じない→使い方の見直し／別商品　　一時的な事情→休止",
     PORANGE, ORANGE),
    ("③ それでも解約なら　止めない。すぐ解約導線を出す",
     "★ここで引き止めると特商法違反のリスク。これは引き止め施策ではありません",
     PRED, RED),
]
# 高さの合計＝1.5+1.5+2.3+2.3（箱）＋0.55×3（矢印の間）＝9.25cm。
# CY0+0.2 から始めて 13.75 で終わり、帯（14.05〜15.95）と重ならない。
y = CY0 + 0.2
for i, (head, body, fill, col) in enumerate(flow):
    h = 2.3 if i >= 2 else 1.5
    box(s, CX0, y, CW, h, fill=fill, line=col, lw=1.6 if i >= 2 else 1.0)
    T(s, CX0 + 0.5, y + 0.12, CW - 1.0, h - 0.24,
      [one(head, 14 if i == 3 else 12.5, True, col, ls=1.2, sa=5),
       one(body, 10.5, True if i == 3 else None, INK, ls=1.3)], anchor="m")
    y += h
    if i < len(flow) - 1:
        arrow_down(s, CX0 + CW / 2, y + 0.07, 0.41)
        y += 0.55
box(s, CX0, CY0 + 9.75, CW, 1.9, fill=NAVY)
T(s, CX0 + 0.6, CY0 + 9.9, CW - 1.2, 1.6,
  [one("解約アンケートの設置と、理由別オファーの自動化だけで、"
       "月次チャーン率は1〜2%削減できる（業界水準）。", 12.5, True, WHITE, ls=1.3)],
  anchor="m")
foot(s, "法令｜特定商取引法は契約解除の妨害行為に罰則を定めている。"
        "本シナリオは解約導線を塞ぐものではなく、選択肢を先に提示するもの")

# ============================================================
# S25 年間の企画投稿カレンダー
# ============================================================
s = slides[24]
frame(s, "配信設計⑤｜年間の企画投稿カレンダー",
      ["S10の季節性データが入り次第、「◯月：企画（参照データ）」の形で確定させる。"])
cal = [
    ("1月", "新年の習慣づくり企画", "年始は継続の意思が固まりやすい"),
    ("3月", "新生活・環境変化のフォロー", "生活リズムが変わり離脱が出やすい"),
    ("5月", "母の日ギフト（LINEギフト）", "贈答需要。受け取り側の新規友だち化"),
    ("6月", "梅雨・体調の変化に合わせた使い方", "使用継続のフォロー"),
    ("8月", "夏の使い方・保管方法", "高温期は保管の相談が増える"),
    ("11月", "年末に向けた買い置き提案", "まとめ買い・プラン変更の好機"),
]
rows = [[m, t, note] for m, t, note in cal]
simple_table(s, CX0, CY0 + 0.2, CW, 5.0, ["月", "企画", "参照データ／理由"],
             rows, col_w=[3.0, 10.0, 12.12], align=["c", "l", "l"], row_h=0.8)
box(s, CX0, CY0 + 5.7, CW, 2.6, fill=PORANGE, line=ORANGE, lw=1.3)
T(s, CX0 + 0.6, CY0 + 5.9, CW - 1.2, 2.2,
  [one("★上の企画は、季節性データが入るまでの叩き台", 12, True, ORANGE, sa=4),
   one("Googleトレンド（S10）で山が確定したら、「参照データ」の列に実データを入れて確定させる。"
       "根拠のない企画カレンダーは作らない。", 10.5, None, INK, ls=1.3)],
  anchor="m")
box(s, CX0, CY0 + 8.7, CW, 2.7, fill=PALE)
T(s, CX0 + 0.6, CY0 + 8.9, CW - 1.2, 2.3,
  [one("企画投稿の役割は「売ること」ではない", 12, True, NAVY, sa=4),
   one("配信を止めるとアカウントは忘れられる。忘れられたアカウントからの"
       "スキップ提案や解約受付は開かれない。接点を絶やさないことが目的。",
       10.5, None, INK, ls=1.3)],
  anchor="m")
foot(s, "★S10のGoogleトレンド取得後に確定する。現時点の企画は根拠列が未確定の叩き台")

# ============================================================
# S26 通知メッセージ
# ============================================================
s = slides[25]
frame(s, "配信設計⑥｜通知メッセージ（友だちでなくても届く）",
      ["電話番号ベースで届く。＝ まだ友だちになっていない既存顧客にも到達できる。"])
uses = [
    ("発送完了・配送状況", "問い合わせが一番多い用件。自動化で受電が減る"),
    ("★定期の次回発送前", "スキップ・間隔変更を選べる導線を、友だち以外にも届けられる"),
    ("決済エラー", "カード期限切れ等。放置すると意図しない解約になる"),
    ("再入荷・欠品のお詫び", "待っている人にだけ届く"),
]
for i, (h, b) in enumerate(uses):
    y = CY0 + 0.2 + i * 1.9
    hot = h.startswith("★")
    box(s, CX0, y, CW, 1.7, fill=PORANGE if hot else PALE,
        line=ORANGE if hot else None, lw=1.4)
    T(s, CX0 + 0.4, y + 0.15, 8.4, 1.4,
      [one(h, 12, True, ORANGE if hot else NAVY)], anchor="m")
    T(s, CX0 + 9.2, y + 0.15, CW - 9.6, 1.4, [one(b, 10.5, None, INK, ls=1.3)],
      anchor="m")
box(s, CX0, CY0 + 8.0, CW, 3.4, fill=NAVY)
T(s, CX0 + 0.6, CY0 + 8.2, CW - 1.2, 3.0,
  [one("★決済エラーの放置は「意図しない解約」を生む", 13, True, ORANGE, sa=5),
   one("カードの期限切れで止まった定期は、顧客が辞めたわけではない。"
       "気づかないまま関係が切れる。通知メッセージで拾えば、そのまま継続する。",
       11.5, None, WHITE, ls=1.32)],
  anchor="m")
foot(s, "業界水準｜通知メッセージは電話番号をキーに配信できるため、"
        "未友だちの既存顧客にも到達する。利用には所定の申請と要件確認が必要")

# ============================================================
# S27 月次解約率の改善（★本資料の山）
# ============================================================
s = slides[26]
frame(s, "改善①｜LTVを決めているのは、解約率",
      ["月次解約率が4ポイント違うだけで、顧客の寿命は2倍以上変わる。"])
cw2 = (CW - 0.5) / 2
box(s, CX0, CY0 + 0.2, cw2, 3.6, fill=PALE)
T(s, CX0 + 0.4, CY0 + 0.4, cw2 - 0.8, 0.7,
  [one("月次解約率 3% なら", 12, True, NAVY)], anchor="m")
T(s, CX0 + 0.4, CY0 + 1.15, cw2 - 0.8, 1.8,
  [one("約33ヶ月", 34, True, NAVY)], anchor="m")
T(s, CX0 + 0.4, CY0 + 3.0, cw2 - 0.8, 0.6,
  [one("平均継続期間", 9.5, None, MUT)], anchor="m")
box(s, CX0 + cw2 + 0.5, CY0 + 0.2, cw2, 3.6, fill=PRED, line=RED, lw=1.5)
T(s, CX0 + cw2 + 0.9, CY0 + 0.4, cw2 - 0.8, 0.7,
  [one("月次解約率 7% なら", 12, True, RED)], anchor="m")
T(s, CX0 + cw2 + 0.9, CY0 + 1.15, cw2 - 0.8, 1.8,
  [one("約14ヶ月", 34, True, RED)], anchor="m")
T(s, CX0 + cw2 + 0.9, CY0 + 3.0, cw2 - 0.8, 0.6,
  [one("平均継続期間", 9.5, None, MUT)], anchor="m")
box(s, CX0, CY0 + 4.2, CW, 2.4, fill=PORANGE, line=ORANGE, lw=1.5)
T(s, CX0 + 0.6, CY0 + 4.4, CW - 1.2, 2.0,
  [one("そして、解約率は下げられる。", 14, True, ORANGE, sa=4),
   one("解約アンケートの設置と、リテンションオファー（スキップ／プラン変更／割引）の"
       "自動化だけで、月次チャーン率は 1〜2% 削減できる（業界水準）。",
       11.5, None, INK, ls=1.32)],
  anchor="m")
levers = [
    ("レバー① F2転換率", "30〜40%", "初回離脱50〜60%超を削る。使用フォローとF2オファー（S22）"),
    ("レバー② 月次解約率", "1〜2%削減", "解約の受け皿とスキップ提示（S23・S24）"),
    ("休眠の掘り起こし", "―", "解約後もブロックされなければ再開提案ができる"),
]
for i, (h, v, b) in enumerate(levers):
    y = CY0 + 7.0 + i * 1.5
    box(s, CX0, y, CW, 1.3, fill=PALE)
    T(s, CX0 + 0.4, y + 0.1, 6.2, 1.1, [one(h, 11.5, True, NAVY)], anchor="m")
    T(s, CX0 + 6.8, y + 0.1, 3.4, 1.1, [one(v, 13, True, ORANGE)], anchor="m")
    T(s, CX0 + 10.6, y + 0.1, CW - 11.0, 1.1, [one(b, 10, None, INK, ls=1.25)],
      anchor="m")
foot(s, "業界水準｜継続期間・チャーン削減幅は定期通販CRM各社の公開値。"
        "初回離脱率・F2転換率は美容／健康食品の単品通販の水準")

# ============================================================
# S28 工数削減 ＋ 改善モデル
# ============================================================
s = slides[27]
frame(s, "改善②｜工数削減と、改善が積み上がる仕組み",
      ["解約の電話をLINEに逃がすと、受電が減り、同時に解約も減る。"])
box(s, CX0, CY0 + 0.2, CW * 0.48, 5.4, fill=PALE)
T(s, CX0 + 0.35, CY0 + 0.4, CW * 0.48 - 0.7, 0.7,
  [one("工数が減るところ", 13, True, NAVY)], anchor="m")
red = ["解約・休止の受付（電話 → LINE）",
       "配送状況・お届け日の問い合わせ",
       "お届け間隔・数量の変更受付",
       "使い方・保管方法のよくある質問",
       "決済エラーの個別連絡"]
for j, it in enumerate(red):
    T(s, CX0 + 0.35, CY0 + 1.3 + j * 0.85, CW * 0.48 - 0.7, 0.8,
      [one("・" + it, 10.5, None, INK, ls=1.25)], anchor="m")
box(s, CX0 + CW * 0.48 + 0.3, CY0 + 0.2, CW * 0.52 - 0.3, 5.4, fill=PORANGE,
    line=ORANGE, lw=1.4)
T(s, CX0 + CW * 0.48 + 0.65, CY0 + 0.4, CW * 0.52 - 0.95, 5.0,
  [one("★人件費の削減額では語らない", 12.5, True, ORANGE, sa=5),
   one("金額にすると小さく見えるし、論点もズレる。", 10.5, None, INK, ls=1.3, sa=4),
   one("問題は「ピーク時に手が止まること」。解約の電話が集中する時間帯に"
       "受電が詰まると、繋がらない → 苦情 → 行政の目、という順に悪化する。",
       10.5, None, INK, ls=1.3, sa=4),
   one("LINEなら24時間、同時に何人でも受けられる。", 11, True, INK, ls=1.3)],
  anchor="m")
box(s, CX0, CY0 + 6.3, CW, 5.1, fill=NAVY)
T(s, CX0 + 0.6, CY0 + 6.5, CW - 1.2, 4.7,
  [one("改善モデル｜なぜ運用するほど良くなるのか", 12.5, True, ORANGE, sa=6),
   one("① 解約アンケートの回答が溜まる　→　理由の分布が分かる", 12, None, WHITE, ls=1.32, sa=3),
   one("② 一番多い理由に対して、先に出すオファーを変える", 12, None, WHITE, ls=1.32, sa=3),
   one("③ 効いたオファーを、次回発送前フォロー（S23）に前倒しする",
       12, None, WHITE, ls=1.32, sa=3),
   one("④ 解約に至る前に片づく件数が増える　→　月次解約率が下がる",
       12, True, WHITE, ls=1.32, sa=5),
   one("★ 顧客が辞める理由を、こちらが知っている状態をつくることが本質。",
       12, True, ORANGE, ls=1.32)],
  anchor="m")
foot(s, "意見｜アンケートの回答は、配信改善だけでなく商品開発・同梱物の改善にも使える")

# ============================================================
# S29 効果測定の設計
# ============================================================
s = slides[28]
frame(s, "成果｜効果測定の設計",
      ["握る指標と、その先の指標を分ける。LTVは6ヶ月では動ききらないので約束しない。"])
simple_table(s, CX0, CY0 + 0.2, CW, 3.3,
             ["", "指標", "測れる時期", "扱い"],
             [["握る", "F2転換率", "1〜2ヶ月で動く", "成果指標"],
              ["握る", "月次解約率（チャーン）", "3〜6ヶ月で動く", "成果指標"],
              ["その先", "LTV → 限界CPO", "12ヶ月以上", "構造で示すだけ。約束しない"]],
             col_w=[3.2, 9.0, 6.4, 6.52], align=["c", "l", "l", "l"], row_h=0.78)
box(s, CX0, CY0 + 4.0, CW, 4.2, fill=PALE)
T(s, CX0 + 0.6, CY0 + 4.2, CW - 1.2, 3.8,
  [one("Before（広告のみ）", 12, True, MUT, sa=3),
   one("1回買って終わる／解約は電話でしか受けられない　"
       "→　LTVが伸びない　→　限界CPOが低い　→　広告で勝てる範囲が狭い",
       10.5, None, INK, ls=1.3, sa=6),
   one("After（広告 × LINE）── 広告費も新規獲得数も据え置き", 12, True, NAVY, sa=3),
   one("① 使用フォロー・F2オファー　→　平均購入回数が増える", 10.5, None, INK, ls=1.3, sa=2),
   one("② 解約の手前で受け止める　→　平均継続期間が伸びる", 10.5, None, INK, ls=1.3, sa=2),
   one("③ ①②の結果　→　LTVが上がる　→　限界CPOが上がる　"
       "→　同じ広告枠でも勝てる範囲が広がる", 10.5, True, ORANGE, ls=1.3)],
  anchor="m")
box(s, CX0, CY0 + 8.6, CW, 2.8, fill=NAVY)
T(s, CX0 + 0.6, CY0 + 8.8, CW - 1.2, 2.4,
  [one("本資料では、具体的な数値シミュレーションは載せていません。", 12.5, True, WHITE, sa=4),
   one("御社の実績値（初回獲得数・F2転換率・月次解約率・客単価・利益率）を"
       "いただければ、貴社専用のシミュレーションを作成します。", 11.5, None, WHITE, ls=1.3)],
  anchor="m")
foot(s, "意見｜業界汎用資料のため仮置きの数値を並べていない。"
        "限界CPO ＝ 年間LTV −（広告費を除く年間の総費用 ÷ 顧客総数）")

# ============================================================
# S30 費用プラン
# ============================================================
s = slides[29]
frame(s, "費用プラン（6ヶ月〜／税抜）",
      ["初期費用と月額を分けて記載。継続費は複数箇所に明記している。"])
simple_table(s, CX0, CY0 + 0.2, CW, 4.4,
             ["プラン", "初期費用", "月額", "内容"],
             [["① コンサル基本", "10万円〜", "20万／30万／50万",
               "月3投稿／5投稿／9投稿。設計から運用まで伴走"],
              ["② 初動設計＋運用", "10万円〜", "5万円〜", "構築を中心に、運用は最小限"],
              ["③ 運用代行・効率改善", "10万円〜", "0万円〜", "アカウント費のみ。既存運用の改善"],
              ["④ 成果報酬型", "5万円〜", "単価×成果数＋固定費", "成果地点を合意のうえ設定"]],
             col_w=[6.0, 4.4, 6.2, 8.52], align=["l", "c", "c", "l"], row_h=0.88)
box(s, CX0, CY0 + 5.1, CW * 0.5, 3.4, fill=PALE)
T(s, CX0 + 0.35, CY0 + 5.3, CW * 0.5 - 0.7, 3.0,
  [one("無償で付帯するもの", 12, True, NAVY, sa=4),
   one("アカウント開設／プロフィール／リッチメニュー／あいさつメッセージ／"
       "初期アンケート／キーワード自動応答／ステップ配信／セグメント配信／"
       "タグ管理／GAレポート連携／クリエイティブ／定例会",
       9.5, None, INK, ls=1.3)],
  anchor="m")
box(s, CX0 + CW * 0.5 + 0.3, CY0 + 5.1, CW * 0.5 - 0.3, 3.4, fill=GREY)
T(s, CX0 + CW * 0.5 + 0.65, CY0 + 5.3, CW * 0.5 - 0.95, 3.0,
  [one("別途費用がかかるもの", 12, True, INK, sa=4),
   one("離脱防止ツール（Sitelead）初期10万＋月5万／"
       "API連携／通知メッセージ／配信ツール利用料",
       9.5, None, INK, ls=1.3, sa=3),
   one("※配信ツールは要件に応じて1種を選定します", 9.5, None, MUT, ls=1.3)],
  anchor="m")
box(s, CX0, CY0 + 9.0, CW, 2.4, fill=PORANGE, line=ORANGE, lw=1.3)
T(s, CX0 + 0.6, CY0 + 9.2, CW - 1.2, 2.0,
  [one("継続費のご確認", 12, True, ORANGE, sa=4),
   one("初期費用は初月のみ。2ヶ月目以降は月額のみが継続してかかります。"
       "契約期間は6ヶ月〜。7ヶ月目以降も同額の月額が継続します。",
       11, None, INK, ls=1.3)],
  anchor="m")
foot(s, "すべて税抜。契約期間・成果地点は個別のご要望に応じて調整します")

# ============================================================
# S31 運用スケジュール（6ヶ月）
# ============================================================
s = slides[30]
frame(s, "運用スケジュール（6ヶ月）",
      ["★1〜2ヶ月目のID連携がすべての前提。ここが遅れると全体が後ろ倒しになる。"])
months = [
    ("1ヶ月目", ["キックオフ・要件確認", "★ID連携の可否確認と設計",
                 "アカウント開設・プロフィール"], True),
    ("2ヶ月目", ["★ID連携の実装", "同梱物QR・購入完了画面の導線",
                 "あいさつ／初期アンケート／リッチメニュー"], True),
    ("3ヶ月目", ["シナリオ①②の実装（F2・解約予防）", "通知メッセージの申請・設定",
                 "配信開始"], False),
    ("4ヶ月目", ["★シナリオ③（解約）の実装", "企画投稿の開始",
                 "F2転換率の初回レビュー"], False),
    ("5ヶ月目", ["解約アンケートの回答分析", "オファーの改善",
                 "セグメント配信の追加"], False),
    ("6ヶ月目", ["月次解約率のレビュー", "半期の振り返り",
                 "次期の設計"], False),
]
mw = (CW - 5 * 0.2) / 6
for i, (m, items, hot) in enumerate(months):
    x = CX0 + i * (mw + 0.2)
    box(s, x, CY0 + 0.2, mw, 0.9, fill=ORANGE if hot else NAVY)
    T(s, x, CY0 + 0.2, mw, 0.9, [one(m, 11, True, WHITE, align="c")],
      anchor="m", ml=0.05, mr=0.05)
    box(s, x, CY0 + 1.2, mw, 5.4, fill=PORANGE if hot else PALE)
    for j, it in enumerate(items):
        star = it.startswith("★")
        T(s, x + 0.16, CY0 + 1.4 + j * 1.7, mw - 0.32, 1.6,
          [one(it, 9.5, True if star else None, ORANGE if star else INK, ls=1.28)],
          anchor="t")
box(s, CX0, CY0 + 7.4, CW, 4.0, fill=NAVY)
T(s, CX0 + 0.6, CY0 + 7.6, CW - 1.2, 3.6,
  [one("★解約シナリオ（S24）は4ヶ月目に置いている", 12.5, True, ORANGE, sa=5),
   one("先に「買った直後〜F2」（シナリオ①）と「次回発送前」（シナリオ②）を回して、"
       "解約に至る前に片づく件数を増やしてから、解約シナリオを実装する。",
       11.5, None, WHITE, ls=1.32, sa=4),
   one("順番を逆にすると、解約の受付だけが増える。", 11.5, True, WHITE, ls=1.32)],
  anchor="m")
foot(s, "意見｜ID連携の可否と、カート／定期システムの仕様確認を最優先で行う")

# ============================================================
# S32 サポート体制
# ============================================================
s = slides[31]
frame(s, "サポート体制",
      ["配信を作る人と、法令を見る人を分ける。単品通販はここが要る。"])
roles = [
    ("戦略・設計", "シナリオ設計／KPI設計／定例会。F2転換率と月次解約率の推移を追う"),
    ("配信制作", "文面・クリエイティブ制作／リッチメニュー更新／企画投稿"),
    ("★法令チェック", "薬機法・景表法・特商法の観点で配信文面を確認。"
     "★LINE配信は薬機法上の「広告」に該当するため必須"),
    ("技術・連携", "ID連携／通知メッセージ／API／計測環境の構築と保守"),
]
for i, (h, b) in enumerate(roles):
    y = CY0 + 0.2 + i * 2.4
    hot = h.startswith("★")
    box(s, CX0, y, CW, 2.1, fill=PORANGE if hot else PALE,
        line=ORANGE if hot else None, lw=1.5)
    T(s, CX0 + 0.4, y + 0.15, 6.4, 1.8,
      [one(h, 13, True, ORANGE if hot else NAVY)], anchor="m")
    T(s, CX0 + 7.2, y + 0.15, CW - 7.6, 1.8, [one(b, 11, None, INK, ls=1.3)],
      anchor="m")
box(s, CX0, CY0 + 9.9, CW, 1.5, fill=GREY)
T(s, CX0 + 0.6, CY0 + 9.95, CW - 1.2, 1.4,
  [one("定例会は月1回。F2転換率・月次解約率・解約理由の分布を定点で見ます。",
       11.5, True, INK)], anchor="m")
foot(s, "意見｜法令チェックの体制は単品通販では省略できない。"
        "配信文面が薬機法上の広告に該当するため")

# ============================================================
# S33 飛び道具
# ============================================================
s = slides[32]
frame(s, "締め①｜飛び道具｜解約ボタンを、LINEに置く",
      ["法令を守りながらリテンションできる場所は、いまここしかない。"])
box(s, CX0, CY0 + 0.2, CW, 3.4, fill=ORANGE)
T(s, CX0 + 0.8, CY0 + 0.4, CW - 1.6, 3.0,
  [one("「解約はこちら」を、リッチメニューの1軍に置く。", 20, True, WHITE, ls=1.25, sa=6),
   one("隠さない。むしろ目立たせる。", 14, None, WHITE, ls=1.25)],
  anchor="m")
merits = [
    ("① 法令リスクが下がる", "解約導線を分かりにくくすること自体が、"
     "特商法上の指摘対象になり得る。目立たせるほど安全になる"),
    ("② 苦情が減る", "「電話がつながらない」が定期購入の苦情の中心。"
     "24時間受けられる窓口があれば、そもそも発生しない"),
    ("③ ★リテンションできる", "解約したい人が必ずLINEに来る。"
     "＝ スキップ・周期変更・休止を提示できる場所に、自分から来てくれる"),
    ("④ 理由が溜まる", "解約アンケートが自動で溜まる。"
     "商品開発・同梱物・配信の改善材料になる"),
]
for i, (h, b) in enumerate(merits):
    y = CY0 + 4.0 + i * 1.9
    hot = "★" in h
    box(s, CX0, y, CW, 1.7, fill=PORANGE if hot else PALE,
        line=ORANGE if hot else None, lw=1.4)
    T(s, CX0 + 0.4, y + 0.12, 6.6, 1.45,
      [one(h, 12, True, ORANGE if hot else NAVY)], anchor="m")
    T(s, CX0 + 7.4, y + 0.12, CW - 7.8, 1.45, [one(b, 10.5, None, INK, ls=1.3)],
      anchor="m")
foot(s, "意見｜「解約を隠す」運用は、法令・苦情・改善のすべてで不利になる。"
        "本提案は解約を止めるものではない")

# ============================================================
# S34 第2の提案軸｜通知メッセージ
# ============================================================
s = slides[33]
frame(s, "締め②｜第2の提案軸｜通知メッセージで既存顧客を全部拾う",
      ["友だちでなくても、電話番号で届く。＝ いま繋がっていない既存顧客に到達できる。"])
box(s, CX0, CY0 + 0.2, CW, 3.0, fill=NAVY)
T(s, CX0 + 0.8, CY0 + 0.4, CW - 1.6, 2.6,
  [one("既存顧客のうち、LINEの友だちになっているのは一部にすぎない。",
       15, True, WHITE, ls=1.28, sa=5),
   one("通知メッセージは電話番号をキーに配信できるため、"
       "友だち未追加の既存顧客にも届く。", 12, None, WHITE, ls=1.3)],
  anchor="m")
cw2 = (CW - 0.4) / 2
axis = [
    ("いま到達できていない層に届く",
     ["購入したが友だち追加していない人", "解約したが連絡先は残っている人",
      "決済エラーで止まっている人"]),
    ("そこから友だち化につなげる",
     ["通知の中に友だち追加の導線を置く", "一度届けば、以降は通常配信で追える",
      "同梱物QRと合わせて到達率を底上げする"]),
]
for i, (h, items) in enumerate(axis):
    x = CX0 + i * (cw2 + 0.4)
    box(s, x, CY0 + 3.7, cw2, 5.2, fill=PALE)
    T(s, x + 0.35, CY0 + 3.9, cw2 - 0.7, 0.8, [one(h, 13, True, NAVY, ls=1.2)],
      anchor="m")
    for j, it in enumerate(items):
        T(s, x + 0.35, CY0 + 4.9 + j * 1.3, cw2 - 0.7, 1.2,
          [one("・" + it, 10.5, None, INK, ls=1.3)], anchor="t")
box(s, CX0, CY0 + 9.3, CW, 2.1, fill=PORANGE, line=ORANGE, lw=1.3)
T(s, CX0 + 0.6, CY0 + 9.45, CW - 1.2, 1.8,
  [one("本提案（S16-26）のID連携基盤の上に、そのまま乗せられます。",
       12, True, ORANGE, ls=1.3)],
  anchor="m")
foot(s, "業界水準｜通知メッセージの利用には所定の申請と要件確認が必要。"
        "配信できる内容にも制限がある")

# ============================================================
# S35 LINEOA実績
# ============================================================
s = slides[34]
frame(s, "締め③｜LINEOA実績",
      ["LINEヤフー公式データと、D2C・単品通販の導入事例を並べる。"])
cw2 = (CW - 0.4) / 2
nums = [("LINE国内MAU", "1億人突破", "2025年12月末時点"),
        ("メッセージ開封率（当日中）", "約8割", "受信直後 約2割／3〜6時間で約5割")]
for i, (label, val, sub) in enumerate(nums):
    stat(s, CX0 + i * (cw2 + 0.4), CY0 + 0.2, cw2, label, val, sub, h=2.6, vsz=24)
placeholder(s, CX0, CY0 + 3.2, CW, 4.4,
            "★差込枠｜LINEヤフー公式 D2C・単品通販の導入事例（未取得）",
            "取得先：lycbiz.com（LINE公式アカウントの導入事例）\n"
            "掲載前に各事例ページの原典で数値を再確認すること。"
            "公式事例が存在しない場合は「該当なし」と正直に報告する")
box(s, CX0, CY0 + 8.1, CW, 3.3, fill=GREY)
T(s, CX0 + 0.6, CY0 + 8.3, CW - 1.2, 2.9,
  [one("数値の出所について", 11.5, True, INK, sa=4),
   one("本資料では、LINEヤフー公式値と業界水準を明確に分けて記載しています。"
       "「開封率はメルマガの◯倍」といった倍率表現は公式値ではないため使用していません。",
       10.5, None, INK, ls=1.3)],
  anchor="m")
foot(s, "LY公式｜LINE国内MAU・開封率はLINEヤフー公式値。★導入事例は未取得")


# ============================================================
# 自動チェック①：業種残骸の正規表現スキャン
# ============================================================
RESIDUE_PATTERNS = [
    r"○○業界", r"〇〇業界", r"クロスセル", r"アップセル", r"転職", r"求人",
    r"賃貸", r"入居", r"オーナー", r"注文住宅", r"来店予約", r"カゴ落ち",
    r"利用状況タグ", r"誕生日.{0,3}記念日タグ",
]
residue_hits = []
for i, sl in enumerate(slides, 1):
    for sh in sl.shapes:
        if sh.has_text_frame:
            txt = sh.text_frame.text
            for pat in RESIDUE_PATTERNS:
                if re.search(pat, txt):
                    residue_hits.append((i, pat, txt[:40]))
if residue_hits:
    print("★業種残骸が残っています：")
    for i, pat, txt in residue_hits:
        print(f"  slide {i}: /{pat}/ -> {txt!r}")
else:
    print("業種残骸チェック：残存0")

# ============================================================
# 自動チェック②：★GUARD｜「解約を止める」表現が入っていないか
#   特商法で契約解除の妨害には罰則がある。引き止め表現は差し戻し。
# ============================================================
BANNED = [r"解約を止め", r"引き止め", r"引き留め", r"解約を防[ぐぎ]",
          r"解約させな", r"解約を思いとどま"]
# 否定・打消しの文脈で出るのは正しい用法（「引き止めるのではなく」等）。
# 許可リストを育てるのではなく、直後に否定語が来るかで機械的に判定する。
NEGATORS = [r"ではな", r"ではあり", r"のではなく", r"のではありません", r"ません",
            r"ない", r"とは", r"できな", r"リスク", r"違反", r"禁じ", r"禁止"]
banned_hits = []
for i, sl in enumerate(slides, 1):
    for sh in sl.shapes:
        if sh.has_text_frame:
            txt = sh.text_frame.text
            for pat in BANNED:
                for m in re.finditer(pat, txt):
                    after = txt[m.end(): m.end() + 24]
                    if any(re.search(n, after) for n in NEGATORS):
                        continue
                    seg = txt[max(0, m.start() - 30): m.end() + 30]
                    banned_hits.append((i, pat, seg.replace("\n", " ")[:60]))
if banned_hits:
    print("★禁止表現（引き止め）が入っています：")
    for i, pat, seg in banned_hits:
        print(f"  slide {i}: /{pat}/ -> {seg!r}")
else:
    print("引き止め表現チェック：残存0")

# ============================================================
# 自動チェック③：S24に「止めない」が明記されているか
# ============================================================
s24 = "\n".join(sh.text_frame.text for sh in slides[23].shapes if sh.has_text_frame)
if "止めない" in s24 and "引き止め施策ではありません" in s24:
    print("S24「止めない」明記チェック：OK")
else:
    print("★S24に「止めない」の明記がありません。差し戻し対象です。")

# ============================================================
# 自動チェック④：LY非公式の倍率表現が入っていないか
# ============================================================
#   ※「1.7倍」（CPC）や「2倍以上変わる」（継続期間）は別の話なので、
#     LINE・メルマガ・開封の語が近くにある倍率だけを見る。
RATIO = [r"メルマガの[0-9０-９]", r"[0-9０-９.０-９]+倍", r"60〜80%", r"60～80%"]
LINE_WORDS = ["LINE", "メルマガ", "開封", "クリック率"]
ratio_hits = []
for i, sl in enumerate(slides, 1):
    for sh in sl.shapes:
        if sh.has_text_frame:
            txt = sh.text_frame.text
            for pat in RATIO:
                for m in re.finditer(pat, txt):
                    seg = txt[max(0, m.start() - 40): m.end() + 40]
                    if not any(w in seg for w in LINE_WORDS):
                        continue
                    if "倍率表現は" in seg or "倍率ではなく" in seg:
                        continue
                    ratio_hits.append((i, pat, seg.replace("\n", " ")[:60]))
if ratio_hits:
    print("★LY非公式の倍率表現の疑いがあります：")
    for i, pat, seg in ratio_hits:
        print(f"  slide {i}: /{pat}/ -> {seg!r}")
else:
    print("倍率表現チェック：残存0")

prs.save(OUT)
print("saved:", OUT)
