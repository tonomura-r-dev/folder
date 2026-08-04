# -*- coding: utf-8 -*-
"""ビイサイド(vivical) LINEOA提案 33枚
クラフトマンv02(44枚)をコピー→11枚削除＋並べ替え→中盤をvivical用に再構築"""
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
SRC = str(ROOT / "_templates" / "DYM_LINEOA_33p_base.pptx")
OUT = str(ROOT / "20260804_株式会社ビイサイドプランニング御中_公式LINEのご提案.pptx")

TNAVY = "002060"   # タイトル
NAVY = "1F285A"    # 打ち手・カード見出し
RED = "C00000"     # 課題
INK = "333333"     # 本文
MUT = "808080"
WHITE = "FFFFFF"
PALE = "F4F7FF"
GREY = "F2F2F2"
BORDER = "D9D9D9"

# 新順序（クラフトマンの1-based番号）
ORDER = [1, 2, 3, 4, 5, 6, 7, 8, 9, 10, 11, 12,
         13, 14, 15, 16, 17, 18, 29, 20, 31, 21, 19, 37, 22,
         23, 24, 28, 34, 33, 27, 35, 44]

shutil.copyfile(SRC, OUT)
prs = Presentation(OUT)

# ---------- 構造作業（削除・並べ替え）を先に完了させる ----------
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
assert len(ids) == 44, len(ids)
keep_rids = [ids[c - 1].rId for c in ORDER]
keep_els = [ids[c - 1] for c in ORDER]
for el in ids:
    if el.rId not in keep_rids:
        prs.part.drop_rel(el.rId)
for el in ids:
    sldIdLst.remove(el)
for el in keep_els:
    sldIdLst.append(el)
slides = list(prs.slides)
assert len(slides) == 33, len(slides)

# ---------- helpers ----------
def set_font(run, size, bold=None, color=INK, name="メイリオ"):
    f = run.font
    f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", name)
    f.color.rgb = RGBColor.from_string(color)

ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
ANCH = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}

def reset_tf(tf):
    """既存の段落・runを全消去（追記事故の防止）"""
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)

def put_text(tf, paras, anchor="t", ml=0.08, mr=0.08, mt=0.04, mb=0.04, wrap=True):
    reset_tf(tf)
    tf.word_wrap = wrap
    tf.margin_left = Inches(ml); tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt); tf.margin_bottom = Inches(mb)
    tf.vertical_anchor = ANCH[anchor]
    first = True
    for p in paras:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = ALIGN[p.get("align", "l")]
        if p.get("sa") is not None:
            para.space_after = Pt(p["sa"])
        if p.get("ls") is not None:
            para.line_spacing = p["ls"]
        for (t, sz, b, c) in p["runs"]:
            r = para.add_run(); r.text = t
            set_font(r, sz, b, c)
    return tf

def add_text(slide, x, y, w, h, paras, anchor="t", **kw):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    put_text(box.text_frame, paras, anchor=anchor, **kw)
    return box

def add_box(slide, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor.from_string(line); sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius if radius is not None else 0.08
        except Exception:
            pass
    return sp

def card(slide, x, y, w, h, head, body, hcol=NAVY, fill=PALE, hsz=15, bsz=12, line=None):
    sp = add_box(slide, x, y, w, h, fill=fill, line=line, lw=1.0)
    paras = [{"runs": [(head, hsz, True, hcol)], "sa": 4}]
    if body:
        paras.append({"runs": [(body, bsz, None, INK)], "ls": 1.15})
    put_text(sp.text_frame, paras, anchor="m", ml=0.16, mr=0.14, mt=0.08, mb=0.08)
    return sp

def clear_slide(slide):
    spTree = slide.shapes._spTree
    for el in list(spTree):
        if el.tag.split("}")[-1] in ("sp", "cxnSp", "pic", "graphicFrame", "grpSp"):
            spTree.remove(el)

def frame(slide, title, lead):
    """川本スタイル：タイトル20pt紺 ＋ リード14pt太字"""
    clear_slide(slide)
    add_text(slide, 0.60, 0.13, 9.60, 0.45,
             [{"runs": [(title, 20, True, TNAVY)]}], anchor="m", ml=0.0, mr=0.0)
    add_text(slide, 0.60, 0.70, 9.60, 0.58,
             [{"runs": [(lead, 14, True, INK)], "ls": 1.1}], anchor="m", ml=0.0, mr=0.0)

def replace_first_text(shape, text):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    if not p.runs:
        return
    p.runs[0].text = text
    for r in p.runs[1:]:
        r._r.getparent().remove(r._r)
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)

def find_shape_by_text(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle in sh.text_frame.text:
            return sh
    return None

# ================= Slide 1: 表紙 =================
s = slides[0]
sh = find_shape_by_text(s, "クラフトマン")
if sh:
    replace_first_text(sh, "株式会社ビイサイドプランニング 御中")

# ================= Slide 5 / 12 / 26: アジェンダ =================
s = slides[4]
sh = find_shape_by_text(s, "LINE公式アカウント運用のメリット")
if sh:
    replace_first_text(sh, "LINE公式アカウント運用のメリット")
s = slides[25]
sh = find_shape_by_text(s, "プラン・シミュレーション")
if sh:
    replace_first_text(sh, "プラン・シミュレーション／サポート体制")

# ================= Slide 6: 業種文言の差し替え（EC向け残骸→フェア向け） =================
s = slides[5]
for sh in s.shapes:
    if not sh.has_text_frame:
        continue
    txt = sh.text_frame.text
    if "定期購入" in txt:
        put_text(sh.text_frame,
                 [{"runs": [("次回フェア・個別求人", 9, True, WHITE)], "align": "c"},
                  {"runs": [("への動線設計", 9, True, WHITE)], "align": "c"}],
                 anchor="m", ml=0.02, mr=0.02, mt=0.0, mb=0.0)
    elif "ヒアリングの上" in txt:
        put_text(sh.text_frame,
                 [{"runs": [("Meta広告のみ", 12, True, INK)], "align": "c", "sa": 4},
                  {"runs": [("予約後の接点なし", 12, True, INK)], "align": "c"}],
                 anchor="m", ml=0.06, mr=0.06, mt=0.0, mb=0.0)

# ================= Slide 13: 現状分析 =================
s = slides[12]
frame(s, "現状分析｜求人vivical様の現状とボトルネック",
      "目標予約数は達成している一方、予約者の3〜4割が直前辞退し、実質来場単価が16,000〜18,000円に高騰しています。")
add_text(s, 0.60, 1.42, 4.60, 0.32, [{"runs": [("現状の実績（Meta広告）", 13, True, NAVY)]}], ml=0.0)
facts = [("月額広告予算", "80万円"), ("事前予約単価（CPA）", "10,000円"), ("月間事前予約数", "80名")]
for i, (lab, val) in enumerate(facts):
    y = 1.82 + i * 0.86
    sp = add_box(s, 0.60, y, 4.60, 0.74, fill=PALE)
    put_text(sp.text_frame,
             [{"runs": [(lab + "　", 12, None, INK), (val, 17, True, NAVY)]}],
             anchor="m", ml=0.20, mr=0.10, mt=0.0, mb=0.0)
add_text(s, 0.60, 4.46, 4.60, 0.40,
         [{"runs": [("目標予約数は達成（各会場合計）。獲得面の課題はない。", 11.5, None, INK)]}], ml=0.0)
add_text(s, 5.55, 1.42, 4.65, 0.32, [{"runs": [("ボトルネック（来場フェーズ）", 13, True, RED)]}], ml=0.0)
card(s, 5.55, 1.82, 4.65, 1.28,
     "課題① 予約後の直前辞退",
     "予約者の3〜4割が当日辞退。予約→来場の引き上げ率は60〜70%にとどまる。",
     hcol=RED, fill=WHITE, line=RED, hsz=14, bsz=12)
card(s, 5.55, 3.22, 4.65, 1.28,
     "課題② 実質来場単価の高騰",
     "予約単価は10,000円だが、来場ベースでは16,000〜18,000円まで上昇。",
     hcol=RED, fill=WHITE, line=RED, hsz=14, bsz=12)
card(s, 5.55, 4.62, 4.65, 1.28,
     "課題③ 未応募層が資産化されない",
     "来場後・未応募の求職者と接点が切れ、次回フェアへ再案内できていない。",
     hcol=RED, fill=WHITE, line=RED, hsz=14, bsz=12)
bar = add_box(s, 0.60, 6.10, 9.60, 0.52, fill=NAVY)
put_text(bar.text_frame,
         [{"runs": [("課題は「予約の獲得」ではなく「予約後の引き上げ」にある", 14, True, WHITE)], "align": "c"}],
         anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
add_text(s, 0.60, 6.70, 9.60, 0.30,
         [{"runs": [("参照元：貴社ご提供の実績数値（2026年6〜7月開催／3会場）", 10, None, MUT)]}], ml=0.0)

# ================= Slide 14: 解決の方向性 =================
s = slides[13]
frame(s, "解決の方向性",
      "「予約後の接点」をLINEで設計し、引き上げ率60〜70%→80%、実質来場単価16,667円→10,000円を目指します。")
rows = [
    ("課題①", "直前辞退", "① Meta広告×LINE連携（オフラインCV最適化）",
     "来場・未来場データをMeta広告へ還元し、来場につながる層へ配信を最適化。"),
    ("課題②", "不安・迷い", "② 直前リマインド・安心コンテンツ配信",
     "前日・当日朝のリマインドと「服装自由・履歴書不要」等のQ&Aで参加辞退を抑止。"),
    ("課題③", "未応募層の離脱", "③ 属性セグメント配信による再アプローチ",
     "アンケート取得した希望職種・年代をもとに、次回フェア・個別求人へ再案内。"),
]
for i, (num, issue, act, desc) in enumerate(rows):
    y = 1.50 + i * 1.32
    sp = add_box(s, 0.60, y, 2.60, 1.15, fill=WHITE, line=RED, lw=1.0)
    put_text(sp.text_frame,
             [{"runs": [(num, 14, True, RED)], "align": "c", "sa": 2},
              {"runs": [(issue, 14, True, RED)], "align": "c"}],
             anchor="m", ml=0.05, mr=0.05, mt=0.0, mb=0.0)
    add_box(s, 3.32, y + 0.36, 0.52, 0.44, fill=NAVY, shape=MSO_SHAPE.RIGHT_ARROW)
    card(s, 3.95, y, 6.25, 1.15, act, desc, hcol=NAVY, fill=PALE, hsz=15, bsz=12)
tb = add_box(s, 0.60, 5.52, 9.60, 1.10, fill=WHITE, line=NAVY, lw=1.5)
put_text(tb.text_frame,
         [{"runs": [("目標効果", 12, True, NAVY)], "align": "c", "sa": 4},
          {"runs": [("予約→来場の引き上げ率　", 12, None, INK), ("60〜70%", 14, True, INK),
                    (" → ", 12, None, MUT), ("80%", 20, True, NAVY),
                    ("　／　実質来場単価　", 12, None, INK), ("16,667円", 14, True, INK),
                    (" → ", 12, None, MUT), ("10,000円", 20, True, NAVY)], "align": "c"}],
         anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
add_text(s, 0.60, 6.70, 9.60, 0.30,
         [{"runs": [("※提案後の数値は本施策の実行を前提とした目標値です（成果を保証するものではありません）", 10, None, MUT)]}], ml=0.0)

# ================= Slide 15: メッセージ設計 Before/After =================
s = slides[14]
frame(s, "施策① メッセージ設計 Before / After",
      "「予約して終わり」から「予約後に接点を持ち続ける」へ。同じ予約数でも、来場率が変わります。")
bf = add_box(s, 0.60, 1.42, 4.60, 2.05, fill=WHITE, line=BORDER, lw=1.25)
put_text(bf.text_frame,
         [{"runs": [("Before（現状）", 14, True, RED)], "sa": 6},
          {"runs": [("・予約完了後はメールでの一方的な案内のみ", 12, None, INK)], "sa": 4, "ls": 1.15},
          {"runs": [("・当日までの数週間、接点が実質ゼロ", 12, None, INK)], "sa": 4, "ls": 1.15},
          {"runs": [("・服装や持ち物への不安が解消されないまま当日を迎える", 12, None, INK)], "ls": 1.15}],
         anchor="t", ml=0.20, mr=0.16, mt=0.14)
add_box(s, 5.30, 2.18, 0.62, 0.52, fill=NAVY, shape=MSO_SHAPE.RIGHT_ARROW)
af = add_box(s, 6.00, 1.42, 4.20, 2.05, fill=PALE, line=NAVY, lw=1.5)
put_text(af.text_frame,
         [{"runs": [("After（改善後）", 14, True, NAVY)], "sa": 6},
          {"runs": [("・友だち追加直後にアンケート（希望職種・年代）", 12, None, INK)], "sa": 4, "ls": 1.15},
          {"runs": [("・「服装自由・履歴書不要」等の安心Q&Aを自動送付", 12, None, INK)], "sa": 4, "ls": 1.15},
          {"runs": [("・前日・当日朝にリマインドで来場を後押し", 12, None, INK)], "ls": 1.15}],
         anchor="t", ml=0.20, mr=0.16, mt=0.14)
add_text(s, 0.60, 3.66, 9.60, 0.34,
         [{"runs": [("設計の考え方｜「情報を送る」のではなく「不安を消す」", 13, True, NAVY)]}], ml=0.0)
cols = [
    ("① 価値を先に渡す", "登録直後に診断・特典を提供し、\n「登録してよかった」を作る"),
    ("② 不安を先回りで消す", "服装・持ち物・回り方など、\n当日を具体的に想像できる状態に"),
    ("③ 決断を後押しする", "前日・当日朝のリマインドで\n「行く」を最後に確定させる"),
]
for i, (h, b) in enumerate(cols):
    x = 0.60 + i * 3.28
    card(s, x, 4.06, 3.05, 1.60, h, b.replace("\n", ""), hcol=NAVY, fill=PALE, hsz=13, bsz=12)
bar = add_box(s, 0.60, 5.90, 9.60, 0.52, fill=NAVY)
put_text(bar.text_frame,
         [{"runs": [("予約後の空白期間を埋めることが、来場率を左右する", 14, True, WHITE)], "align": "c"}],
         anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)

# ================= Slide 16: ツナグバ事例（タイトルのみ調整） =================
s = slides[15]
sh = find_shape_by_text(s, "参考：人材業界の運用例")
if sh:
    replace_first_text(sh, "参考：人材・イベント業界の運用事例（ツナグバ）")

# ================= Slide 18: 属性取得アンケート =================
s = slides[17]
frame(s, "具体施策① 友だち追加直後の属性取得アンケート",
      "「30秒でわかる希望条件診断」で回答ハードルを下げ、取得した属性を来場後の追客まで活用します。")
add_text(s, 0.60, 1.40, 4.70, 0.32, [{"runs": [("アンケート設問（案）", 13, True, NAVY)]}], ml=0.0)
qs = [
    ("Q1", "希望職種は？", "事務／製造・軽作業／販売・サービス／営業／その他"),
    ("Q2", "就職・転職の希望時期は？", "すぐにでも／3ヶ月以内／半年以内／情報収集中"),
    ("Q3", "重視する条件は？", "給与／勤務地／休日／未経験歓迎"),
]
for i, (no, q, opt) in enumerate(qs):
    y = 1.78 + i * 1.16
    sp = add_box(s, 0.60, y, 4.70, 1.00, fill=PALE)
    put_text(sp.text_frame,
             [{"runs": [(no + "　", 12, True, NAVY), (q, 13, True, INK)], "sa": 4},
              {"runs": [(opt, 11, None, INK)], "ls": 1.15}],
             anchor="m", ml=0.18, mr=0.12, mt=0.06, mb=0.06)
add_text(s, 0.60, 5.30, 4.70, 0.34,
         [{"runs": [("いずれもボタン選択式。所要30秒で完了します。", 11.5, None, INK)]}], ml=0.0)
add_text(s, 5.60, 1.40, 4.60, 0.32, [{"runs": [("取得した属性の活用", 13, True, NAVY)]}], ml=0.0)
uses = [
    ("① タグ付けによる自動仕分け", "回答内容をタグとして自動付与し、希望職種・時期別にリスト化。"),
    ("② セグメント配信", "「事務希望×すぐにでも」など、条件に合う出展企業だけを案内。"),
    ("③ イベント後の追客", "未応募の求職者へ、次回フェア・個別求人をタグ別に再案内。"),
]
for i, (h, b) in enumerate(uses):
    y = 1.78 + i * 1.16
    card(s, 5.60, y, 4.60, 1.00, h, b, hcol=NAVY, fill=WHITE, hsz=13, bsz=11.5, line=NAVY)
bar = add_box(s, 0.60, 5.86, 9.60, 0.52, fill=NAVY)
put_text(bar.text_frame,
         [{"runs": [("アンケートは「聞く」ためではなく、配信を当てるための設計", 14, True, WHITE)], "align": "c"}],
         anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)

# ================= Slide 19: リッチメニュー設計 =================
s = slides[18]
frame(s, "具体施策② リッチメニュー設計（曜日・時間帯別の出し分け）",
      "求める情報へ1タップで遷移できる6エリア構成とし、週末は「決断を促す」構成へ切り替えます。")
add_box(s, 0.60, 1.45, 4.55, 3.55, fill=WHITE, line=NAVY, lw=1.75, shape=MSO_SHAPE.RECTANGLE)
cells_top = [["出展企業一覧"], ["会場MAP", "アクセス"], ["開催日程"]]
cells_btm = [["事前予約", "特典受取"], ["キャリア相談", "（1on1）"], ["よくある質問"]]
gap = 0.05
cw = (4.55 - gap * 4) / 3
ch = (3.55 - gap * 3) / 2
for r, row in enumerate([cells_top, cells_btm]):
    for c, label in enumerate(row):
        x = 0.60 + gap + c * (cw + gap)
        y = 1.45 + gap + r * (ch + gap)
        fill = PALE if r == 0 else NAVY
        tcol = NAVY if r == 0 else WHITE
        cell = add_box(s, x, y, cw, ch, fill=fill, shape=MSO_SHAPE.RECTANGLE)
        put_text(cell.text_frame,
                 [{"runs": [(ln, 11, True, tcol)], "align": "c", "ls": 1.15} for ln in label],
                 anchor="m", ml=0.03, mr=0.03, mt=0.0, mb=0.0)
add_text(s, 0.60, 5.06, 4.55, 0.30,
         [{"runs": [("上段＝情報提供／下段＝アクション誘導", 11, None, MUT), ("", 11, None, MUT)], "align": "c"}], ml=0.0, mr=0.0)
modes = [
    ("通常時", "情報収集フェーズ。出展企業・会場MAP・開催日程を上段に配置し、比較検討を促す。"),
    ("週末（土曜夜・日曜夜）", "決断フェーズ。「事前予約」「特典受取」を最上段に入れ替え、申込への最短動線に切り替える。"),
]
for i, (h, b) in enumerate(modes):
    y = 1.45 + i * 1.30
    card(s, 5.45, y, 4.75, 1.15, h, b, hcol=NAVY, fill=PALE, hsz=14, bsz=12)
card(s, 5.45, 4.05, 4.75, 0.95,
     "フェア直前（3日前〜当日）",
     "「会場MAP・アクセス」を最上段へ。当日の迷いを解消し、来場を後押しする。",
     hcol=NAVY, fill=PALE, hsz=14, bsz=12)
bar = add_box(s, 0.60, 5.52, 9.60, 0.52, fill=NAVY)
put_text(bar.text_frame,
         [{"runs": [("メニューは固定せず、求職者の検討フェーズに合わせて入れ替える", 14, True, WHITE)], "align": "c"}],
         anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
add_text(s, 0.60, 6.14, 9.60, 0.30,
         [{"runs": [("※メニュー項目・文言は貴社と協議のうえ確定します", 10, None, MUT)]}], ml=0.0)

# ================= Slide 20: 曜日別コンテンツ配信 =================
s = slides[19]
frame(s, "具体施策③ 曜日別コンテンツ配信方針（転職熱量のピーク活用）",
      "転職意欲のピークは日曜夜。週末2段構えで「種まき」→「決断の後押し」を設計します。")
days = [
    ("土曜 夜", "種まき｜じっくり見る", "注目企業TOP5",
     "腰を据えて情報収集する時間帯。出展企業の魅力を比較形式で提示し、興味を醸成する。", PALE, NAVY),
    ("日曜 夜", "後押し｜決断させる", "背中を押す配信",
     "「明日からまた仕事…」の心理が最も高まるタイミング。事前予約への最短導線を提示する。", NAVY, WHITE),
]
for i, (d, role, head, body, bg, tc) in enumerate(days):
    x = 0.60 + i * 4.95
    hd = add_box(s, x, 1.45, 4.65, 0.52, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    put_text(hd.text_frame,
             [{"runs": [(d + "　｜　" + role, 14, True, WHITE)], "align": "c"}],
             anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
    bd = add_box(s, x, 2.05, 4.65, 2.30, fill=bg, line=NAVY if bg == PALE else None, lw=1.0)
    put_text(bd.text_frame,
             [{"runs": [(head, 17, True, tc if bg == NAVY else NAVY)], "sa": 8},
              {"runs": [(body, 12, None, tc if bg == NAVY else INK)], "ls": 1.2}],
             anchor="m", ml=0.20, mr=0.16, mt=0.10, mb=0.10)
add_text(s, 0.60, 4.58, 9.60, 0.34,
         [{"runs": [("配信を「しない」日の設計", 13, True, NAVY)]}], ml=0.0)
nos = [
    ("月曜朝は配信しない", "意欲は高いが、出勤直後で行動に移せない。開封されても流れて終わるため配信対象外とする。"),
    ("平日日中は最小限", "勤務中で反応が取れない。お役立ちコラム等の軽い接点にとどめる。"),
]
for i, (h, b) in enumerate(nos):
    x = 0.60 + i * 4.95
    card(s, x, 4.96, 4.65, 1.10, h, b, hcol=NAVY, fill=WHITE, hsz=13, bsz=11.5, line=BORDER)
bar = add_box(s, 0.60, 6.22, 9.60, 0.48, fill=NAVY)
put_text(bar.text_frame,
         [{"runs": [("「いつ送るか」で反応は変わる。週末2段設計で来場予約を積み上げる", 13.5, True, WHITE)], "align": "c"}],
         anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)

# ================= Slide 21: 時間帯別コンテンツ配信 =================
s = slides[20]
frame(s, "具体施策④ 時間帯別コンテンツ配信方針（朝・昼・夜）",
      "1日の中でも求職者の状態は変わります。時間帯ごとに「読める内容」を出し分けます。")
times = [
    ("朝　7時台", "通勤前のスキマ時間", "新着参加企業・会場アクセス情報",
     "短時間で読める情報を配信。「今日も仕事か」という気持ちが動くタイミングで接点を持つ。"),
    ("昼　12時台", "休憩中のリラックス時間", "お役立ちコラム（面接対策・履歴書の書き方）",
     "すぐの行動は求めず、有益情報の提供でアカウントの価値を高め、ブロックを防ぐ。"),
    ("夜　19〜21時", "帰宅後の検討時間", "イベント事前予約・特典受け取りオファー",
     "最も反応が取れる時間帯。予約フォームへの誘導と限定特典で、その場での申込を促す。"),
]
for i, (t, state, head, body) in enumerate(times):
    x = 0.60 + i * 3.28
    hd = add_box(s, x, 1.45, 3.05, 0.72, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    put_text(hd.text_frame,
             [{"runs": [(t, 15, True, WHITE)], "align": "c", "sa": 2},
              {"runs": [(state, 10.5, None, WHITE)], "align": "c"}],
             anchor="m", ml=0.05, mr=0.05, mt=0.0, mb=0.0)
    bd = add_box(s, x, 2.25, 3.05, 2.85, fill=PALE)
    put_text(bd.text_frame,
             [{"runs": [(head, 13, True, NAVY)], "sa": 8, "ls": 1.15},
              {"runs": [(body, 11.5, None, INK)], "ls": 1.2}],
             anchor="t", ml=0.18, mr=0.14, mt=0.16)
add_text(s, 0.60, 5.30, 9.60, 0.34,
         [{"runs": [("配信量の目安", 13, True, NAVY)]}], ml=0.0)
add_text(s, 0.60, 5.66, 9.60, 0.56,
         [{"runs": [("通常週は週2〜3通、フェア直前週は最大週5通まで。開封率・ブロック率を見ながら通数を調整し、「送りすぎ」による離脱を防ぎます。", 12, None, INK)], "ls": 1.2}], ml=0.0)
bar = add_box(s, 0.60, 6.28, 9.60, 0.48, fill=NAVY)
put_text(bar.text_frame,
         [{"runs": [("夜19〜21時が予約獲得の主戦場。朝・昼は関係構築に使う", 13.5, True, WHITE)], "align": "c"}],
         anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)

# ================= Slide 22: 季節・開催サイクル別 =================
s = slides[21]
frame(s, "具体施策⑤ 季節・イベント開催サイクル別のコンテンツ方針",
      "フェア開催サイクルに合わせ、「直前の追い込み」と「開催後の資産化」を年間で繰り返します。")
phases = [
    ("① 告知期", "開催3〜4週間前", "受付開始のお知らせ・出展企業の先行公開で、予約の母数をつくる。"),
    ("② 追い込み期", "開催1週間前〜前日", "リマインドと安心Q&Aを集中配信。予約者の来場を確実にする。"),
    ("③ 当日", "開催当日 朝〜", "会場MAP付きの最終通知。当日の迷いをなくし、来場を最大化する。"),
    ("④ 資産化期", "開催後〜次回告知まで", "お礼と個別相談の案内。未応募層を次回フェアへ引き継ぐ。"),
]
for i, (ph, when, body) in enumerate(phases):
    x = 0.60 + i * 2.46
    hd = add_box(s, x, 1.45, 2.28, 0.68, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    put_text(hd.text_frame,
             [{"runs": [(ph, 13.5, True, WHITE)], "align": "c", "sa": 2},
              {"runs": [(when, 10, None, WHITE)], "align": "c"}],
             anchor="m", ml=0.04, mr=0.04, mt=0.0, mb=0.0)
    if i < 3:
        add_box(s, x + 2.30, 1.64, 0.14, 0.30, fill=MUT, shape=MSO_SHAPE.RIGHT_ARROW)
    bd = add_box(s, x, 2.21, 2.28, 1.70, fill=PALE)
    put_text(bd.text_frame, [{"runs": [(body, 11.5, None, INK)], "ls": 1.2}],
             anchor="t", ml=0.16, mr=0.12, mt=0.14)
add_text(s, 0.60, 4.10, 9.60, 0.34,
         [{"runs": [("2026年 秋クールの開催スケジュール", 13, True, NAVY)]}], ml=0.0)
sched = [("8/28", "受付開始"), ("9/26", "東近江"), ("10/3", "草津"), ("10/31", "水口"), ("11/8", "彦根")]
for i, (d, place) in enumerate(sched):
    x = 0.60 + i * 1.95
    sp = add_box(s, x, 4.50, 1.80, 0.78, fill=WHITE, line=NAVY, lw=1.0)
    put_text(sp.text_frame,
             [{"runs": [(d, 14, True, NAVY)], "align": "c", "sa": 2},
              {"runs": [(place, 11.5, None, INK)], "align": "c"}],
             anchor="m", ml=0.04, mr=0.04, mt=0.0, mb=0.0)
add_text(s, 0.60, 5.42, 9.60, 0.56,
         [{"runs": [("4会場を約6週間で回るサイクルのため、「告知→追い込み→当日→資産化」を会場ごとに重ねて運用します。前会場の来場データを次会場の配信・広告最適化に反映できる点が、複数会場開催の強みです。", 12, None, INK)], "ls": 1.2}], ml=0.0)
bar = add_box(s, 0.60, 6.14, 9.60, 0.48, fill=NAVY)
put_text(bar.text_frame,
         [{"runs": [("4会場のデータを次の会場へ引き継ぎ、開催ごとに精度を高める", 13.5, True, WHITE)], "align": "c"}],
         anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
add_text(s, 0.60, 6.70, 9.60, 0.30,
         [{"runs": [("出典：しが就職・転職フェア 公式サイト（2026年秋開催情報）", 10, None, MUT)]}], ml=0.0)

# ================= Slide 23: ステップ配信 Day0-14 =================
s = slides[22]
frame(s, "具体施策⑥ ステップ配信設計（Day 0 〜 Day 14）",
      "アクティブ率が最も高い追加直後に価値を渡し、開催前後の6ステップで来場と再接点を確保します。")
steps = [
    ("Day 0", "友だち追加・アンケート", "サンクスメッセージ＋希望条件診断（30秒）"),
    ("Day 1", "診断結果＆企業提示", "回答に合う出展企業をピックアップして提示"),
    ("Day 3", "不安解消コンテンツ", "服装・持ち物・ブースの回り方をQ&A形式で"),
    ("Day 6", "開催前日リマインド", "アクセスMAP付きで来場を後押し"),
    ("Day 7", "当日朝の最終通知", "開催時間・会場を再通知し、当日の迷いを解消"),
    ("Day 14", "イベント後フォロー", "お礼＋無料個別キャリア相談への誘導"),
]
for i, (d, title, body) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = 0.60 + col * 3.28
    y = 1.45 + row * 2.05
    hd = add_box(s, x, y, 3.05, 0.50, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    put_text(hd.text_frame,
             [{"runs": [(d + "　", 13, True, WHITE), (title, 12, True, WHITE)], "align": "c"}],
             anchor="m", ml=0.06, mr=0.06, mt=0.0, mb=0.0)
    bd = add_box(s, x, y + 0.56, 3.05, 1.20, fill=PALE)
    put_text(bd.text_frame, [{"runs": [(body, 11.5, None, INK)], "ls": 1.2}],
             anchor="m", ml=0.16, mr=0.12, mt=0.06, mb=0.06)
    if col < 2:
        add_box(s, x + 3.07, y + 0.13, 0.16, 0.26, fill=MUT, shape=MSO_SHAPE.RIGHT_ARROW)
add_text(s, 0.60, 5.62, 9.60, 0.56,
         [{"runs": [("Day 0〜3で「登録してよかった」を作り、Day 6〜7で来場を確定させ、Day 14で次の接点につなげます。開催日程に合わせてDay数は自動調整します。", 12, None, INK)], "ls": 1.2}], ml=0.0)
bar = add_box(s, 0.60, 6.24, 9.60, 0.48, fill=NAVY)
put_text(bar.text_frame,
         [{"runs": [("6ステップすべて自動配信。運用工数を増やさずに来場率を高める", 13.5, True, WHITE)], "align": "c"}],
         anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)

# ================= Slide 24: Meta広告×LINE連携 =================
s = slides[23]
frame(s, "具体施策⑦ Meta広告×LINE連携（オフラインCV最適化）",
      "「未来場者データの自動除外」により、予約からの来場引き上げ率を8割まで改善します。")
flow = [
    ("Meta広告", "予約・来場につながりやすいユーザーを獲得", GREY, INK),
    ("LINE公式アカウント", "予約フォーム送信後にLINEへ自動登録・リマインド配信", NAVY, WHITE),
    ("LINE公式アカウント", "当日の来場判定をデータ化（オフラインCV）", NAVY, WHITE),
    ("Meta広告へ連携", "非来場ユーザーを除外し、配信精度がさらに向上", GREY, INK),
]
for i, (head, body, cf, ct) in enumerate(flow):
    x = 0.60 + i * 2.46
    hd = add_box(s, x, 1.45, 2.28, 0.46, fill=cf, shape=MSO_SHAPE.RECTANGLE)
    put_text(hd.text_frame, [{"runs": [(head, 11.5, True, ct)], "align": "c"}],
             anchor="m", ml=0.03, mr=0.03, mt=0.0, mb=0.0)
    bd = add_box(s, x, 1.97, 2.28, 1.30, fill=WHITE, line=BORDER, lw=1.0, shape=MSO_SHAPE.RECTANGLE)
    put_text(bd.text_frame, [{"runs": [(body, 11.5, None, INK)], "ls": 1.2}],
             anchor="t", ml=0.14, mr=0.10, mt=0.12)
    if i < 3:
        add_box(s, x + 2.30, 2.44, 0.14, 0.30, fill=NAVY, shape=MSO_SHAPE.RIGHT_ARROW)
loop = add_box(s, 1.00, 3.42, 8.80, 0.44, fill=PALE, line=NAVY, lw=1.25, shape=MSO_SHAPE.LEFT_ARROW)
put_text(loop.text_frame,
         [{"runs": [("来場・非来場データをMeta広告へ連携 → 「来場する人」へ配信を最適化（ドタキャン層を除外）", 12, True, NAVY)], "align": "c"}],
         anchor="m", ml=0.55, mr=0.10, mt=0.0, mb=0.0)
eff = add_box(s, 0.60, 4.10, 5.60, 1.45, fill=PALE, line=NAVY, lw=1.25)
put_text(eff.text_frame,
         [{"runs": [("期待効果", 12, True, NAVY)], "sa": 5},
          {"runs": [("予約→来場の引き上げ率　", 12, None, INK), ("60〜70%", 14, True, INK),
                    (" → ", 12, None, MUT), ("80%", 20, True, NAVY)], "sa": 4},
          {"runs": [("実質来場コスト　", 12, None, INK), ("最大30〜40%削減", 14, True, NAVY)], "sa": 4},
          {"runs": [("他クライアントの実績値ベース（参照元：社内配信実績）", 10, None, MUT)]}],
         anchor="m", ml=0.20, mr=0.14, mt=0.06, mb=0.06)
qa = add_box(s, 6.45, 4.10, 3.75, 1.45, fill=WHITE, line=BORDER, lw=1.0)
put_text(qa.text_frame,
         [{"runs": [("「知らぬ間の登録」へのご配慮", 12, True, NAVY)], "sa": 5},
          {"runs": [("申込フォーム送信後、予約完了通知・会場案内をLINEでお届けする自然な遷移導線のため、ユーザーの違和感が生じにくい設計です。", 11, None, INK)], "ls": 1.2}],
         anchor="m", ml=0.18, mr=0.14, mt=0.06, mb=0.06)
bar = add_box(s, 0.60, 5.75, 9.60, 0.52, fill=NAVY)
put_text(bar.text_frame,
         [{"runs": [("この導線の中核はLINE公式アカウント ─ LINE運用の質が、広告配信の精度まで左右する", 14, True, WHITE)], "align": "c"}],
         anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)

# ================= Slide 25: LP離脱防止ポップアップ =================
s = slides[24]
frame(s, "具体施策⑧ LP離脱防止ポップアップ・LINE誘導導線",
      "特典をフックにLINE登録率を底上げし、「予約に至らなかった訪問者」も友だち資産に変えます。")
bf = add_box(s, 0.60, 1.45, 4.35, 2.30, fill=WHITE, line=BORDER, lw=1.25)
put_text(bf.text_frame,
         [{"runs": [("Before（現状）", 14, True, RED)], "sa": 6},
          {"runs": [("・LP内のLINE導線は「最新情報配信中」の告知のみ", 12, None, INK)], "sa": 4, "ls": 1.15},
          {"runs": [("・登録するメリットが提示されていない", 12, None, INK)], "sa": 4, "ls": 1.15},
          {"runs": [("・予約に至らなかった訪問者は、そのまま離脱", 12, None, INK)], "ls": 1.15}],
         anchor="t", ml=0.20, mr=0.16, mt=0.14)
add_box(s, 5.05, 2.36, 0.62, 0.50, fill=NAVY, shape=MSO_SHAPE.RIGHT_ARROW)
af = add_box(s, 5.85, 1.45, 4.35, 2.30, fill=PALE, line=NAVY, lw=1.5)
put_text(af.text_frame,
         [{"runs": [("After（ご提案）", 14, True, NAVY)], "sa": 6},
          {"runs": [("・離脱の兆候を検知してポップアップを表示", 12, None, INK)], "sa": 4, "ls": 1.15},
          {"runs": [("・『滋賀県内 業界別平均年収＆面接対策シート』を特典化", 12, None, INK)], "sa": 4, "ls": 1.15},
          {"runs": [("・登録メリットを明示し、友だち追加率を底上げ", 12, None, INK)], "ls": 1.15}],
         anchor="t", ml=0.20, mr=0.16, mt=0.14)
add_text(s, 0.60, 3.92, 9.60, 0.34,
         [{"runs": [("ご提供特典（案）", 13, True, NAVY)]}], ml=0.0)
gifts = [
    ("① 滋賀県内 業界別平均年収シート", "地元の相場観を提示し、転職検討の材料として活用いただく。"),
    ("② 面接・自己PR対策チェックシート", "当日のブース面談で使える実用コンテンツとして提供。"),
]
for i, (h, b) in enumerate(gifts):
    x = 0.60 + i * 4.95
    card(s, x, 4.30, 4.65, 1.05, h, b, hcol=NAVY, fill=WHITE, hsz=13, bsz=11.5, line=NAVY)
bar = add_box(s, 0.60, 5.62, 9.60, 0.52, fill=NAVY)
put_text(bar.text_frame,
         [{"runs": [("予約に至らなかった訪問者こそ、次回フェアの見込み客になる", 14, True, WHITE)], "align": "c"}],
         anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
add_text(s, 0.60, 6.24, 9.60, 0.30,
         [{"runs": [("※特典の内容・名称は貴社と協議のうえ確定します", 10, None, MUT)]}], ml=0.0)

# ================= Slide 27: 費用対効果サマリー =================
s = slides[26]
frame(s, "想定の費用対効果（サマリー）",
      "広告予算80万円は据え置いたまま、実来場者数を月48名→80名へ拡大することを目指します。")
kpis = [("初期費用", "¥215,000", "初期構築20万＋離脱防止1.5万"),
        ("月次固定費", "¥80,000", "コンサル5万＋離脱防止1.5万＋アカウント費1.5万"),
        ("契約期間", "3ヶ月〜", "秋開催4会場（9月末〜11月頭）を1クールとして検証")]
for i, (lab, val, note) in enumerate(kpis):
    x = 0.60 + i * 3.28
    sp = add_box(s, x, 1.42, 3.05, 1.62, fill=PALE)
    put_text(sp.text_frame,
             [{"runs": [(lab, 12, True, NAVY)], "align": "c", "sa": 4},
              {"runs": [(val, 20, True, NAVY)], "align": "c", "sa": 4},
              {"runs": [(note, 10, None, INK)], "align": "c", "ls": 1.15}],
             anchor="m", ml=0.12, mr=0.12, mt=0.06, mb=0.06)
rows = [
    ("月額広告予算", "80万円", "80万円", "±0円（予算据え置き）"),
    ("事前予約単価（CPA）", "10,000円", "8,000円", "20%削減：LINE内予約導線・CVR改善"),
    ("月間事前予約数", "80名", "100名", "+20名（+25%）"),
    ("予約からの来場引き上げ率", "60%", "80%", "+20pt：リマインド＋Meta除外最適化"),
    ("月間イベント実来場者数", "48名", "80名", "+32名（+66%）"),
    ("実質来場単価", "16,667円", "10,000円", "40%削減"),
]
gf = s.shapes.add_table(7, 5, Inches(0.60), Inches(3.20), Inches(9.60), Inches(2.86))
tbl = gf.table
tbl.first_row = False
tbl.horz_banding = False
for i, w in enumerate([2.55, 1.45, 0.45, 1.65, 3.50]):
    tbl.columns[i].width = Inches(w)
tbl.rows[0].height = Inches(0.36)
for r in range(1, 7):
    tbl.rows[r].height = Inches(0.41)
for c, htext in enumerate(["指標", "現状の実績", "", "提案後（目標）", "改善インパクト"]):
    cell = tbl.cell(0, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor.from_string(NAVY)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
    put_text(cell.text_frame, [{"runs": [(htext, 11, True, WHITE)], "align": "c"}],
             anchor="m", ml=0.02, mr=0.02, mt=0.0, mb=0.0)
for r, (met, cur, prop, drv) in enumerate(rows, start=1):
    band = WHITE if r % 2 == 1 else PALE
    hl = r in (5, 6)
    vals = [(met, 10.5, True, INK, "l"), (cur, 11, None, INK, "c"), ("→", 11, None, MUT, "c"),
            (prop, 14 if hl else 12, True, NAVY, "c"), (drv, 10, None, INK, "l")]
    for c, (t, sz, b, col, al) in enumerate(vals):
        cell = tbl.cell(r, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor.from_string(band)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.04)
        cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
        put_text(cell.text_frame, [{"runs": [(t, sz, b, col)], "align": al}],
                 anchor="m", ml=0.02, mr=0.02, mt=0.0, mb=0.0)
add_text(s, 0.60, 6.16, 9.60, 0.72,
         [{"runs": [("※現状値は貴社ヒアリングに基づく実情数値、提案後はLINE×Meta連携運用に基づく試算目標値です。", 10, None, MUT)], "sa": 2},
          {"runs": [("※複数プランにおける詳細な収支シミュレーションは、別添のExcel資料（SIM）をご参照ください。", 10, None, MUT)]}], ml=0.0)

# ================= Slide 29: サポート体制 =================
s = slides[28]
frame(s, "サポート体制（チーム体制・コミュニケーション）",
      "営業2名＋運用担当＋運用統括の4名体制で、企画から運用改善までを一気通貫で担当します。")
members = [
    ("営業担当", "高橋", "窓口・企画提案", "貴社の課題整理と施策の企画、\nご提案・進行管理を担当。"),
    ("営業担当", "武者", "窓口・企画提案", "会場ごとの運用要件を整理し、\n社内制作チームへ連携。"),
    ("運用担当", "―", "配信・分析の実務", "配信設計、シナリオ実装、\n数値分析と改善提案を担当。"),
    ("運用統括", "―", "品質管理・戦略設計", "配信戦略の監修と品質担保、\n広告連携の全体最適を担当。"),
]
for i, (role, name, duty, body) in enumerate(members):
    x = 0.60 + i * 2.46
    hd = add_box(s, x, 1.45, 2.28, 0.68, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    put_text(hd.text_frame,
             [{"runs": [(role, 13, True, WHITE)], "align": "c", "sa": 2},
              {"runs": [(name if name != "―" else "専任担当", 10.5, None, WHITE)], "align": "c"}],
             anchor="m", ml=0.04, mr=0.04, mt=0.0, mb=0.0)
    bd = add_box(s, x, 2.21, 2.28, 1.85, fill=PALE)
    put_text(bd.text_frame,
             [{"runs": [(duty, 12, True, NAVY)], "sa": 6, "align": "c"},
              {"runs": [(body.replace("\n", ""), 11, None, INK)], "ls": 1.2}],
             anchor="t", ml=0.16, mr=0.12, mt=0.14)
add_text(s, 0.60, 4.25, 9.60, 0.34,
         [{"runs": [("体制のポイント", 13, True, NAVY)]}], ml=0.0)
points = [
    ("案件数を絞った専任制", "1担当あたり5〜8件に案件数を限定。担当が貴社の状況を把握した状態で、質の高い提案を継続します。"),
    ("スピーディな社内連携", "営業・運用・制作が同一チーム。配信内容の変更や会場ごとの出し分けにも即時対応します。"),
]
for i, (h, b) in enumerate(points):
    x = 0.60 + i * 4.95
    card(s, x, 4.63, 4.65, 1.22, h, b, hcol=NAVY, fill=WHITE, hsz=13, bsz=11.5, line=NAVY)
bar = add_box(s, 0.60, 6.10, 9.60, 0.52, fill=NAVY)
put_text(bar.text_frame,
         [{"runs": [("4名体制・案件数を絞った専任制で、秋クールの4会場を伴走支援します", 14, True, WHITE)], "align": "c"}],
         anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)

# ================= Slide 31: 計測環境の構築 =================
s = slides[30]
frame(s, "計測環境の構築（オフラインCV連携＆GA計測）",
      "LINE上のアクションと、LP・Meta広告側の成果を紐づけて計測できる環境を構築します。")
layers = [
    ("① LP／申込フォーム", "GA4パラメータ管理", "LINE経由の流入をパラメータで判別し、予約フォーム到達・完了までを計測。"),
    ("② LINE公式アカウント", "友だち・配信・反応の計測", "友だち数、ブロック率、開封率、クリック率、タグ別のセグメント反応を管理画面で取得。"),
    ("③ フェア当日（オフライン）", "来場判定のデータ化", "受付での来場確認をデータ化し、予約者のうち誰が来場したかを紐づけ。"),
    ("④ Meta広告", "オフラインCVの還元", "来場・非来場データを広告側へ連携し、配信対象の最適化に活用。"),
]
for i, (title, sub, body) in enumerate(layers):
    y = 1.42 + i * 1.10
    hd = add_box(s, 0.60, y, 2.75, 0.94, fill=NAVY)
    put_text(hd.text_frame,
             [{"runs": [(title, 12.5, True, WHITE)], "align": "c", "sa": 3},
              {"runs": [(sub, 10.5, None, WHITE)], "align": "c"}],
             anchor="m", ml=0.08, mr=0.08, mt=0.0, mb=0.0)
    add_box(s, 3.47, y + 0.30, 0.42, 0.34, fill=MUT, shape=MSO_SHAPE.RIGHT_ARROW)
    bd = add_box(s, 4.00, y, 6.20, 0.94, fill=PALE)
    put_text(bd.text_frame, [{"runs": [(body, 12, None, INK)], "ls": 1.2}],
             anchor="m", ml=0.18, mr=0.14, mt=0.06, mb=0.06)
add_text(s, 0.60, 5.92, 9.60, 0.56,
         [{"runs": [("これにより「広告→予約→来場」までを一本の数値で追跡でき、会場ごとの費用対効果を開催のたびに検証できます。", 12, None, INK)], "ls": 1.2}], ml=0.0)
bar = add_box(s, 0.60, 6.30, 9.60, 0.48, fill=NAVY)
put_text(bar.text_frame,
         [{"runs": [("「予約まで」ではなく「来場まで」を計測できる状態をつくる", 13.5, True, WHITE)], "align": "c"}],
         anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)

# ================= Slide 32: 運用スケジュール =================
s = slides[31]
frame(s, "運用スケジュール（キックオフから運用開始まで約1ヶ月）",
      "8月中に構築を完了させ、8/28の受付開始から9/26の初回フェアに間に合わせます。")
steps = [
    ("STEP 1", "お申込み", "3ヶ月契約（秋開催4会場分）の書類締結"),
    ("STEP 2", "キックオフMTG", "目標数値・役割分担・スケジュールの確定"),
    ("STEP 3", "連携設定・タグ確認", "GA4パラメータ設定、Metaデータ連携の設計"),
    ("STEP 4", "初期構築", "リッチメニュー・アンケート・シナリオ制作、動画クリエイティブ確認"),
    ("STEP 5", "運用開始・検証", "8/28受付開始に合わせて配信開始、9/26フェアで初回検証"),
]
for i, (st, title, body) in enumerate(steps):
    y = 1.45 + i * 0.92
    sp = add_box(s, 0.60, y, 1.55, 0.76, fill=NAVY)
    put_text(sp.text_frame, [{"runs": [(st, 13, True, WHITE)], "align": "c"}],
             anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
    bd = add_box(s, 2.30, y, 7.90, 0.76, fill=PALE)
    put_text(bd.text_frame,
             [{"runs": [(title + "　", 13, True, NAVY), (body, 11.5, None, INK)]}],
             anchor="m", ml=0.20, mr=0.14, mt=0.04, mb=0.04)
bar = add_box(s, 0.60, 6.10, 9.60, 0.52, fill=NAVY)
put_text(bar.text_frame,
         [{"runs": [("8月中の構築完了により、秋クール4会場すべてで施策を稼働できます", 14, True, WHITE)], "align": "c"}],
         anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
add_text(s, 0.60, 6.70, 9.60, 0.30,
         [{"runs": [("※貴社作業・確認にかかる時間により前後する可能性がございます。", 10, None, MUT)]}], ml=0.0)

# ================= Slide 33: 空のフッター枠にCopyrightを復元 =================
s = slides[32]
for sh in s.shapes:
    if sh.has_text_frame and sh.name == "正方形/長方形 5" and not sh.text_frame.text.strip():
        put_text(sh.text_frame,
                 [{"runs": [("Copyright(c) DYM Co., Ltd.  All Rights Reserved.", 10, True, WHITE)]}],
                 anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)

prs.save(OUT)
print("saved:", OUT)
print("slides:", len(list(Presentation(OUT).slides)))
