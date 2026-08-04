# -*- coding: utf-8 -*-
"""しが就職・転職フェア LINEOA提案 8枚ビルド（FMT(3)コピー→既存編集経路）v2: 本文フォント拡大＋余白圧縮"""
import shutil
from copy import deepcopy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
SRC = str(ROOT / "_templates" / "DYM_LINEOA_FMT.pptx")
OUT = str(ROOT / "20260804_求人vivical『しが就職・転職フェア』様_LINE公式アカウント活用のご提案_ver2.pptx")

NAVY = "002060"
INK = "343434"
BLUE = "1565C0"
PALE = "F4F7FF"
GREYBG = "ECEDF3"
BORDER = "D9D9D9"
MUT = "808080"
WHITE = "FFFFFF"

shutil.copyfile(SRC, OUT)
prs = Presentation(OUT)
slides = list(prs.slides)

# ---- 素材キャプチャ：スライド2の区切り線（y=1390675） ----
div_el = None
for sh in slides[1].shapes:
    if sh._element.tag.endswith('}cxnSp') and abs(sh.top - 1390675) < 60000:
        div_el = deepcopy(sh._element)
        break
assert div_el is not None, "divider not found"

# ---- スライド10以降を削除（9枚構成） ----
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst)[9:]:
    prs.part.drop_rel(sldId.rId)
    sldIdLst.remove(sldId)
slides = list(prs.slides)
assert len(slides) == 9

# ---------- helpers ----------
def set_font(run, size, bold=None, color=INK, name="メイリオ"):
    f = run.font
    f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", name)
    f.color.rgb = RGBColor.from_string(color)

ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
ANCH = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}

def put_text(tf, paras, anchor="t", ml=0.06, mr=0.06, mt=0.03, mb=0.03, wrap=True):
    tf.word_wrap = wrap
    tf.margin_left = Inches(ml)
    tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt)
    tf.margin_bottom = Inches(mb)
    tf.vertical_anchor = ANCH[anchor]
    first = True
    for p in paras:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = ALIGN[p.get("align", "l")]
        if p.get("sa") is not None:
            para.space_after = Pt(p["sa"])
        if p.get("ls") is not None:
            para.line_spacing = p["ls"]
        for (t, sz, b, c) in p["runs"]:
            r = para.add_run()
            r.text = t
            set_font(r, sz, b, c)
    return tf

def add_text(slide, x, y, w, h, paras, anchor="t", **kw):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    put_text(box.text_frame, paras, anchor=anchor, **kw)
    return box

def add_box(slide, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor.from_string(line)
        sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    return sp

def clear_slide(slide):
    spTree = slide.shapes._spTree
    for el in list(spTree):
        if el.tag.split("}")[-1] in ("sp", "cxnSp", "pic", "graphicFrame", "grpSp"):
            spTree.remove(el)

def frame(slide, header, band):
    """タイトル＋結論バンド＋区切り線（FMT実測準拠）"""
    clear_slide(slide)
    slide.shapes._spTree.append(deepcopy(div_el))
    add_text(slide, 0.60, 0.13, 8.10, 0.40,
             [{"runs": [(header, 16, True, NAVY)], "align": "l"}],
             anchor="m", ml=0.0, mr=0.0)
    add_text(slide, 0.42, 0.66, 10.02, 0.80,
             [{"runs": [(band, 14, None, INK)], "align": "c"}],
             anchor="m", ml=0.0, mr=0.0)

def replace_text(shape, text):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    runs = p.runs
    runs[0].text = text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)

# ---------- Slide 1: 表紙（既存シェイプのテキスト差し替え） ----------
s = slides[0]
for sh in s.shapes:
    nm = sh.name
    if nm == "正方形/長方形 5":
        replace_text(sh, "求人vivical『しが就職・転職フェア』様向け")
    elif nm == "正方形/長方形 1":
        replace_text(sh, "LINE公式アカウント活用のご提案")
    elif nm == "正方形/長方形 2":
        replace_text(sh, "〜「直前辞退の防止」と「次回イベントへの継続囲い込み」による採用成果の最大化〜")
        sh.left = Inches(1.07)
        sh.width = Inches(8.70)
    elif nm == "正方形/長方形 6":
        replace_text(sh, "2026年8月｜株式会社DYM")
        sh.left = Inches(3.72)
        sh.width = Inches(3.40)

# ---------- Slide 2: 背景・目的 ----------
s = slides[1]
frame(s, "ご提案の背景と目的",
      "転職フェアの成果最大化には「予約後の引き上げ率（来場率）向上」が不可欠である")
for (y, chip, txt) in [
    (1.78, "現状", "月額広告費80万円で目標予約数は達成（予約単価10,000円）。一方で予約者の3〜4割が直前辞退し、来場引き上げ率は60〜70%にとどまる"),
    (2.26, "課題", "その結果、実質的な「来場獲得単価」は16,000〜18,000円まで高騰している"),
]:
    c = add_box(s, 0.50, y, 1.00, 0.38, fill=NAVY)
    put_text(c.text_frame, [{"runs": [(chip, 11, True, WHITE)], "align": "c"}], anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
    add_text(s, 1.65, y - 0.06, 8.68, 0.50,
             [{"runs": [(txt, 11.5, None, INK)]}], anchor="m", ml=0.0)
add_text(s, 0.50, 2.80, 6.0, 0.34,
         [{"runs": [("本施策で目指す3つのゴール", 13, True, NAVY)]}], anchor="m", ml=0.0)
goals = [
    ("GOAL 1｜来場率の向上", "予約後の直前辞退を防ぎ、実来場数を底上げする",
     "予約→来場の引き上げ率", "60〜70%→80%", "LINEリマインド＋Meta連携による改善"),
    ("GOAL 2｜来場単価の削減", "ドタキャン層への無駄な広告費を排除する",
     "実質来場単価", "16,667円→10,000円", "予算据え置きのまま実来場数を拡大"),
    ("GOAL 3｜LTVの最大化", "今回未応募の求職者をストックし、継続的に接点を持つ",
     "20代後半〜30代前半を中心に", "再案内", "次回フェア・個別求人へセグメント配信"),
]
for i, (g, desc, lab, val, mech) in enumerate(goals):
    x = 0.50 + i * 3.36
    add_box(s, x, 3.18, 3.11, 3.28, fill=PALE)
    hd = add_box(s, x, 3.18, 3.11, 0.48, fill=NAVY)
    put_text(hd.text_frame, [{"runs": [(g, 12.5, True, WHITE)], "align": "c"}], anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
    add_text(s, x + 0.15, 3.80, 2.81, 0.80, [{"runs": [(desc, 11, None, INK)], "ls": 1.12}])
    add_text(s, x + 0.15, 4.76, 2.81, 0.30, [{"runs": [(lab, 10, None, MUT)], "align": "c"}], ml=0.0, mr=0.0)
    add_text(s, x + 0.15, 5.06, 2.81, 0.60,
             [{"runs": [(val, 19 if i < 2 else 26, True, BLUE)], "align": "c"}], anchor="m", ml=0.0, mr=0.0)
    add_text(s, x + 0.15, 5.80, 2.81, 0.62, [{"runs": [(mech, 10.5, None, INK)], "align": "c", "ls": 1.1}], ml=0.0, mr=0.0)
add_text(s, 0.50, 6.56, 9.83, 0.30,
         [{"runs": [("参照元：貴社ご提供の実績数値（現状）／提案後は本施策の実行を前提とした目標値です（成果を保証するものではありません）", 9.5, None, MUT)]}], ml=0.0)

# ---------- Slide 3: 3つの離脱ポイント ----------
s = slides[2]
frame(s, "課題の整理｜求職者の3つの離脱ポイント",
      "「予約後のフォロー不足」により、広告予算の約30〜40%が来場に繋がらず流出している")
cols = [
    ("① イベント前", "引き上げ率60〜70%（3〜4割が直前辞退）。実質来場単価は16,000〜18,000円に高騰",
     "前日・当日朝のLINEリマインド配信と安心Q&A（服装自由・履歴書不要）の送付により、引き上げ率80%へ引き上げ"),
    ("② イベント当日", "会場内の回遊不足により、企業ブース訪問が伸びない",
     "リッチメニューで会場MAP・出展企業一覧を即時提供"),
    ("③ イベント後", "未応募層が離脱し、フェアの記憶が風化する",
     "アンケート属性に基づくセグメント配信で次回フェア・個別求人を案内"),
]
for i, (ph, issue, action) in enumerate(cols):
    x = 0.50 + i * 3.36
    hd = add_box(s, x, 1.78, 3.11, 0.46, fill=NAVY)
    put_text(hd.text_frame, [{"runs": [(ph, 13, True, WHITE)], "align": "c"}], anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
    if i < 2:
        add_box(s, x + 3.14, 1.88, 0.19, 0.26, fill=MUT, shape=MSO_SHAPE.RIGHT_ARROW)
    kb = add_box(s, x, 2.40, 3.11, 1.30, fill=GREYBG)
    put_text(kb.text_frame,
             [{"runs": [("課題", 10.5, True, MUT)], "sa": 4},
              {"runs": [(issue, 11.5 if i == 0 else 12, None, INK)], "ls": 1.12}],
             anchor="t", ml=0.12, mr=0.12, mt=0.09)
    add_box(s, x + 1.43, 3.82, 0.25, 0.28, fill=NAVY, shape=MSO_SHAPE.DOWN_ARROW)
    ab = add_box(s, x, 4.20, 3.11, 1.80, fill=PALE, line=NAVY, lw=1.0)
    put_text(ab.text_frame,
             [{"runs": [("LINEでの打ち手", 10.5, True, NAVY)], "sa": 4},
              {"runs": [(action, 11.5 if i == 0 else 12, None, INK)], "ls": 1.12}],
             anchor="t", ml=0.12, mr=0.12, mt=0.09)
bar = add_box(s, 0.50, 6.35, 9.83, 0.44, fill=NAVY)
put_text(bar.text_frame,
         [{"runs": [("3つの離脱ポイントをLINE公式アカウントで一気通貫にカバーする", 12, True, WHITE)], "align": "c"}],
         anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)

# ---------- Slide 4: ステップ配信 ----------
s = slides[3]
frame(s, "施策①｜ステップ配信・リマインド設計",
      "属性取得アンケートと適切なタイミングのリマインドにより、来場率を大幅に改善する")
steps = [
    ("① 友だち追加直後",
     ["・サンクスメッセージ", "・フェアの回り方ガイド", "・簡単アンケート（希望職種などの属性取得）"],
     "属性取得・不安解消"),
    ("② 開催3日前〜前日",
     ["・会場アクセス案内", "・「服装自由・履歴書不要」などの安心Q&A"],
     "来場準備・安心醸成"),
    ("③ 開催当日 朝",
     ["・開催リマインド通知", "・会場MAPリンク付きで当日の迷いを解消"],
     "来場の最終後押し"),
    ("④ イベント終了翌日",
     ["・ご来場お礼メッセージ", "・無料個別キャリア相談へ誘導"],
     "関係継続・再接点化"),
]
for i, (tm, items, aim) in enumerate(steps):
    x = 0.50 + i * 2.515
    hd = add_box(s, x, 1.76, 2.28, 0.46, fill=NAVY)
    put_text(hd.text_frame, [{"runs": [(tm, 11.5, True, WHITE)], "align": "c"}], anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
    if i < 3:
        add_box(s, x + 2.30, 1.87, 0.19, 0.24, fill=MUT, shape=MSO_SHAPE.RIGHT_ARROW)
    bd = add_box(s, x, 2.36, 2.28, 2.30, fill=WHITE, line=BORDER, lw=1.0)
    put_text(bd.text_frame,
             [{"runs": [(t, 11, None, INK)], "ls": 1.12, "sa": 5} for t in items],
             anchor="t", ml=0.10, mr=0.08, mt=0.10)
    add_text(s, x, 4.76, 2.28, 0.66,
             [{"runs": [("狙い", 9.5, None, MUT)], "align": "c", "sa": 2},
              {"runs": [(aim, 12, True, BLUE)], "align": "c"}],
             ml=0.0, mr=0.0)
pb = add_box(s, 0.50, 5.95, 9.83, 0.70, fill=PALE)
chip = add_box(s, 0.65, 6.12, 1.15, 0.36, fill=NAVY)
put_text(chip.text_frame, [{"runs": [("ポイント", 11, True, WHITE)], "align": "c"}], anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
add_text(s, 1.95, 5.95, 8.38, 0.70,
         [{"runs": [("取得した属性（希望職種・年代など）は、イベント後のセグメント配信（次回フェア・個別求人のご案内）に活用します", 11.5, None, INK)]}],
         anchor="m", ml=0.0)

# ---------- Slide 5: リッチメニュー ----------
s = slides[4]
frame(s, "施策②｜リッチメニュー設計（6エリア構成）",
      "求める情報へ1タップで遷移できる6エリア構成により、利用頻度と来場意欲を高める")
add_box(s, 0.50, 1.95, 4.72, 4.20, fill=WHITE, line=NAVY, lw=1.75)
cells_top = [["次回フェア日程", "会場案内"], ["出展企業一覧を見る"], ["会場MAP・アクセス"]]
cells_btm = [["事前参加予約", "特典受取"], ["無料キャリア相談", "（1on1）"], ["よくある質問"]]
gap = 0.05
cw = (4.72 - gap * 4) / 3
ch = (4.20 - gap * 3) / 2
for r, row in enumerate([cells_top, cells_btm]):
    for c, label in enumerate(row):
        x = 0.50 + gap + c * (cw + gap)
        y = 1.95 + gap + r * (ch + gap)
        fill = PALE if r == 0 else NAVY
        tcol = NAVY if r == 0 else WHITE
        cell = add_box(s, x, y, cw, ch, fill=fill)
        put_text(cell.text_frame,
                 [{"runs": [(line, 11, True, tcol)], "align": "c", "ls": 1.15} for line in label],
                 anchor="m", ml=0.03, mr=0.03, mt=0.0, mb=0.0)
add_text(s, 0.50, 6.25, 4.72, 0.30,
         [{"runs": [("リッチメニュー構成イメージ（トーク画面下部に常時表示）", 9.5, None, MUT)], "align": "c"}], ml=0.0, mr=0.0)
panels = [
    (1.95, "上段：情報提供エリア", "イベント概要や出展企業を常時可視化し、当日の会場回遊率を高めます（フェア日程／出展企業一覧／会場MAP）"),
    (4.20, "下段：アクション誘導エリア", "事前予約・個別相談・疑問解消へ1タップで誘導し、予約・相談のCVRを改善します（参加予約／相談／FAQ）"),
]
for (y, hd, body) in panels:
    add_box(s, 5.42, y, 4.91, 1.95, fill=PALE)
    add_text(s, 5.60, y + 0.12, 4.55, 0.36, [{"runs": [(hd, 13, True, NAVY)]}], ml=0.0)
    add_text(s, 5.60, y + 0.56, 4.55, 1.25, [{"runs": [(body, 12, None, INK)], "ls": 1.2}], ml=0.0)
add_text(s, 5.42, 6.25, 4.91, 0.30,
         [{"runs": [("※メニュー項目・文言は貴社と協議のうえ確定します", 9.5, None, MUT)]}], ml=0.0)

# ---------- Slide 6: LP導線 Before/After ----------
s = slides[5]
frame(s, "施策③｜友だち追加導線の改善（LP）",
      "LP上のLINE導線に「求職者限定特典」を紐付けることで、友だち追加率を底上げする")
# Before
add_box(s, 0.50, 1.80, 4.42, 3.15, fill=WHITE, line=BORDER, lw=1.25)
hb = add_box(s, 0.50, 1.80, 4.42, 0.50, fill=GREYBG)
put_text(hb.text_frame, [{"runs": [("Before｜現状", 12.5, True, INK)], "align": "c"}], anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
btn = add_box(s, 1.36, 2.76, 2.70, 0.55, fill=GREYBG, line=MUT, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
put_text(btn.text_frame, [{"runs": [("LINE友だち追加はこちら", 11.5, None, INK)], "align": "c"}], anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
add_text(s, 0.68, 3.52, 4.06, 1.30,
         [{"runs": [("・登録メリットの提示がない", 11.5, None, INK)], "sa": 5, "ls": 1.15},
          {"runs": [("・「追加する理由」が弱く、登録前に離脱が発生", 11.5, None, INK)], "ls": 1.15}], ml=0.0)
# arrow
add_box(s, 5.01, 3.10, 0.81, 0.55, fill=NAVY, shape=MSO_SHAPE.RIGHT_ARROW)
# After
add_box(s, 5.91, 1.80, 4.42, 3.15, fill=PALE, line=NAVY, lw=1.5)
ha = add_box(s, 5.91, 1.80, 4.42, 0.50, fill=NAVY)
put_text(ha.text_frame, [{"runs": [("After｜ご提案", 12.5, True, WHITE)], "align": "c"}], anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
btn2 = add_box(s, 6.77, 2.76, 2.70, 0.55, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
put_text(btn2.text_frame, [{"runs": [("限定特典を受け取って友だち追加", 11, True, WHITE)], "align": "c"}], anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
add_text(s, 6.09, 3.52, 4.06, 1.30,
         [{"runs": [("・登録者限定！『滋賀県内 業界別平均年収＆面接対策チェックシート』をプレゼント", 11.5, None, INK)], "sa": 5, "ls": 1.15},
          {"runs": [("・登録メリットを明確化し、友だち追加率を底上げ", 11.5, None, INK)], "ls": 1.15}], ml=0.0)
add_text(s, 0.50, 5.15, 6.0, 0.32, [{"runs": [("ご提供特典（案）", 12, True, NAVY)]}], ml=0.0)
for i, t in enumerate(["① 滋賀県転職市場レポート", "② 面接・自己PR対策チェックシート（PDF）"]):
    cbox = add_box(s, 0.50 + i * 5.0, 5.52, 4.83, 0.78, fill=WHITE, line=NAVY, lw=1.0)
    put_text(cbox.text_frame, [{"runs": [(t, 12, True, INK)], "align": "c"}], anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
add_text(s, 0.50, 6.46, 9.83, 0.30,
         [{"runs": [("※特典の内容・名称は貴社と協議のうえ確定します", 9.5, None, MUT)]}], ml=0.0)

# ---------- Slide 7: Meta広告×LINE連携（オフラインCV最適化） ----------
s = slides[6]
frame(s, "施策④｜Meta広告×LINE連携（オフラインCV最適化）",
      "「未来場者データの自動除外」により、予約からの来場引き上げ率を8割まで改善する")
flow = [
    ("Meta広告", GREYBG, INK, "予約・来場につながりやすいユーザーを獲得"),
    ("LINE公式アカウント", NAVY, WHITE, "友だち追加＆離脱防止リマインドを配信"),
    ("LINE公式アカウント", NAVY, WHITE, "当日の来場判定をデータ化（オフラインCV）"),
    ("Meta広告へ連携", GREYBG, INK, "非来場ユーザーを除外し、配信精度がさらに向上"),
]
for i, (chip, cf, ct, body) in enumerate(flow):
    x = 0.50 + i * 2.56
    card = add_box(s, x, 1.85, 2.16, 1.70, fill=WHITE, line=BORDER, lw=1.0)
    hd = add_box(s, x, 1.85, 2.16, 0.42, fill=cf)
    put_text(hd.text_frame, [{"runs": [(chip, 10.5, True, ct)], "align": "c"}], anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
    add_text(s, x + 0.10, 2.38, 1.96, 1.05,
             [{"runs": [(body, 11, None, INK)], "ls": 1.15}], ml=0.0, mr=0.0)
    if i < 3:
        add_box(s, x + 2.22, 2.58, 0.28, 0.24, fill=NAVY, shape=MSO_SHAPE.RIGHT_ARROW)
loop = add_box(s, 0.90, 3.80, 8.93, 0.42, fill=PALE, line=NAVY, lw=1.25, shape=MSO_SHAPE.LEFT_ARROW)
put_text(loop.text_frame,
         [{"runs": [("来場・非来場データをMeta広告へ連携 → 「来場する人」へ配信を最適化（ドタキャン層を除外）", 11, True, NAVY)], "align": "c"}],
         anchor="m", ml=0.55, mr=0.10, mt=0.0, mb=0.0)
eff = add_box(s, 0.50, 4.50, 5.90, 1.35, fill=PALE, line=NAVY, lw=1.25)
put_text(eff.text_frame,
         [{"runs": [("期待効果", 11, True, NAVY)], "sa": 4},
          {"runs": [("予約→来場の引き上げ率　", 11, None, INK), ("60〜70%", 13, True, INK),
                    (" → ", 11.5, None, MUT), ("80%", 19, True, BLUE)], "sa": 3},
          {"runs": [("実質来場コスト　", 11, None, INK), ("最大30〜40%削減", 13, True, BLUE)], "sa": 3},
          {"runs": [("他クライアントの実績値ベース（参照元：社内配信実績）", 9, None, MUT)]}],
         anchor="m", ml=0.18, mr=0.12, mt=0.05, mb=0.05)
qa = add_box(s, 6.60, 4.50, 3.73, 1.35, fill=WHITE, line=BORDER, lw=1.0)
put_text(qa.text_frame,
         [{"runs": [("「知らぬ間の登録」へのご配慮", 11, True, NAVY)], "sa": 5},
          {"runs": [("申込フォーム送信後、予約完了通知・会場案内をLINEでお届けする自然な遷移導線のため、ユーザーの違和感が生じにくい設計です", 10, None, INK)], "ls": 1.15}],
         anchor="m", ml=0.16, mr=0.12, mt=0.06, mb=0.06)
bar = add_box(s, 0.50, 6.15, 9.83, 0.48, fill=NAVY)
put_text(bar.text_frame,
         [{"runs": [("この導線の中核はLINE公式アカウント ─ LINE運用の質が、広告配信の精度まで左右する", 12, True, WHITE)], "align": "c"}],
         anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)

# ---------- Slide 8: シミュレーション ----------
s = slides[7]
frame(s, "効果シミュレーション",
      "広告予算（80万円）は据え置きのまま、実来場者数を「月48名 → 80名（1.66倍）」へ拡大する")
rows = [
    ("月額広告予算", "80万円", "80万円", "±0円（予算据え置き）"),
    ("事前予約単価（CPA）", "10,000円", "8,000円", "20%削減：LINE内予約導線・CVR改善"),
    ("月間事前予約数", "80名", "100名", "+20名（+25%）：予約単価改善に伴う獲得数増"),
    ("予約からの来場引き上げ率", "60%", "80%", "+20pt：LINEリマインド＋Meta除外最適化"),
    ("月間イベント実来場者数", "48名", "80名", "+32名（+66%）：予約数増×来場率80%の相乗効果"),
    ("実質来場単価", "16,667円", "10,000円", "40%削減：ドタキャン層への広告費を排除"),
]
gf = s.shapes.add_table(7, 5, Inches(0.50), Inches(1.78), Inches(9.83), Inches(3.66))
tbl = gf.table
tbl.first_row = False
tbl.horz_banding = False
widths = [2.50, 1.55, 0.50, 1.80, 3.48]
for i, w in enumerate(widths):
    tbl.columns[i].width = Inches(w)
tbl.rows[0].height = Inches(0.42)
for r in range(1, 7):
    tbl.rows[r].height = Inches(0.54)
headers = ["指標", "現状の実績", "", "提案後（目標）", "改善インパクト・ドライバー"]
for c, htext in enumerate(headers):
    cell = tbl.cell(0, c)
    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor.from_string(NAVY)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.04)
    cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
    put_text(cell.text_frame, [{"runs": [(htext, 11.5, True, WHITE)], "align": "c"}], anchor="m", ml=0.02, mr=0.02, mt=0.0, mb=0.0)
for r, (met, cur, prop, drv) in enumerate(rows, start=1):
    band = WHITE if r % 2 == 1 else PALE
    hl = r in (5, 6)
    vals = [(met, 11, True, INK, "l"), (cur, 11.5, None, INK, "c"),
            ("→", 11.5, None, MUT, "c"),
            (prop, 15 if hl else 13, True, BLUE, "c"),
            (drv, 10, None, INK, "l")]
    for c, (t, sz, b, col, al) in enumerate(vals):
        cell = tbl.cell(r, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor.from_string(band)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Inches(0.08); cell.margin_right = Inches(0.04)
        cell.margin_top = Inches(0.01); cell.margin_bottom = Inches(0.01)
        put_text(cell.text_frame, [{"runs": [(t, sz, b, col)], "align": al}], anchor="m", ml=0.02, mr=0.02, mt=0.0, mb=0.0)
bar = add_box(s, 0.50, 5.62, 9.83, 0.52, fill=NAVY)
put_text(bar.text_frame,
         [{"runs": [("広告費を1円も増やさずに、実来場者数 月48名 → 80名（1.66倍）・来場単価40%削減を実現する", 12.5, True, WHITE)], "align": "c"}],
         anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
add_text(s, 0.50, 6.26, 9.83, 0.70,
         [{"runs": [("※現状値は貴社ヒアリングに基づく実情数値、提案後はLINE×Meta連携運用に基づく試算目標値です。", 9.5, None, MUT)], "sa": 2},
          {"runs": [("※複数プランにおける詳細な収支シミュレーションは、別添のExcel資料（SIM）をご参照ください。", 9.5, None, MUT)]}], ml=0.0)

# ---------- Slide 9: ロードマップ ----------
s = slides[8]
frame(s, "導入までの進め方（3フェーズ）",
      "9〜11月開催（4会場分）の3ヶ月運用からスタートし、成果検証を踏まえて継続判断が可能")
phases = [
    ("Phase 1｜設計・構築", "（8月）",
     ["・要件定義", "・配信シナリオ作成", "・リッチメニューデザイン", "・アンケート設計"], False),
    ("Phase 2｜テスト・実装", "（8月下旬〜9月）",
     ["・アカウント設定", "・LP導線改修（特典訴求の実装）", "・Meta広告とのデータ連携設定", "・縦型動画クリエイティブ制作"], False),
    ("Phase 3｜実運用・改善", "（9月末〜11月）",
     ["・4会場での実運用（9/26・10月2会場・11月）", "・アナリティクス分析", "・成果検証と改善施策の実行"], True),
]
for i, (ph, mo, items, hl) in enumerate(phases):
    x = 0.50 + i * 3.32
    cv = add_box(s, x, 1.78, 3.19, 0.72, fill=NAVY, shape=MSO_SHAPE.CHEVRON)
    put_text(cv.text_frame,
             [{"runs": [(ph, 12, True, WHITE)], "align": "c", "sa": 1},
              {"runs": [(mo, 9.5, None, WHITE)], "align": "c"}],
             anchor="m", ml=0.22, mr=0.12, mt=0.0, mb=0.0)
    card = add_box(s, x, 2.66, 3.19, 1.95,
                   fill=PALE if hl else WHITE,
                   line=NAVY if hl else BORDER, lw=1.5 if hl else 1.0)
    put_text(card.text_frame,
             [{"runs": [(t, 11, None, INK)], "ls": 1.15, "sa": 7} for t in items],
             anchor="t", ml=0.16, mr=0.10, mt=0.14)
add_text(s, 0.50, 4.78, 6.0, 0.32, [{"runs": [("運用のポイント", 12, True, NAVY)]}], ml=0.0)
points = [
    ("契約形態", "9月末〜11月頭開催の4会場分に向けた3ヶ月契約を想定"),
    ("クリエイティブ", "秋開催に向け、Meta広告用の縦型動画クリエイティブを導入・運用"),
    ("チーム体制", "営業担当2名＋運用担当＋運用統括の4名体制で手厚くフォロー"),
]
for i, (lab, body) in enumerate(points):
    x = 0.50 + i * 3.32
    pc = add_box(s, x, 5.14, 3.19, 1.00, fill=PALE)
    put_text(pc.text_frame,
             [{"runs": [(lab, 10.5, True, NAVY)], "sa": 3},
              {"runs": [(body, 10.5, None, INK)], "ls": 1.15}],
             anchor="m", ml=0.14, mr=0.10, mt=0.06, mb=0.06)
bar = add_box(s, 0.50, 6.30, 9.83, 0.48, fill=NAVY)
put_text(bar.text_frame,
         [{"runs": [("秋開催4会場で実運用 → 成果検証 → 次回（冬）開催へ改善を反映", 12, True, WHITE)], "align": "c"}],
         anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)

prs.save(OUT)
print("saved:", OUT)
print("slides:", len(list(Presentation(OUT).slides)))
