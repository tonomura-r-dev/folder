# -*- coding: utf-8 -*-
"""賃貸業界 Ver.B（オーナー開拓版）LINEOA施策提案（36枚）

_templates/DYM_LINEOA_FMT.pptx（業界汎用FMT・36枚）をコピーし、
全36枚を _drafts/賃貸業界_VerB_スライド作成プロンプト.md の指定内容で構築する。

  python _build/build_chintai_verB.py
  python _build/qa_render.py 賃貸業界_LINEOA施策提案_VerB_オーナー開拓.pptx

【この資料の背骨】
オーナーは営業されたくないが、情報は欲しい。しかも検討は1〜3年。
訪問・電話・DMは「その瞬間に検討中の人」しか拾えない。
やる意思はあるが今すぐではない層を育て続ければ、面談獲得の単価（CPA）が下がる。

【実データは取得済み】
- Googleトレンド：_data/trends/souzoku_5year.csv（相続）・verB_5year.csv（確定申告）
  → _images/verB_trend_*.png（実測。土地活用・アパート経営・空室は測定不能と正直に明記）
- LINEヤフー前後検索：取得済み（S08に反映。「アパート経営の落とし穴」が起点より前）
- LINEOA実績：オーナー向け公式事例は0件（2回探索済み）→ 近接事例で代替＋正直に明記

【未取得データ（★破線の差込枠）】
- S09 競合6社（大東建託・東建コーポレーション等）の友だち数

【落とし穴メモ（_build/README.md・deck-apply スキルより）】
- put_text() は必ず reset_tf() を通す（既存テキストへの追記事故を防ぐ）
- スライドの新規追加はしない。ベース36枚を clear_slide() して作り直す
- 一括置換をしない。数値は行・列を特定して書く
"""
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
SRC = str(ROOT / "_templates" / "DYM_LINEOA_FMT.pptx")
OUT = str(ROOT / "賃貸業界_LINEOA施策提案_VerB_オーナー開拓.pptx")
IMG = ROOT / "_images"

# ---- 配色（DYM 3色：ネイビー・炭黒・オレンジ。LINEグリーンはUI部分のみ）----
TNAVY = "002060"   # タイトル
NAVY = "1F285A"    # 打ち手・カード見出し
ORANGE = "ED7D31"  # 強調
RED = "C00000"     # 課題・警告
INK = "333333"     # 本文
MUT = "7F7F7F"
WHITE = "FFFFFF"
PALE = "F4F7FF"    # 淡ネイビー
PORANGE = "FCE4D6"  # 淡オレンジ
GREY = "F2F2F2"
BORDER = "D9D9D9"
GREEN = "06C755"   # LINE UI 専用
BEZEL = "2B2B2B"
SCREEN = "EFF2F7"
PRED = "FDF2F2"

# ---- FMTのレイアウト座標（cm）----
SW, SH = 27.52, 19.05
TITLE_XY = (1.52, 0.38, 24.4, 0.94)
LEAD_XY = (1.20, 1.80, 25.1, 1.90)
DIV_Y = 3.86
CX0, CW = 1.20, 25.12
CY0, CY1 = 4.30, 17.00
FOOT_Y = 17.35

shutil.copyfile(SRC, OUT)
prs = Presentation(OUT)
slides = list(prs.slides)
assert len(slides) == 36, len(slides)
assert (prs.slide_width, prs.slide_height) == (9906000, 6858000)


# ================= helpers（build_chumon_jutaku.py と共通の作法）=================
def set_font(run, size, bold=None, color=INK, name="メイリオ"):
    f = run.font
    f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", name)
    f.color.rgb = RGBColor.from_string(color)


ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
ANCH = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}


def reset_tf(tf):
    """既存の段落・runを全消去（追記事故の防止）"""
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)


def put_text(tf, paras, anchor="t", ml=0.14, mr=0.14, mt=0.06, mb=0.06, wrap=True):
    reset_tf(tf)
    tf.word_wrap = wrap
    tf.margin_left = Cm(ml)
    tf.margin_right = Cm(mr)
    tf.margin_top = Cm(mt)
    tf.margin_bottom = Cm(mb)
    tf.vertical_anchor = ANCH[anchor]
    first = True
    for p in paras:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = ALIGN[p.get("align", "l")]
        if p.get("sa") is not None:
            para.space_after = Pt(p["sa"])
        if p.get("ls") is not None:
            para.line_spacing = p["ls"]
        for item in p["runs"]:
            t, sz, b, c = item
            r = para.add_run()
            r.text = t
            set_font(r, sz, b, c)
    return tf


def T(slide, x, y, w, h, paras, anchor="t", **kw):
    box_ = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    put_text(box_.text_frame, paras, anchor=anchor, **kw)
    return box_


def one(text, sz, b=None, c=INK, align="l", sa=None, ls=None):
    d = {"runs": [(text, sz, b, c)], "align": align}
    if sa is not None:
        d["sa"] = sa
    if ls is not None:
        d["ls"] = ls
    return d


def multi(runs, align="l", sa=None, ls=None):
    """1段落に複数run（色・太さを混在させたいとき）"""
    d = {"runs": runs, "align": align}
    if sa is not None:
        d["sa"] = sa
    if ls is not None:
        d["ls"] = ls
    return d


def box(slide, x, y, w, h, fill=None, line=None, lw=1.0,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06, dash=None):
    sp = slide.shapes.add_shape(shape, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor.from_string(line)
        sp.line.width = Pt(lw)
        if dash:
            sp.line.dash_style = dash
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    reset_tf(sp.text_frame)
    return sp


def card(slide, x, y, w, h, head, body, hcol=NAVY, fill=PALE,
         hsz=11, bsz=9, line=None, anchor="t", ls=1.18):
    sp = box(slide, x, y, w, h, fill=fill, line=line)
    paras = [one(head, hsz, True, hcol, sa=3)]
    for b in (body if isinstance(body, list) else [body]):
        if b:
            paras.append(one(b, bsz, None, INK, ls=ls, sa=1))
    put_text(sp.text_frame, paras, anchor=anchor, ml=0.22, mr=0.18, mt=0.14, mb=0.10)
    return sp


def clear_slide(slide):
    spTree = slide.shapes._spTree
    for el in list(spTree):
        if el.tag.split("}")[-1] in ("sp", "cxnSp", "pic", "graphicFrame", "grpSp"):
            spTree.remove(el)


def frame(slide, title, lead):
    """FMT準拠：タイトル16pt紺（y=0.38）＋リード12pt（y=1.80・2行以内）＋区切り線"""
    clear_slide(slide)
    T(slide, *TITLE_XY, [one(title, 16, True, TNAVY)], anchor="m", ml=0, mr=0)
    T(slide, *LEAD_XY, [one(l, 12, None, INK, ls=1.28) for l in lead],
      anchor="m", ml=0, mr=0)
    ln = slide.shapes.add_connector(1, Cm(0), Cm(DIV_Y), Cm(SW), Cm(DIV_Y))
    ln.line.color.rgb = RGBColor.from_string(BORDER)
    ln.line.width = Pt(1.0)


def foot(slide, text):
    T(slide, CX0, FOOT_Y, CW, 0.9, [one(text, 7.5, None, MUT, ls=1.15)], ml=0, mr=0)


def badge(slide, x, y, w, h, text, fill=PORANGE, col=ORANGE, sz=8.5, dash=None):
    sp = box(slide, x, y, w, h, fill=fill, line=col, lw=1.0, dash=dash)
    put_text(sp.text_frame, [one(text, sz, True, col, align="c")],
             anchor="m", ml=0.06, mr=0.06, mt=0, mb=0)
    return sp


def band(slide, y, text, fill=NAVY, col=WHITE, sz=11.5, h=0.86, x=CX0, w=CW):
    sp = box(slide, x, y, w, h, fill=fill, radius=0.10)
    put_text(sp.text_frame, [one(text, sz, True, col, align="c")],
             anchor="m", ml=0.2, mr=0.2, mt=0, mb=0)
    return sp


def placeholder(slide, x, y, w, h, label, note):
    """★未取得データの差込枠（破線）。捏造しない。"""
    sp = box(slide, x, y, w, h, fill=WHITE, line=MUT, lw=1.25,
              dash=MSO_LINE_DASH_STYLE.DASH)
    put_text(sp.text_frame,
              [one(label, 11, True, MUT, align="c", sa=4),
               one(note, 8.5, None, MUT, align="c", ls=1.25)],
              anchor="m", ml=0.3, mr=0.3, mt=0.1, mb=0.1)
    return sp


def phone(slide, x, y, w, h, header, lines):
    """LINEトーク画面のモック。lines = [(kind, text)]
    kind: 'in'（BOT吹き出し）/ 'chip'（選択肢）/ 'btn'（CTA）/ 'note'（注記）"""
    box(slide, x, y, w, h, fill=BEZEL, radius=0.10)
    box(slide, x + 0.14, y + 0.42, w - 0.28, h - 0.56, fill=SCREEN, radius=0.02,
        shape=MSO_SHAPE.RECTANGLE)
    hd = box(slide, x + 0.14, y + 0.42, w - 0.28, 0.52, fill=GREEN,
             shape=MSO_SHAPE.RECTANGLE)
    put_text(hd.text_frame, [one(header, 7, True, WHITE, align="c")],
             anchor="m", ml=0.05, mr=0.05, mt=0, mb=0)
    cy = y + 1.05
    iw = w - 0.62
    for kind, text in lines:
        nlines = text.count("\n") + 1
        bh = 0.30 + 0.30 * nlines
        if kind == "in":
            sp = box(slide, x + 0.31, cy, iw, bh, fill=WHITE, line=BORDER, lw=0.75)
            put_text(sp.text_frame, [one(l, 6.5, None, INK, ls=1.18)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.14, mr=0.10, mt=0.05, mb=0.05)
        elif kind == "chip":
            sp = box(slide, x + 0.31, cy, iw, bh, fill=WHITE, line=GREEN, lw=0.9)
            put_text(sp.text_frame, [one(l, 6.5, None, "0B7A3B", align="c", ls=1.15)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.08, mr=0.08, mt=0.03, mb=0.03)
        elif kind == "btn":
            sp = box(slide, x + 0.31, cy, iw, bh, fill=GREEN)
            put_text(sp.text_frame, [one(l, 6.8, True, WHITE, align="c", ls=1.15)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.08, mr=0.08, mt=0.03, mb=0.03)
        else:  # note
            sp = slide.shapes.add_textbox(Cm(x + 0.31), Cm(cy), Cm(iw), Cm(bh))
            put_text(sp.text_frame, [one(l, 6.2, None, MUT, ls=1.15)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.05, mr=0.05, mt=0, mb=0)
        cy += bh + 0.12
    return cy


def richmenu(slide, x, y, w, h, tabs):
    """リッチメニュー3タブ×6ボタンのモック。
    tabs = [(タブ名, [ボタン文言6個]), ...]。最初のタブをアクティブ表示。"""
    tab_h = 0.55
    box(slide, x, y, w, h, fill=BEZEL, radius=0.06)
    tw = (w - 0.12) / len(tabs)
    for i, (name, _) in enumerate(tabs):
        active = (i == 0)
        tb = box(slide, x + 0.06 + i * tw, y + 0.06, tw - 0.03, tab_h,
                  fill=GREEN if active else "3F3F3F",
                  shape=MSO_SHAPE.RECTANGLE)
        put_text(tb.text_frame, [one(name, 8, True, WHITE, align="c")],
                 anchor="m", ml=0, mr=0, mt=0, mb=0)
    _, btns = tabs[0]
    gx, gy = x + 0.06, y + 0.06 + tab_h + 0.06
    gw, gh = w - 0.12, h - tab_h - 0.18
    cols, rows = 3, 2
    bw, bh = (gw - 0.06 * (cols - 1)) / cols, (gh - 0.06 * (rows - 1)) / rows
    for i, label in enumerate(btns):
        r, c = divmod(i, cols)
        bb = box(slide, gx + c * (bw + 0.06), gy + r * (bh + 0.06), bw, bh,
                  fill=SCREEN, line=BORDER, lw=0.75, radius=0.04)
        put_text(bb.text_frame, [one(l, 7.2, True, NAVY, align="c", ls=1.1)
                                 for l in label.split("\n")],
                 anchor="m", ml=0.05, mr=0.05, mt=0, mb=0)


def simple_table(slide, x, y, w, h, headers, rows, col_w=None,
                  hsz=9, bsz=8.5, header_fill=NAVY, zebra=PALE,
                  align=None, row_h=None):
    """罫線テーブル（ネイティブpptxテーブル）。headers=[str], rows=[[str]]"""
    n_r, n_c = len(rows) + 1, len(headers)
    gf = slide.shapes.add_table(n_r, n_c, Cm(x), Cm(y), Cm(w), Cm(h))
    tbl = gf.table
    if col_w:
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Cm(cw)
    if row_h:
        for r in tbl.rows:
            r.height = Cm(row_h)
    for j, htext in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor.from_string(header_fill)
        c.margin_left = c.margin_right = Cm(0.12)
        c.margin_top = c.margin_bottom = Cm(0.05)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        put_text(c.text_frame, [one(htext, hsz, True, WHITE,
                                    align=(align[j] if align else "c"))],
                 anchor="m", ml=0.12, mr=0.12, mt=0.02, mb=0.02)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor.from_string(
                zebra if (zebra and i % 2 == 0) else WHITE)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            put_text(c.text_frame, [one(str(val), bsz, None, INK,
                                        align=(align[j] if align else "l"))],
                     anchor="m", ml=0.12, mr=0.12, mt=0.02, mb=0.02)
    return gf


def find_shape(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle in sh.text_frame.text:
            return sh
    return None


# ============================================================
# S01 表紙
# ============================================================
s = slides[0]
sh = find_shape(s, "業界：")
if sh:
    put_text(sh.text_frame, [one("業界：賃貸業界（オーナー開拓・管理受託主）", 12, True, TNAVY)], anchor="m")
sh = find_shape(s, "最上位パートナーの知見")
if sh:
    put_text(sh.text_frame,
             [one("やる意思はあるが今すぐではないオーナーを、時期が来るまで育て続ける", 13, True, INK)],
             anchor="m")

# ============================================================
# S02 資料アジェンダ
# ============================================================
s = slides[1]
frame(s, "資料アジェンダ", ["この提案は「オーナーを増やす」話ではなく、「やる意思があるオーナーを、",
                          "時期が来るまで育て続ける」話である。"])
AGENDA = [
    ("1", "なぜLINEか", "S03-06"), ("2", "実態調査", "S07-10"),
    ("3", "全体設計", "S11-15"), ("4", "構築", "S16-20"),
    ("5", "配信設計", "S21-26"), ("6", "歩留まり・工数", "S27-29"),
    ("7", "成果と体制", "S30-33"), ("8", "締め", "S34-36"),
]
gx, gy, gw, gh = CX0, CY0 + 0.3, CW, 11.6
cw = gw / 4
ch = gh / 2
for i, (n, t, r) in enumerate(AGENDA):
    row, col = divmod(i, 4)
    x = gx + col * cw
    y = gy + row * ch
    box(s, x + 0.12, y + 0.12, cw - 0.24, ch - 0.24, fill=PALE)
    T(s, x + 0.12, y + 0.22, cw - 0.24, 0.9, [one(n, 22, True, ORANGE, align="c")], anchor="t")
    T(s, x + 0.12, y + 0.95, cw - 0.24, 1.3, [one(t, 12.5, True, NAVY, align="c")], anchor="t")
    T(s, x + 0.12, y + ch - 0.55, cw - 0.24, 0.5, [one(r, 8.5, None, MUT, align="c")], anchor="t")
foot(s, "全8章・36枚｜主語は「LINE公式アカウントの運用」。広告運用の提案書ではない")

# ============================================================
# S03 検討背景
# ============================================================
s = slides[2]
frame(s, "検討背景",
      ["メールは開かれない、電話は出ない、訪問は嫌われる。残るのがLINE。",
       "※「オーナーもLINEを使う」と積極的に言い切れる根拠はない。消去法で書く。"])
nums = [("LINE国内MAU", "1億人突破", "2025年12月末時点"),
        ("メッセージ開封率", "約55%", "メールは約20%｜2022年6月時点")]
cw2 = (CW - 0.4) / 2
for i, (label, val, sub) in enumerate(nums):
    x = CX0 + i * (cw2 + 0.4)
    box(s, x, CY0 + 0.2, cw2, 2.6, fill=PALE)
    T(s, x + 0.3, CY0 + 0.35, cw2 - 0.6, 0.6, [one(label, 12, True, NAVY)], anchor="m")
    T(s, x + 0.3, CY0 + 1.0, cw2 - 0.6, 1.1, [one(val, 26, True, ORANGE)], anchor="m")
    T(s, x + 0.3, CY0 + 2.05, cw2 - 0.6, 0.6, [one(sub, 10, None, INK)], anchor="m")
box(s, CX0, CY0 + 3.1, CW, 2.3, fill=PRED, line=RED, lw=1.0)
T(s, CX0 + 0.3, CY0 + 3.25, CW - 0.6, 2.1,
  [one("★要注意：不動産オーナーが経営に活用したいツールは「メール」が首位（GMO調査・第三者）", 11, True, RED, ls=1.3, sa=4),
   one("根拠が出せないので「オーナーもLINEを使う」と言い切らない。メールは開かれない／電話は出ない／訪問は嫌われる、という消去法の論で構成する。", 10.5, None, INK, ls=1.3)],
  anchor="m")
foot(s, "業界水準｜出典：LINEヤフー公式（国内MAU・開封率）。GMO調査は第三者・要一次確認")

# ============================================================
# S04 管理受託のWebマーケ構造
# ============================================================
s = slides[3]
frame(s, "管理受託のWebマーケ構造",
      ["今すぐ検討中の人しか拾っていない。意思はあるが時期が来ていない層を毎期捨てている。"])
box(s, CX0, CY0 + 0.3, CW, 2.6, fill=NAVY)
T(s, CX0 + 0.4, CY0 + 0.5, CW - 0.8, 2.2,
  [one("訪問・電話・DMは「その瞬間に検討中の人」しか拾えない。", 14, True, WHITE, ls=1.35, sa=6),
   one("検討期間は数ヶ月〜3年。今すぐでない層は、次にオーナーが動くまで接点がゼロになる。", 12, None, "E9EDF7", ls=1.35)],
  anchor="m")
pts = [
    ("検討期間が長い", "数ヶ月〜1年（管理）\n1〜3年（建築請負）"),
    ("受注が年に数件〜数十件", "分母が小さく\n月次で動かない"),
    ("今すぐでない層は毎期リセット", "接点が無いまま\n他社に流れる"),
]
cw3 = (CW - 0.6) / 3
for i, (h, b) in enumerate(pts):
    x = CX0 + i * (cw3 + 0.3)
    card(s, x, CY0 + 3.2, cw3, 2.4, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=11.5, bsz=9.5, ls=1.3)
foot(s, "意見｜DYM提案の構造整理")

# ============================================================
# S05 オーナーの心理
# ============================================================
s = slides[4]
frame(s, "オーナーの心理",
      ["今の管理会社に不満はあるが、替えるのが面倒。",
       "営業が来ると断れないから、そもそも会いたくない。"])
psy = [
    ("不満はあるが\n動くのが面倒", "空室・家賃下落・滞納・修繕費への\n不満はあっても、乗り換え手続きの\n手間で動かない"),
    ("営業が来ると\n断れない", "訪問・電話営業を受けると\n断りきれず、そもそも\n会うこと自体を避ける"),
    ("情報は\n欲しい", "税金・相続・修繕の情報は\n知りたい。営業されずに\n情報だけ欲しい"),
]
cw3 = (CW - 0.6) / 3
for i, (h, b) in enumerate(psy):
    x = CX0 + i * (cw3 + 0.3)
    card(s, x, CY0 + 0.2, cw3, 3.8, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=13, bsz=10, ls=1.35)
box(s, CX0, CY0 + 4.4, CW, 1.6, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 4.4, CW - 0.6, 1.6,
  [one("「営業されたくないが、情報は欲しい」を満たせるのがLINE", 13, True, WHITE, align="c", ls=1.3)],
  anchor="m")
foot(s, "意見｜DYM仮説。上司確認で裏付けを取ること")

# ============================================================
# S06 できること／できないこと ＋ スコープ宣言
# ============================================================
s = slides[5]
frame(s, "Webマーケで「できること／できないこと」",
      ["管理の実務品質は変えられない。変えられるのは接点の作り方と育て方だけ。"])
box(s, CX0, CY0 + 0.2, 8.0, 6.0, fill=GREY)
T(s, CX0 + 0.3, CY0 + 0.4, 7.4, 0.6, [one("できないこと", 13, True, MUT)], anchor="m")
for i, t in enumerate(["管理の実務品質（対応の速さ・工事の質）", "管理料率・保証条件", "営業担当の提案力そのもの"]):
    T(s, CX0 + 0.3, CY0 + 1.1 + i * 1.5, 7.4, 1.3, [one("✕ " + t, 11, None, INK, ls=1.25)], anchor="m")
bx2 = CX0 + 8.3
box(s, bx2, CY0 + 0.2, CW - 8.3, 6.0, fill=PALE)
T(s, bx2 + 0.3, CY0 + 0.4, CW - 8.3 - 0.6, 0.6, [one("できること（＝本提案の範囲）", 13, True, NAVY)], anchor="m")
can = [
    "訪問も電話もせずに、痛みを可視化する接点をつくる（診断）",
    "今すぐでない層を捨てずに育て続ける（モーメント配信）",
    "相談・セミナーの予約と当日実施率を上げる",
    "既存オーナーに言われる前に出す（空室アラート・修繕提案）",
    "訪問・架電を減らして、営業を増員せずに面談数を増やす",
]
for i, t in enumerate(can):
    T(s, bx2 + 0.3, CY0 + 1.0 + i * 1.0, CW - 8.3 - 0.6, 0.95, [one(f"{i+1}. {t}", 10, None, INK, ls=1.2)], anchor="m")
box(s, CX0, CY0 + 6.4, CW, 1.0, fill=PORANGE)
T(s, CX0 + 0.3, CY0 + 6.4, CW - 0.6, 1.0,
  [one("「寝かせる」と言わない。「やる意思はあるが今すぐではない層を育てる」で統一", 11, True, ORANGE, align="c")],
  anchor="m")
foot(s, "意見｜DYMの提案スコープ宣言")

# ============================================================
# S07 モーメント分析（★実データ：Googleトレンド）
# ============================================================
s = slides[6]
frame(s, "実態調査｜モーメント分析（★重点）",
      ["オーナーが動くのは季節ではなくイベント。相続が最大のトリガー。",
       "「土地活用」「アパート経営」はGoogleトレンドで測定不能なほど検索Volが小さい。"])
pic_w = (CW - 0.4) / 2
s.shapes.add_picture(str(IMG / "verB_trend_souzoku5y.png"), Cm(CX0), Cm(CY0 + 0.2), width=Cm(pic_w))
s.shapes.add_picture(str(IMG / "verB_trend_kakutei5y.png"), Cm(CX0 + pic_w + 0.4), Cm(CY0 + 0.2), width=Cm(pic_w))
box(s, CX0, CY0 + 7.0, CW, 1.0, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 7.0, CW - 0.6, 1.0,
  [one("前後検索で「相続 → 土地活用」の順序が確認できている（S08）。相続がトリガー", 12, True, WHITE, align="c", ls=1.25)],
  anchor="m")
foot(s, "業界水準｜出典：Googleトレンド（2026年8月取得・実データ）。「土地活用」「アパート経営」「空室」は本グラフでは測定不能")

# ============================================================
# S08 前後検索
# ============================================================
s = slides[7]
frame(s, "実態調査｜前後検索（★重点）",
      ["前は相続と税金、後は他社の口コミ。「アパート経営の落とし穴」が起点より前＝不安が先。"])
bands3 = [
    ("検索“前”15〜7日", "相続がトリガー\n（遺産相続／贈与税／\n相続登記義務化）＋\n★アパート経営の落とし穴", MUT),
    ("起点0日", "大東建託／土地活用相談／\n土地活用 ランキング", NAVY),
    ("検索“後”7〜15日", "他社評判＋自宅建築へ流出\n（住宅ローン／住宅展示場／\n大東建託 口コミ）", RED),
]
cw3b = (CW - 0.6) / 3
for i, (h, sub, c) in enumerate(bands3):
    x = CX0 + i * (cw3b + 0.3)
    bb = box(s, x, CY0 + 0.2, cw3b, 1.1, fill=c)
    put_text(bb.text_frame, [one(h, 11.5, True, WHITE, align="c")], anchor="m")
    T(s, x, CY0 + 1.4, cw3b, 2.6, [one(l, 9.5, None, INK, align="c", ls=1.3) for l in sub.split("\n")], anchor="t")
box(s, CX0, CY0 + 4.3, CW, 2.6, fill=PRED, line=RED, lw=1.0)
T(s, CX0 + 0.3, CY0 + 4.45, CW - 0.6, 2.4,
  [one("★「アパート経営の落とし穴」が起点より前にある＝不安のほうが先に来ている", 12, True, RED, ls=1.3, sa=4),
   one("オーナーは「アパート」単体ではなく、駐車場・トランクルーム・自動販売機・レンタルスペース・民泊・ガレージ・太陽光も並行して調べている。同時に市街化調整区域・農地転用・公図も検索＝そもそも建てられるかを気にしている（S18の診断設計に反映）。", 10, None, INK, ls=1.3)],
  anchor="t")
foot(s, "事実｜出典：LINEヤフー媒体資料（取得済み・反映済み）。起点＝「土地活用」")

# ============================================================
# S09 他社分析①
# ============================================================
s = slides[8]
frame(s, "他社分析①｜競合のLINE運用ステータス",
      ["友だち数は page.line.me から取得日つきで実測する。憶測で書かない。"])
comps = ["大東建託", "東建コーポレーション", "生和コーポレーション", "シノケン", "積水ハウス不動産", "大和ハウス"]
rows = [[c, "★未取得", "★未取得"] for c in comps]
simple_table(s, CX0, CY0 + 0.2, CW, 3.3,
             ["アカウント", "友だち数", "オーナー向け別アカウントの有無"],
             rows, col_w=[CW * 0.36, CW * 0.28, CW * 0.36], row_h=0.5)
placeholder(s, CX0, CY0 + 3.8, CW, 3.1, "★差込枠｜友だち数・配信内容の実測（未取得）",
            "page.line.me から取得日つきで取得。オーナー向け別アカウントの有無も記録すること")
foot(s, "選定は知名度優先（スキル §3 で事前承認済み）")

# ============================================================
# S10 他社分析②
# ============================================================
s = slides[9]
frame(s, "他社分析②｜運用の型4分類",
      ["オーナー向けLINE運用は4型。診断型はほぼ空いている。"])
types = [
    ("診断ドリブン型", "診断で属性取得→\n検討段階別に出し分け", "★ほぼ空いている"),
    ("相談カウンター型", "有人相談・店舗送客が主。\n配信していないことが多い", "競合多数"),
    ("特典・会員型", "限定特典で\n友だち化", "一部あり"),
    ("カタログ型", "物件・事例一覧・\n予約が中心", "競合多数"),
]
cw4 = (CW - 0.3 * 3) / 4
for i, (h, b, tag) in enumerate(types):
    x = CX0 + i * (cw4 + 0.3)
    highlight = "空いている" in tag
    card(s, x, CY0 + 0.3, cw4, 3.6, h, [b, tag], hcol=(ORANGE if highlight else NAVY),
         fill=(PORANGE if highlight else PALE), hsz=11, bsz=9, ls=1.25)
box(s, CX0, CY0 + 4.2, CW, 1.6, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 4.2, CW - 0.6, 1.6,
  [one("診断ドリブン型がほぼ空いている＝S18（土地活用タイプ診断）の差別化余地", 12.5, True, WHITE, align="c", ls=1.3)],
  anchor="m")
foot(s, "★未取得｜4分類は一般的な運用パターン。各社の実際の型は実測で確認（要取得）")


# ============================================================
# S11 カスタマージャーニー
# ============================================================
s = slides[10]
frame(s, "全体設計｜カスタマージャーニー",
      ["8フェーズのうち接点が無いのは「意思はあるが今すぐでない」期間そのもの。",
       "ここが1〜3年ある。"])
PH = ["①接触\n（認知・流入）", "②離脱\n（対策）", "③育成\n（興味喚起）", "④リード獲得\n（web/LINE）",
      "⑤リード有効化", "⑥有効リードの\n再育成", "⑦マネタイズ\n（管理受託）", "⑧LTV最大化\n・紹介"]
n = len(PH)
cwp = CW / n
for i, label in enumerate(PH):
    x = CX0 + i * cwp
    gap = (i in (2, 3, 4, 5))
    bb = box(s, x + 0.08, CY0 + 0.4, cwp - 0.16, 2.2, fill=(PRED if gap else PALE), line=(RED if gap else None), lw=1.25)
    put_text(bb.text_frame, [one(l, 8.5, True, (RED if gap else NAVY), align="c", ls=1.15) for l in label.split("\n")], anchor="m")
box(s, CX0 + 2 * cwp, CY0 + 2.9, 4 * cwp, 0.7, fill=RED)
T(s, CX0 + 2 * cwp, CY0 + 2.9, 4 * cwp, 0.7, [one("★接点が無い区間＝1〜3年ある", 10, True, WHITE, align="c")], anchor="m")
box(s, CX0, CY0 + 4.1, CW, 1.6, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 4.1, CW - 0.6, 1.6,
  [one("この長い空白区間をLINEのモーメント配信で埋める。DM・チラシ・訪問には手を入れない。", 12.5, True, WHITE, align="c", ls=1.3)],
  anchor="m")
foot(s, "意見｜DYM標準のカスタマージャーニーフレームにオーナー開拓を当てはめ")

# ============================================================
# S12 施策全体像
# ============================================================
s = slides[11]
frame(s, "全体設計｜施策全体像",
      ["DM・チラシ・訪問には手を入れない。その後ろの区間をLINEで埋める。"])
flow = ["DM・チラシ\n・訪問", "→", "★LINE友だち化\n（診断で入口）", "→", "育成\n（1〜3年）", "→", "相談\n・セミナー", "→", "管理受託\n・請負"]
mw = CW / len(flow)
for i, f in enumerate(flow):
    x = CX0 + i * mw
    if f == "→":
        T(s, x, CY0 + 1.6, mw, 1.2, [one(f, 16, True, ORANGE, align="c")], anchor="m")
    else:
        highlight = "LINE" in f
        bb = box(s, x + 0.1, CY0 + 1.3, mw - 0.2, 1.8, fill=(PORANGE if highlight else PALE),
                  line=(ORANGE if highlight else None), lw=1.25)
        put_text(bb.text_frame, [one(l, 9, True, (ORANGE if highlight else NAVY), align="c", ls=1.15) for l in f.split("\n")], anchor="m")
box(s, CX0, CY0 + 3.6, CW, 1.8, fill=PALE)
T(s, CX0 + 0.3, CY0 + 3.6, CW - 0.6, 1.8,
  [one("触るのはLINE友だち化から相談・セミナーまでの区間だけ。DM・チラシ・訪問の設計は変えない。", 12.5, True, TNAVY, ls=1.3)],
  anchor="m")
foot(s, "意見｜DYM提案の施策全体像")

# ============================================================
# S13 効果を最大化する2軸
# ============================================================
s = slides[12]
frame(s, "全体設計｜効果を最大化する2軸",
      ["新規オーナーの追加動線と、既存オーナーの育成・CV動線を別々に設計する。"])
axes = [
    ("① 新規オーナー追加動線", "DM・チラシ・セミナーからの\n友だち化を最大化", NAVY, PALE),
    ("② 既存オーナー育成・CV動線", "すでに友だちのオーナーを\n相談・セミナーへ引き上げ", ORANGE, PORANGE),
]
cw2 = (CW - 0.4) / 2
for i, (h, b, c, fill) in enumerate(axes):
    x = CX0 + i * (cw2 + 0.4)
    box(s, x, CY0 + 0.2, cw2, 3.0, fill=fill, line=c, lw=1.25)
    T(s, x + 0.3, CY0 + 0.4, cw2 - 0.6, 0.7, [one(h, 13, True, c)], anchor="m")
    T(s, x + 0.3, CY0 + 1.2, cw2 - 0.6, 1.6, [one(l, 10.5, None, INK, ls=1.35) for l in b.split("\n")], anchor="t")
box(s, CX0, CY0 + 3.5, CW, 1.6, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 3.5, CW - 0.6, 1.6,
  [one("2軸を分けないと、新規獲得の忙しさで既存の育成が止まる。別トラックで運用する", 12, True, WHITE, align="c", ls=1.3)],
  anchor="m")
foot(s, "意見｜DYM提案の設計方針")

# ============================================================
# S14 対策領域マップ
# ============================================================
s = slides[13]
frame(s, "全体設計｜対策領域マップ",
      ["新規開拓と既存維持の2領域。既存側がLTVを作る。"])
areas = [
    ("① 新規開拓領域", "友だち化率・診断完了率", [
        "DM・チラシのQRから友だち化",
        "セミナー申込からの友だち化",
        "土地活用タイプ診断で入口を作る",
    ], NAVY, PALE),
    ("② 既存維持領域", "相談・セミナー予約CVR・LTV", [
        "既存オーナーの一斉登録（最速で貯まる）",
        "年次モーメント配信で関係維持",
        "空室アラート・修繕提案で先回り",
    ], ORANGE, PORANGE),
]
cw2 = (CW - 0.4) / 2
for i, (h, sub, items, c, fill) in enumerate(areas):
    x = CX0 + i * (cw2 + 0.4)
    box(s, x, CY0 + 0.2, cw2, 6.2, fill=fill)
    T(s, x + 0.3, CY0 + 0.4, cw2 - 0.6, 0.7, [one(h, 15, True, c)], anchor="m")
    T(s, x + 0.3, CY0 + 1.15, cw2 - 0.6, 0.7, [one(sub, 10.5, None, MUT)], anchor="m")
    for j, it in enumerate(items):
        T(s, x + 0.3, CY0 + 2.1 + j * 1.5, cw2 - 0.6, 1.4, [one(f"・{it}", 11, None, INK, ls=1.3)], anchor="t")
foot(s, "意見｜②既存維持領域がLTVを作る（月次フィーで長期化）")

# ============================================================
# S15 施策展開図（初期・月次）
# ============================================================
s = slides[14]
frame(s, "全体設計｜施策展開図（初期・月次）",
      ["初期構築は1ヶ月、以降は月次。配信の中身は年間で決め打ちできる。"])
band(s, CY0 + 0.2, "初期（初動・構築）", fill=NAVY, h=0.75)
init_items = ["あいさつメッセージ設計（訪問しません宣言）", "土地活用タイプ診断設計（4問）",
              "リッチメニュー3タブ18ボタン", "年次モーメントカレンダー設計", "初動30日ステップの設計"]
for i, it in enumerate(init_items):
    card(s, CX0 + (i % 3) * (CW / 3), CY0 + 1.2 + (i // 3) * 1.9, CW / 3 - 0.2, 1.7,
         f"{i+1}", it, hcol=NAVY, fill=PALE, hsz=13, bsz=10, ls=1.25)
band(s, CY0 + 5.0, "月次（定例運用）", fill=ORANGE, col=WHITE, h=0.75)
month_items = ["企画投稿（今月の数字を出し続ける）", "モーメント配信（確定申告・相続等）",
               "通知メッセージ運用（空室・修繕）", "定例レポート・改善提案"]
for i, it in enumerate(month_items):
    card(s, CX0 + (i % 4) * (CW / 4), CY0 + 6.0, CW / 4 - 0.2, 1.4,
         f"{i+1}", it, hcol=ORANGE, fill=PORANGE, hsz=12, bsz=9, ls=1.2)
foot(s, "意見｜配信の中身は年間で決め打ちできる（S21年次モーメントカレンダー）")

# ============================================================
# S16 友だち追加動線
# ============================================================
s = slides[15]
frame(s, "構築｜友だち追加動線",
      ["DM・チラシのQR／セミナー申込／既存オーナーの一斉登録。既存が最速で貯まる。"])
routes = [
    ("① 既存オーナー一斉登録", "管理受託中のオーナーへ\n案内し一斉に友だち化"),
    ("② DM・チラシQR", "既存の販促物に\nQRを追加するだけ"),
    ("③ セミナー申込", "セミナー申込フォームから\n友だち化を促す"),
    ("④ 土地活用タイプ診断", "「まず診断する」を\n入口にした友だち化"),
]
cw4 = (CW - 0.3 * 3) / 4
for i, (h, b) in enumerate(routes):
    x = CX0 + i * (cw4 + 0.3)
    card(s, x, CY0 + 0.3, cw4, 4.6, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=11, bsz=9.5, ls=1.3)
box(s, CX0, CY0 + 5.2, CW, 1.6, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 5.2, CW - 0.6, 1.6,
  [one("①既存オーナーの一斉登録が最速で貯まる。ここから始める。", 12.5, True, WHITE, align="c", ls=1.25)],
  anchor="m")
foot(s, "意見｜DYM提案の友だち追加動線設計")

# ============================================================
# S17 あいさつメッセージ
# ============================================================
s = slides[16]
frame(s, "構築｜あいさつメッセージ",
      ["「訪問はいたしません」を最初に書く。それが警戒を解く。"])
phone(s, CX0 + 0.3, CY0 + 0.2, 8.6, 6.6, "●●管理",
      [("in", "友だち追加ありがとうございます。\n担当の佐藤です。"),
       ("in", "✅ こちらから訪問はいたしません\n✅ 営業のお電話もいたしません"),
       ("in", "土地活用・空室でお困りの際に\n役立つ情報だけお届けします。"),
       ("chip", "まず土地活用タイプを診断する"),
       ("note", "※ 訪問・お電話はいたしません")])
box(s, CX0 + 9.6, CY0 + 0.2, CW - 9.6, 6.6, fill=PALE)
T(s, CX0 + 9.9, CY0 + 0.4, CW - 9.9 - 0.3, 0.6, [one("設計の要点", 12.5, True, NAVY)], anchor="m")
pts = [
    ("「訪問はいたしません」を最初に", "S05の心理（会いたくない）\nへ直接答える"),
    ("営業色を出さない", "情報提供に徹する。\n提案は診断結果から入る"),
    ("診断への誘導で終える", "低ハードルの入口として\n土地活用タイプ診断へ"),
]
for i, (h, b) in enumerate(pts):
    T(s, CX0 + 9.9, CY0 + 1.2 + i * 1.8, CW - 9.9 - 0.3, 1.7,
      [one(f"・{h}", 11, True, NAVY, ls=1.25, sa=2)] +
      [one(l, 9.5, None, INK, ls=1.2) for l in b.split("\n")], anchor="t")
foot(s, "意見｜配信文面は初稿")

# ============================================================
# S18 土地活用タイプ診断（★設計変更あり）
# ============================================================
s = slides[17]
frame(s, "構築｜土地活用タイプ診断（★重点・設計変更あり）",
      ["「アパートを建てますか」ではなく「この土地に何が向くか」から入る。",
       "順番が重要。Q1で「そもそも建てられるか」を最初に判定する。"])
Q = [
    ("Q1（最初に判定）そもそも建てられるか", "市街化区域／市街化調整区域／農地／わからない"),
    ("Q2 土地の状況", "更地／アパートが建っている／自宅の隣／相続で取得した"),
    ("Q3 いま困っていること", "空室／家賃が下がっている／修繕費／税金／まだ何も"),
    ("Q4 検討時期", "1年以内／1〜3年／まだ先／情報だけ"),
]
cw2 = (CW - 0.4) / 2
for i, (h, opt) in enumerate(Q):
    r, c = divmod(i, 2)
    x = CX0 + c * (cw2 + 0.4)
    y = CY0 + 0.2 + r * 2.1
    card(s, x, y, cw2, 1.9, h, [opt], hcol=NAVY, fill=PALE, hsz=10.5, bsz=9, ls=1.25)
box(s, CX0, CY0 + 4.5, CW, 1.5, fill=PORANGE)
T(s, CX0 + 0.3, CY0 + 4.5, CW - 0.6, 1.5,
  [one("選択肢には駐車場・トランクルーム・レンタルスペース・太陽光も並べる（S08の検索実態と一致）", 10.5, True, RED, ls=1.3)],
  anchor="m")
box(s, CX0, CY0 + 6.1, CW, 1.0, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 6.1, CW - 0.6, 1.0,
  [one("診断結果は 共感 → 実例 → 収支 → CTA の順。収支を先に出さない", 11.5, True, WHITE, align="c", ls=1.2)],
  anchor="m")
foot(s, "意見｜DYM提案の診断設計。S08前後検索データに基づき順番を変更")

# ============================================================
# S19 リッチメニュー
# ============================================================
s = slides[18]
frame(s, "構築｜リッチメニュー",
      ["18ボタンすべてに役割がある。左上は「まず診断する」。"])
richmenu(s, CX0 + 1.5, CY0 + 0.3, 12.0, 6.5, [
    ("まず診断", ["土地活用\nタイプ診断", "空室・家賃\n査定", "建てられるか\n調べる",
              "活用事例を\n見る", "よくある\n失敗", "相談する"]),
    ("収支・税金", ["収支を\n試算する", "確定申告の\n備え", "相続税対策",
               "固定資産税\nの見方", "修繕費の\n相場", "FPに相談"]),
    ("事例・セミナー", ["セミナーに\n申し込む", "オーナーの声", "エリア相場",
                 "法改正の\nまとめ", "オンライン\n相談", "資料を\n受け取る"]),
])
T(s, CX0 + 14.2, CY0 + 0.5, CW - 14.2, 6.0,
  [one("設計の要点", 12.5, True, NAVY, sa=6),
   one("・左上＝「まず診断する」固定", 10.5, None, INK, ls=1.3, sa=4),
   one("・収支・税金＝S07のモーメント\n　（確定申告・相続）に対応", 10.5, None, INK, ls=1.3, sa=4),
   one("・事例・セミナー＝1対1が重い\n　オーナーへの1対多の入口", 10.5, None, INK, ls=1.3, sa=4),
   one("・タブ切替で3倍の情報量を\n　1画面に収める", 10.5, None, INK, ls=1.3)],
  anchor="t")
foot(s, "意見｜DYM提案のリッチメニュー設計")

# ============================================================
# S20 長期育成の設計思想
# ============================================================
s = slides[19]
frame(s, "構築｜長期育成の設計思想",
      ["追いかけずに育て続けると決める。それがこの業種の正解。"])
box(s, CX0, CY0 + 0.2, CW, 2.2, fill=NAVY)
T(s, CX0 + 0.4, CY0 + 0.4, CW - 0.8, 1.8,
  [one("検討期間は数ヶ月〜3年。短期の追客では届かない。", 13, True, WHITE, ls=1.35, sa=4),
   one("「今すぐ客」だけを追うのではなく、「今すぐでない客」を捨てずに育て続ける設計に切り替える。", 11.5, None, "E9EDF7", ls=1.35)],
  anchor="m")
pts = [
    ("モーメント配信が主役", "年次（相続・確定申告・修繕周期）。\nステップ配信は初動30日だけ"),
    ("企画投稿で思い出される", "毎月「今月の数字」を出し続ける\n会社が選ばれる"),
    ("急かさない設計", "相談オファーは焦らせず、\n情報提供の延長線上に置く"),
]
cw3 = (CW - 0.6) / 3
for i, (h, b) in enumerate(pts):
    x = CX0 + i * (cw3 + 0.3)
    card(s, x, CY0 + 2.7, cw3, 2.4, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=11.5, bsz=9.5, ls=1.3)
foot(s, "意見｜DYM提案の設計思想")


# ============================================================
# S21 年次モーメントカレンダー
# ============================================================
s = slides[20]
frame(s, "配信設計｜年次モーメントカレンダー（★重点）",
      ["確定申告・相続・修繕周期・更新期。いつ何を出すかを1年分決めておく。"])
CAL = [
    ("1〜3月", "確定申告", "トレンドで33倍の山（実データ）", "経費・減価償却・青色申告"),
    ("4〜6月", "固定資産税", "納付通知のタイミング", "評価額の見方・軽減措置"),
    ("7〜9月", "大規模修繕", "検討期", "修繕積立・工事相場"),
    ("10〜12月", "年末調整・相続", "相続の相談が増える", "相続税対策・生前贈与"),
]
cw4 = (CW - 0.3 * 3) / 4
for i, (m, tag, ref, b) in enumerate(CAL):
    x = CX0 + i * (cw4 + 0.3)
    tb = box(s, x, CY0 + 0.2, cw4, 1.1, fill=NAVY)
    put_text(tb.text_frame, [one(f"{m}｜{tag}", 10.5, True, WHITE, align="c")], anchor="m")
    card(s, x, CY0 + 1.4, cw4, 3.0, ref, [b], hcol=ORANGE, fill=PALE, hsz=9.5, bsz=9, ls=1.3)
box(s, CX0, CY0 + 4.7, CW, 1.3, fill=PORANGE)
T(s, CX0 + 0.3, CY0 + 4.7, CW - 0.6, 1.3,
  [one("「◯月：モーメント（参照データ）」形式で、根拠を1行ずつ添える", 11.5, True, RED, align="c", ls=1.25)],
  anchor="m")
foot(s, "業界水準｜確定申告のピークはGoogleトレンド実データ（S07：3→100・33倍）")

# ============================================================
# S22 企画投稿案（月次）
# ============================================================
s = slides[21]
frame(s, "配信設計｜企画投稿案（月次）",
      ["毎月「今月の数字」を出し続ける会社が思い出される。"])
pts = [
    ("今月の成約家賃", "近隣エリアの\n成約家賃を毎月配信"),
    ("今月の空室率", "エリア別の空室動向を\n毎月レポート"),
    ("法改正・税制のまとめ", "オーナーに関係する\n法改正を月次で解説"),
    ("オーナーの声", "既存オーナーの\n体験談を紹介"),
]
cw4 = (CW - 0.3 * 3) / 4
for i, (h, b) in enumerate(pts):
    x = CX0 + i * (cw4 + 0.3)
    card(s, x, CY0 + 0.3, cw4, 3.0, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=11, bsz=9.5, ls=1.3)
box(s, CX0, CY0 + 3.7, CW, 1.6, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 3.7, CW - 0.6, 1.6,
  [one("「今月の数字」を出し続ける会社が、いざという時に思い出される", 12.5, True, WHITE, align="c", ls=1.25)],
  anchor="m")
foot(s, "意見｜DYM提案の企画投稿設計")

# ============================================================
# S23 初動30日のステップ
# ============================================================
s = slides[22]
frame(s, "配信設計｜初動30日のステップ",
      ["短期集中はここだけ。診断結果→事例→収支→相談オファー。以降は月次に合流。"])
sc = [
    ("① 診断結果", "Day0", "共感→実例→収支のCTA。ここでは収支を先に出さない"),
    ("② 事例紹介", "Day7", "似た属性のオーナー事例を配信"),
    ("③ 収支の目安", "Day14〜21", "「◯◯タイプなら収支はこの水準」を提示"),
    ("④ 相談オファー", "Day30", "「今すぐでなくてOK」の低ハードルな相談オファー"),
]
cw4 = CW / 4
for i, (h, d, b) in enumerate(sc):
    x = CX0 + i * cw4
    tb = box(s, x + 0.1, CY0 + 0.3, cw4 - 0.2, 0.8, fill=(ORANGE if i == 3 else NAVY))
    put_text(tb.text_frame, [one(f"{h}｜{d}", 11, True, WHITE, align="c")], anchor="m")
    card(s, x + 0.1, CY0 + 1.2, cw4 - 0.2, 2.8, "", [b], hcol=NAVY, fill=(PORANGE if i == 3 else PALE), hsz=1, bsz=9.5, ls=1.3, anchor="m")
box(s, CX0, CY0 + 4.4, CW, 1.4, fill=PALE)
T(s, CX0 + 0.3, CY0 + 4.4, CW - 0.6, 1.4,
  [one("30日を過ぎたら年次モーメント配信（S21）に合流する。以降は日数固定にしない", 11.5, True, TNAVY, ls=1.3)],
  anchor="m")
foot(s, "意見｜DYM提案のステップ配信設計")

# ============================================================
# S24 実文面
# ============================================================
s = slides[23]
frame(s, "配信設計｜実文面",
      ["Day0（診断結果）・Day7（事例）・Day30（相談オファー）。実際に届く文面はこれ。"])
phone(s, CX0 + 0.6, CY0 + 0.2, 8.0, 6.7, "Day0",
      [("in", "診断結果です。\nこの土地には「駐車場」\n「トランクルーム」も\n候補に入ります"),
       ("btn", "3つの活用案を見る"),
       ("note", "※ 訪問・お電話はいたしません")])
phone(s, CX0 + 8.9, CY0 + 0.2, 8.0, 6.7, "Day7",
      [("in", "似た土地状況の\nオーナー様の事例です"),
       ("in", "「まさか駐車場が\n向いているとは」という\nお声をいただきました"),
       ("chip", "事例を詳しく見る")])
phone(s, CX0 + 17.2, CY0 + 0.2, 8.0, 6.7, "Day30",
      [("in", "ここまでお読みいただき\nありがとうございます"),
       ("in", "今すぐでなくて構いません。\n一度お話だけでも"),
       ("chip", "相談してみる"),
       ("note", "※ 押し売りしない低ハードルなオファー")])
foot(s, "意見｜配信文面は初稿。絵文字・改行込みで実際の配信に近い形")

# ============================================================
# S25 セミナー送客
# ============================================================
s = slides[24]
frame(s, "配信設計｜セミナー送客",
      ["1対1の商談が重いなら、1対多で始める。"])
pts = [
    ("会場セミナー", "地域のオーナー向け\n少人数セミナー"),
    ("オンラインセミナー", "移動不要で\n参加ハードルを下げる"),
    ("リッチメニューから申込", "「セミナーに申し込む」を\n常設ボタンに"),
]
cw3 = (CW - 0.6) / 3
for i, (h, b) in enumerate(pts):
    x = CX0 + i * (cw3 + 0.3)
    card(s, x, CY0 + 0.3, cw3, 2.6, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=12, bsz=10, ls=1.3)
box(s, CX0, CY0 + 3.2, CW, 1.6, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 3.2, CW - 0.6, 1.6,
  [one("1対1の商談が重いなら、1対多のセミナーで温度を上げてから個別相談へ", 12, True, WHITE, align="c", ls=1.25)],
  anchor="m")
foot(s, "意見｜DYM提案のセミナー送客設計")

# ============================================================
# S26 通知メッセージ
# ============================================================
s = slides[25]
frame(s, "配信設計｜通知メッセージ",
      ["言われる前に出すのが管理会社の信用。空室アラート・修繕提案。"])
notif = [
    ("空室アラート", "管理物件の空室発生を\n即時通知"), ("修繕提案", "修繕時期が近づいたら\n事前に提案"),
    ("確定申告リマインド", "1〜2月に\n事前案内"), ("契約更新案内", "更新期の\n2〜3ヶ月前に通知"),
]
cw4 = (CW - 0.3 * 3) / 4
for i, (h, b) in enumerate(notif):
    x = CX0 + i * (cw4 + 0.3)
    card(s, x, CY0 + 0.3, cw4, 3.2, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=12, bsz=10, ls=1.3)
box(s, CX0, CY0 + 3.9, CW, 1.8, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 3.9, CW - 0.6, 1.8,
  [one("★言われる前に出すのが管理会社の信用", 13, True, WHITE, ls=1.3, sa=4),
   one("空室・修繕を先回りで案内すると、「ちゃんと見てくれている」という信頼につながる。", 10.5, None, "E9EDF7", ls=1.3)],
  anchor="m")
foot(s, "意見｜通知メッセージはAPI連携（Messaging API）が前提。別途費用（S31参照）")

# ============================================================
# S27 相談の歩留まり改善
# ============================================================
s = slides[26]
frame(s, "歩留まり・工数｜相談の歩留まり改善",
      ["前日リマインドとオンライン相談でドタキャンが減る。移動時間も消える。"])
pts = [
    ("前日リマインド", "相談予約の前日に\n自動送信"), ("オンライン相談の選択肢", "訪問しなくても\n相談できる導線"),
    ("当日朝の再通知", "当日の流れを\n再案内"),
]
cw3 = (CW - 0.6) / 3
for i, (h, b) in enumerate(pts):
    x = CX0 + i * (cw3 + 0.3)
    card(s, x, CY0 + 0.3, cw3, 2.6, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=12, bsz=10, ls=1.3)
box(s, CX0, CY0 + 3.2, CW, 1.6, fill=PORANGE)
T(s, CX0 + 0.3, CY0 + 3.2, CW - 0.6, 1.6,
  [one("オンライン相談は移動時間もゼロになる。営業1人あたりの面談可能数が増える", 11.5, True, RED, align="c", ls=1.3)],
  anchor="m")
foot(s, "意見｜DYM提案の歩留まり改善設計")

# ============================================================
# S28 工数削減
# ============================================================
s = slides[27]
frame(s, "歩留まり・工数｜工数削減",
      ["訪問と架電をやめると、営業は増員せずに面談数を増やせる。"])
pts = [
    ("よくある質問の自動応答", "「相続税はどうなる？」等を\n24時間自動応答"),
    ("日程調整の自動化", "相談・セミナーの\n日程調整をLINEで完結"),
    ("配信の自動化", "年次モーメント・企画投稿は\n事前設計で自動配信"),
]
cw3 = (CW - 0.6) / 3
for i, (h, b) in enumerate(pts):
    x = CX0 + i * (cw3 + 0.3)
    card(s, x, CY0 + 0.3, cw3, 2.6, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=11.5, bsz=9.5, ls=1.3)
box(s, CX0, CY0 + 3.2, CW, 1.6, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 3.2, CW - 0.6, 1.6,
  [one("訪問・架電をやめて自動化に置き換えると、営業は面談（本業）に集中できる", 12, True, WHITE, align="c", ls=1.25)],
  anchor="m")
foot(s, "意見｜DYM提案の工数削減設計")

# ============================================================
# S29 改善モデル
# ============================================================
s = slides[28]
frame(s, "歩留まり・工数｜改善モデル",
      ["リード有効化→相談→受託→修繕・追加。各段に打ち手が1対1で対応。"])
model = ["リード\n有効化", "→", "相談\n・セミナー", "→", "管理\n受託", "→", "修繕\n・追加棟"]
mw = CW / len(model)
for i, m in enumerate(model):
    x = CX0 + i * mw
    if m == "→":
        T(s, x, CY0 + 1.0, mw, 1.4, [one(m, 18, True, ORANGE, align="c")], anchor="m")
    else:
        bb = box(s, x + 0.1, CY0 + 0.9, mw - 0.2, 1.6, fill=PALE)
        put_text(bb.text_frame, [one(l, 10, True, NAVY, align="c", ls=1.15) for l in m.split("\n")], anchor="m")
moves = [
    ("リード有効化", "土地活用タイプ診断＋モーメント配信で育成"), ("相談・セミナー", "低ハードルな相談オファー＋オンライン選択肢"),
    ("管理受託", "1件の単価が大きいので受注効果が大きい"), ("修繕・追加棟", "既存オーナーへの先回り提案でLTV拡張"),
]
for i, (h, b) in enumerate(moves):
    x = CX0 + i * (CW / 4)
    card(s, x, CY0 + 3.0, CW / 4 - 0.2, 2.6, h, [b], hcol=NAVY, fill=WHITE, line=BORDER, hsz=10.5, bsz=9, ls=1.25)
foot(s, "意見｜DYM提案の改善モデル")


# ============================================================
# S30 効果測定の設計
# ============================================================
s = slides[29]
frame(s, "成果と体制｜効果測定の設計",
      ["各段で何を測るかを先に決める。主KPI＝面談CPA・CVR。CPOは約束しない。"])
simple_table(s, CX0, CY0 + 0.2, CW, 2.2,
             ["階層", "指標", "測れる時期"],
             [["① 先行", "友だち追加数／診断完了率", "1〜2ヶ月目"],
              ["② 主KPI", "相談・セミナー予約数／予約CPA／友だち→予約CVR", "3〜6ヶ月目"],
              ["③ 最終（参考）", "管理受託・受注件数／CPO", "1〜3年"]],
             col_w=[CW * 0.16, CW * 0.56, CW * 0.28], row_h=0.6)
box(s, CX0, CY0 + 2.7, CW, 3.4, fill=PALE)
T(s, CX0 + 0.3, CY0 + 2.85, CW - 0.6, 3.2,
  [one("★1件の単価が大きいので、育てていたリードが年1件起きるだけで投資は回収される。", 11.5, True, TNAVY, ls=1.3, sa=4),
   one("ただし「6ヶ月で証明できるのは面談CPAとCVR」であり、CPOは約束しない。", 11, None, INK, ls=1.3, sa=4),
   one("CPOは3年累積の「構造」だけを図示：1年目は受注が微増しCPOはほぼ横ばい。2年目以降、育てたリードが起き始め、追加販促費ゼロで受注が乗るためCPOが下がる。", 10, None, INK, ls=1.3, sa=4),
   one("締め：御社の実績をいただければ、この構造に数字を入れてSIMを作ります。", 11.5, True, TNAVY, ls=1.3)],
  anchor="t")
foot(s, "意見｜数値は入れない（業界汎用のためSIMは作らない）。主KPI＝面談CPA・CVR")

# ============================================================
# S31 費用プラン
# ============================================================
s = slides[30]
frame(s, "成果と体制｜費用プラン",
      ["無償付帯と別途費用を分けて示す（6ヶ月〜・税抜）。"])
plans = [
    ("① コンサル基本", "初期10万〜", "月20万（3投稿）\n月30万（5投稿）\n月50万（9投稿）"),
    ("② 初動設計＋運用", "初期10万〜", "月5万〜"),
    ("③ 運用代行・効率改善", "初期10万〜", "月0万〜\n（アカウント費のみ）"),
    ("④ 成果報酬型", "初期5万〜", "月＝単価×成果数\n＋固定費"),
]
cw4 = (CW - 0.3 * 3) / 4
for i, (h, init, m) in enumerate(plans):
    x = CX0 + i * (cw4 + 0.3)
    box(s, x, CY0 + 0.2, cw4, 3.3, fill=PALE)
    T(s, x + 0.2, CY0 + 0.35, cw4 - 0.4, 0.9, [one(h, 11, True, NAVY, ls=1.2)], anchor="t")
    T(s, x + 0.2, CY0 + 1.3, cw4 - 0.4, 0.6, [one(init, 13, True, ORANGE)], anchor="t")
    T(s, x + 0.2, CY0 + 1.95, cw4 - 0.4, 1.4, [one(l, 9, None, INK, ls=1.25) for l in m.split("\n")], anchor="t")
T(s, CX0, CY0 + 3.8, CW, 0.5, [one("無償付帯", 11.5, True, NAVY)], anchor="m")
T(s, CX0, CY0 + 4.3, CW, 0.9,
  [one("アカウント開設／プロフィール／リッチメニュー／あいさつ／初期アンケート／KW自動応答／ステップ配信／セグメント配信／タグ管理／GAレポート連携／クリエイティブ／定例会", 9.5, None, INK, ls=1.3)],
  anchor="t")
T(s, CX0, CY0 + 5.3, CW, 0.5, [one("別途費用", 11.5, True, RED)], anchor="m")
T(s, CX0, CY0 + 5.8, CW, 0.9,
  [one("離脱防止（Sitelead）初期10万＋月5万／API連携／通知メッセージ 等。配信ツールはLstep／Hachidoriのいずれかを要件で選定", 9.5, None, INK, ls=1.3)],
  anchor="t")
foot(s, "意見｜DYM標準プラン（6ヶ月〜・税抜）")

# ============================================================
# S32 運用スケジュール
# ============================================================
s = slides[31]
frame(s, "成果と体制｜運用スケジュール",
      ["1ヶ月目に構築、2ヶ月目に稼働、3ヶ月目から改善。6ヶ月で年間配信の型ができる。"])
sched = [
    ("1ヶ月目", "構築", "あいさつMSG・診断・リッチメニュー・年次カレンダー設計"),
    ("2ヶ月目", "稼働", "友だち追加動線を稼働。初動30日ステップ開始"),
    ("3〜4ヶ月目", "改善①", "モーメント配信・セミナー送客を開始。効果測定"),
    ("5〜6ヶ月目", "改善②", "年間配信の型を確立"),
]
cw4 = CW / 4
for i, (h, tag, b) in enumerate(sched):
    x = CX0 + i * cw4
    tb = box(s, x + 0.1, CY0 + 0.3, cw4 - 0.2, 0.8, fill=NAVY)
    put_text(tb.text_frame, [one(h, 11.5, True, WHITE, align="c")], anchor="m")
    T(s, x + 0.1, CY0 + 1.25, cw4 - 0.2, 0.6, [one(tag, 12, True, ORANGE, align="c")], anchor="m")
    card(s, x + 0.1, CY0 + 1.95, cw4 - 0.2, 2.7, "", b, hcol=NAVY, fill=PALE, hsz=1, bsz=9.5, ls=1.3, anchor="m")
    if i < 3:
        T(s, x + cw4 - 0.15, CY0 + 1.9, 0.4, 1.0, [one("▶", 14, True, MUT, align="c")], anchor="m")
foot(s, "意見｜6ヶ月で年間配信の型を確立。以降は月次運用フェーズへ")

# ============================================================
# S33 サポート体制
# ============================================================
s = slides[32]
frame(s, "成果と体制｜サポート体制",
      ["定例会・レポート・クリエイティブ制作まで含む。"])
supp = [
    ("定例会", "月次で振り返り・\n翌月の企画をすり合わせ"),
    ("レポーティング", "GA連携で友だち数・\n予約数を数値化"),
    ("クリエイティブ制作", "配信文面・リッチメニュー\n画像の制作を内包"),
    ("運用代行", "日々の配信設計・\n自動応答の調整"),
]
cw4 = (CW - 0.3 * 3) / 4
for i, (h, b) in enumerate(supp):
    x = CX0 + i * (cw4 + 0.3)
    card(s, x, CY0 + 0.3, cw4, 3.0, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=11.5, bsz=9.5, ls=1.3)
foot(s, "意見｜DYM標準のサポート体制")

# ============================================================
# S34 飛び道具（4案）
# ============================================================
s = slides[33]
frame(s, "締め｜飛び道具",
      ["不利を先に言う会社が選ばれる。"])
tobi = [
    ("A", "サブリースの落とし穴を自社から開示する配信", "★前後検索で「アパート経営の落とし穴」が起点より前に確認済み。不利を先に言う会社が選ばれる", True),
    ("B", "近隣の成約家賃を毎月配信", "自分の物件が高いか安いかが毎月分かる。思い出される理由になり続ける", False),
    ("C", "決算期に合わせた収支レポート配信", "税理士に持っていく資料になる。第三者経由で信用が増える", False),
    ("D", "自宅建築への送客（Ver.Aとの相互送客）", "★前後検索で後7〜15日に住宅ローン・住宅展示場・タマホーム 平屋。土地活用検討者の一部が自宅建築に流れている", True),
]
cw2 = (CW - 0.4) / 2
for i, (tag, h, b, verified) in enumerate(tobi):
    r, c = divmod(i, 2)
    x = CX0 + c * (cw2 + 0.4)
    y = CY0 + 0.2 + r * 3.4
    card(s, x, y, cw2, 3.1, f"{tag}　{h}", [b], hcol=NAVY, fill=PALE, hsz=11.5, bsz=9.5, ls=1.35)
    if verified:
        badge(s, x + cw2 - 2.5, y + 0.15, 2.2, 0.5, "実データ確認済み", fill=PORANGE, col=ORANGE, sz=7.5)
foot(s, "事実｜A・D案は前後検索データ（S08）で裏付け済み。B・Cは意見（施策アイデア）")

# ============================================================
# S35 LINEOA実績
# ============================================================
s = slides[34]
frame(s, "締め｜LINEOA実績",
      ["賃貸管理・オーナー向けの公式事例は、2回探して1件も見つかっていない。",
       "近接事例で代替しつつ、正直に「先行者になれる」と書く。"])
box(s, CX0, CY0 + 0.2, CW, 1.3, fill=PRED, line=RED, lw=1.0)
T(s, CX0 + 0.3, CY0 + 0.2, CW - 0.6, 1.3,
  [one("★オーナー向けの公式事例は存在しない＝先行者になれる", 13, True, RED, align="c", ls=1.25)],
  anchor="m")
cases = [
    ("お部屋探しのハートサポート（賃貸仲介・近接事例）", "LINE経由の問い合わせから30〜40件成約（2020年1〜7月）"),
    ("LIFULL HOME'S（不動産情報サイト・近接事例）", "ステップ配信でCPC約6割改善・運用工数1/4"),
]
cw2 = (CW - 0.4) / 2
for i, (h, b) in enumerate(cases):
    x = CX0 + i * (cw2 + 0.4)
    card(s, x, CY0 + 1.7, cw2, 2.4, h, [b], hcol=NAVY, fill=PALE, hsz=11, bsz=9.5, ls=1.3)
foot(s, "事実｜オーナー向け公式事例なし（2回探索済み）。近接事例はLINEヤフー公式。掲載前に原典で再確認すること")


# ============================================================
# 業種残骸の最終チェック（正規表現スキャン）
# ============================================================
import re

RESIDUE_PATTERNS = [
    r"○○業界", r"〇〇業界", r"クロスセル", r"アップセル", r"転職", r"求人",
    r"利用状況タグ", r"誕生日.{0,3}記念日タグ", r"人材\(転職\)業界",
]
residue_hits = []
for i, sl in enumerate(slides, 1):
    for sh in sl.shapes:
        if sh.has_text_frame:
            txt = sh.text_frame.text
            for pat in RESIDUE_PATTERNS:
                if re.search(pat, txt):
                    residue_hits.append((i, pat, txt[:40]))
if residue_hits:
    print("★業種残骸が残っています：")
    for i, pat, txt in residue_hits:
        print(f"  slide {i}: /{pat}/ -> {txt!r}")
else:
    print("業種残骸チェック：残存0")

prs.save(OUT)
print("saved:", OUT)
