# -*- coding: utf-8 -*-
"""ジム・フィットネス業界資料：S39/S40/S41 のレイアウトをS42に合わせる。

背景：S39〜41は3面カードが縦に詰まって見づらく、文字サイズもS42より小さかった。
文言・列・横位置・シェイプ自体は変更せず、以下の2点だけをS42基準に揃える。
  1) 各カードの .top（縦位置） … S39〜41はFMT由来で行間が詰まっていたため、
     S42のクラスタ（行の代表position）に合わせて再配置。
  2) 各シェイプの文字サイズ（run.font.size） … S39〜41は一回り小さいサイズが
     入っていたため、S42の対応シェイプのサイズをコピー。

シェイプの対応づけは「文字が入っているシェイプ」を対象に、topクラスタ（許容誤差0.18cm）
→ その中をleft（横位置）昇順、の2段階で行う。空文字の装飾シェイプ（背景等）は
同じ行のズレ幅で一緒に動かすだけで、サイズ変更の対象にはしない。

実行前提：SRC には元の43枚デッキ（S1=表紙, S39〜41=対象, S42=お手本, S43=背面）を指定。
"""
from pptx import Presentation

SRC = "元ファイル.pptx"
OUT = "ジム・フィットネス業界_LINEOA施策提案.pptx"
TOL = 0.18  # cm
TARGETS = [(38, "S39"), (39, "S40"), (40, "S41")]
REF_IDX = 41  # S42


def cm(emu):
    return emu / 360000


def cm_to_emu(v):
    return int(round(v * 360000))


def walk(shapes):
    out = []
    for s in shapes:
        if s.shape_type == 6:
            out.extend(walk(s.shapes))
        else:
            out.append(s)
    return out


def col_of(sh):
    return 1 if cm(sh.left) < 13 else 2


def cluster_by_top(shs, tol=TOL):
    shs = sorted(shs, key=lambda sh: cm(sh.top))
    clusters, cur, cur_top = [], [], None
    for sh in shs:
        t = cm(sh.top)
        if cur and (t - cur_top) > tol:
            clusters.append(cur)
            cur = []
        cur.append(sh)
        cur_top = t
    if cur:
        clusters.append(cur)
    return [(sum(cm(sh.top) for sh in c) / len(c), c) for c in clusters]


def classify(slide):
    """戻り値: content_header[col], content_rows[col][row]  … 非空テキストのみ
               deco_header[col],   deco_rows[col][row]      … 空テキストのみ（背景等）
    """
    shapes = [sh for sh in walk(slide.shapes) if sh.has_text_frame]
    body = [sh for sh in shapes if not (cm(sh.left) < 2.0 and cm(sh.width) > 15)]

    numbers = [sh for sh in body
               if sh.text_frame.text.strip() in ("1", "2", "3") and cm(sh.width) < 1.5]
    row_top = {1: {}, 2: {}}
    for sh in numbers:
        row_top[col_of(sh)][sh.text_frame.text.strip()] = cm(sh.top)

    def row_of(sh):
        c = col_of(sh)
        tops = sorted(row_top[c].items(), key=lambda kv: kv[1])
        t = cm(sh.top)
        for i, (rn, rt) in enumerate(tops):
            lo = rt - 0.3
            hi = tops[i + 1][1] - 0.3 if i + 1 < len(tops) else 9999
            if lo <= t < hi:
                return int(rn)
        return None

    content_h, deco_h = {1: [], 2: []}, {1: [], 2: []}
    content_r = {1: {1: [], 2: [], 3: []}, 2: {1: [], 2: [], 3: []}}
    deco_r = {1: {1: [], 2: [], 3: []}, 2: {1: [], 2: [], 3: []}}
    for sh in body:
        c = col_of(sh)
        r = row_of(sh)
        is_content = sh.text_frame.text.strip() != ""
        if r is None:
            (content_h if is_content else deco_h)[c].append(sh)
        else:
            (content_r if is_content else deco_r)[c][r].append(sh)

    ch = {c: cluster_by_top(content_h[c]) for c in (1, 2)}
    cr = {c: {r: cluster_by_top(content_r[c][r]) for r in (1, 2, 3)} for c in (1, 2)}
    return ch, cr, deco_h, deco_r


def all_runs(sh):
    runs = []
    for p in sh.text_frame.paragraphs:
        runs.extend(p.runs)
    return runs


def build_plan(ref_h, ref_r, hdr, rows):
    ok = True
    plan = {"header": {}, "rows": {}, "row_delta": {1: {}, 2: {}}, "header_delta": {}}
    for c in (1, 2):
        if len(hdr[c]) != len(ref_h[c]):
            ok = False
        else:
            plan["header"][c] = list(zip(hdr[c], ref_h[c]))
            plan["header_delta"][c] = ref_h[c][0][0] - hdr[c][0][0]
        for r in (1, 2, 3):
            if len(rows[c][r]) != len(ref_r[c][r]):
                ok = False
            else:
                plan["rows"].setdefault(c, {})[r] = list(zip(rows[c][r], ref_r[c][r]))
                plan["row_delta"].setdefault(c, {})[r] = [
                    (old_top, new_top - old_top)
                    for (old_top, _), (new_top, _) in zip(rows[c][r], ref_r[c][r])
                ]
    return ok, plan


def apply_position(plan, deco_h, deco_r):
    for c in (1, 2):
        for old_cluster, new_cluster in plan["header"][c]:
            _, shs = old_cluster
            new_top, _ = new_cluster
            for sh in shs:
                sh.top = cm_to_emu(new_top)
        for r in (1, 2, 3):
            for old_cluster, new_cluster in plan["rows"][c][r]:
                _, shs = old_cluster
                new_top, _ = new_cluster
                for sh in shs:
                    sh.top = cm_to_emu(new_top)
    for c in (1, 2):
        for sh in deco_h[c]:
            sh.top = cm_to_emu(cm(sh.top) + plan["header_delta"][c])
        for r in (1, 2, 3):
            band_deltas = plan["row_delta"][c][r]
            for sh in deco_r[c][r]:
                t = cm(sh.top)
                _, delta = min(band_deltas, key=lambda bd: abs(bd[0] - t))
                sh.top = cm_to_emu(t + delta)


def apply_font(plan):
    for c in (1, 2):
        for old_cluster, new_cluster in plan["header"][c]:
            _apply_cluster_font(old_cluster, new_cluster)
        for r in (1, 2, 3):
            for old_cluster, new_cluster in plan["rows"][c][r]:
                _apply_cluster_font(old_cluster, new_cluster)


def _apply_cluster_font(old_cluster, new_cluster):
    _, old_shapes = old_cluster
    _, new_shapes = new_cluster
    # left昇順。同一left（FMT由来の完全重複シェイプ等）はshape_idで決定的にタイブレーク。
    old_sorted = sorted(old_shapes, key=lambda sh: (cm(sh.left), sh.shape_id))
    new_sorted = sorted(new_shapes, key=lambda sh: (cm(sh.left), sh.shape_id))
    for old_sh, new_sh in zip(old_sorted, new_sorted):
        new_runs = all_runs(new_sh)
        if not new_runs or new_runs[0].font.size is None:
            continue
        ref_size = new_runs[0].font.size
        for r in all_runs(old_sh):
            r.font.size = ref_size


def main():
    prs = Presentation(SRC)
    ref_h, ref_r, _, _ = classify(prs.slides[REF_IDX])

    plans = {}
    for idx, name in TARGETS:
        hdr, rows, deco_h, deco_r = classify(prs.slides[idx])
        ok, plan = build_plan(ref_h, ref_r, hdr, rows)
        plan["deco_header"], plan["deco_rows"] = deco_h, deco_r
        if not ok:
            raise SystemExit(f"{name}: クラスタ件数がS42と一致しません。手動確認してください。")
        plans[idx] = plan
        print(f"{name}: クラスタ対応OK")

    for idx, name in TARGETS:
        plan = plans[idx]
        apply_position(plan, plan["deco_header"], plan["deco_rows"])
        apply_font(plan)
        print(f"{name}: 位置・文字サイズをS42基準に更新")

    prs.save(OUT)
    print("saved:", OUT)


if __name__ == "__main__":
    main()
