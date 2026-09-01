# -*- coding: utf-8 -*-
"""注文住宅業界 LINEOA施策提案（36枚）

_templates/DYM_LINEOA_FMT.pptx（業界汎用FMT・36枚）をコピーし、
14枚を全面再構築＋残りを注文住宅語に置換して完成させる。

  python _build/build_chumon_jutaku.py
  python _build/qa_render.py 注文住宅業界_LINEOA施策提案.pptx

【落とし穴メモ（_build/README.md より）】
- put_text() は必ず reset_tf() を通す（既存テキストへの追記事故を防ぐ）
- スライドの新規追加はしない。ベース36枚を clear_slide() して作り直す
- ベースは別業種（EC・人材）の残骸あり。REPLACE で一括掃除している
"""
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
SRC = str(ROOT / "_templates" / "DYM_LINEOA_FMT.pptx")
OUT = str(ROOT / "注文住宅業界_LINEOA施策提案.pptx")

# ---- 配色（DYM 3色：ネイビー・炭黒・オレンジ。LINEグリーンはUI部分のみ）----
TNAVY = "002060"   # タイトル
NAVY = "1F285A"    # 打ち手・カード見出し
ORANGE = "ED7D31"  # 強調
RED = "C00000"     # 課題・警告
INK = "333333"     # 本文
MUT = "7F7F7F"
WHITE = "FFFFFF"
PALE = "F4F7FF"    # 淡ネイビー
PORANGE = "FCE4D6"  # 淡オレンジ
GREY = "F2F2F2"
BORDER = "D9D9D9"
GREEN = "06C755"   # LINE UI 専用
BEZEL = "2B2B2B"
SCREEN = "EFF2F7"
PRED = "FDF2F2"

# ---- FMTのレイアウト座標（cm）----
SW, SH = 27.52, 19.05
TITLE_XY = (1.52, 0.38, 24.4, 0.94)
LEAD_XY = (1.20, 1.80, 25.1, 1.90)
DIV_Y = 3.86
CX0, CW = 1.20, 25.12
CY0, CY1 = 4.30, 17.00
FOOT_Y = 17.35

shutil.copyfile(SRC, OUT)
prs = Presentation(OUT)
slides = list(prs.slides)
assert len(slides) == 36, len(slides)
assert (prs.slide_width, prs.slide_height) == (9906000, 6858000)


# ================= helpers =================
def set_font(run, size, bold=None, color=INK, name="メイリオ"):
    f = run.font
    f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", name)
    f.color.rgb = RGBColor.from_string(color)


ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
ANCH = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}


def reset_tf(tf):
    """既存の段落・runを全消去（追記事故の防止）"""
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)


def put_text(tf, paras, anchor="t", ml=0.14, mr=0.14, mt=0.06, mb=0.06, wrap=True):
    reset_tf(tf)
    tf.word_wrap = wrap
    tf.margin_left = Cm(ml)
    tf.margin_right = Cm(mr)
    tf.margin_top = Cm(mt)
    tf.margin_bottom = Cm(mb)
    tf.vertical_anchor = ANCH[anchor]
    first = True
    for p in paras:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = ALIGN[p.get("align", "l")]
        if p.get("sa") is not None:
            para.space_after = Pt(p["sa"])
        if p.get("ls") is not None:
            para.line_spacing = p["ls"]
        for item in p["runs"]:
            t, sz, b, c = item
            r = para.add_run()
            r.text = t
            set_font(r, sz, b, c)
    return tf


def T(slide, x, y, w, h, paras, anchor="t", **kw):
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    put_text(box.text_frame, paras, anchor=anchor, **kw)
    return box


def one(text, sz, b=None, c=INK, align="l", sa=None, ls=None):
    d = {"runs": [(text, sz, b, c)], "align": align}
    if sa is not None:
        d["sa"] = sa
    if ls is not None:
        d["ls"] = ls
    return d


def box(slide, x, y, w, h, fill=None, line=None, lw=1.0,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06, dash=None):
    sp = slide.shapes.add_shape(shape, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor.from_string(line)
        sp.line.width = Pt(lw)
        if dash:
            sp.line.dash_style = dash
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    reset_tf(sp.text_frame)
    return sp


def card(slide, x, y, w, h, head, body, hcol=NAVY, fill=PALE,
         hsz=11, bsz=9, line=None, anchor="t", ls=1.18):
    sp = box(slide, x, y, w, h, fill=fill, line=line)
    paras = [one(head, hsz, True, hcol, sa=3)]
    for b in (body if isinstance(body, list) else [body]):
        if b:
            paras.append(one(b, bsz, None, INK, ls=ls, sa=1))
    put_text(sp.text_frame, paras, anchor=anchor, ml=0.22, mr=0.18, mt=0.14, mb=0.10)
    return sp


def clear_slide(slide):
    spTree = slide.shapes._spTree
    for el in list(spTree):
        if el.tag.split("}")[-1] in ("sp", "cxnSp", "pic", "graphicFrame", "grpSp"):
            spTree.remove(el)


def frame(slide, title, lead):
    """FMT準拠：タイトル16pt紺（y=0.38）＋リード14pt（y=1.80・2行以内）＋区切り線"""
    clear_slide(slide)
    T(slide, *TITLE_XY, [one(title, 16, True, TNAVY)], anchor="m", ml=0, mr=0)
    T(slide, *LEAD_XY, [one(l, 12, None, INK, ls=1.28) for l in lead],
      anchor="m", ml=0, mr=0)
    ln = slide.shapes.add_connector(1, Cm(0), Cm(DIV_Y), Cm(SW), Cm(DIV_Y))
    ln.line.color.rgb = RGBColor.from_string(BORDER)
    ln.line.width = Pt(1.0)


def foot(slide, text):
    T(slide, CX0, FOOT_Y, CW, 0.9, [one(text, 7.5, None, MUT, ls=1.15)], ml=0, mr=0)


def badge(slide, x, y, w, h, text, fill=PORANGE, col=ORANGE, sz=8.5, dash=None):
    sp = box(slide, x, y, w, h, fill=fill, line=col, lw=1.0, dash=dash)
    put_text(sp.text_frame, [one(text, sz, True, col, align="c")],
             anchor="m", ml=0.06, mr=0.06, mt=0, mb=0)
    return sp


def band(slide, y, text, fill=NAVY, col=WHITE, sz=11.5, h=0.86, x=CX0, w=CW):
    sp = box(slide, x, y, w, h, fill=fill, radius=0.10)
    put_text(sp.text_frame, [one(text, sz, True, col, align="c")],
             anchor="m", ml=0.2, mr=0.2, mt=0, mb=0)
    return sp


def phone(slide, x, y, w, h, header, lines):
    """LINEトーク画面のモック。lines = [(kind, text)]
    kind: 'in'（BOT吹き出し）/ 'chip'（選択肢）/ 'btn'（CTA）/ 'note'（注記）"""
    box(slide, x, y, w, h, fill=BEZEL, radius=0.10)
    box(slide, x + 0.14, y + 0.42, w - 0.28, h - 0.56, fill=SCREEN, radius=0.02,
        shape=MSO_SHAPE.RECTANGLE)
    hd = box(slide, x + 0.14, y + 0.42, w - 0.28, 0.52, fill=GREEN,
             shape=MSO_SHAPE.RECTANGLE)
    put_text(hd.text_frame, [one(header, 7, True, WHITE, align="c")],
             anchor="m", ml=0.05, mr=0.05, mt=0, mb=0)
    cy = y + 1.05
    iw = w - 0.62
    for kind, text in lines:
        nlines = text.count("\n") + 1
        bh = 0.30 + 0.30 * nlines
        if kind == "in":
            sp = box(slide, x + 0.31, cy, iw, bh, fill=WHITE, line=BORDER, lw=0.75)
            put_text(sp.text_frame, [one(l, 6.5, None, INK, ls=1.18)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.14, mr=0.10, mt=0.05, mb=0.05)
        elif kind == "chip":
            sp = box(slide, x + 0.31, cy, iw, bh, fill=WHITE, line=GREEN, lw=0.9)
            put_text(sp.text_frame, [one(l, 6.5, None, "0B7A3B", align="c", ls=1.15)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.08, mr=0.08, mt=0.03, mb=0.03)
        elif kind == "btn":
            sp = box(slide, x + 0.31, cy, iw, bh, fill=GREEN)
            put_text(sp.text_frame, [one(l, 6.8, True, WHITE, align="c", ls=1.15)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.08, mr=0.08, mt=0.03, mb=0.03)
        else:  # note
            sp = slide.shapes.add_textbox(Cm(x + 0.31), Cm(cy), Cm(iw), Cm(bh))
            put_text(sp.text_frame, [one(l, 6.2, None, MUT, ls=1.15)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.05, mr=0.05, mt=0, mb=0)
        cy += bh + 0.12
    return cy


def find_shape(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle in sh.text_frame.text:
            return sh
    return None


# ============================================================
# S01 表紙
# ============================================================
s = slides[0]
sh = find_shape(s, "業界：")
if sh:
    put_text(sh.text_frame, [one("業界：注文住宅", 12, True, TNAVY)], anchor="m")
sh = find_shape(s, "最上位パートナーの知見")
if sh:
    put_text(sh.text_frame,
             [one("検討期間10ヶ月を伴走し、来場率と契約棟数を引き上げる動線を構築する", 13, True, INK)],
             anchor="m")

# ============================================================
# S05 CPC推移
# ============================================================
s = slides[4]
frame(s, "市場環境｜検索広告におけるCPCの推移",
      ["注文住宅の主要KWは3年でCPCが約1.45倍。一方で検索需要はほぼ横ばいのため、",
       "同一の契約棟数を維持するコストが構造的に増え続けている状態。"])

# 左：CPC縦棒グラフ
GX, GY, GW, GH = 1.60, 5.10, 12.6, 9.40
base_y = GY + GH - 1.35
T(s, GX, GY - 0.62, GW, 0.55,
  [one("平均CPC（円）｜「注文住宅＋地域名」系キーワード", 9.5, True, NAVY)], ml=0)
cpc = [("2023年", 420), ("2024年", 480), ("2025年", 545), ("2026年", 610)]
maxv, bar_w, gap = 700.0, 1.85, 1.05
plot_h = GH - 2.15
for i, (lab, v) in enumerate(cpc):
    bx = GX + 1.30 + i * (bar_w + gap)
    bh = plot_h * (v / maxv)
    last = (i == len(cpc) - 1)
    box(s, bx, base_y - bh, bar_w, bh, fill=ORANGE if last else "F4B183",
        shape=MSO_SHAPE.RECTANGLE)
    T(s, bx - 0.42, base_y - bh - 0.66, bar_w + 0.84, 0.56,
      [one(f"{v}円", 10 if last else 9, True, ORANGE if last else INK, align="c")], ml=0, mr=0)
    T(s, bx - 0.42, base_y + 0.10, bar_w + 0.84, 0.52,
      [one(lab, 8.5, None, INK, align="c")], ml=0, mr=0)
ln = s.shapes.add_connector(1, Cm(GX + 0.70), Cm(base_y), Cm(GX + GW - 0.40), Cm(base_y))
ln.line.color.rgb = RGBColor.from_string(MUT)
ln.line.width = Pt(1.0)
badge(s, GX + 8.60, GY + 0.30, 3.30, 0.80, "3年で 1.45倍")

# 右：3つの要因カード
RX, RW = 15.10, 11.20
T(s, RX, GY - 0.62, RW, 0.55, [one("CPCを押し上げている3つの層", 9.5, True, NAVY)], ml=0)
factors = [
    ("① 市場競争要因", ["着工棟数は減少局面。同じパイを大手・ローコスト・",
                   "地場ビルダーで奪い合う → 同一KWの入札社数が増加"]),
    ("② マクロ経済要因", ["建築資材・人件費の上昇で1棟あたり原価が上昇 →",
                    "各社が許容CPAを引き上げてでも獲りに来る"]),
    ("③ 媒体・技術要因", ["検索結果最上部のAI要約（AIO）でサイトに来ずに解決 →",
                    "「相場」「後悔」など情報収集KWの流入が構造的に減少"]),
]
for i, (h, b) in enumerate(factors):
    card(s, RX, GY + i * 2.55, RW, 2.25, h, b, hsz=10.5, bsz=8.5, fill=PALE)
band(s, GY + 7.85, "「広告を増やす」以外の打ち手が必要な局面に入っている", fill=ORANGE,
     x=RX, w=RW, h=1.05, sz=11)
foot(s, "※ CPCは業界水準に基づく推計レンジ。検索Vol指数は年平均=100。"
        "提案前に貴社の実績CPC／Googleトレンド実データに差し替えます。")

# ============================================================
# S06 広告実績CVR設計 ＋ 広告審査の留意点
# ============================================================
s = slides[5]
frame(s, "広告与件｜CV地点別の設計マトリクスと、広告審査の留意点",
      ["数値の差込だらけの実績表ではなく「時期×対象×訴求×CV地点」の設計として提示。",
       "実績値は貴社データで置換し、脚注へ退避する。"])

hdr = ["時期", "対象（誰に）", "訴求（何を）", "CV地点", "想定CVR"]
colw = [3.20, 5.60, 8.20, 5.10, 3.02]
rows = [
    ["通年", "検討初期・土地未定", "無理のない予算がわかる60秒診断", "LINE友だち追加", "―"],
    ["通年", "検討初期・予算不明", "総額の内訳（外構込み）を先に開示", "資料請求", "2.0%"],
    ["4〜6月", "翌年入居を狙う層", "土地探しの正しい順番セミナー", "個別相談予約", "0.8%"],
    ["10〜12月", "検討中期・比較段階", "冬の断熱体感／構造見学会", "見学会予約", "0.6%"],
]
ty, rh = 5.05, 1.02
x = CX0
for i, h in enumerate(hdr):
    sp = box(s, x, ty, colw[i], 0.86, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    put_text(sp.text_frame, [one(h, 9, True, WHITE, align="c")], anchor="m",
             ml=0.06, mr=0.06, mt=0, mb=0)
    x += colw[i]
for r, row in enumerate(rows):
    x = CX0
    bg = WHITE if r % 2 == 0 else GREY
    for i, cellv in enumerate(row):
        sp = box(s, x, ty + 0.86 + r * rh, colw[i], rh, fill=bg,
                 line=BORDER, lw=0.75, shape=MSO_SHAPE.RECTANGLE)
        emph = (i == 3)
        put_text(sp.text_frame,
                 [one(cellv, 8.5, True if emph else None,
                      ORANGE if emph else INK, align="c" if i != 2 else "l", ls=1.15)],
                 anchor="m", ml=0.16, mr=0.10, mt=0, mb=0)
        x += colw[i]
T(s, CX0, ty + 0.86 + 4 * rh + 0.16, CW, 0.55,
  [one("※ CVRは業界水準（CV地点＝各行に明記）。貴社の実績CPC/CVR/CPAが入り次第、全行を実測値に置換します。",
       8, None, MUT)], ml=0)

ay = 10.90
T(s, CX0, ay, CW, 0.55, [one("広告審査・表示上の留意点（先に潰しておく3点）", 11, True, RED)], ml=0)
alerts = [
    ("① 断定・最上級表現", ["「必ず」「No.1」「日本一」は根拠資料がないと審査落ち。",
                     "調査主体・時点・範囲を明記できるものだけ使う"]),
    ("② 価格・ローン表記", ["「月々○万円で家が建つ」は前提（借入額・金利・年数・",
                     "自己資金）の併記が必須。総額表示との整合も取る"]),
    ("③ 施工事例・体験談", ["ビフォーアフター写真は撮影条件を揃える。",
                     "お客様の声は実在・許諾が必要。イメージなら明記"]),
]
for i, (h, b) in enumerate(alerts):
    card(s, CX0 + i * 8.44, ay + 0.62, 8.10, 2.55, h, b, hcol=RED, fill=PRED,
         line=RED, hsz=10, bsz=8.5)
foot(s, "※ 広告審査の運用は媒体・時期で変わります。入稿前に最新のガイドラインで再確認してください。")

# ============================================================
# S07 シーズナリティ
# ============================================================
s = slides[6]
frame(s, "ニーズ調査（シーズナリティ）｜需要の波と9ヶ月のズレ",
      ["契約のピークは1〜3月だが、その層が情報収集を始めるのは前年の4〜6月。",
       "＝ 繁忙期に広告を集中しても、その時点で既に他社と接触済み。"])

CAL_X, CAL_Y, CAL_W = CX0, 5.15, CW
months = ["1月", "2月", "3月", "4月", "5月", "6月", "7月", "8月", "9月", "10月", "11月", "12月"]
keiyaku = [115, 130, 145, 90, 80, 75, 85, 90, 95, 105, 110, 100]
kentou = [85, 90, 95, 110, 125, 135, 120, 115, 100, 95, 88, 80]
cellw = CAL_W / 12.0
row_h = 3.05
for i, m in enumerate(months):
    mx = CAL_X + i * cellw
    sp = box(s, mx, CAL_Y, cellw - 0.06, 0.62, fill=GREY, line=BORDER, lw=0.6,
             shape=MSO_SHAPE.RECTANGLE)
    put_text(sp.text_frame, [one(m, 8, True, INK, align="c")], anchor="m",
             ml=0, mr=0, mt=0, mb=0)
# 上段：契約数指数（ネイビー）
T(s, CAL_X, CAL_Y + 0.72, 6.0, 0.5, [one("契約数指数（ネイビー）", 8.5, True, NAVY)], ml=0)
base1 = CAL_Y + 0.72 + row_h + 0.55
for i, v in enumerate(keiyaku):
    mx = CAL_X + i * cellw
    bh = row_h * (v / 160.0)
    box(s, mx + 0.22, base1 - bh, cellw - 0.50, bh, fill=NAVY if v >= 115 else "8496B0",
        shape=MSO_SHAPE.RECTANGLE)
# 下段：検討開始数指数（オレンジ）
T(s, CAL_X, base1 + 0.28, 8.0, 0.5, [one("検討開始数指数（オレンジ）", 8.5, True, ORANGE)], ml=0)
base2 = base1 + 0.28 + row_h + 0.55
for i, v in enumerate(kentou):
    mx = CAL_X + i * cellw
    bh = row_h * (v / 160.0)
    box(s, mx + 0.22, base2 - bh, cellw - 0.50, bh,
        fill=ORANGE if v >= 120 else "F8CBAD", shape=MSO_SHAPE.RECTANGLE)
badge(s, CAL_X + 7.30, base2 + 0.22, 8.40, 0.86, "ピークのズレ ≒ 9ヶ月", sz=11)

cy = base2 + 1.35
card(s, CX0, cy, 12.30, 2.05, "今の打ち手",
     ["繁忙期（1〜3月）に広告を集中させる。",
      "→ その時点で顧客は既に3社と接触済み。比較表の1行にしかならない"],
     hcol=RED, fill=PRED, line=RED, hsz=10.5, bsz=9)
card(s, CX0 + 12.82, cy, 12.30, 2.05, "あるべき打ち手",
     ["閑散期（4〜6月）に低ハードルで接点を作る。",
      "→ 9ヶ月伴走して育てた会社が、翌年の契約を持っていく"],
     hcol=ORANGE, fill=PORANGE, line=ORANGE, hsz=10.5, bsz=9)
foot(s, "※ 指数は業界水準に基づく推計（年平均=100）。リードタイムは情報収集〜契約で8〜12ヶ月（業界水準）。"
        "Googleトレンド（注文住宅／ハウスメーカー／注文住宅 相場／土地 探し方・5年＋1年）の実データで差し替えます。")

# ============================================================
# S08 前後検索
# ============================================================
s = slides[7]
frame(s, "ニーズ調査（前後検索）｜検索の15日前と15日後で、不安は姿を変える",
      ["検索前15日は「会社」ではなく「不安」を調べている。検索後15日には他社の指名検索と",
       "ネガティブKWが同時に立つ。＝ 友だち追加を仕掛けるのは「前15日〜起点」。"])

BY, BH2 = 5.05, 8.30
w1, w2, w3 = 9.10, 6.30, 9.30
gap2 = 0.20
x1 = CX0
x2 = x1 + w1 + gap2
x3 = x2 + w2 + gap2

box(s, x1, BY, w1, BH2, fill=GREY, line=BORDER, lw=0.75)
T(s, x1 + 0.30, BY + 0.22, w1 - 0.60, 0.62,
  [one("検索 −15日｜漠然とした不安", 10.5, True, NAVY)], ml=0)
kw_before = ["注文住宅 相場", "注文住宅 3000万 どんな家", "ハウスメーカー ランキング",
             "注文住宅 後悔", "間取り 失敗", "土地 探し方",
             "住宅ローン いくらまで借りられる", "建売 注文住宅 違い",
             "注文住宅 総額 内訳", "家 建てる 何から"]
T(s, x1 + 0.30, BY + 0.95, w1 - 0.60, BH2 - 1.20,
  [one(k, 9, True if i < 3 else None, ORANGE if i < 3 else INK, ls=1.55)
   for i, k in enumerate(kw_before)], ml=0)

box(s, x2, BY, w2, BH2, fill=NAVY)
T(s, x2 + 0.25, BY + 0.22, w2 - 0.50, 0.62,
  [one("起点｜検討が具体化", 10.5, True, WHITE, align="c")], ml=0, mr=0)
T(s, x2 + 0.25, BY + 1.15, w2 - 0.50, 2.6,
  [one(k, 9.5, True, WHITE, align="c", ls=1.70)
   for k in ["注文住宅 ●●市", "ハウスメーカー 比較", "工務店 ハウスメーカー どっち"]],
  ml=0, mr=0)
badge(s, x2 + 0.35, BY + 4.55, w2 - 0.70, 1.55,
      "★ 友だち追加を\n仕掛けるのはココ", sz=10)
T(s, x2 + 0.25, BY + 6.35, w2 - 0.50, 1.6,
  [one("フォームは押せない段階。\nLINEなら押せる。", 8.5, None, WHITE, align="c", ls=1.35)],
  ml=0, mr=0)

box(s, x3, BY, w3, BH2, fill=PRED, line=RED, lw=0.9)
T(s, x3 + 0.30, BY + 0.22, w3 - 0.60, 0.62,
  [one("検索 +15日｜比較と疑心", 10.5, True, RED)], ml=0)
kw_after = [("積水ハウス 坪単価", 0), ("一条工務店 評判", 0), ("●●ホーム 欠陥", 1),
            ("ハウスメーカー 断り方", 1), ("住宅展示場 しつこい", 1),
            ("注文住宅 契約後 解約 手付金", 1), ("打ち合わせ 疲れた", 1),
            ("値引き 交渉 タイミング", 0), ("相見積もり 失礼", 0),
            ("ハウスメーカー 営業 うざい", 1)]
T(s, x3 + 0.30, BY + 0.95, w3 - 0.60, BH2 - 1.20,
  [one(k, 9, True if neg else None, RED if neg else INK, ls=1.55)
   for k, neg in kw_after], ml=0)
badge(s, x3 + w3 - 3.60, BY + 0.18, 3.30, 0.70, "他社と接触済み", fill=RED, col=WHITE, sz=8.5)

cy = BY + BH2 + 0.30
pts = ["① 前15日は「会社」ではなく「不安」を検索している。会社名を出しても刺さらない",
       "② 後15日は他社の指名検索とネガティブKWが同時に立つ ＝ 営業されること自体への嫌悪",
       "③ だから、比較が始まる前（前15日〜起点）に低ハードルの接点を置く"]
T(s, CX0, cy, CW, 2.0, [one(p, 9.5, None, INK, ls=1.45, sa=2) for p in pts], ml=0)
foot(s, "※ 現時点は業界知見に基づく（仮説）。LINEヤフーの前後検索データ"
        "（対象KW：注文住宅 相場／ハウスメーカー 比較／注文住宅 ●●市）で検証・差し替えを行います。")

# ============================================================
# S09 他社配信アカウント（実測差込枠）
# ============================================================
s = slides[8]
frame(s, "他社分析①｜主要8社のLINE運用ステータス（実測差込）",
      ["友だち数を推測で載せることはしません。公式（page.line.me）で取得日つきの実測に",
       "差し替える前提で、選定リストと取得手順を先に確定させています。"])

hdr = ["#", "社名", "LINE公式アカウント", "友だち数", "運用の型"]
colw = [1.10, 6.40, 7.60, 5.30, 4.72]
comp = [
    ("1", "積水ハウス", "あり（page.line.me/891uhetr）"),
    ("2", "住友林業", "確認中"),
    ("3", "一条工務店", "あり（page.line.me/507jdzdn）"),
    ("4", "ヘーベルハウス（旭化成ホームズ）", "確認中"),
    ("5", "三井ホーム", "確認中"),
    ("6", "タマホーム", "確認中"),
    ("7", "アイ工務店（中堅・伸長）", "確認中"),
    ("8", "クレバリーホーム（中堅・FC）", "確認中"),
]
ty, rh = 5.05, 0.83
x = CX0
for i, h in enumerate(hdr):
    sp = box(s, x, ty, colw[i], 0.76, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    put_text(sp.text_frame, [one(h, 9, True, WHITE, align="c")], anchor="m",
             ml=0.04, mr=0.04, mt=0, mb=0)
    x += colw[i]
from pptx.enum.dml import MSO_LINE_DASH_STYLE as DASH
for r, (num, name, acc) in enumerate(comp):
    x = CX0
    bg = WHITE if r % 2 == 0 else GREY
    for i, cellv in enumerate([num, name, acc]):
        sp = box(s, x, ty + 0.76 + r * rh, colw[i], rh, fill=bg, line=BORDER,
                 lw=0.7, shape=MSO_SHAPE.RECTANGLE)
        put_text(sp.text_frame,
                 [one(cellv, 8.5, None, INK, align="c" if i == 0 else "l")],
                 anchor="m", ml=0.16, mr=0.08, mt=0, mb=0)
        x += colw[i]
    for i in (3, 4):
        sp = box(s, x, ty + 0.76 + r * rh, colw[i], rh, fill=WHITE, line=ORANGE,
                 lw=0.9, shape=MSO_SHAPE.RECTANGLE, dash=DASH.DASH)
        put_text(sp.text_frame, [one("［実測差込］", 8, None, ORANGE, align="c")],
                 anchor="m", ml=0.04, mr=0.04, mt=0, mb=0)
        x += colw[i]

sy = ty + 0.76 + 8 * rh + 0.35
T(s, CX0, sy, CW, 0.55,
  [one("友だち数の取得手順（提案前に実施・DYM側で対応）", 11, True, NAVY)], ml=0)
steps = [
    ("STEP 1", ["page.line.me/<アカウントID> を開く。", "スライドに載せる数値は公式のみ"]),
    ("STEP 2", ["_accountHeadSubText 直後の数値を取得。", "「おすすめ」欄の他社数値は拾わない"]),
    ("STEP 3", ["取得日を必ず併記（例：2026年●月●日時点）。", "古ければ提案直前に再取得"]),
    ("STEP 4", ["実機スクショは調査用の別アカウントで取得。", "本人アカウントで追加しない"]),
]
for i, (h, b) in enumerate(steps):
    card(s, CX0 + i * 6.34, sy + 0.62, 6.00, 2.30, h, b, hsz=10, bsz=8.5, fill=PALE)
foot(s, "※ 競合8社は知名度優先で自動選定（大手中心＋運用が特徴的な中堅2社）。"
        "実名掲載の最終可否は上司確認事項です。第三者集計サイトの数値はスライドに掲載しません。")

# ============================================================
# S10 他社分析（運用の型4分類）
# ============================================================
s = slides[9]
frame(s, "他社分析②｜運用の型は4つ。診断ドリブン型が空いている",
      ["規模で勝っている社と、伸び率で勝っている社は一致しない。この2軸で読むと示唆が出る。",
       "現時点で「検討初期の不安を診断で受け止める」ポジションは取られていない。"])

MX, MY, MW, MH = CX0, 5.05, 15.40, 10.10
box(s, MX, MY, MW, MH, fill=WHITE, line=BORDER, lw=1.0, shape=MSO_SHAPE.RECTANGLE)
lnv = s.shapes.add_connector(1, Cm(MX + MW / 2), Cm(MY), Cm(MX + MW / 2), Cm(MY + MH))
lnv.line.color.rgb = RGBColor.from_string(BORDER)
lnh = s.shapes.add_connector(1, Cm(MX), Cm(MY + MH / 2), Cm(MX + MW), Cm(MY + MH / 2))
lnh.line.color.rgb = RGBColor.from_string(BORDER)
T(s, MX, MY - 0.58, MW, 0.5, [one("配信頻度　少 ←──────────→ 多", 8.5, None, MUT, align="c")], ml=0)
T(s, MX - 0.10, MY + MH + 0.10, MW, 0.5,
  [one("縦軸：友だち増加率（低↓ / 高↑）", 8.5, None, MUT, align="c")], ml=0)

quad = [
    (0, 0, "特典・会員型", ["限定特典で友だち化。ローコスト系に多い。",
                       "増加率は出るが友だちの質が薄くなりがち"], PALE, NAVY),
    (1, 0, "診断ドリブン型", ["診断で属性を取得し、検討段階別に出し分け。",
                        "★ この型が最も少ない ＝ 空いている"], PORANGE, ORANGE),
    (0, 1, "相談カウンター型", ["有人相談・展示場送客が主。配信していないことが多い。",
                         "地場ビルダー・FC加盟店に多い"], GREY, NAVY),
    (1, 1, "カタログ型", ["施工事例・見学会一覧の閲覧が中心。大手に多い。",
                     "友だち数は多いが増加率は鈍い"], GREY, NAVY),
]
for cx_i, cy_i, h, b, fl, hc in quad:
    qx = MX + 0.28 + cx_i * (MW / 2)
    qy = MY + 0.28 + cy_i * (MH / 2)
    card(s, qx, qy, MW / 2 - 0.56, MH / 2 - 0.56, h, b, hcol=hc, fill=fl,
         hsz=11, bsz=8.5, line=hc if hc == ORANGE else None)
badge(s, MX + MW / 2 + 0.60, MY + MH / 2 - 1.35, 6.20, 0.78,
      "★ 御社が取りに行くポジション", sz=9)
badge(s, MX + MW - 3.10, MY + 0.10, 2.90, 0.62, "仮説｜実測後に確定",
      fill=WHITE, col=ORANGE, sz=7.5, dash=DASH.DASH)

RX2 = MX + MW + 0.50
T(s, RX2, MY, CW - MW - 0.50, 0.55, [one("読み方の2軸", 11, True, NAVY)], ml=0)
card(s, RX2, MY + 0.62, 9.22, 2.05, "規模で勝っているのは誰か",
     ["＝ 友だち数。長く運用している大手が強い"], hsz=10, bsz=8.5)
card(s, RX2, MY + 2.90, 9.22, 2.05, "伸び率で勝っているのは誰か",
     ["＝ 友だち増加率。伸び率トップは友だち数最少の", "後発社であることが多い"],
     hsz=10, bsz=8.5, fill=PORANGE, hcol=ORANGE)
card(s, RX2, MY + 5.18, 9.22, 2.35, "御社にとっての示唆",
     ["大手はカタログ型に寄っている。検討初期の",
      "「不安」に応える型が空いている。ここを取る"],
     hsz=10, bsz=8.5, fill=NAVY, hcol=WHITE)
sp = find_shape(s, "御社にとっての示唆")
if sp:
    put_text(sp.text_frame,
             [one("御社にとっての示唆", 10, True, WHITE, sa=3),
              one("大手はカタログ型に寄っている。検討初期の", 8.5, None, WHITE, ls=1.18),
              one("「不安」に応える型が空いている。ここを取る", 8.5, None, WHITE, ls=1.18)],
             anchor="t", ml=0.22, mr=0.18, mt=0.14, mb=0.10)
foot(s, "※ 型の分類はDYM見解。配信頻度と増加率の関係は相関であり、因果ではありません。"
        "各社のプロットは実測後に確定します。")

# ============================================================
# S12 ステップ配信案（14日・前後クエリ連動）
# ============================================================
s = slides[11]
frame(s, "具体企画①｜ステップ配信案（14日）— 前後クエリと連動させる",
      ["根拠は2つ。①友だち追加から15日前後でアクティブ率が急落する ②起点KW検索の後15日で",
       "他社の指名検索が立つ。読まれる15日と候補が決まる15日が重なるため14日に寄せる。"])

TL_Y = 5.05
days = ["Day0", "Day1", "Day3", "Day6", "Day9", "Day11", "Day14"]
cvd = {3, 6}
step = CW / 7.0
for i, d in enumerate(days):
    dx = CX0 + i * step
    is_cv = i in cvd
    sp = box(s, dx + 0.30, TL_Y, step - 0.60, 0.82,
             fill=ORANGE if is_cv else NAVY, radius=0.20)
    put_text(sp.text_frame, [one(d, 10, True, WHITE, align="c")], anchor="m",
             ml=0, mr=0, mt=0, mb=0)
    if is_cv:
        T(s, dx + 0.10, TL_Y + 0.88, step - 0.20, 0.5,
          [one("★CVオファー", 8, True, ORANGE, align="c")], ml=0, mr=0)

ty2 = TL_Y + 1.55
hdr = ["配信日", "配信内容（方向性）", "根拠（前後検索データとの連動）"]
colw = [2.90, 9.40, 12.82]
plan = [
    ("Day0", "あいさつ＋60秒の予算診断（4問・選択式）",
     "前15日の最多KWが「相場」。金額の不安に最初に答えるため"),
    ("Day1", "無理のない返済額の考え方（手取りの20〜25%）",
     "「住宅ローン いくらまで借りられる」が前15日に立つため"),
    ("Day3", "よくある後悔TOP5と、契約前に潰す方法",
     "「注文住宅 後悔」「間取り 失敗」が前15日の上位にあるため"),
    ("Day6", "★CVオファー① 完成見学会（貸切・営業電話なし）",
     "「住宅展示場 しつこい」への先回りが押される条件になるため"),
    ("Day9", "土地探しの正しい順番（土地から探さない）",
     "「土地 探し方」が前15日に立ち、土地未定が離脱理由の上位のため"),
    ("Day11", "OB施主インタビュー（3社比較・9ヶ月の検討）",
     "後15日に「評判」「口コミ」が立つ。第三者の声を先に渡すため"),
    ("Day14", "★CVオファー② 個別相談30分（オンライン可）",
     "後15日で他社の指名検索が始まる。その前に判断軸を渡すため"),
]
rh = 1.12
x = CX0
for i, h in enumerate(hdr):
    sp = box(s, x, ty2, colw[i], 0.78, fill=NAVY, shape=MSO_SHAPE.RECTANGLE)
    put_text(sp.text_frame, [one(h, 9, True, WHITE, align="c")], anchor="m",
             ml=0.06, mr=0.06, mt=0, mb=0)
    x += colw[i]
for r, row in enumerate(plan):
    x = CX0
    is_cv = row[0] in ("Day6", "Day14")
    bg = PORANGE if is_cv else (WHITE if r % 2 == 0 else GREY)
    for i, cellv in enumerate(row):
        sp = box(s, x, ty2 + 0.78 + r * rh, colw[i], rh, fill=bg, line=BORDER,
                 lw=0.7, shape=MSO_SHAPE.RECTANGLE)
        put_text(sp.text_frame,
                 [one(cellv, 8.5, True if (i == 0 or is_cv and i == 1) else None,
                      ORANGE if is_cv and i <= 1 else INK,
                      align="c" if i == 0 else "l", ls=1.18)],
                 anchor="m", ml=0.16, mr=0.10, mt=0, mb=0)
        x += colw[i]
foot(s, "※ CVオファーは6日目前後と14日目の2回置く（1回で終わらせない）。"
        "出口はCV者＝予約確定・前日リマインドへ分岐／未CV者＝タグ別の定常配信へ合流。")

# ============================================================
# S13 ステップ配信 実文面
# ============================================================
s = slides[12]
frame(s, "具体企画②｜実際にスマホへ届く文面（Day0 / Day3 / Day6 / Day14）",
      ["概要ではなく、絵文字・改行まで含めた実文を確定させる。",
       "開封は通知プレビューの冒頭15文字で決まるため、全通で設計する。"])

PY0, PW, PH = 5.30, 5.85, 10.30
gapx = (CW - PW * 4) / 3.0
previews = ["はじめまして。●●ホ", "【後悔TOP5】1位", "今週末、実際に建", "2週間、おつきあ"]
labels = ["Day0｜あいさつ＋診断", "Day3｜後悔TOP5", "Day6｜★CVオファー①", "Day14｜★CVオファー②"]
contents = [
    [("in", "はじめまして。\n●●ホームの田中と申します🏠"),
     ("in", "✅ こちらからお電話は一切しません\n✅ 営業訪問もいたしません\n✅ 不要になればブロックでOK"),
     ("in", "有人対応は平日 9:00〜18:00 です😊\n4問・60秒の診断で、無理なく\n建てられる予算をお出しします。"),
     ("chip", "診断をはじめる"),
     ("chip", "あとで")],
    [("in", "【建ててから気づく後悔 TOP5】\n先に言います。契約前に潰せます。"),
     ("in", "🥇1位 コンセントの位置と数\n🥈2位 収納の\"量\"より\"場所\"\n🥉3位 窓が多すぎた\n4位 スイッチの高さ\n5位 外構費を見ていなかった"),
     ("in", "5位が最多です。外構150〜250万円は\n建物見積もりに入っていません。\nうちは最初から入れて出します。"),
     ("btn", "資金計画の相談をする")],
    [("in", "今週末、実際に建てられた\nお家を見ていただけます🏠"),
     ("in", "📍●●市●●町\n📅●月●日(土)・●日(日)\n🕙10:00〜17:00 完全予約制\n※1組60分・貸切です"),
     ("in", "・貸切なので他のお客様と会いません\n・その場で契約の話はしません\n・見学後の営業電話もしません"),
     ("btn", "見学を予約する"),
     ("chip", "今回は都合が合わない")],
    [("in", "2週間、おつきあいいただき\nありがとうございました😊"),
     ("in", "ここから先はご家庭ごとに\n答えが変わります。年収・自己資金・\n土地・お子さまの年齢・親の援助。"),
     ("in", "30分の個別相談をご用意しました。\n①専用の資金計画書\n②総予算の上限額（根拠つき）\n③ご希望エリアの土地相場"),
     ("btn", "相談日程を選ぶ"),
     ("chip", "まだ情報だけ受け取りたい")],
]
for i in range(4):
    px = CX0 + i * (PW + gapx)
    is_cv = i >= 2
    T(s, px, PY0 - 1.30, PW, 0.52,
      [one(labels[i], 9.5, True, ORANGE if is_cv else NAVY, align="c")], ml=0, mr=0)
    pv = box(s, px, PY0 - 0.72, PW, 0.62, fill=WHITE, line=ORANGE if is_cv else BORDER, lw=0.9)
    put_text(pv.text_frame,
             [one("通知冒頭15字：" + previews[i], 6.8, True, ORANGE if is_cv else MUT, align="c")],
             anchor="m", ml=0.04, mr=0.04, mt=0, mb=0)
    phone(s, px, PY0, PW, PH, "●●ホーム", contents[i])
foot(s, "※ 文面は業界知見に基づく初稿です。実績数値・お客様の声は貴社の実データに差し替えます。"
        "「今回は都合が合わない」を必ず置き、押した方には次回開催を自動案内してリストから外しません。")

# ============================================================
# S14 通知メッセージ 利用シーン想定
# ============================================================
s = slides[13]
frame(s, "具体企画③｜LINE通知メッセージの利用シーン想定",
      ["通知メッセージは、電話番号を保有していれば友だち以外にも届く。",
       "＝ 資料請求だけで止まっている過去リストに、もう一度リーチできる。"])
scenes = [
    ("① 見学会・相談の前日リマインド",
     ["予約後〜当日までの3日間にキャンセルは起きる。",
      "前日18:00に地図・持ち物・服装・担当者名を配信。",
      "「日程を変更する」ボタンを堂々と置き、",
      "無断キャンセルをリスト内に残す。"], ORANGE),
    ("② 資料請求だけで止まった過去リストの掘り起こし",
     ["電話が繋がらず追客が止まった層に、番号ベースで再接触。",
      "「無理のない予算がわかる60秒診断」を入口にする。",
      "営業電話ではないため心理的な抵抗が小さい。"], NAVY),
    ("③ 着工〜引渡しの工程通知（契約後）",
     ["基礎・上棟・configure検査・引渡しの節目を自動通知。",
      "「今どうなっているか」の問い合わせ対応が消える。",
      "施主の満足度が上がり、紹介につながる。"], NAVY),
    ("④ 引渡し後の定期点検・メンテ案内（LTV）",
     ["6ヶ月・1年・2年・10年点検の案内を自動化。",
      "外構・リフォーム・紹介キャンペーンの入口になる。",
      "1棟で終わらせず、LTVを伸ばす動線。"], NAVY),
]
cw2 = (CW - 0.60) / 2
ch2 = 5.05
for i, (h, b, hc) in enumerate(scenes):
    cxi = CX0 + (i % 2) * (cw2 + 0.60)
    cyi = 5.05 + (i // 2) * (ch2 + 0.55)
    card(s, cxi, cyi, cw2, ch2, h, b,
         hcol=hc, fill=PORANGE if hc == ORANGE else PALE,
         line=hc if hc == ORANGE else None, hsz=11, bsz=9.5, ls=1.35)
foot(s, "※ 通知メッセージは別途費用。利用にはLINEヤフーの審査と、電話番号の取得経路に関する"
        "同意設計が必要です。導入前に対象リストの取得同意状況を確認します。")

# ============================================================
# S15 飛び道具提案
# ============================================================
s = slides[14]
frame(s, "具体企画④｜飛び道具（他社がやっていない企画軸）",
      ["定番施策だけでは「どこも同じ」になる。注文住宅特有の意思決定構造を突く3案。",
       "いずれも追加の広告費をほぼ必要としない。"])
gimmicks = [
    ("A｜「親に転送される」コンテンツ",
     ["注文住宅は決裁者が3人（ご夫婦＋親）。",
      "ところが各社のLINEはご夫婦向けしかない。",
      "",
      "▶ リッチメニューに「親からの援助・贈与の話」を常設",
      "▶ 相続時精算課税・住宅取得等資金の非課税枠を図解",
      "▶ 奥様が親に転送した瞬間、御社は親にもリーチする",
      "",
      "他社が置いていない枠。転送数はタグで計測できる。"]),
    ("B｜「不利を先に言う」シリーズ",
     ["お客様が最も怖いのは、騙されること。",
      "",
      "▶ Day3で外構費150〜250万円の未計上を自分から暴露",
      "▶ 「うちで建てないほうがいい人」を明記した配信",
      "▶ 高低差・地盤改良で+200〜300万円になる土地の例",
      "",
      "「不利なことを先に言う会社は信用できる」。",
      "これが注文住宅で最強の差別化になる。"]),
    ("C｜構造見学会のLINE限定公開",
     ["完成見学会は各社やっている。差がつかない。",
      "",
      "▶ 建築中の現場（断熱材・耐力壁・配管）をLINE限定で公開",
      "▶ 「見せられる会社」であること自体が証明になる",
      "▶ 検討後期の決め手になり、比較で優位に立てる",
      "",
      "枠が限られるため、友だち限定の希少性が効く。"]),
]
gw = (CW - 1.00) / 3
for i, (h, b) in enumerate(gimmicks):
    card(s, CX0 + i * (gw + 0.50), 5.05, gw, 10.90, h, b,
         hcol=ORANGE, fill=PORANGE, line=ORANGE, hsz=11.5, bsz=9, ls=1.32)
foot(s, "※ B案は「不利の開示」を伴うため、掲載範囲と表現は事前に社内合意が必要です。"
        "C案は現場の安全管理・施主様の許諾が前提になります。")

# ============================================================
# S16 LINEOA実績（LINEヤフー公式事例）
# ============================================================
s = slides[15]
frame(s, "LINEOA実績｜LINEヤフー公式の導入事例（住宅・不動産）",
      ["注文住宅の請負を主業とする会社の「LINE公式アカウント」事例は、公式サイトに掲載がない。",
       "＝ 御社が最初の事例になれる。以下は公式に掲載のある近接事例。"])
cases = [
    ("auka（ギバーテイクオール）", "注文住宅の無料相談窓口",
     ["導線：広告 → LP → 友だち追加 → 家づくりアンケート",
      "2021年3月〜LINE広告／6月〜友だち追加広告（CPF）",
      "クリエイティブ4〜5種を同時運用して効果を分析"],
     ["友だち　約24,000人", "CV　週1〜2件で安定", "CPA　約10,000円（CPF単体）"]),
    ("オープンハウスグループ", "分譲・注文住宅／LINE広告",
     ["2021年10月〜LINE広告をインハウス化",
      "エリア別の顧客データ分析 ＋ 営業担当へのヒアリング",
      "データと「現場の声」をクリエイティブに反映"],
     ["CVR　前年比 158%改善", "CPA　前年比 68%抑制", "資料請求数　142%増"]),
    ("ミサワホーム北越", "住宅建築・リフォーム（新潟）",
     ["LINE公式アカウントで初回リフォーム代金オフを配布",
      "2020年12月〜LINE広告「友だち追加」を利用",
      "興味関心に住宅＋「旅行」「美容」を混ぜてセグメント"],
     ["友だち数　5倍に増加", "友だち　6.6万人", "（2022年8月時点）"]),
    ("LIFULL HOME'S", "不動産情報サイト／LINE公式アカウント",
     ["アンケートでユーザー情報を一元管理しOne to One配信",
      "手動配信をステップ配信に置き換え",
      "通知メッセージを問い合わせ確認に利用"],
     ["CPC　約6割改善", "運用工数　1/4に削減", "通知メッセのCTR　メール比+15%"]),
]
cw3 = (CW - 0.60) / 2
ch3 = 5.35
for i, (name, sub, acts, kpis) in enumerate(cases):
    cxi = CX0 + (i % 2) * (cw3 + 0.60)
    cyi = 5.05 + (i // 2) * (ch3 + 0.45)
    box(s, cxi, cyi, cw3, ch3, fill=WHITE, line=BORDER, lw=1.0)
    T(s, cxi + 0.30, cyi + 0.22, cw3 - 0.60, 1.05,
      [one(name, 11.5, True, NAVY, sa=2), one(sub, 8.5, None, MUT)], ml=0)
    T(s, cxi + 0.30, cyi + 1.42, cw3 * 0.56, 2.6,
      [one("■ 施策", 8.5, True, INK, sa=2)] +
      [one("・" + a, 8, None, INK, ls=1.30, sa=1) for a in acts], ml=0)
    kx = cxi + cw3 * 0.60
    kbox = box(s, kx, cyi + 1.42, cw3 * 0.36, 3.35, fill=PORANGE, line=ORANGE, lw=0.9)
    put_text(kbox.text_frame,
             [one("成果", 8.5, True, ORANGE, align="c", sa=3)] +
             [one(k, 9, True, ORANGE, align="c", ls=1.30, sa=2) for k in kpis],
             anchor="m", ml=0.10, mr=0.10, mt=0.08, mb=0.08)
    T(s, cxi + 0.30, cyi + ch3 - 0.78, cw3 - 0.60, 0.62,
      [one("出典：LINEヤフー for Business 導入事例（lycbiz.com/jp/case-study）", 7, None, MUT)], ml=0)
foot(s, "※ 数値は各事例ページの記載に基づきます（盛り・丸めなし）。資料化にあたっては掲載前に"
        "各ページで再確認してください。ミサワホーム（検討初期層への配信で申込の75%が新規）は"
        "LINE広告ではなくYahoo!広告ディスプレイ広告の事例のため、本ページには収録していません。")

# ============================================================
# S33 診断設問（注文住宅・4問）
# ============================================================
s = slides[32]
frame(s, "具体施策（友だち追加〜短期）｜60秒 予算診断の設問設計",
      ["4問すべて選択式・自由入力ゼロ。設問はそのまま配信セグメントのタグになるものだけに絞る。",
       "アンケートは1通目に置く（参加率が最大化する）。"])
qs = [
    ("Q1", "ご入居はいつ頃をお考えですか？",
     "［1年以内］［1〜2年後］［2年より先］［まだ決めていない］", "検討時期"),
    ("Q2", "土地はお決まりですか？",
     "［所有している］［親の土地を使う予定］［探している最中］［これから探す］", "土地有無"),
    ("Q3", "ご予算のイメージに近いものは？",
     "［〜3,000万円］［3,000〜4,000万円］［4,000〜5,000万円］［まだ分からない］", "予算帯"),
    ("Q4", "いま、いちばん知りたいことは？",
     "［費用・総額］［間取り・デザイン］［土地］［会社の選び方］", "関心テーマ"),
]
qy, qh = 5.05, 1.95
for i, (qn_, q, opts, tag) in enumerate(qs):
    y = qy + i * (qh + 0.30)
    box(s, CX0, y, 16.60, qh, fill=WHITE, line=BORDER, lw=1.0)
    T(s, CX0 + 0.35, y + 0.20, 1.60, 0.6, [one(qn_, 12, True, ORANGE)], ml=0)
    T(s, CX0 + 1.95, y + 0.18, 14.30, 1.60,
      [one(q, 10.5, True, NAVY, sa=3), one(opts, 9, None, INK, ls=1.25)], ml=0)
    tb = box(s, CX0 + 17.20, y + 0.35, 3.60, 1.20, fill=PORANGE, line=ORANGE, lw=0.9)
    put_text(tb.text_frame,
             [one("取れるタグ", 7.5, None, ORANGE, align="c", sa=2),
              one(tag, 10, True, ORANGE, align="c")],
             anchor="m", ml=0.06, mr=0.06, mt=0, mb=0)
    T(s, CX0 + 21.10, y + 0.35, 4.00, 1.20,
      [one({"検討時期": "ホット／ウォーム\n／コールドの判定",
            "土地有無": "土地なし層へ\n土地情報を優先",
            "予算帯": "施工事例の\n価格帯出し分け",
            "関心テーマ": "ステップ配信の\n初回内容を切替"}[tag], 8, None, INK, ls=1.25)],
      anchor="m", ml=0, mr=0)

ny = qy + 4 * (qh + 0.30) + 0.25
card(s, CX0, ny, 12.30, 2.70, "Q3に「まだ分からない」を必ず用意する",
     ["実際に最も多く選ばれるのがこの選択肢。そして、",
      "この層が最も育つ。予算が決まっている人は、",
      "すでに他社と話が進んでいる。"],
     hcol=ORANGE, fill=PORANGE, line=ORANGE, hsz=11, bsz=9.5, ls=1.32)
card(s, CX0 + 12.82, ny, 12.30, 2.70, "診断結果は「共感 → 実例 → 費用 → CTA」の順",
     ["費用を先に出すと逃げられる。まず受け止め、",
      "実例を見せ、そのあとで金額。CTAは2つ用意し、",
      "「もう少し情報だけ」の逃げ道を必ず残す。"],
     hcol=NAVY, fill=PALE, hsz=11, bsz=9.5, ls=1.32)
foot(s, "※ あいさつメッセージでは①担当者名を出す ②有人対応の時間帯を明記 "
        "③アンケートを1通目に置く ④「こちらからお電話はしません」を先に書く、の4点を守ります。")

# ============================================================
# S34 季節企画カレンダー（注文住宅）
# ============================================================
s = slides[33]
frame(s, "具体施策（長期）｜注文住宅の年間サイクルに合わせた企画投稿案",
      ["契約ピーク（1〜3月）ではなく、検討開始ピーク（4〜6月）に仕込むのが要点。",
       "各期の企画は、その時期に立つ検索ニーズを根拠に設計する。"])
season = [
    ("1〜3月", "契約の最繁忙期", "「年度内ご入居・新生活応援フェア」",
     ["4月入学・入社に合わせた入居希望が最大化する時期。",
      "施策：LINEで「年度内着工の締切カレンダー」を配信。",
      "残枠のカウントダウンで意思決定を後押しする。"],
     "KPI：見学会予約数／契約棟数"),
    ("4〜6月", "翌年組の情報収集スタート", "「家づくり はじめの一歩・土地探しセミナー」",
     ["入居ピークと同時に、翌年組が動き始める最重要期。",
      "施策：「土地から探さない」順番を教える無料セミナー。",
      "友だち追加の最大の仕込み期として広告も厚くする。"],
     "KPI：友だち追加数／診断完了率"),
    ("7〜9月", "夏休みの来場が動く", "「夏休み・親子で構造見学会」",
     ["家族で動ける時期。子どもの学区を意識し始める層も多い。",
      "施策：建築中の現場をLINE限定で公開。断熱の体感も。",
      "お子さま同伴前提でキッズスペースを事前案内。"],
     "KPI：来場率／同伴者数"),
    ("10〜12月", "見学会ラッシュ・年内契約", "「冬のあたたかさ体感フェア」",
     ["寒さで断熱性能への関心が最も高まる時期。",
      "施策：実測データ（室温・光熱費）を配信で先出し。",
      "年明けの相談枠を先に押さえる導線を張る。"],
     "KPI：個別相談予約数／成約率"),
]
sw2 = (CW - 1.05) / 4
for i, (per, theme, title, body, kpi) in enumerate(season):
    sx = CX0 + i * (sw2 + 0.35)
    hb = box(s, sx, 5.05, sw2, 1.35, fill=NAVY, radius=0.10)
    put_text(hb.text_frame,
             [one(per, 13, True, WHITE, align="c", sa=2),
              one(theme, 8.5, None, WHITE, align="c")],
             anchor="m", ml=0.08, mr=0.08, mt=0, mb=0)
    bb = box(s, sx, 6.55, sw2, 8.10, fill=PALE, line=BORDER, lw=0.8)
    put_text(bb.text_frame,
             [one(title, 10, True, ORANGE, sa=5)] +
             [one(b, 8.5, None, INK, ls=1.32, sa=2) for b in body],
             anchor="t", ml=0.24, mr=0.20, mt=0.20, mb=0.10)
    kb = box(s, sx, 14.80, sw2, 0.95, fill=PORANGE, line=ORANGE, lw=0.8)
    put_text(kb.text_frame, [one(kpi, 8, True, ORANGE, align="c")],
             anchor="m", ml=0.08, mr=0.08, mt=0, mb=0)
foot(s, "※ 企画の時期設定は業界水準の年間サイクルに基づく初稿です。"
        "Googleトレンドの実データ（5年＋直近1年）が入り次第、ピーク月を実測に合わせて調整します。")

# ============================================================
# S35 診断タイプ別の追客分岐
# ============================================================
s = slides[34]
frame(s, "具体施策（長期）｜診断タイプ別に追客の強弱を決める",
      ["すぐCVしない層に一律配信すると、ブロック率が上がるだけ。",
       "診断で取った4タグから4タイプに分け、配信頻度とオファーの強さを変える。"])
types = [
    ("① 今すぐ建てたい型", "1年以内 × 土地あり",
     ["最短でCVできる層。ここに配信リソースを寄せる。",
      "週2回配信／見学会・個別相談を即オファー。",
      "有人チャットを優先的に割り当てる。"], ORANGE, "強｜週2回＋即オファー"),
    ("② 土地探し型", "土地なし／これから探す",
     ["土地が決まらない限り前に進まない層。",
      "週1回配信／エリア別の土地情報を主コンテンツに。",
      "「土地から探さない」順番の教育を継続。"], NAVY, "中｜週1回＋土地情報"),
    ("③ 予算不安型", "予算「まだ分からない」",
     ["最も人数が多く、最も育つ層。焦らせない。",
      "週1回配信／総額の内訳・返済シミュレーションを軸に。",
      "6日目・14日目のCVオファーは必ず通す。"], NAVY, "中｜週1回＋費用教育"),
    ("④ 情報収集型", "2年より先／まだ決めていない",
     ["今オファーを出すと確実にブロックされる層。",
      "月2回・お役立ちのみ。オファーを外す。",
      "3ヶ月後に再オファー、または再診断で温度感を再取得。"], MUT, "弱｜月2回・オファーなし"),
]
tw = (CW - 1.05) / 4
for i, (name, cond, body, col, freq) in enumerate(types):
    tx = CX0 + i * (tw + 0.35)
    hb = box(s, tx, 5.05, tw, 1.45, fill=col, radius=0.10)
    put_text(hb.text_frame,
             [one(name, 11, True, WHITE, align="c", sa=2),
              one(cond, 8, None, WHITE, align="c")],
             anchor="m", ml=0.08, mr=0.08, mt=0, mb=0)
    bb = box(s, tx, 6.60, tw, 7.00,
             fill=PORANGE if col == ORANGE else PALE, line=col if col == ORANGE else BORDER, lw=0.8)
    put_text(bb.text_frame, [one(b, 9, None, INK, ls=1.38, sa=3) for b in body],
             anchor="t", ml=0.24, mr=0.20, mt=0.22, mb=0.10)
    fb = box(s, tx, 13.75, tw, 1.00, fill=WHITE, line=col, lw=1.0)
    put_text(fb.text_frame, [one(freq, 8.5, True, col, align="c")],
             anchor="m", ml=0.08, mr=0.08, mt=0, mb=0)
band(s, 15.05, "④に「オファーを出さない」と決めるのが、①〜③の配信効果を守る",
     fill=ORANGE, h=0.95, sz=11.5)
foot(s, "※ タイプ判定はQ1（検討時期）×Q2（土地有無）×Q3（予算帯）のタグ組み合わせで自動化します。"
        "配信頻度は運用2ヶ月目のブロック率実測を見て調整します。")

# ============================================================
# S21 施策展開図（8ブロックを注文住宅語に）
# ============================================================
s = slides[20]
blocks21 = {
    "◯LINE友だち追加によるメリット提示":
        "◯LINE友だち追加のメリット提示\n✓60秒予算診断　✓総額の内訳資料\n\n"
        "◯LINE上でのコミュニケーション誘導\n✓LINEで資金相談　✓LINEで土地相談",
    "◯LINEでのナーチャリング":
        "◯LINEでのナーチャリング\n✓家づくりアンケート回答\n\n"
        "◯LINEで予算診断\n✓無理のない総予算レンジを提示",
    "◯企画投稿（月間投稿数）":
        "◯企画投稿（月5本）\n✓年度内入居、夏の構造見学会 等\n\n"
        "ライフイベント→よくある後悔\nアンケ回答→費用・土地の学習配信",
    "◯ステップ配信":
        "◯ステップ配信\n✓友だち追加後14日で設計\n（前後検索ニーズと連動）\n\n"
        "✓特典＞不安解消＞共感＞信頼＞オファー",
    "◯あいさつメッセージ":
        "◯あいさつメッセージ\n（担当者名・営業電話なしを明記）\n\n"
        "◯初期アンケート（4問）\n配信セグメントに活用\n\n◯プロフ・リッチメニュー3タブ",
    "◯(アフター)フォロー":
        "◯(アフター)フォロー\n✓よくあるQAの自動応答\n✓前日リマインド・事前ヒアリング\n\n"
        "◯ご紹介の斡旋\n✓引渡し後の紹介CPへ誘導",
    "◯アクション済":
        "◯見学・相談 実施済\n✓事後フォローで満足度↑\n\n"
        "◯外構・メンテ・追加工事\n✓点検時期に合わせて自動案内",
    "◯再・簡易アンケート(診断)":
        "◯再診断（簡易アンケート）\n✓検討が止まった層に対して\n住まいの軸と異なる診断を用意\n\n"
        "◯再オファー\n✓3ヶ月後に見学会を再提示",
}
for needle, newtext in blocks21.items():
    sh = find_shape(s, needle)
    if sh:
        put_text(sh.text_frame,
                 [one(l, 8, True if l.startswith("◯") else None,
                      NAVY if l.startswith("◯") else INK, ls=1.20)
                  for l in newtext.split("\n")],
                 anchor="t", ml=0.10, mr=0.08, mt=0.06, mb=0.04)

# ============================================================
# S03 / S04 カスタマージャーニー：リード定義を注文住宅のCV3段に
# ============================================================
for idx in (2, 3):
    s = slides[idx]
    sh = find_shape(s, "カスタマージャーニーマップにあわせて")
    if sh:
        put_text(sh.text_frame,
                 [one("カスタマージャーニーマップにあわせてLINE投稿を配置。", 13, True, INK, ls=1.20),
                  one("※CV3段＝【リード】友だち追加・予算診断／【有効リード】見学会・個別相談の予約／"
                      "【マネタイズ】ご契約。", 11, True, INK, ls=1.20)],
                 anchor="m", ml=0.10, mr=0.10, mt=0.04, mb=0.04)

# ============================================================
# 全スライド一括：業種残骸の掃除
# ============================================================
REPLACE = [
    # --- サブスク／EC前提の本文を注文住宅のLTV文脈へ ---
    ("基礎・上棟・configure検査・引渡しの節目を自動通知。",
     "基礎・上棟・竣工検査・引渡しの節目を自動通知。"),
    ("「年末の〇〇に向けて、今から準備を始めるお客様が増えています」という業界固有モーメントに乗せた再接触配信。"
     "悩みのピーク2〜3週前に配信することがCVR最大化のコツ。「売り込み」ではなく「タイミングの案内」として届ける。",
     "「来年4月のご入居に向けて、今から動き始めるお客様が増えています」という業界固有モーメントに乗せた再接触配信。"
     "検討ピークの2〜3ヶ月前に配信することがCVR最大化のコツ。「売り込み」ではなく「タイミングの案内」として届ける。"),
    ("申込を迷っているユーザーに「〇〇を始めて3ヶ月の方の今」「先輩ユーザーの1週間スケジュール」など"
     "リアルな体験談をLINEで定期配信。「買った後どうなるか」の不安を、実際のユーザー事例で解消する。",
     "ご契約を迷っているお客様に「入居して3ヶ月のお住まい」「OB施主の1日の暮らし」など"
     "リアルな体験談をLINEで定期配信。「建てた後どうなるか」の不安を、実際の施主事例で解消する。"),
    ("利用量・頻度タグを活用して「よく使っている層」に上位プラン・関連サービスを自動配信。"
     "「今のプランより月〇〇円で〇〇が追加できます」という具体的メリット提示がCVRを高める。"
     "満足度が高い購入後3ヶ月が最適タイミング。",
     "入居後の経過年数タグを活用し、点検・メンテの時期に合わせて外構・リフォームを自動提案。"
     "「お引渡しから1年、そろそろ外構の追加をご検討の時期です」という具体提案がCVRを高める。"
     "満足度が高い入居後3ヶ月〜1年が最適タイミング。"),
    ("「お友達をご紹介いただいた方にAmazonギフト○○円プレゼント」をLINEで配信。"
     "満足度が最も高い利用開始3ヶ月後が最適タイミング。紹介ユーザーのLTVは通常ユーザーの1.5〜2倍で広告比CPA1/5〜1/10。",
     "「ご友人・ご親族をご紹介いただいた方に特典」をLINEで配信。"
     "満足度が最も高い入居後3ヶ月が最適タイミング。紹介経由のお客様は成約率が高く、獲得コストを大きく抑えられる。"),
    ("誕生日・入会記念日を記念日タグとして登録し「○○様、今日はご誕生日おめでとうございます。"
     "特別クーポンをお贈りします」を自動配信。「自分のことを覚えてくれている」体験がロイヤリティを高め年間LTVを平均20%向上。",
     "お引渡し日を記念日タグとして登録し「お引渡しから1年、おめでとうございます。"
     "1年点検のご案内です」を自動配信。「覚えてくれている」体験が信頼を高め、点検受診率とご紹介につながる。"),
    ("「〇〇診断」でゲーム感覚の登録動線を設計", "「60秒 予算診断」でゲーム感覚の登録動線を設計"),
    ("「あなたの〇〇タイプ診断」など自分に関係する情報を得られる診断コンテンツをLINE登録の特典に設定。"
     "結果表示→そのまま育成配信へ移行。",
     "「あなたの家づくりタイプ診断」など自分に関係する情報を得られる診断コンテンツをLINE登録の特典に設定。"
     "結果表示→そのまま育成配信へ移行。"),
    ("利用状況タグを活用して「よく使っている層」に上位プランの提案を自動配信。"
     "「今のプランより月1,000円で〇〇が使い放題に」という具体的メリット提示がCVRを高める。",
     "入居後の経過年数タグを活用し、点検時期に合わせた外構・リフォーム提案を自動配信。"
     "「1年点検と一緒に、お庭のご相談も承ります」という具体提案がCVRを高める。"),
    ("サービス利用後30日・3ヶ月・6ヶ月のタイミングで自動配信。「その後いかがですか？」"
     "「定期メンテナンスのご案内」など適切なタイミングで次回利用を自然に促す。",
     "お引渡し後30日・3ヶ月・6ヶ月のタイミングで自動配信。「その後お住まいはいかがですか？」"
     "「定期点検のご案内」など適切なタイミングで次のご相談を自然に促す。"),
    ("誕生日や初回購入日を記念日として登録し、該当日に特別クーポン・優待を自動配信。"
     "「あなただけへの特別オファー」感が再購入率を高める。",
     "お引渡し日を記念日として登録し、該当日に点検案内・特別優待を自動配信。"
     "「覚えてくれている」感が、ご紹介と追加工事につながる。"),
    # --- 短い見出し・チップ（長い文字列を先に置換すること）---
    # ※ この行は「リピート率」「クロスセル率」の個別置換より前に置くこと（先勝ち）
    ("リピート率、クロスセル率、客単価を改善し、LTVを最大化する。",
     "ご紹介率・追加工事率・1棟あたり粗利を改善し、LTVを最大化する。"),
    ("利用状況タグ × アップグレード提案で客単価を向上", "経過年数タグ × 点検連動提案で追加工事を獲得"),
    ("利用量・頻度に応じたアップグレード提案", "入居後年数に応じた外構・リフォーム提案"),
    ("誕生日記念日タグ × パーソナライズ配信でLTV向上", "お引渡し記念日タグ × 点検案内でLTV向上"),
    ("誕生日・契約記念日に特別オファーを自動配信", "お引渡し記念日・点検時期に特別オファーを自動配信"),
    ("カゴ落ち・問い合わせのみ層へのリマインド配信", "資料請求のみ・見学未来場層へのリマインド配信"),
    ("関連サービス・上位プランへの誘導設計", "外構・リフォーム・ご紹介への誘導設計"),
    ("→ 商談・再購入誘導", "→ 商談・ご紹介誘導"),
    ("利用状況タグの設定", "入居後年数タグの設定"),
    ("利用状況タグ付与", "経過年数タグ付与"),
    ("利用状況連動", "経過年数連動"),
    ("高頻度ユーザー特定", "入居後年数で特定"),
    ("アップグレード提案", "外構・リフォーム提案"),
    ("アップセル提案", "点検連動提案"),
    ("客単価 -- %向上", "追加工事の受注"),
    ("客単価 ── %向上", "追加工事の受注"),
    ("再購入率10%前後", "ご紹介の発生"),
    ("誕生日・入会日等", "お引渡し日・点検時期"),
    ("プラン変更CV", "リフォームCV"),
    ("再購入CV", "追加工事CV"),
    ("購入完了", "お引渡し"),
    ("定期購入・クロスセル", "外構・メンテ・ご紹介"),
    ("クロスセル・アップセル", "外構・リフォーム提案"),
    ("再購入・クロスセル", "追加工事・ご紹介"),
    ("クロスセル配信", "外構・メンテ提案配信"),
    ("クロスセルCV", "外構・リフォームCV"),
    ("クロスセル施策", "LTV拡張施策"),
    ("クロスセル率", "紹介・追加工事率"),
    ("クロスセル", "追加工事・紹介"),
    ("〇〇業界", "注文住宅業界"),
    ("○○業界", "注文住宅業界"),
    ("人材(転職)業界の例", "注文住宅業界の設計"),
    ("即転職型", "今すぐ建てたい型"),
    ("転職成功事例・年収比較データ", "施工事例・総額の内訳データ"),
    ("転職", "家づくり"),
    ("求人", "施工事例"),
    ("リピート率", "紹介率"),
    ("再購入時期、併売商品を設定", "点検時期・追加工事の提案時期を設定"),
    ("商材への動線設計", "見学会予約・ご契約への動線設計"),
]


def sweep(slide):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            if not para.runs:
                continue
            txt = "".join(r.text for r in para.runs)
            new = txt
            for a, b in REPLACE:
                new = new.replace(a, b)
            if new != txt:
                para.runs[0].text = new
                for r in para.runs[1:]:
                    r._r.getparent().remove(r._r)


for sl in slides:
    sweep(sl)

prs.save(OUT)
print("saved:", OUT)
print("slides:", len(Presentation(OUT).slides))
