# -*- coding: utf-8 -*-
"""八十吉 パフォーマー指名予約 ご予約の流れ（確認用）"""
import sys, os
SKILL = r'C:\Users\tonomura-r\.claude\skills\dym-format'
sys.path.insert(0, os.path.join(SKILL, 'scripts'))

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from helpers import copy_template, safe_save, delete_slide, duplicate_slide, \
    replace_text_in_slide, COLORS

OUT = sys.argv[1]
prs = copy_template(os.path.join(SKILL, 'assets', 'template.pptx'),
                    os.path.join(os.path.dirname(__file__), '_flow.pptx'))

BLANK = 3
N_ORIG = len(prs.slides._sldIdLst)
for _ in range(1):
    duplicate_slide(prs, BLANK)
for idx in range(N_ORIG - 1, -1, -1):
    delete_slide(prs, idx)
S = list(prs.slides)

NAVY = COLORS['navy_700']
ACCENT = COLORS['accent']
INK = COLORS['ink']
MUTED = COLORS['ink_muted']
SW = 10.8333

EXIST_FILL = 'E8EBF2'      # 既存の仕組み
EXIST_LINE = NAVY
LINE_FILL = 'FDF0E6'       # 今回LINEで作る
LINE_LINE = ACCENT


def set_title(slide, text):
    replace_text_in_slide(slide, {'［スライドタイトルを入力］': text})
    sh = list(slide.shapes)[0]
    for p in sh.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(28)
            r.font.bold = True
            r.font.name = 'メイリオ'
            r.font.color.rgb = RGBColor.from_string(NAVY)


def text(slide, left, top, width, height, lines, size=16, bold=False,
         color=INK, space=6, align=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space)
        if align is not None:
            p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = 'メイリオ'
        r.font.color.rgb = RGBColor.from_string(color)
    return tb


def arrow(slide, cx, cy, w=0.30):
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(cx - w / 2),
                               Inches(cy - 0.10), Inches(w), Inches(0.20))
    a.fill.solid()
    a.fill.fore_color.rgb = RGBColor.from_string(COLORS['navy_300'])
    a.line.fill.background()
    a.shadow.inherit = False


def box(slide, x, y, w, h, fill, edge, num, head, where, detail):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = RGBColor.from_string(fill)
    sh.line.color.rgb = RGBColor.from_string(edge)
    sh.line.width = Pt(1.5)
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = Pt(4)
    tf.margin_top = Pt(9)
    # (テキスト, サイズ, 太字, 色, ブロック末の余白)
    rows = [(num, 14, True, edge, 3),
            (head, 15, True, INK, 6),
            (where, 11.5, False, edge, 6),
            (detail, 11, False, MUTED, 0)]
    first = True
    for t, sz, bd, col, sp in rows:
        lines = t.split('\n')
        for j, ln in enumerate(lines):
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(sp if j == len(lines) - 1 else 0)
            r = p.add_run()
            r.text = ln
            r.font.size = Pt(sz)
            r.font.bold = bd
            r.font.name = 'メイリオ'
            r.font.color.rgb = RGBColor.from_string(col)
    return sh


# =====================================================================
# 1枚目：全体の道順
# =====================================================================
s = S[0]
set_title(s, 'ご予約の流れ（ご確認用）')

STEPS = [
    ('①', 'お客様が\nご予約', 'ぐるなび・食べログ', 'パフォーマーを\nご指名', False),
    ('②', 'パフォーマー\nへ連絡', 'LINE', '指名が入ったことを\n通知', True),
    ('③', 'パフォーマー\nが了承', 'LINE', '受けられるかを\n返信', True),
    ('④', 'お席を\n確保', 'サイトコントローラー', '既存の在庫を\n押さえる', False),
    ('⑤', '店長へ\nお知らせ', 'LINE', '成立したことを\n共有', True),
]

BW, BH, GAP = 1.94, 2.15, 0.23
total = BW * 5 + GAP * 4
x0 = (SW - total) / 2
y0 = 1.55

for i, (num, head, where, detail, is_line) in enumerate(STEPS):
    x = x0 + i * (BW + GAP)
    box(s, x, y0, BW, BH,
        LINE_FILL if is_line else EXIST_FILL,
        LINE_LINE if is_line else EXIST_LINE,
        num, head, where, detail)
    if i < 4:
        arrow(s, x + BW + GAP / 2, y0 + BH / 2)

# 凡例
ly = y0 + BH + 0.42
lx = x0
for label, fill, edge in [('いまお使いの仕組み', EXIST_FILL, EXIST_LINE),
                          ('今回LINEで作るところ', LINE_FILL, LINE_LINE)]:
    sq = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(lx), Inches(ly),
                            Inches(0.28), Inches(0.28))
    sq.fill.solid()
    sq.fill.fore_color.rgb = RGBColor.from_string(fill)
    sq.line.color.rgb = RGBColor.from_string(edge)
    sq.line.width = Pt(1.5)
    sq.shadow.inherit = False
    text(s, lx + 0.36, ly - 0.06, 2.6, 0.4, [label], size=13, color=INK)
    lx += 3.1

bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x0), Inches(ly + 0.62),
                         Inches(total), Inches(0.72))
bar.fill.solid()
bar.fill.fore_color.rgb = RGBColor.from_string('FDF0E6')
bar.line.fill.background()
bar.shadow.inherit = False
tf = bar.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = 'ご予約の受付とお席の管理は、いまの仕組みのまま。LINEは「連絡」だけを引き受けます。'
r.font.size = Pt(17)
r.font.bold = True
r.font.name = 'メイリオ'
r.font.color.rgb = RGBColor.from_string(ACCENT)

safe_save(prs, OUT)
print('SAVED:', OUT, len(prs.slides._sldIdLst), 'slides')
