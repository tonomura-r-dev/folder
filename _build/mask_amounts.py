# -*- coding: utf-8 -*-
"""クライアント固有の金額・通数を伏せて、横展開できる雛形にする。

半透明の白をかぶせる方式ではなく、**テキスト自体を差し替える**。
かぶせるだけだと、コピペやテキスト抽出で元の数字が読めてしまうため。

伏せるもの：
  ・金額（¥/\\ 付き、カンマ区切り）        → ¥---
  ・通数など カンマ区切りの数値            → ---
  ・単価など 小数（¥1.40, ¥2.5 …）         → ¥-.--
  ・弊社の還元率（従量課金分の◯%をご提供） → 〇%
  ・本文中の「4500万弱」などの概算          → 〇〇〇〇万弱
  ・クライアント名                          → 〇〇

残すもの：
  ・期間（202508 などの YYYYMM）
  ・¥0 や -- などの構造上の値
  ・表の見出し行

使い方:
    python _build/mask_amounts.py <入力.pptx> [出力.pptx] [クライアント名] [--skip=2,10]

--skip には、媒体の公開料金だけを載せているページなど、
マスクすると資料として成立しなくなるページ番号を渡す。
"""
import re
import shutil
import sys
from pathlib import Path
from pptx import Presentation

SIGN = r'[-+±]?'
MONEY = re.compile(rf'^{SIGN}[¥\\]\s?{SIGN}\d{{1,3}}(?:,\d{{3}})+(?:\.\d+)?$')
COUNT = re.compile(rf'^{SIGN}\d{{1,3}}(?:,\d{{3}})+$')
UNIT = re.compile(rf'^{SIGN}[¥\\]?\s?{SIGN}\d+\.\d+$')
YYYYMM = re.compile(r'^\d{6}')
KEEP = {"¥0", "\\0", "--", "-", ""}

MASK_MONEY, MASK_COUNT, MASK_UNIT = "¥---", "---", "¥-.--"


def replace_runs(tf, new_text):
    """1つ目の run の書式を残したままテキストを差し替える。"""
    for para in list(tf.paragraphs)[1:]:
        para._p.getparent().remove(para._p)
    para = tf.paragraphs[0]
    runs = list(para.runs)
    if not runs:
        return False
    runs[0].text = new_text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)
    return True


def mask_value(text):
    """マスク後の文字列を返す。対象外なら None。"""
    t = text.strip()
    if t in KEEP or YYYYMM.match(t):
        return None
    if MONEY.match(t):
        return MASK_MONEY
    if COUNT.match(t):
        return MASK_COUNT
    if UNIT.match(t):
        return MASK_UNIT
    return None


# 文中の数字を潰すパターン（順に適用）
INLINE = [
    (re.compile(r'\d{3,4}万弱'), '〇〇〇〇万弱'),
    (re.compile(r'\d{3,4}万'), '〇〇〇〇万'),
    (re.compile(r'[-+±]?[¥\\]\s?[-+±]?\d{1,3}(?:,\d{3})+'), MASK_MONEY),
    (re.compile(r'(従量課金分の)\s*[０-９0-9]+\s*(%|％)'), r'\1〇\2'),
]


def mask_inline(text):
    out = text
    for pat, rep in INLINE:
        out = pat.sub(rep, out)
    return out if out != text else None


def mask_deck(src, out, client=None, skip=()):
    shutil.copyfile(src, out)
    prs = Presentation(out)
    report = []

    for idx, s in enumerate(prs.slides, 1):
        if idx in skip:
            continue
        for sh in s.shapes:
            if sh.has_table:
                for ri, row in enumerate(sh.table.rows):
                    for ci, cell in enumerate(row.cells):
                        before = cell.text
                        new = mask_value(before)
                        if new and replace_runs(cell.text_frame, new):
                            report.append((idx, f"表 r{ri}c{ci}", before, new))
            if not sh.has_text_frame:
                continue
            original = sh.text_frame.text
            if not original.strip():
                continue
            new = mask_inline(original)
            if client and client in (new or original):
                new = (new or original).replace(client, "〇〇")
            if new and new != original and replace_runs(sh.text_frame, new):
                report.append((idx, sh.name[:20], original[:40], new[:40]))

    prs.save(out)
    return out, report


if __name__ == "__main__":
    if len(sys.argv) < 2:
        sys.exit(__doc__)
    args = [a for a in sys.argv[1:] if not a.startswith("--skip=")]
    skip = {int(n) for a in sys.argv[1:] if a.startswith("--skip=")
            for n in a.split("=", 1)[1].split(",") if n}
    src = Path(args[0])
    out = Path(args[1]) if len(args) > 1 else src.with_name(src.stem + "_横展開用.pptx")
    client = args[2] if len(args) > 2 else None
    out, report = mask_deck(str(src), str(out), client, skip)
    for slide, where, before, after in report:
        print(f"S{slide:>2} {where:<22} {before!r} -> {after!r}")
    print(f"\n{len(report)} 箇所をマスクしました -> {out}")
