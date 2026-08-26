# -*- coding: utf-8 -*-
"""既存デッキの S9「特別プランのご提供について」を作り直す。

S8（LINEヤフー社の料金改定で調整プランが一度停止）を受けた続きのページとして、
「10月＝一度停止 → 11月以降＝弊社で調整」の時系列で組み直す。

使い方:
    python _build/rebuild_s9_special_plan.py <入力.pptx> [出力.pptx]

※ 金額・率（0.5%確約 など）は一切載せない。S9の元テキストにあった
   「（0.5%確約）」は意図的に落としている。
"""
import shutil
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

TARGET_SLIDE = 9

NAVY = "1F285A"    # 打ち手・カード見出し
RED = "FF0000"     # このデッキの強調色（S4/S8 に準拠）
INK = "333333"     # 本文
MUT = "808080"
GREY = "8C8C8C"    # 済んだ話・トーンを落とすチップ
WHITE = "FFFFFF"
PALE = "F4F7FF"
FADE = "F2F2F2"
BORDER = "D9D9D9"


def set_font(run, size, bold=None, color=INK, name="メイリオ"):
    f = run.font
    f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", name)
    f.color.rgb = RGBColor.from_string(color)


ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
ANCH = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}


def put_text(tf, paras, anchor="t", ml=0.06, mr=0.06, mt=0.03, mb=0.03, wrap=True):
    tf.word_wrap = wrap
    tf.margin_left = Inches(ml)
    tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt)
    tf.margin_bottom = Inches(mb)
    tf.vertical_anchor = ANCH[anchor]
    first = True
    for p in paras:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = ALIGN[p.get("align", "l")]
        if p.get("sa") is not None:
            para.space_after = Pt(p["sa"])
        if p.get("ls") is not None:
            para.line_spacing = p["ls"]
        for (t, sz, b, c) in p["runs"]:
            r = para.add_run()
            r.text = t
            set_font(r, sz, b, c)
    return tf


def add_text(slide, x, y, w, h, paras, anchor="t", **kw):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    put_text(box.text_frame, paras, anchor=anchor, **kw)
    return box


def add_box(slide, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor.from_string(line)
        sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    sp.text_frame.word_wrap = True
    return sp


def chip(slide, x, y, w, h, text, fill, size=14):
    c = add_box(slide, x, y, w, h, fill=fill)
    put_text(c.text_frame, [{"runs": [(text, size, True, WHITE)], "align": "c"}],
             anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
    return c


def clear_slide(slide):
    spTree = slide.shapes._spTree
    for el in list(spTree):
        if el.tag.split("}")[-1] in ("sp", "cxnSp", "pic", "graphicFrame", "grpSp"):
            spTree.remove(el)


def build(src, out):
    shutil.copyfile(src, out)
    prs = Presentation(out)
    s = prs.slides[TARGET_SLIDE - 1]
    clear_slide(s)

    # タイトル（このデッキの実測位置・16pt メイリオ）
    add_text(s, 0.47, 0.10, 8.22, 0.40,
             [{"runs": [("料金改定へのご対応について", 16, None, NAVY)]}],
             anchor="m", ml=0.1, mr=0.0)

    # リード（S8 と同じ位置・16pt bold）＝このページの結論
    add_text(s, 0.16, 0.88, 10.58, 0.64,
             [{"runs": [("プランのご改定でご対応いたします。開始は10月からではなく、", 16, True, INK),
                        ("11月から", 16, True, RED),
                        ("のご案内です。", 16, True, INK)], "align": "c"}],
             anchor="m", ml=0.0, mr=0.0)

    def step(y, h, label, label_fill, paras, fill, line, lw=1.0, chip_size=14):
        add_box(s, 0.50, y, 9.83, h, fill=fill, line=line, lw=lw)
        chip(s, 0.85, y + (h - 0.44) / 2, 1.45, 0.44, label, label_fill, size=chip_size)
        add_text(s, 2.65, y, 7.40, h, paras, anchor="m", ml=0.0)

    def down_arrow(y):
        a = add_box(s, 5.19, y, 0.45, 0.32, fill=BORDER, shape=MSO_SHAPE.DOWN_ARROW)
        a.line.fill.background()

    # ===== ① 10月〜：料金改定。放置するとコストが上振れ =====
    step(1.72, 1.10, "10月〜", GREY,
         [{"runs": [("LINEヤフー社の料金プランが改定されます。", 15, None, INK)], "ls": 1.2},
          {"runs": [("このままの運用では、コストが上振れします。", 15, True, RED)], "ls": 1.2}],
         fill=FADE, line=BORDER)

    down_arrow(2.90)

    # ===== ② ご対応：プランの改定 =====
    step(3.30, 0.80, "ご対応", NAVY,
         [{"runs": [("プランのご改定により、コストを抑えます。", 16, True, INK)]}],
         fill=WHITE, line=NAVY)

    down_arrow(4.18)

    # ===== ③ 11月〜：対応開始時期（このページの主役） =====
    step(4.58, 1.90, "11月〜", NAVY,
         [{"runs": [("ご対応の開始は、10月からではなく11月からとなります", 12.5, None, MUT)], "sa": 6},
          {"runs": [("各社様の配信金額と媒体インセンティブの兼ね合いから、", 18, True, INK)], "ls": 1.25},
          {"runs": [("弊社にて特別プランの調整をいたします。", 18, True, INK)], "ls": 1.25}],
         fill=PALE, line=NAVY, lw=1.5)

    prs.save(out)
    return out


if __name__ == "__main__":
    if len(sys.argv) < 2:
        sys.exit(__doc__)
    src = Path(sys.argv[1])
    out = Path(sys.argv[2]) if len(sys.argv) > 2 else src.with_name(src.stem + "_S9修正.pptx")
    print("saved:", build(str(src), str(out)))
