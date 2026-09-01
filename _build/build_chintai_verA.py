# -*- coding: utf-8 -*-
"""賃貸業界 Ver.A（入居者集客版）LINEOA施策提案（36枚）

_templates/DYM_LINEOA_FMT.pptx（業界汎用FMT・36枚）をコピーし、
全36枚を _drafts/賃貸業界_VerA_スライド作成プロンプト.md の指定内容で構築する。

  python _build/build_chintai_verA.py
  python _build/qa_render.py 賃貸業界_LINEOA施策提案_VerA_入居者集客.pptx

【この資料の背骨】
ポータル反響は、いま「使い捨て」になっている。Webマーケでできるのは、
反響を資産に変えてCPOを下げること。接客の質ではなく、反響の再利用率で勝つ。

【実データは取得済み（EC版と違い★差込枠が少ない）】
- Googleトレンド：_data/trends/verA_1year.csv（実測）→ _images/verA_trend_*.png
- LINEヤフー前後検索：取得済み（S08に反映）
- LINEOA実績：ハートサポート／LIFULL HOME'Sの実数値（LY公式）

【未取得データ（★破線の差込枠）】
- S09 競合2社（いい部屋ネット本部／エイブルAGENT）の友だち数
- S10 大手FC店舗別アカウントの友だち数分散状況

【落とし穴メモ（_build/README.md・deck-apply スキルより）】
- put_text() は必ず reset_tf() を通す（既存テキストへの追記事故を防ぐ）
- スライドの新規追加はしない。ベース36枚を clear_slide() して作り直す
- 一括置換をしない。数値は行・列を特定して書く
"""
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
SRC = str(ROOT / "_templates" / "DYM_LINEOA_FMT.pptx")
OUT = str(ROOT / "賃貸業界_LINEOA施策提案_VerA_入居者集客.pptx")
IMG = ROOT / "_images"

# ---- 配色（DYM 3色：ネイビー・炭黒・オレンジ。LINEグリーンはUI部分のみ）----
TNAVY = "002060"   # タイトル
NAVY = "1F285A"    # 打ち手・カード見出し
ORANGE = "ED7D31"  # 強調
RED = "C00000"     # 課題・警告
INK = "333333"     # 本文
MUT = "7F7F7F"
WHITE = "FFFFFF"
PALE = "F4F7FF"    # 淡ネイビー
PORANGE = "FCE4D6"  # 淡オレンジ
GREY = "F2F2F2"
BORDER = "D9D9D9"
GREEN = "06C755"   # LINE UI 専用
BEZEL = "2B2B2B"
SCREEN = "EFF2F7"
PRED = "FDF2F2"

# ---- FMTのレイアウト座標（cm）----
SW, SH = 27.52, 19.05
TITLE_XY = (1.52, 0.38, 24.4, 0.94)
LEAD_XY = (1.20, 1.80, 25.1, 1.90)
DIV_Y = 3.86
CX0, CW = 1.20, 25.12
CY0, CY1 = 4.30, 17.00
FOOT_Y = 17.35

shutil.copyfile(SRC, OUT)
prs = Presentation(OUT)
slides = list(prs.slides)
assert len(slides) == 36, len(slides)
assert (prs.slide_width, prs.slide_height) == (9906000, 6858000)


# ================= helpers（build_chumon_jutaku.py と共通の作法）=================
def set_font(run, size, bold=None, color=INK, name="メイリオ"):
    f = run.font
    f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", name)
    f.color.rgb = RGBColor.from_string(color)


ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
ANCH = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}


def reset_tf(tf):
    """既存の段落・runを全消去（追記事故の防止）"""
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)


def put_text(tf, paras, anchor="t", ml=0.14, mr=0.14, mt=0.06, mb=0.06, wrap=True):
    reset_tf(tf)
    tf.word_wrap = wrap
    tf.margin_left = Cm(ml)
    tf.margin_right = Cm(mr)
    tf.margin_top = Cm(mt)
    tf.margin_bottom = Cm(mb)
    tf.vertical_anchor = ANCH[anchor]
    first = True
    for p in paras:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = ALIGN[p.get("align", "l")]
        if p.get("sa") is not None:
            para.space_after = Pt(p["sa"])
        if p.get("ls") is not None:
            para.line_spacing = p["ls"]
        for item in p["runs"]:
            t, sz, b, c = item
            r = para.add_run()
            r.text = t
            set_font(r, sz, b, c)
    return tf


def T(slide, x, y, w, h, paras, anchor="t", **kw):
    box_ = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    put_text(box_.text_frame, paras, anchor=anchor, **kw)
    return box_


def one(text, sz, b=None, c=INK, align="l", sa=None, ls=None):
    d = {"runs": [(text, sz, b, c)], "align": align}
    if sa is not None:
        d["sa"] = sa
    if ls is not None:
        d["ls"] = ls
    return d


def multi(runs, align="l", sa=None, ls=None):
    """1段落に複数run（色・太さを混在させたいとき）"""
    d = {"runs": runs, "align": align}
    if sa is not None:
        d["sa"] = sa
    if ls is not None:
        d["ls"] = ls
    return d


def box(slide, x, y, w, h, fill=None, line=None, lw=1.0,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06, dash=None):
    sp = slide.shapes.add_shape(shape, Cm(x), Cm(y), Cm(w), Cm(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor.from_string(line)
        sp.line.width = Pt(lw)
        if dash:
            sp.line.dash_style = dash
    sp.shadow.inherit = False
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    reset_tf(sp.text_frame)
    return sp


def card(slide, x, y, w, h, head, body, hcol=NAVY, fill=PALE,
         hsz=11, bsz=9, line=None, anchor="t", ls=1.18):
    sp = box(slide, x, y, w, h, fill=fill, line=line)
    paras = [one(head, hsz, True, hcol, sa=3)]
    for b in (body if isinstance(body, list) else [body]):
        if b:
            paras.append(one(b, bsz, None, INK, ls=ls, sa=1))
    put_text(sp.text_frame, paras, anchor=anchor, ml=0.22, mr=0.18, mt=0.14, mb=0.10)
    return sp


def clear_slide(slide):
    spTree = slide.shapes._spTree
    for el in list(spTree):
        if el.tag.split("}")[-1] in ("sp", "cxnSp", "pic", "graphicFrame", "grpSp"):
            spTree.remove(el)


def frame(slide, title, lead):
    """FMT準拠：タイトル16pt紺（y=0.38）＋リード12pt（y=1.80・2行以内）＋区切り線"""
    clear_slide(slide)
    T(slide, *TITLE_XY, [one(title, 16, True, TNAVY)], anchor="m", ml=0, mr=0)
    T(slide, *LEAD_XY, [one(l, 12, None, INK, ls=1.28) for l in lead],
      anchor="m", ml=0, mr=0)
    ln = slide.shapes.add_connector(1, Cm(0), Cm(DIV_Y), Cm(SW), Cm(DIV_Y))
    ln.line.color.rgb = RGBColor.from_string(BORDER)
    ln.line.width = Pt(1.0)


def foot(slide, text):
    T(slide, CX0, FOOT_Y, CW, 0.9, [one(text, 7.5, None, MUT, ls=1.15)], ml=0, mr=0)


def badge(slide, x, y, w, h, text, fill=PORANGE, col=ORANGE, sz=8.5, dash=None):
    sp = box(slide, x, y, w, h, fill=fill, line=col, lw=1.0, dash=dash)
    put_text(sp.text_frame, [one(text, sz, True, col, align="c")],
             anchor="m", ml=0.06, mr=0.06, mt=0, mb=0)
    return sp


def band(slide, y, text, fill=NAVY, col=WHITE, sz=11.5, h=0.86, x=CX0, w=CW):
    sp = box(slide, x, y, w, h, fill=fill, radius=0.10)
    put_text(sp.text_frame, [one(text, sz, True, col, align="c")],
             anchor="m", ml=0.2, mr=0.2, mt=0, mb=0)
    return sp


def placeholder(slide, x, y, w, h, label, note):
    """★未取得データの差込枠（破線）。捏造しない。"""
    sp = box(slide, x, y, w, h, fill=WHITE, line=MUT, lw=1.25,
              dash=MSO_LINE_DASH_STYLE.DASH)
    put_text(sp.text_frame,
              [one(label, 11, True, MUT, align="c", sa=4),
               one(note, 8.5, None, MUT, align="c", ls=1.25)],
              anchor="m", ml=0.3, mr=0.3, mt=0.1, mb=0.1)
    return sp


def phone(slide, x, y, w, h, header, lines):
    """LINEトーク画面のモック。lines = [(kind, text)]
    kind: 'in'（BOT吹き出し）/ 'chip'（選択肢）/ 'btn'（CTA）/ 'note'（注記）"""
    box(slide, x, y, w, h, fill=BEZEL, radius=0.10)
    box(slide, x + 0.14, y + 0.42, w - 0.28, h - 0.56, fill=SCREEN, radius=0.02,
        shape=MSO_SHAPE.RECTANGLE)
    hd = box(slide, x + 0.14, y + 0.42, w - 0.28, 0.52, fill=GREEN,
             shape=MSO_SHAPE.RECTANGLE)
    put_text(hd.text_frame, [one(header, 7, True, WHITE, align="c")],
             anchor="m", ml=0.05, mr=0.05, mt=0, mb=0)
    cy = y + 1.05
    iw = w - 0.62
    for kind, text in lines:
        nlines = text.count("\n") + 1
        bh = 0.30 + 0.30 * nlines
        if kind == "in":
            sp = box(slide, x + 0.31, cy, iw, bh, fill=WHITE, line=BORDER, lw=0.75)
            put_text(sp.text_frame, [one(l, 6.5, None, INK, ls=1.18)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.14, mr=0.10, mt=0.05, mb=0.05)
        elif kind == "chip":
            sp = box(slide, x + 0.31, cy, iw, bh, fill=WHITE, line=GREEN, lw=0.9)
            put_text(sp.text_frame, [one(l, 6.5, None, "0B7A3B", align="c", ls=1.15)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.08, mr=0.08, mt=0.03, mb=0.03)
        elif kind == "btn":
            sp = box(slide, x + 0.31, cy, iw, bh, fill=GREEN)
            put_text(sp.text_frame, [one(l, 6.8, True, WHITE, align="c", ls=1.15)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.08, mr=0.08, mt=0.03, mb=0.03)
        else:  # note
            sp = slide.shapes.add_textbox(Cm(x + 0.31), Cm(cy), Cm(iw), Cm(bh))
            put_text(sp.text_frame, [one(l, 6.2, None, MUT, ls=1.15)
                                     for l in text.split("\n")],
                     anchor="m", ml=0.05, mr=0.05, mt=0, mb=0)
        cy += bh + 0.12
    return cy


def richmenu(slide, x, y, w, h, tabs):
    """リッチメニュー3タブ×6ボタンのモック。
    tabs = [(タブ名, [ボタン文言6個]), ...]。最初のタブをアクティブ表示。"""
    tab_h = 0.55
    box(slide, x, y, w, h, fill=BEZEL, radius=0.06)
    tw = (w - 0.12) / len(tabs)
    for i, (name, _) in enumerate(tabs):
        active = (i == 0)
        tb = box(slide, x + 0.06 + i * tw, y + 0.06, tw - 0.03, tab_h,
                  fill=GREEN if active else "3F3F3F",
                  shape=MSO_SHAPE.RECTANGLE)
        put_text(tb.text_frame, [one(name, 8, True, WHITE, align="c")],
                 anchor="m", ml=0, mr=0, mt=0, mb=0)
    _, btns = tabs[0]
    gx, gy = x + 0.06, y + 0.06 + tab_h + 0.06
    gw, gh = w - 0.12, h - tab_h - 0.18
    cols, rows = 3, 2
    bw, bh = (gw - 0.06 * (cols - 1)) / cols, (gh - 0.06 * (rows - 1)) / rows
    for i, label in enumerate(btns):
        r, c = divmod(i, cols)
        bb = box(slide, gx + c * (bw + 0.06), gy + r * (bh + 0.06), bw, bh,
                  fill=SCREEN, line=BORDER, lw=0.75, radius=0.04)
        put_text(bb.text_frame, [one(l, 7.2, True, NAVY, align="c", ls=1.1)
                                 for l in label.split("\n")],
                 anchor="m", ml=0.05, mr=0.05, mt=0, mb=0)


def simple_table(slide, x, y, w, h, headers, rows, col_w=None,
                  hsz=9, bsz=8.5, header_fill=NAVY, zebra=PALE,
                  align=None, row_h=None):
    """罫線テーブル（ネイティブpptxテーブル）。headers=[str], rows=[[str]]"""
    n_r, n_c = len(rows) + 1, len(headers)
    gf = slide.shapes.add_table(n_r, n_c, Cm(x), Cm(y), Cm(w), Cm(h))
    tbl = gf.table
    if col_w:
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Cm(cw)
    if row_h:
        for r in tbl.rows:
            r.height = Cm(row_h)
    for j, htext in enumerate(headers):
        c = tbl.cell(0, j)
        c.fill.solid()
        c.fill.fore_color.rgb = RGBColor.from_string(header_fill)
        c.margin_left = c.margin_right = Cm(0.12)
        c.margin_top = c.margin_bottom = Cm(0.05)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        put_text(c.text_frame, [one(htext, hsz, True, WHITE,
                                    align=(align[j] if align else "c"))],
                 anchor="m", ml=0.12, mr=0.12, mt=0.02, mb=0.02)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = tbl.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor.from_string(
                zebra if (zebra and i % 2 == 0) else WHITE)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            put_text(c.text_frame, [one(str(val), bsz, None, INK,
                                        align=(align[j] if align else "l"))],
                     anchor="m", ml=0.12, mr=0.12, mt=0.02, mb=0.02)
    return gf


def find_shape(slide, needle):
    for sh in slide.shapes:
        if sh.has_text_frame and needle in sh.text_frame.text:
            return sh
    return None


# ============================================================
# S01 表紙
# ============================================================
s = slides[0]
sh = find_shape(s, "業界：")
if sh:
    put_text(sh.text_frame, [one("業界：賃貸業界（入居者集客・仲介＋管理併営）", 12, True, TNAVY)], anchor="m")
sh = find_shape(s, "最上位パートナーの知見")
if sh:
    put_text(sh.text_frame,
             [one("反響を「資産」に変えて、CPOを下げる。接客の質ではなく再利用率で勝つ", 13, True, INK)],
             anchor="m")

# ============================================================
# S02 資料アジェンダ
# ============================================================
s = slides[1]
frame(s, "資料アジェンダ", ["この資料が言っていること：ポータル反響は、いま「使い捨て」になっている。",
                          "Webマーケでできるのは、反響を資産に変えてCPOを下げること。"])
AGENDA = [
    ("1", "なぜLINEか", "S03-06"), ("2", "実態調査", "S07-10"),
    ("3", "全体設計", "S11-14"), ("4", "構築", "S15-19"),
    ("5", "配信設計", "S20-25"), ("6", "来店改善・工数", "S26-28"),
    ("7", "成果と体制", "S29-32"), ("8", "締め", "S33-36"),
]
gx, gy, gw, gh = CX0, CY0 + 0.3, CW, 11.6
cw = gw / 4
ch = gh / 2
for i, (n, t, r) in enumerate(AGENDA):
    row, col = divmod(i, 4)
    x = gx + col * cw
    y = gy + row * ch
    box(s, x + 0.12, y + 0.12, cw - 0.24, ch - 0.24, fill=PALE)
    T(s, x + 0.12, y + 0.22, cw - 0.24, 0.9, [one(n, 22, True, ORANGE, align="c")], anchor="t")
    T(s, x + 0.12, y + 0.95, cw - 0.24, 1.3, [one(t, 12.5, True, NAVY, align="c")], anchor="t")
    T(s, x + 0.12, y + ch - 0.55, cw - 0.24, 0.5, [one(r, 8.5, None, MUT, align="c")], anchor="t")
foot(s, "全8章・36枚｜主語は「LINE公式アカウントの運用」。広告運用の提案書ではない")

# ============================================================
# S03 検討背景
# ============================================================
s = slides[2]
frame(s, "検討背景",
      ["ユーザーはアプリ登録も会員登録もしたがらない。唯一ハードルを越えられるのがLINE。",
       "国内MAU1億人突破、メッセージは受信直後から高い割合で読まれる媒体。"])
nums = [("LINE国内MAU", "1億人突破", "2025年12月末時点"),
        ("メッセージ開封率", "約55%", "メールは約20%｜2022年6月時点")]
cw2 = (CW - 0.4) / 2
for i, (label, val, sub) in enumerate(nums):
    x = CX0 + i * (cw2 + 0.4)
    box(s, x, CY0 + 0.2, cw2, 3.2, fill=PALE)
    T(s, x + 0.3, CY0 + 0.4, cw2 - 0.6, 0.6, [one(label, 12, True, NAVY)], anchor="m")
    T(s, x + 0.3, CY0 + 1.1, cw2 - 0.6, 1.3, [one(val, 28, True, ORANGE)], anchor="m")
    T(s, x + 0.3, CY0 + 2.4, cw2 - 0.6, 0.7, [one(sub, 10, None, INK)], anchor="m")
box(s, CX0, CY0 + 3.7, CW, 2.0, fill=NAVY)
T(s, CX0 + 0.4, CY0 + 3.7, CW - 0.8, 2.0,
  [one("アプリ登録・会員登録は「ハードルが高い」と感じるユーザーが増えている。", 12, None, "E9EDF7", ls=1.3, sa=3),
   one("フォーム入力よりも低ハードルなLINEなら、唯一そのハードルを越えられる。", 12, True, WHITE, ls=1.3)],
  anchor="m")
foot(s, "業界水準｜出典：LINEヤフー公式（国内MAU・開封率）。倍率表現は使用していない")

# ============================================================
# S04 賃貸仲介のWebマーケ構造（3つの壁）
# ============================================================
s = slides[3]
frame(s, "賃貸仲介のWebマーケ構造｜3つの壁",
      ["物件は全社共通。増やせるのは反響数ではなく、反響あたりの成約数だけ。",
       "②③を潰すと、同じ反響数・同じ出稿費のままCPOが下がる。"])
walls = [
    ("① ポータル依存", "反響単価は言い値。\n下げられない", MUT, "触れない"),
    ("② 取りこぼし", "初動が遅い／\n電話が繋がらない", ORANGE, "触れる（自動化）"),
    ("③ 使い捨て", "決めなかった人・更新期の人が\n資産化されていない", NAVY, "触れる（リスト化）"),
]
cw3 = (CW - 0.6) / 3
for i, (h, b, c, tag) in enumerate(walls):
    x = CX0 + i * (cw3 + 0.3)
    wb_ = box(s, x, CY0 + 0.2, cw3, 3.1, fill=WHITE, line=c, lw=1.5)
    put_text(wb_.text_frame,
             [one(h, 13, True, c, align="c", sa=6)] +
             [one(l, 10, None, INK, align="c", ls=1.25) for l in b.split("\n")] +
             [one(tag, 10.5, True, (MUT if i == 0 else c), align="c", sa=8)],
             anchor="m")
box(s, CX0, CY0 + 3.6, CW, 2.1, fill=PALE)
T(s, CX0 + 0.4, CY0 + 3.6, CW - 0.8, 2.1,
  [one("②取りこぼし・③使い捨てを潰すと、同じ反響数・同じ出稿費のままCPOが下がる。", 13, True, TNAVY, ls=1.3)],
  anchor="m")
foot(s, "意見｜DYM提案の構造整理")

# ============================================================
# S05 部屋探しユーザーの心理
# ============================================================
s = slides[4]
frame(s, "部屋探しユーザーの心理",
      ["電話に出ないのは冷たいからではない。知らない番号に出ないだけ。",
       "追いかけるほど逃げる。"])
psy = [
    ("知らない番号には\n出ない", "平均内見物件数は2.7件\n（減少傾向）"),
    ("自分で\n決めたい", "オンライン内見実施率\n32.5%（2年連続）"),
    ("リスクは\n自分で調べる", "ハザードマップを自分で\n確認 48.3%"),
]
cw3 = (CW - 0.6) / 3
for i, (h, b) in enumerate(psy):
    x = CX0 + i * (cw3 + 0.3)
    card(s, x, CY0 + 0.2, cw3, 3.6, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=13, bsz=10.5, ls=1.35)
box(s, CX0, CY0 + 4.2, CW, 1.8, fill=NAVY)
T(s, CX0 + 0.4, CY0 + 4.2, CW - 0.8, 1.8,
  [one("部屋探しから契約までは約3週間〜1ヶ月。この短期間に「追いかけずに寄り添う」設計が要る。", 12.5, True, WHITE, ls=1.3)],
  anchor="m")
foot(s, "業界水準｜出典：SUUMOリサーチセンター 賃貸契約者動向調査（2022〜2024年度・首都圏）。"
        "※検索スニペット経由のため提案前に一次ソースで再確認すること")

# ============================================================
# S06 できること／できないこと ＋ スコープ宣言
# ============================================================
s = slides[5]
frame(s, "Webマーケで「できること／できないこと」",
      ["接客の質は変えられない。変えられるのは反響の初動と再利用の6項目だけ。"])
box(s, CX0, CY0 + 0.2, 8.0, 6.3, fill=GREY)
T(s, CX0 + 0.3, CY0 + 0.4, 7.4, 0.6, [one("できないこと", 13, True, MUT)], anchor="m")
for i, t in enumerate(["担当者の接客トークの質そのもの", "物件の仕入れ力・管理戸数", "ポータルの掲載単価（言い値）"]):
    T(s, CX0 + 0.3, CY0 + 1.1 + i * 1.5, 7.4, 1.3, [one("✕ " + t, 11, None, INK, ls=1.25)], anchor="m")
bx2 = CX0 + 8.3
box(s, bx2, CY0 + 0.2, CW - 8.3, 6.3, fill=PALE)
T(s, bx2 + 0.3, CY0 + 0.4, CW - 8.3 - 0.6, 0.6, [one("できること（＝本提案の範囲）", 13, True, NAVY)], anchor="m")
can = [
    ("反響の初動を自動化する（5分以内の自動返信）", "反響→来店率"),
    ("反響を友だち化してリスト資産にする", "翌期のCPO"),
    ("条件タグで物件レコメンドを自動配信する", "来店率・工数"),
    ("来店予約のリマインドでドタキャンを減らす", "来店→成約率"),
    ("決めなかった人・更新期が近い人を再アプローチする", "追加成約・CPO"),
    ("ポータルを通さない自社反響をつくる", "ポータル依存度"),
]
for i, (t, kpi) in enumerate(can):
    y = CY0 + 1.0 + i * 0.85
    T(s, bx2 + 0.3, y, CW - 8.3 - 3.3, 0.82, [one(f"{i+1}. {t}", 9.5, None, INK, ls=1.15)], anchor="m")
    badge(s, bx2 + CW - 8.3 - 2.8, y + 0.1, 2.5, 0.55, kpi, fill=WHITE, col=ORANGE, sz=8)
foot(s, "意見｜DYMの提案スコープ宣言")


# ============================================================
# S07 シーズナリティ（★実データ：Googleトレンド）
# ============================================================
s = slides[6]
frame(s, "ニーズ調査｜シーズナリティ",
      ["検索は年中フラットなのに、契約は1〜3月に集中している。",
       "＝業界は1〜3月しか刈っていない。閑散期に探している人を、まるごと取りこぼしている。"])
pic_w = (CW - 0.4) / 2
s.shapes.add_picture(str(IMG / "verA_trend_chintai.png"), Cm(CX0), Cm(CY0 + 0.2), width=Cm(pic_w))
s.shapes.add_picture(str(IMG / "verA_trend_hikkoshi.png"), Cm(CX0 + pic_w + 0.4), Cm(CY0 + 0.2), width=Cm(pic_w))
box(s, CX0, CY0 + 7.0, CW, 1.0, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 7.0, CW - 0.6, 1.0,
  [one("閑散期こそ仕込み時。S24（年間企画投稿カレンダー）の設計図になる。", 12, True, WHITE, align="c", ls=1.25)],
  anchor="m")
foot(s, "業界水準｜出典：Googleトレンド（2026年8月取得・実データ）。「内見」は検索Vol不足のため測定不能")

# ============================================================
# S08 前後検索
# ============================================================
s = slides[7]
frame(s, "ニーズ調査｜前後検索",
      ["前はエリア名と相場、後は「審査に落ちる不安」。不安の中身が2週間で入れ替わる。",
       "友だち追加を仕掛けるのは前15日〜起点（比較が始まる前）。"])
bands3 = [
    ("検索“前”15〜7日", "エリア名で探し始める\n（大分 賃貸／千葉市 賃貸 等）＋\n家賃相場・敷金礼金ゼロ物件", MUT),
    ("起点0日", "ポータル名で検索\n（スーモ／ホームズ／athome）", NAVY),
    ("検索“後”7〜15日", "★審査に落ちる不安が集中\n（賃貸 審査 落ちる確率／\n保証人なし 賃貸 等）", RED),
]
cw3b = (CW - 0.6) / 3
for i, (h, sub, c) in enumerate(bands3):
    x = CX0 + i * (cw3b + 0.3)
    bb = box(s, x, CY0 + 0.2, cw3b, 1.1, fill=c)
    put_text(bb.text_frame, [one(h, 11.5, True, WHITE, align="c")], anchor="m")
    T(s, x, CY0 + 1.4, cw3b, 2.4, [one(l, 9.5, None, INK, align="c", ls=1.3) for l in sub.split("\n")], anchor="t")
box(s, CX0, CY0 + 4.2, CW, 2.7, fill=PALE)
T(s, CX0 + 0.3, CY0 + 4.35, CW - 0.6, 2.5,
  [one("★新発見：外国人入居者セグメント", 11.5, True, NAVY, sa=3),
   one("在留カード更新・永住権・帰化条件・技能実習等のKWが検索前後に多数確認された。保証審査と日本語対応が壁になっている層が実在する（S33 D案の根拠）。", 10, None, INK, ls=1.3, sa=4),
   one("ライフイベント連動：「仕事探し」（前1日）→「バイト探し」（後3日）", 9.5, None, MUT, ls=1.25)],
  anchor="t")
foot(s, "事実｜出典：LINEヤフー媒体資料（取得済み・反映済み）")

# ============================================================
# S09 他社分析①
# ============================================================
s = slides[8]
frame(s, "他社分析①｜競合のLINE運用ステータス",
      ["競合は友だち数を持っていても、使っているとは限らない。",
       "友だち数は page.line.me から取得日つきで実測する。"])
simple_table(s, CX0, CY0 + 0.2, CW, 1.6,
             ["アカウント", "ID", "友だち数", "運用の型"],
             [["いい部屋ネット（本部）", "883ubdur", "★未取得", "要実測"],
              ["エイブルAGENT（単身）", "869oufcj", "★未取得", "要実測"]],
             col_w=[CW * 0.32, CW * 0.22, CW * 0.20, CW * 0.26], row_h=0.55)
placeholder(s, CX0, CY0 + 2.4, CW, 4.5, "★差込枠｜友だち数・配信内容の実測（未取得）",
            "page.line.me から取得日つきで取得。「おすすめ」欄の数値は拾わない。\n"
            "実機スクリーンショット（あいさつMSG・リッチメニュー・直近配信）も併せて取得")
foot(s, "選定は知名度優先（スキル §3 で事前承認済み）")

# ============================================================
# S10 他社分析②
# ============================================================
s = slides[9]
frame(s, "他社分析②｜大手FCは店舗ごとにアカウントが乱立",
      ["本部で運用が統一されておらず、友だちが店舗ごとに分散している。",
       "数は持っていても、リストが会社の資産になっていない。"])
FC = [
    ("アパマンショップ", "高崎/中央/十条/北浦和/東中野/中野の6店舗"),
    ("エイブル", "豊中/横浜/小田急相模原/寝屋川の4店舗"),
    ("いい部屋ネット", "三宮駅前/西荻窪の2店舗"),
]
cw3 = (CW - 0.6) / 3
for i, (h, b) in enumerate(FC):
    x = CX0 + i * (cw3 + 0.3)
    card(s, x, CY0 + 0.2, cw3, 2.6, h, [b, "友だち数：★未取得"], hcol=NAVY, fill=PALE, hsz=12, bsz=10, ls=1.3)
box(s, CX0, CY0 + 3.1, CW, 2.0, fill=PRED, line=RED, lw=1.0)
T(s, CX0 + 0.3, CY0 + 3.1, CW - 0.6, 2.0,
  [one("本部で運用が統一されておらず、友だちが店舗ごとに分散している。", 12, True, RED, ls=1.3, sa=3),
   one("数は持っていても、リストが会社の資産になっていない＝併営会社（管理併営）なら差別化できる余地", 10.5, None, INK, ls=1.3)],
  anchor="m")
foot(s, "★未取得｜店舗別アカウントの存在は確認済み。友だち数は page.line.me から実測（要取得）")

# ============================================================
# S11 カスタマージャーニー
# ============================================================
s = slides[10]
frame(s, "全体設計｜カスタマージャーニー",
      ["8フェーズのうち接点が無いのは②サイト離脱〜④リード獲得。ここが失注の正体。"])
PH = ["①接触\n（認知・流入）", "②サイト離脱\n（対策）", "③育成\n（興味喚起）", "④リード獲得\n（web/LINE）",
      "⑤リード有効化", "⑥有効リードの\n再育成", "⑦マネタイズ\n（成約/購入）", "⑧LTV最大化\n・紹介"]
n = len(PH)
cwp = CW / n
for i, label in enumerate(PH):
    x = CX0 + i * cwp
    gap = (i in (1, 2, 3))
    bb = box(s, x + 0.08, CY0 + 0.4, cwp - 0.16, 2.2, fill=(PRED if gap else PALE), line=(RED if gap else None), lw=1.25)
    put_text(bb.text_frame, [one(l, 8.5, True, (RED if gap else NAVY), align="c", ls=1.15) for l in label.split("\n")], anchor="m")
box(s, CX0 + 1 * cwp, CY0 + 2.9, 3 * cwp, 0.7, fill=RED)
T(s, CX0 + 1 * cwp, CY0 + 2.9, 3 * cwp, 0.7, [one("★接点が無い区間＝失注の正体", 10, True, WHITE, align="c")], anchor="m")
box(s, CX0, CY0 + 4.1, CW, 1.6, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 4.1, CW - 0.6, 1.6,
  [one("この空白区間をLINEで埋める。ポータルと店舗の動線には手を入れない。", 12.5, True, WHITE, align="c", ls=1.3)],
  anchor="m")
foot(s, "意見｜DYM標準のカスタマージャーニーフレームに賃貸仲介を当てはめ")

# ============================================================
# S12 施策全体像
# ============================================================
s = slides[11]
frame(s, "全体設計｜施策全体像",
      ["ポータルと店舗には手を入れない。その間の「反響〜来店」だけを埋める。"])
flow = ["ポータル\n反響", "→", "★LINE友だち化\n（初動自動化）", "→", "来店予約\n・内見", "→", "成約\n・入居", "→", "更新\n・住み替え"]
mw = CW / len(flow)
for i, f in enumerate(flow):
    x = CX0 + i * mw
    if f == "→":
        T(s, x, CY0 + 1.6, mw, 1.2, [one(f, 18, True, ORANGE, align="c")], anchor="m")
    else:
        highlight = "LINE" in f
        bb = box(s, x + 0.1, CY0 + 1.3, mw - 0.2, 1.8, fill=(PORANGE if highlight else PALE),
                  line=(ORANGE if highlight else None), lw=1.25)
        put_text(bb.text_frame, [one(l, 9.5, True, (ORANGE if highlight else NAVY), align="c", ls=1.15) for l in f.split("\n")], anchor="m")
box(s, CX0, CY0 + 3.6, CW, 1.8, fill=PALE)
T(s, CX0 + 0.3, CY0 + 3.6, CW - 0.6, 1.8,
  [one("触るのはLINE友だち化から成約までの区間だけ。ポータル出稿費もLPも触らない。", 12.5, True, TNAVY, ls=1.3)],
  anchor="m")
foot(s, "意見｜DYM提案の施策全体像")

# ============================================================
# S13 対策領域マップ
# ============================================================
s = slides[12]
frame(s, "全体設計｜対策領域マップ",
      ["打ち手は新規反響と既存入居者の2領域。併営（仲介＋管理）なら両方取れる。"])
areas = [
    ("① 新規反響領域", "取りこぼしの回収（来店率・CPO）", [
        "反響の初動を自動化（5分以内）",
        "希望条件アンケートで物件レコメンド",
        "来店リマインドでドタキャン防止",
    ], NAVY, PALE),
    ("② 既存入居者領域", "使い捨ての解消（更新・紹介・LTV）", [
        "更新期の住み替え提案",
        "決めなかった人の再アプローチ",
        "退去予告後の引き止め・紹介導線",
    ], ORANGE, PORANGE),
]
cw2 = (CW - 0.4) / 2
for i, (h, sub, items, c, fill) in enumerate(areas):
    x = CX0 + i * (cw2 + 0.4)
    box(s, x, CY0 + 0.2, cw2, 6.2, fill=fill)
    T(s, x + 0.3, CY0 + 0.4, cw2 - 0.6, 0.7, [one(h, 15, True, c)], anchor="m")
    T(s, x + 0.3, CY0 + 1.15, cw2 - 0.6, 0.7, [one(sub, 10.5, None, MUT)], anchor="m")
    for j, it in enumerate(items):
        T(s, x + 0.3, CY0 + 2.1 + j * 1.5, cw2 - 0.6, 1.4, [one(f"・{it}", 11, None, INK, ls=1.3)], anchor="t")
foot(s, "意見｜併営（仲介＋管理）の場合、②既存入居者領域が差別化になる")

# ============================================================
# S14 施策展開図（初期・月次）
# ============================================================
s = slides[13]
frame(s, "全体設計｜施策展開図（初期・月次）",
      ["初期構築は1ヶ月、以降は月次運用。やることは決まっていて増え続けない。"])
band(s, CY0 + 0.2, "初期（初動・構築）", fill=NAVY, h=0.75)
init_items = ["あいさつメッセージ設計（電話しません宣言）", "希望条件アンケート設計（5問）",
              "リッチメニュー3タブ18ボタン", "14日ステップ配信の設計", "反響即時対応（5分以内自動返信）"]
for i, it in enumerate(init_items):
    card(s, CX0 + (i % 3) * (CW / 3), CY0 + 1.2 + (i // 3) * 1.9, CW / 3 - 0.2, 1.7,
         f"{i+1}", it, hcol=NAVY, fill=PALE, hsz=13, bsz=10, ls=1.25)
band(s, CY0 + 5.0, "月次（定例運用）", fill=ORANGE, col=WHITE, h=0.75)
month_items = ["年間企画投稿（閑散期の仕込み）", "物件レコメンド配信の運用",
               "通知メッセージ運用（更新・退去予告）", "定例レポート・改善提案"]
for i, it in enumerate(month_items):
    card(s, CX0 + (i % 4) * (CW / 4), CY0 + 6.0, CW / 4 - 0.2, 1.4,
         f"{i+1}", it, hcol=ORANGE, fill=PORANGE, hsz=12, bsz=9, ls=1.2)
foot(s, "意見｜DYM標準の展開図を賃貸施策に置換")


# ============================================================
# S15 友だち追加動線
# ============================================================
s = slides[14]
frame(s, "構築｜友だち追加動線",
      ["反響が来た瞬間にLINEへ寄せる。ここを外すと以降が全部乗らない。"])
routes = [
    ("① 反響直後の自動誘導", "問い合わせ直後にLINE誘導。\n5分以内の初動に接続"),
    ("② サイト離脱防止", "検索結果一覧の離脱意図を\n検知してLINE追加バナー"),
    ("③ LINE広告（CPF）", "友だち追加動線として\nタイムラインに配信"),
    ("④ 既存リストの活用", "メルマガ・テレマの\n既存リストをLINEへ移行"),
]
cw4 = (CW - 0.3 * 3) / 4
for i, (h, b) in enumerate(routes):
    x = CX0 + i * (cw4 + 0.3)
    card(s, x, CY0 + 0.3, cw4, 4.6, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=11, bsz=9.5, ls=1.3)
box(s, CX0, CY0 + 5.2, CW, 1.6, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 5.2, CW - 0.6, 1.6,
  [one("①反響直後の自動誘導が最速で貯まる。ここを外すと以降が全部乗らない。", 12.5, True, WHITE, align="c", ls=1.25)],
  anchor="m")
foot(s, "意見｜DYM提案の友だち追加動線設計")

# ============================================================
# S16 あいさつメッセージ
# ============================================================
s = slides[15]
frame(s, "構築｜あいさつメッセージ",
      ["「こちらからお電話はしません」を最初に書く。この1行がブロック率を決める。"])
phone(s, CX0 + 0.3, CY0 + 0.2, 8.6, 6.6, "●●不動産",
      [("in", "お問い合わせありがとうございます！\n担当の田中です。"),
       ("in", "✅ こちらからお電話はしません\n✅ 来店を急かすこともしません"),
       ("in", "5つだけ教えてください。\nご希望に合うお部屋をお送りします。"),
       ("chip", "希望条件を入力する"),
       ("note", "※ 反響から5分以内に自動送信")])
box(s, CX0 + 9.6, CY0 + 0.2, CW - 9.6, 6.6, fill=PALE)
T(s, CX0 + 9.9, CY0 + 0.4, CW - 9.9 - 0.3, 0.6, [one("設計の要点", 12.5, True, NAVY)], anchor="m")
pts = [
    ("「お電話はしません」を最初に", "知らない番号に出ない層の\n警戒を解く（S05の心理）"),
    ("来店を急かさない", "追いかけるほど逃げる。\n提案は物件情報から入る"),
    ("5分以内の自動送信", "反響直後の温度が\n最も高い瞬間を逃さない"),
]
for i, (h, b) in enumerate(pts):
    T(s, CX0 + 9.9, CY0 + 1.2 + i * 1.8, CW - 9.9 - 0.3, 1.7,
      [one(f"・{h}", 11, True, NAVY, ls=1.25, sa=2)] +
      [one(l, 9.5, None, INK, ls=1.2) for l in b.split("\n")], anchor="t")
foot(s, "意見｜配信文面は初稿")

# ============================================================
# S17 希望条件アンケート
# ============================================================
s = slides[16]
frame(s, "構築｜希望条件アンケート",
      ["5問すべて選択式・自由入力ゼロ。Q1＝エリアが最初（前後検索データで裏付け済み）。"])
Q = [
    ("Q1 エリアは？", "選択式（沿線・駅）", "前後検索でエリア名が最初に検索されることと一致"),
    ("Q2 家賃の上限は？", "選択式（レンジ）", "物件レコメンドの絞り込み軸"),
    ("Q3 間取りは？", "1K/1LDK/2LDK等", "同上"),
    ("Q4 入居時期は？", "今すぐ／1ヶ月以内／\n2〜3ヶ月／情報収集中", "温度感の把握・ステップ配信の分岐"),
    ("Q5 こだわりは？", "ペット可／初期費用が安い／\n家具付き 等（複数選択）", "レコメンドのタグ付け"),
]
cw2 = (CW - 0.4) / 2
for i, (h, opt, use) in enumerate(Q):
    r, c = divmod(i, 2)
    x = CX0 + c * (cw2 + 0.4)
    y = CY0 + 0.2 + r * 2.5
    if i == 4:
        x = CX0
        y = CY0 + 0.2 + 2 * 2.5
    card(s, x, y, cw2, 2.3, h, [opt, "→ " + use], hcol=NAVY, fill=PALE, hsz=11.5, bsz=9, ls=1.25)
foot(s, "意見｜取れるタグはS23（物件レコメンド配信）で使用")

# ============================================================
# S18 リッチメニュー
# ============================================================
s = slides[17]
frame(s, "構築｜リッチメニュー",
      ["18ボタンすべてに役割がある。左上は常に「今すぐ探す」。"])
richmenu(s, CX0 + 1.5, CY0 + 0.3, 12.0, 6.5, [
    ("探す", ["今すぐ\n物件を探す", "エリアから\n探す", "初期費用が\n安い物件",
            "ペット可・\n家具付き", "新着を\n受け取る", "条件を\n変更する"]),
    ("お金・審査", ["初期費用を\n計算する", "家賃の目安\n（手取り比）", "審査について",
                "保証人なしで\n借りたい", "フリーレント\nとは", "相談する"]),
    ("来店・内見", ["来店を\n予約する", "オンライン\n内見", "内見の流れ",
                "IT重説\nについて", "店舗の場所", "日程を\n変更する"]),
])
T(s, CX0 + 14.2, CY0 + 0.5, CW - 14.2, 6.0,
  [one("設計の要点", 12.5, True, NAVY, sa=6),
   one("・左上＝「今すぐ探す」固定", 10.5, None, INK, ls=1.3, sa=4),
   one("・お金・審査＝「審査に落ちる\n　不安」（後7〜15日の検索）に\n　先回り", 10.5, None, INK, ls=1.3, sa=4),
   one("・来店・内見＝オンライン内見の\n　逃げ道を明示", 10.5, None, INK, ls=1.3, sa=4),
   one("・タブ切替で3倍の情報量を\n　1画面に収める", 10.5, None, INK, ls=1.3)],
  anchor="t")
foot(s, "意見｜DYM提案のリッチメニュー設計")

# ============================================================
# S19 反響即時対応
# ============================================================
s = slides[18]
frame(s, "構築｜反響即時対応（★重点）",
      ["5分以内に返せるかで来店率が変わる。人ではなく自動化で解く。"])
steps = [
    ("① 反響発生", "ポータル・自社サイトから\n問い合わせが入る"),
    ("② 自動応答", "5分以内にLINEで\nあいさつ＋アンケート誘導"),
    ("③ 条件取得", "希望条件アンケートで\nタグ化・レコメンド開始"),
]
cw3 = (CW - 0.6) / 3
for i, (h, b) in enumerate(steps):
    x = CX0 + i * (cw3 + 0.3)
    card(s, x, CY0 + 0.3, cw3, 2.6, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=12.5, bsz=10, ls=1.3)
box(s, CX0, CY0 + 3.2, CW, 2.2, fill=PORANGE)
T(s, CX0 + 0.3, CY0 + 3.35, CW - 0.6, 2.0,
  [one("夜間・繁忙期でも落ちないのが自動化の価値", 13, True, RED, ls=1.3, sa=4),
   one("有人対応では夜間・土日の反響取りこぼしが起きる。自動化なら24時間、初動の速度が落ちない。", 10.5, None, INK, ls=1.3)],
  anchor="m")
foot(s, "意見｜DYM提案の反響即時対応設計")

# ============================================================
# S20 14日ステップ設計
# ============================================================
s = slides[19]
frame(s, "配信設計｜14日ステップ設計（★重点）",
      ["検討2〜4週間に14日を張る。Day10〜14が他社で決まる直前＝最後の勝負どころ。"])
sc = [
    ("① 初動", "Day0", "5分以内の自動返信＋希望条件アンケート"),
    ("② 提案", "Day1・3", "条件タグに合う物件を配信。来店を急がせない"),
    ("③ 誘導", "Day5・7", "内見・来店オファー①（オンライン内見の逃げ道つき）"),
    ("④ クロージング", "Day10・14", "最後の一押し。Day14のテーマは「審査」"),
]
cw4 = CW / 4
for i, (h, d, b) in enumerate(sc):
    x = CX0 + i * cw4
    tb = box(s, x + 0.1, CY0 + 0.3, cw4 - 0.2, 0.8, fill=(ORANGE if i == 3 else NAVY))
    put_text(tb.text_frame, [one(f"{h}｜{d}", 11.5, True, WHITE, align="c")], anchor="m")
    card(s, x + 0.1, CY0 + 1.2, cw4 - 0.2, 2.8, "", [b], hcol=NAVY, fill=(PORANGE if i == 3 else PALE), hsz=1, bsz=10, ls=1.3, anchor="m")
box(s, CX0, CY0 + 4.4, CW, 1.6, fill=PALE)
T(s, CX0 + 0.3, CY0 + 4.4, CW - 0.6, 1.6,
  [one("「◯日後：方向性（〜というデータがあるため）」の形式で、根拠を1行ずつ添える（S08と連動）", 11, True, TNAVY, ls=1.3)],
  anchor="m")
foot(s, "意見｜④のタイミングは前後検索データ（S08：後7〜15日に審査不安が集中）と連動")


# ============================================================
# S21 実文面①
# ============================================================
s = slides[20]
frame(s, "配信設計｜実文面①",
      ["Day0の即レスとDay3の物件提案。実際に届く文面はこれ。"])
phone(s, CX0 + 0.6, CY0 + 0.2, 12.0, 6.7, "Day0",
      [("in", "お問い合わせありがとうございます！\n担当の田中です。"),
       ("in", "✅ こちらからお電話はしません\n✅ 来店を急かすこともしません"),
       ("chip", "希望条件を入力する"),
       ("note", "通知プレビュー：「お問い合わせありがとうございます！担当の")])
phone(s, CX0 + 12.9, CY0 + 0.2, 12.0, 6.7, "Day3",
      [("in", "ご希望条件に合うお部屋が\n3件見つかりました"),
       ("in", "初期費用が安い順に\nご紹介しますね"),
       ("chip", "物件を見てみる")])
foot(s, "意見｜配信文面は初稿。絵文字・改行込みで実際の配信に近い形")

# ============================================================
# S22 実文面②
# ============================================================
s = slides[21]
frame(s, "配信設計｜実文面②",
      ["Day7の内見オファーとDay14。Day14のテーマは「審査」。"])
phone(s, CX0 + 0.6, CY0 + 0.2, 12.0, 6.7, "Day7",
      [("in", "気になる物件は\n見つかりましたか？"),
       ("in", "オンライン内見もできます。\n来店の前にご相談だけでも"),
       ("chip", "来店を予約する"),
       ("chip", "オンライン内見を予約する")])
phone(s, CX0 + 12.9, CY0 + 0.2, 12.0, 6.7, "Day14",
      [("in", "お部屋はお決まりですか？"),
       ("in", "【申込前に確認すべき5つ】\n①審査は保証会社が見ます\n②保証人なしでも通ります\n③必要書類は3点だけ"),
       ("btn", "審査の可否を先に確認する"),
       ("chip", "今回は見送る"),
       ("note", "通知プレビュー：「お部屋はお決まりですか？」")])
foot(s, "意見｜Day14は「今回は見送る」の逃げ道を必ず置く")

# ============================================================
# S23 物件レコメンド配信
# ============================================================
s = slides[22]
frame(s, "配信設計｜物件レコメンド配信",
      ["全員に同じ物件を送るのをやめる。「初期費用が安い」軸を必ず入れる。"])
pts = [
    ("エリアタグで絞る", "S17のQ1回答から\n該当エリアの新着物件のみ配信"),
    ("初期費用が安い順", "「初期費用が安い」を\n軸に含める（S17 Q5と連動）"),
    ("入居時期で優先度を調整", "「今すぐ」層には即時配信、\n「情報収集中」層は週次まとめ"),
]
cw3 = (CW - 0.6) / 3
for i, (h, b) in enumerate(pts):
    x = CX0 + i * (cw3 + 0.3)
    card(s, x, CY0 + 0.3, cw3, 3.0, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=12, bsz=10, ls=1.3)
box(s, CX0, CY0 + 3.7, CW, 1.6, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 3.7, CW - 0.6, 1.6,
  [one("全員に同じ物件を送らない。条件タグに合わないレコメンドはブロック率を上げる。", 12, True, WHITE, align="c", ls=1.25)],
  anchor="m")
foot(s, "意見｜S17（希望条件アンケート）のタグ設計が前提")

# ============================================================
# S24 年間企画投稿カレンダー
# ============================================================
s = slides[23]
frame(s, "配信設計｜年間企画投稿カレンダー",
      ["繁忙期は捌く、閑散期は仕込む。月ごとに何を投稿するかを先に決める。"])
CAL = [
    ("1〜3月", "繁忙期", "捌く", "即レス体制強化・在庫を出し切る配信"),
    ("4〜6月", "新生活後半", "仕込み①", "夏の住み替え需要（7月ピーク）に向けた企画"),
    ("7月", "検索ピーク", "捌く", "「賃貸」検索トレンド年間最高（実データ）"),
    ("8〜12月", "閑散期", "仕込み②", "1〜3月の反響を先取りする認知獲得企画"),
]
cw4 = (CW - 0.3 * 3) / 4
for i, (m, tag, act, b) in enumerate(CAL):
    x = CX0 + i * (cw4 + 0.3)
    tb = box(s, x, CY0 + 0.2, cw4, 0.9, fill=(ORANGE if "捌く" in act else NAVY))
    put_text(tb.text_frame, [one(f"{m}｜{tag}", 10.5, True, WHITE, align="c")], anchor="m")
    card(s, x, CY0 + 1.2, cw4, 3.2, act, [b], hcol=(ORANGE if "捌く" in act else NAVY), fill=(PORANGE if "捌く" in act else PALE), hsz=12, bsz=9, ls=1.3)
foot(s, "意見｜参照データはS07（賃貸検索は年中フラット・7月ピーク／引っ越し検索は1〜3月）")

# ============================================================
# S25 通知メッセージ
# ============================================================
s = slides[24]
frame(s, "配信設計｜通知メッセージ",
      ["友だちでない人にも届く接点がある。内見リマインド／更新案内／退去予告後。"])
notif = [
    ("内見リマインド", "前日・当日朝に\nプッシュ通知"), ("更新案内", "更新期の\n2〜3ヶ月前に事前案内"),
    ("退去予告後", "退去予告を受けた\n入居者へ住み替え提案"), ("空室アラート", "管理物件のオーナー向け\n空室発生を即時通知"),
]
cw4 = (CW - 0.3 * 3) / 4
for i, (h, b) in enumerate(notif):
    x = CX0 + i * (cw4 + 0.3)
    card(s, x, CY0 + 0.3, cw4, 3.2, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=12, bsz=10, ls=1.3)
box(s, CX0, CY0 + 3.9, CW, 1.8, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 3.9, CW - 0.6, 1.8,
  [one("★友だちでなくても届く接点がある", 13, True, WHITE, ls=1.3, sa=4),
   one("ブロックされていても、内見・更新に関わる通知は届く。既存入居者向けはS13②領域に対応。", 10.5, None, "E9EDF7", ls=1.3)],
  anchor="m")
foot(s, "意見｜通知メッセージはAPI連携（Messaging API）が前提。別途費用（S30参照）")

# ============================================================
# S26 来店率改善
# ============================================================
s = slides[25]
frame(s, "来店改善・工数｜来店率改善",
      ["ドタキャンは前日リマインドで減る。来られない人にはオンライン内見の逃げ道。"])
pts = [
    ("前日リマインド", "来店予約の前日夜に\n自動送信"), ("オンライン内見の逃げ道", "来店できない層を\n取りこぼさない"),
    ("当日朝の再通知", "当日の道順・\n持ち物を再案内"),
]
cw3 = (CW - 0.6) / 3
for i, (h, b) in enumerate(pts):
    x = CX0 + i * (cw3 + 0.3)
    card(s, x, CY0 + 0.3, cw3, 2.6, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=12, bsz=10, ls=1.3)
box(s, CX0, CY0 + 3.2, CW, 1.6, fill=PORANGE)
T(s, CX0 + 0.3, CY0 + 3.2, CW - 0.6, 1.6,
  [one("オンライン内見実施率は32.5%（2年連続・SUUMOリサーチセンター）。すでに一定の需要がある選択肢", 11.5, True, RED, ls=1.3)],
  anchor="m")
foot(s, "業界水準｜出典：SUUMOリサーチセンター 賃貸契約者動向調査")

# ============================================================
# S27 工数削減
# ============================================================
s = slides[26]
frame(s, "来店改善・工数｜工数削減",
      ["繁忙期の人手不足は増員では解けない。架電と再送信を自動化して解く。"])
pts = [
    ("反響初動の自動化", "5分以内の自動返信で\n有人対応の待ち行列を解消"),
    ("よくある質問の自動応答", "「保証人なしで借りれる？」等を\n24時間自動応答"),
    ("リマインドの自動化", "前日リマインド・\n更新案内を自動配信"),
]
cw3 = (CW - 0.6) / 3
for i, (h, b) in enumerate(pts):
    x = CX0 + i * (cw3 + 0.3)
    card(s, x, CY0 + 0.3, cw3, 2.6, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=11.5, bsz=9.5, ls=1.3)
box(s, CX0, CY0 + 3.2, CW, 1.6, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 3.2, CW - 0.6, 1.6,
  [one("繁忙期の人手不足は増員では解けない。定型対応を自動化し、接客に集中する設計", 12, True, WHITE, align="c", ls=1.25)],
  anchor="m")
foot(s, "意見｜DYM提案の工数削減設計")

# ============================================================
# S28 改善モデル
# ============================================================
s = slides[27]
frame(s, "来店改善・工数｜改善モデル",
      ["反響の有効化→来店→成約→更新。各段に打ち手が1対1で対応している。"])
model = ["反響\n有効化", "→", "来店\n予約", "→", "内見\n・成約", "→", "更新\n・住み替え"]
mw = CW / len(model)
for i, m in enumerate(model):
    x = CX0 + i * mw
    if m == "→":
        T(s, x, CY0 + 1.0, mw, 1.4, [one(m, 18, True, ORANGE, align="c")], anchor="m")
    else:
        bb = box(s, x + 0.1, CY0 + 0.9, mw - 0.2, 1.6, fill=PALE)
        put_text(bb.text_frame, [one(l, 10, True, NAVY, align="c", ls=1.15) for l in m.split("\n")], anchor="m")
moves = [
    ("反響有効化", "5分以内自動返信＋希望条件アンケート"), ("来店予約", "オンライン内見の逃げ道＋前日リマインド"),
    ("内見・成約", "初期費用が安い順のレコメンド"), ("更新・住み替え", "更新期の事前案内・住み替え提案"),
]
for i, (h, b) in enumerate(moves):
    x = CX0 + i * (CW / 4)
    card(s, x, CY0 + 3.0, CW / 4 - 0.2, 2.6, h, [b], hcol=NAVY, fill=WHITE, line=BORDER, hsz=10.5, bsz=9, ls=1.25)
foot(s, "意見｜DYM提案の改善モデル")


# ============================================================
# S29 効果測定の設計
# ============================================================
s = slides[28]
frame(s, "成果と体制｜効果測定の設計",
      ["各段で何を測るかを先に決める。主KPI＝CPO。数値は入れない。"])
simple_table(s, CX0, CY0 + 0.2, CW, 3.6,
             ["", "Before（ポータルのみ）", "After（1年目）", "After（2年目）"],
             [["反響数", "出稿費÷CPO", "同左（出稿費は据え置き）", "同左"],
              ["来店率", "現状", "① 取りこぼし解消で上昇", "同左"],
              ["翌期の出稿費", "毎期フルで必要", "同左", "② 友だちリスト分だけ圧縮可"],
              ["CPO", "毎期リセット", "改善方向", "★期を追うごとに下がる"]],
             col_w=[CW * 0.16, CW * 0.28, CW * 0.28, CW * 0.28], bsz=8.5, row_h=0.75)
box(s, CX0, CY0 + 4.0, CW, 2.2, fill=PALE)
T(s, CX0 + 0.3, CY0 + 4.15, CW - 0.6, 2.0,
  [one("★肝：ポータル反響は使い捨てなのでCPOは毎期リセットされる。", 11.5, True, RED, ls=1.3, sa=2),
   one("友だち化すると資産になるので、期を追うごとに下がる。", 11.5, None, INK, ls=1.3, sa=4),
   one("締め：御社の実績をいただければ、この構造に数字を入れてSIMを作ります。", 12, True, TNAVY, ls=1.3)],
  anchor="t")
foot(s, "意見｜数値は入れない（業界汎用のためSIMは作らない）。主KPI＝CPO")

# ============================================================
# S30 費用プラン
# ============================================================
s = slides[29]
frame(s, "成果と体制｜費用プラン",
      ["無償付帯と別途費用を分けて示す。追加で払うものを隠さない（6ヶ月〜・税抜）。"])
plans = [
    ("① コンサル基本", "初期10万〜", "月20万（3投稿）\n月30万（5投稿）\n月50万（9投稿）"),
    ("② 初動設計＋運用", "初期10万〜", "月5万〜"),
    ("③ 運用代行・効率改善", "初期10万〜", "月0万〜\n（アカウント費のみ）"),
    ("④ 成果報酬型", "初期5万〜", "月＝単価×成果数\n＋固定費"),
]
cw4 = (CW - 0.3 * 3) / 4
for i, (h, init, m) in enumerate(plans):
    x = CX0 + i * (cw4 + 0.3)
    box(s, x, CY0 + 0.2, cw4, 3.3, fill=PALE)
    T(s, x + 0.2, CY0 + 0.35, cw4 - 0.4, 0.9, [one(h, 11, True, NAVY, ls=1.2)], anchor="t")
    T(s, x + 0.2, CY0 + 1.3, cw4 - 0.4, 0.6, [one(init, 13, True, ORANGE)], anchor="t")
    T(s, x + 0.2, CY0 + 1.95, cw4 - 0.4, 1.4, [one(l, 9, None, INK, ls=1.25) for l in m.split("\n")], anchor="t")
T(s, CX0, CY0 + 3.8, CW, 0.5, [one("無償付帯", 11.5, True, NAVY)], anchor="m")
T(s, CX0, CY0 + 4.3, CW, 0.9,
  [one("アカウント開設／プロフィール／リッチメニュー／あいさつ／初期アンケート／KW自動応答／ステップ配信／セグメント配信／タグ管理／GAレポート連携／クリエイティブ／定例会", 9.5, None, INK, ls=1.3)],
  anchor="t")
T(s, CX0, CY0 + 5.3, CW, 0.5, [one("別途費用", 11.5, True, RED)], anchor="m")
T(s, CX0, CY0 + 5.8, CW, 0.9,
  [one("離脱防止（Sitelead）初期10万＋月5万／API連携／通知メッセージ 等。配信ツールはLstep／Hachidoriのいずれかを要件で選定", 9.5, None, INK, ls=1.3)],
  anchor="t")
foot(s, "意見｜DYM標準プラン（6ヶ月〜・税抜）")

# ============================================================
# S31 運用スケジュール
# ============================================================
s = slides[30]
frame(s, "成果と体制｜運用スケジュール",
      ["1ヶ月目に構築、2ヶ月目に稼働、3ヶ月目から改善。6ヶ月で型ができる。"])
sched = [
    ("1ヶ月目", "構築", "あいさつMSG・アンケート・リッチメニュー・ステップ配信の設計"),
    ("2ヶ月目", "稼働", "友だち追加動線を稼働。14日ステップ配信開始"),
    ("3〜4ヶ月目", "改善①", "反響即時対応・物件レコメンドを開始。効果測定"),
    ("5〜6ヶ月目", "改善②", "年間企画投稿カレンダーの型を確立"),
]
cw4 = CW / 4
for i, (h, tag, b) in enumerate(sched):
    x = CX0 + i * cw4
    tb = box(s, x + 0.1, CY0 + 0.3, cw4 - 0.2, 0.8, fill=NAVY)
    put_text(tb.text_frame, [one(h, 11.5, True, WHITE, align="c")], anchor="m")
    T(s, x + 0.1, CY0 + 1.25, cw4 - 0.2, 0.6, [one(tag, 12, True, ORANGE, align="c")], anchor="m")
    card(s, x + 0.1, CY0 + 1.95, cw4 - 0.2, 2.7, "", b, hcol=NAVY, fill=PALE, hsz=1, bsz=9.5, ls=1.3, anchor="m")
    if i < 3:
        T(s, x + cw4 - 0.15, CY0 + 1.9, 0.4, 1.0, [one("▶", 14, True, MUT, align="c")], anchor="m")
foot(s, "意見｜6ヶ月で年間配信の型を確立。以降は月次運用フェーズへ")

# ============================================================
# S32 サポート体制
# ============================================================
s = slides[31]
frame(s, "成果と体制｜サポート体制",
      ["定例会・レポート・クリエイティブ制作まで含む。運用は丸ごと持つ。"])
supp = [
    ("定例会", "月次で振り返り・\n翌月の企画をすり合わせ"),
    ("レポーティング", "GA連携で友だち数・\n来店率を数値化"),
    ("クリエイティブ制作", "配信文面・リッチメニュー\n画像の制作を内包"),
    ("運用代行", "日々の配信設計・\n自動応答の調整"),
]
cw4 = (CW - 0.3 * 3) / 4
for i, (h, b) in enumerate(supp):
    x = CX0 + i * (cw4 + 0.3)
    card(s, x, CY0 + 0.3, cw4, 3.0, h, b.split("\n"), hcol=NAVY, fill=PALE, hsz=11.5, bsz=9.5, ls=1.3)
foot(s, "意見｜DYM標準のサポート体制。運用は丸ごと持つ")

# ============================================================
# S33 飛び道具（4案）
# ============================================================
s = slides[32]
frame(s, "締め｜飛び道具",
      ["他社がやっていない打ち手が4つある。"])
tobi = [
    ("A", "内見前の360°・動画を友だち限定で配信", "現地に行く前に絞れる。来店時点で温度が高い", False),
    ("B", "決めなかった理由アンケート", "失注をリスト化して次の繁忙期の資産にする", False),
    ("C", "更新3ヶ月前の住み替え提案", "退去を他社に取られる前に拾う。併営会社だけができる", False),
    ("D", "外国人入居者向けの翻訳チャット対応", "★前後検索で在留・ビザ・日本語試験・中国語KWが多数確認。保証審査と日本語対応が壁の層が実在", True),
]
cw2 = (CW - 0.4) / 2
for i, (tag, h, b, verified) in enumerate(tobi):
    r, c = divmod(i, 2)
    x = CX0 + c * (cw2 + 0.4)
    y = CY0 + 0.2 + r * 3.4
    card(s, x, y, cw2, 3.1, f"{tag}　{h}", [b], hcol=NAVY, fill=PALE, hsz=11.5, bsz=9.5, ls=1.35)
    if verified:
        badge(s, x + cw2 - 2.5, y + 0.15, 2.2, 0.5, "実データ確認済み", fill=PORANGE, col=ORANGE, sz=7.5)
foot(s, "事実｜D案は前後検索データ（S08）で裏付け済み。他は意見（施策アイデア）")

# ============================================================
# S34 第2の提案軸｜オーナー向け
# ============================================================
s = slides[33]
frame(s, "締め｜第2の提案軸｜オーナー向け",
      ["この仕組みはそのままオーナー開拓にも使える。"])
box(s, CX0, CY0 + 0.2, CW, 4.0, fill=PALE)
T(s, CX0 + 0.4, CY0 + 0.4, CW - 0.8, 0.7, [one("入居者集客で構築した基盤は、オーナー開拓にも転用できる", 14, True, NAVY)], anchor="m")
pts = [
    "反響対応の自動化ノウハウ → オーナー相談窓口の初動対応に応用",
    "友だち化・タグ設計のノウハウ → 土地活用タイプ診断に応用",
    "管理併営の会社なら、入居者もオーナーも同じLINE基盤で運用できる",
]
for i, t in enumerate(pts):
    T(s, CX0 + 0.4, CY0 + 1.3 + i * 0.9, CW - 0.8, 0.85, [one(f"・{t}", 11, None, INK, ls=1.3)], anchor="m")
box(s, CX0, CY0 + 4.5, CW, 1.5, fill=NAVY)
T(s, CX0 + 0.3, CY0 + 4.5, CW - 0.6, 1.5,
  [one("詳細は「賃貸業界 Ver.B（オーナー開拓版）」でご提案します。", 12.5, True, WHITE, align="c", ls=1.3)],
  anchor="m")
foot(s, "意見｜Ver.Bへの橋渡し")

# ============================================================
# S35 LINEOA実績
# ============================================================
s = slides[34]
frame(s, "締め｜LINEOA実績",
      ["賃貸でLINEが機能した実例がある（LY公式）。"])
cases = [
    ("お部屋探しのハートサポート（賃貸仲介）", "LINE経由の問い合わせから30〜40件成約（2020年1〜7月）。申込〜入居がLINEで完結"),
    ("LIFULL HOME'S（不動産情報サイト）", "ステップ配信でCPC約6割改善・運用工数1/4。通知メッセのCTRがメール比+15%"),
]
cw2 = (CW - 0.4) / 2
for i, (h, b) in enumerate(cases):
    x = CX0 + i * (cw2 + 0.4)
    card(s, x, CY0 + 0.2, cw2, 3.3, h, [b], hcol=NAVY, fill=PALE, hsz=12, bsz=10, ls=1.35)
foot(s, "事実｜出典：LINEヤフー公式事例。数値は検索スニペット経由のため、掲載前に各事例ページの原典で再確認すること")


# ============================================================
# 業種残骸の最終チェック（正規表現スキャン）
# ============================================================
import re

RESIDUE_PATTERNS = [
    r"○○業界", r"〇〇業界", r"クロスセル", r"アップセル", r"転職", r"求人",
    r"利用状況タグ", r"誕生日.{0,3}記念日タグ", r"人材\(転職\)業界",
]
residue_hits = []
for i, sl in enumerate(slides, 1):
    for sh in sl.shapes:
        if sh.has_text_frame:
            txt = sh.text_frame.text
            for pat in RESIDUE_PATTERNS:
                if re.search(pat, txt):
                    residue_hits.append((i, pat, txt[:40]))
if residue_hits:
    print("★業種残骸が残っています：")
    for i, pat, txt in residue_hits:
        print(f"  slide {i}: /{pat}/ -> {txt!r}")
else:
    print("業種残骸チェック：残存0")

prs.save(OUT)
print("saved:", OUT)
