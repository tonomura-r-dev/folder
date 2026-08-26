# -*- coding: utf-8 -*-
"""特別プランのご提供方針 1枚（DYM汎用FMTをコピー→1枚に削減→中身を作り直す）

内容（確定文言）：
  ・10月から原則、特別プランのご提供は難しいです。
  ・11月以降で各社様の配信金額と媒体インセンティブの兼ね合いから
    特別プランの調整をいたします。
"""
import shutil
from copy import deepcopy
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = Path(__file__).resolve().parent.parent
SRC = str(ROOT / "_templates" / "DYM_LINEOA_FMT.pptx")
OUT = str(ROOT / "20260826_特別プランのご提供方針.pptx")

TNAVY = "002060"   # タイトル
NAVY = "1F285A"    # 打ち手・カード見出し
RED = "C00000"     # 制約・注意
INK = "333333"     # 本文
MUT = "808080"
WHITE = "FFFFFF"
PALE = "F4F7FF"
BORDER = "D9D9D9"

shutil.copyfile(SRC, OUT)
prs = Presentation(OUT)
slides = list(prs.slides)

# ---- 素材キャプチャ：スライド2の区切り線（y=1390675） ----
div_el = None
for sh in slides[1].shapes:
    if sh._element.tag.endswith('}cxnSp') and abs(sh.top - 1390675) < 60000:
        div_el = deepcopy(sh._element)
        break
assert div_el is not None, "divider not found"

# ---- 構造作業：スライド2だけ残す（追加はしない） ----
sldIdLst = prs.slides._sldIdLst
ids = list(sldIdLst)
keep = ids[1]
for sldId in ids:
    if sldId is keep:
        continue
    prs.part.drop_rel(sldId.rId)
    sldIdLst.remove(sldId)
slides = list(prs.slides)
assert len(slides) == 1, len(slides)


# ---------- helpers（build_vivical_9p.py と同じ経路） ----------
def set_font(run, size, bold=None, color=INK, name="メイリオ"):
    f = run.font
    f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    f.name = name
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = rPr.find(qn(tag))
        if e is None:
            e = rPr.makeelement(qn(tag), {})
            rPr.append(e)
        e.set("typeface", name)
    f.color.rgb = RGBColor.from_string(color)


ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
ANCH = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}


def put_text(tf, paras, anchor="t", ml=0.06, mr=0.06, mt=0.03, mb=0.03, wrap=True):
    tf.word_wrap = wrap
    tf.margin_left = Inches(ml)
    tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt)
    tf.margin_bottom = Inches(mb)
    tf.vertical_anchor = ANCH[anchor]
    first = True
    for p in paras:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = ALIGN[p.get("align", "l")]
        if p.get("sa") is not None:
            para.space_after = Pt(p["sa"])
        if p.get("ls") is not None:
            para.line_spacing = p["ls"]
        for (t, sz, b, c) in p["runs"]:
            r = para.add_run()
            r.text = t
            set_font(r, sz, b, c)
    return tf


def add_text(slide, x, y, w, h, paras, anchor="t", **kw):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    put_text(box.text_frame, paras, anchor=anchor, **kw)
    return box


def add_box(slide, x, y, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = RGBColor.from_string(line)
        sp.line.width = Pt(lw)
    sp.shadow.inherit = False
    sp.text_frame.word_wrap = True
    return sp


def clear_slide(slide):
    spTree = slide.shapes._spTree
    for el in list(spTree):
        if el.tag.split("}")[-1] in ("sp", "cxnSp", "pic", "graphicFrame", "grpSp"):
            spTree.remove(el)


def frame(slide, header, band):
    """タイトル＋結論バンド＋区切り線（FMT実測準拠）"""
    clear_slide(slide)
    slide.shapes._spTree.append(deepcopy(div_el))
    add_text(slide, 0.60, 0.13, 8.10, 0.40,
             [{"runs": [(header, 16, True, TNAVY)], "align": "l"}],
             anchor="m", ml=0.0, mr=0.0)
    add_text(slide, 0.42, 0.66, 10.02, 0.80,
             [{"runs": band, "align": "c"}],
             anchor="m", ml=0.0, mr=0.0)


def chip(slide, x, y, w, h, text, fill, size=14):
    c = add_box(slide, x, y, w, h, fill=fill)
    put_text(c.text_frame, [{"runs": [(text, size, True, WHITE)], "align": "c"}],
             anchor="m", ml=0.0, mr=0.0, mt=0.0, mb=0.0)
    return c


# ---------- 作図 ----------
s = slides[0]
frame(s, "特別プランのご提供方針",
      [("10月から原則、特別プランのご提供は難しく、", 14, True, INK),
       ("11月以降", 14, True, RED),
       ("に各社様ごとの調整をいたします", 14, True, INK)])

# ===== 10月：原則ご提供なし =====
add_box(s, 0.50, 1.90, 9.83, 1.24, fill=PALE, line=BORDER)
chip(s, 0.82, 2.29, 1.34, 0.46, "10月〜", RED)
add_text(s, 2.42, 1.94, 7.60, 1.16,
         [{"runs": [("原則、特別プランのご提供は難しいです。", 21, True, INK)]}],
         anchor="m", ml=0.0)

arrow = add_box(s, 5.19, 3.28, 0.45, 0.42, fill=BORDER, shape=MSO_SHAPE.DOWN_ARROW)
arrow.line.fill.background()

# ===== 11月以降：配信金額×媒体インセンティブで調整 =====
add_box(s, 0.50, 3.86, 9.83, 2.70, fill=PALE, line=BORDER)
chip(s, 0.82, 4.16, 1.34, 0.46, "11月以降", NAVY)
add_text(s, 2.42, 4.04, 7.60, 0.72,
         [{"runs": [("各社様の配信金額と媒体インセンティブの兼ね合いから、", 15.5, True, INK)]},
          {"runs": [("特別プランの調整をいたします。", 15.5, True, INK)]}],
         anchor="m", ml=0.0)

add_text(s, 0.90, 4.90, 4.00, 0.30,
         [{"runs": [("調整のポイント", 12, True, NAVY)]}], anchor="m", ml=0.0)

TILE_Y, TILE_H = 5.26, 1.00
t1 = add_box(s, 0.90, TILE_Y, 2.75, TILE_H, fill=WHITE, line=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
put_text(t1.text_frame,
         [{"runs": [("各社様の", 10.5, None, MUT)], "align": "c", "sa": 2},
          {"runs": [("配信金額", 16, True, NAVY)], "align": "c"}],
         anchor="m", ml=0.06, mr=0.06, mt=0.0, mb=0.0)

add_text(s, 3.65, TILE_Y, 0.60, TILE_H,
         [{"runs": [("×", 20, True, MUT)], "align": "c"}], anchor="m", ml=0.0, mr=0.0)

t2 = add_box(s, 4.25, TILE_Y, 2.75, TILE_H, fill=WHITE, line=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
put_text(t2.text_frame,
         [{"runs": [("媒体側の", 10.5, None, MUT)], "align": "c", "sa": 2},
          {"runs": [("媒体インセンティブ", 16, True, NAVY)], "align": "c"}],
         anchor="m", ml=0.06, mr=0.06, mt=0.0, mb=0.0)

add_text(s, 7.00, TILE_Y, 0.60, TILE_H,
         [{"runs": [("→", 20, True, MUT)], "align": "c"}], anchor="m", ml=0.0, mr=0.0)

t3 = add_box(s, 7.60, TILE_Y, 2.35, TILE_H, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
put_text(t3.text_frame,
         [{"runs": [("各社様ごとに", 10.5, None, "C8D0E8")], "align": "c", "sa": 2},
          {"runs": [("特別プランを調整", 15, True, WHITE)], "align": "c"}],
         anchor="m", ml=0.06, mr=0.06, mt=0.0, mb=0.0)

prs.save(OUT)
print("saved:", OUT)
