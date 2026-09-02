# -*- coding: utf-8 -*-
"""LINE公式アカウント運用のご提案【業界汎用版】のビルド

DYM汎用FMT（_templates/DYM_LINEOA_FMT.pptx・36枚）を骨格そのままに、
業種依存の記述だけを外した36枚の汎用デッキを生成する。

FMTからの変更点：
  1. 転職業界に特化していた3ページ（33-35）を業種非依存の「型」に書き換え
     （登録直後アンケート／年間の企画カレンダー／診断コンテンツ）
  2. 表紙の未記入プレースホルダ「業界：○○」を「業界汎用版」に
  3. 業種を限定する語を中立化（定期購入→リピート、査定→見積、商圏→エリア）

**ページは1枚も削っていない**（殿村さん指示：FMTの構成をそのまま保つ）。
そのぶん、商談前に手当てが要るページが残る。ビルド後に出る
「差し替えチェックリスト」がその一覧なので、出す前に必ず確認すること。

使い方:
    pip install python-pptx
    python _build/build_generic_lineoa.py
    python _build/qa_render.py "DYM_LINEOA_業界汎用_ご提案.pptx"
"""
from copy import deepcopy
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.text.text import _Paragraph

ROOT = Path(__file__).resolve().parent.parent
SRC = ROOT / "_templates" / "DYM_LINEOA_FMT.pptx"
OUT = ROOT / "DYM_LINEOA_業界汎用_ご提案.pptx"

# 削除するページ（FMTでの1始まり番号）。
# **空 ＝ FMTの36枚を1枚も削らない**（殿村さん指示）。
#
# 以前は下の11枚を削って25枚版にしていた。軽い版が要るときはここに戻す：
#   5, 6              グラフ未挿入の差し込み枠（本文がほぼ空）
#   7, 8              ニーズ調査：人材/転職業界のGoogleトレンド画像＋「あああ」
#   9, 10             他社分析：「すうじ」「ロゴ」等のダミーで全面構成
#   12, 13, 14, 15, 16  具体企画のブレスト：社内メモのまま
DROP = []

# 商談前に手当てが要るページ（ビルド後のチェックリストに出す）
NEEDS_WORK = {
    5: "検索広告CPCの推移：グラフが未挿入。業界のCPCデータを貼る",
    6: "広告実績CVR：グラフが未挿入。実績CVRを貼る",
    7: "ニーズ調査（シーズナリティ）：転職KWのGoogleトレンド画像。業界KWで撮り直す",
    8: "ニーズ調査（前後検索）：転職の前後検索データ＋「あああ」。業界KWに差し替え",
    9: "他社配信アカウント：「すうじ」「ロゴ」等のダミー。競合の実数値を記入",
    10: "他社分析：ダミーのアカウント画面5枚。競合のキャプチャに差し替え",
    12: "具体企画のブレスト：社内メモのまま。企画案を書くか、商談では飛ばす",
    13: "同上",
    14: "同上",
    15: "同上（見出しは「他社分析」だが中身はブレスメモ）",
    16: "同上",
}

# 業種を限定してしまう語の中立化（全ページ対象・run単位で置換）
# FMTは1文が複数runに分かれているため、パターンはrunをまたがない短さにする
NEUTRALIZE = [
    ("定期購入", "リピート"),   # EC限定の語を汎用化
    ("査定", "見積"),           # 不動産・中古車限定の語を汎用化
    ("商圏", "エリア"),         # 店舗業限定の語を汎用化
    ("XXX", "〇〇"),            # FMTに残っていたダミー
]

BADGE = "業種別にカスタマイズ"  # 元の「人材(転職)業界の例」を置き換える


# ---------- helpers ----------
def iter_shapes(shapes):
    """グループ内も含めて全シェイプを辿る"""
    for sh in shapes:
        yield sh
        if sh.shape_type == 6:  # GROUP
            yield from iter_shapes(sh.shapes)


def shape_by_id(slide, shape_id):
    for sh in iter_shapes(slide.shapes):
        if sh.shape_id == shape_id:
            return sh
    raise KeyError(f"shape_id={shape_id} not found")


def set_lines(slide, shape_id, lines):
    """先頭runの書式を引き継いだまま本文を差し替える。

    put_text相当の作り直しはせず、既存の段落・runを使い回す。
    FMTのシェイプは書式が細かく作り込まれているため、
    新規に組み直すと見た目が崩れる。
    """
    if isinstance(lines, str):
        lines = [lines]
    tf = shape_by_id(slide, shape_id).text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs:
        raise ValueError(f"shape_id={shape_id} に書式継承元のrunが無い")

    # 2段落目以降を削除（残すと追記になる）
    for para in tf.paragraphs[1:]:
        para._p.getparent().remove(para._p)
    # 1段落目の2run目以降を削除
    for r in p0.runs[1:]:
        r._r.getparent().remove(r._r)

    p0.runs[0].text = lines[0]
    template = deepcopy(p0._p)
    for line in lines[1:]:
        new_p = deepcopy(template)
        p0._p.getparent().append(new_p)
        _Paragraph(new_p, tf).runs[0].text = line


def neutralize(prs):
    """業種を限定する語を全ページから中立化する"""
    hits = 0
    for slide in prs.slides:
        for sh in iter_shapes(slide.shapes):
            if not sh.has_text_frame:
                continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    for before, after in NEUTRALIZE:
                        if before in run.text:
                            run.text = run.text.replace(before, after)
                            hits += 1
    return hits


def drop_slides(prs, numbers):
    """1始まりのページ番号を指定して削除する"""
    sld_id_lst = prs.slides._sldIdLst
    ids = list(sld_id_lst)
    for n in sorted(numbers, reverse=True):
        el = ids[n - 1]
        prs.part.drop_rel(el.rId)
        sld_id_lst.remove(el)


# ---------- build ----------
shutil.copyfile(SRC, OUT)
prs = Presentation(OUT)
slides = list(prs.slides)
assert len(slides) == 36, f"FMTは36枚のはず（実際: {len(slides)}）"

# =====================================================================
# Slide 1: 表紙
#   FMTの「業界：○○」は未記入に見えてしまうので、汎用版であることを明示する。
#   業界別に出すときはここを「業界：不動産」のように書き換える。
# =====================================================================
set_lines(slides[0], 7, "業界汎用版")

# =====================================================================
# Slide 33: 具体施策（友だち追加～短期施策）／登録直後アンケート
#   転職者向けの5問 → 業種を問わず成立する「検討状態を測る」5問へ
# =====================================================================
s = slides[32]
set_lines(s, 6, BADGE)
QUESTIONS = [
    # (設問ID, 選択肢ID, 狙いID, 設問, 選択肢, 狙い)
    (46, 47, 49,
     "Q1. 今回、ご検討を始めたきっかけは？",
     "A) 費用を抑えたい　B) 品質で選びたい　C) 悩みを解決したい　D) 情報収集のみ　E) その他",
     "セグメント分けの最重要設問。回答別に配信内容を変える"),
    (53, 54, 56,
     "Q2. ご利用・ご導入の希望時期はいつ頃ですか？",
     "A) 1ヶ月以内　B) 3ヶ月以内　C) 半年以内　D) まだ決めていない",
     "「1〜3ヶ月以内」をホットリードとして即オファー対象に"),
    (60, 61, 63,
     "Q3. ご予算・ご希望の価格帯を教えてください",
     "A) 〜〇万円　B) 〇〜〇万円　C) 〇〜〇万円　D) まだ未定",
     "提案プランの絞り込みと、訴求する価格帯の設定に活用"),
    (67, 68, 70,
     "Q4. 今、検討にあたって一番不安なことは？",
     "A) 費用対効果が不安　B) 手間や時間がかかる　C) 品質・実績が不明　D) 社内の合意が必要",
     "ナーチャリングコンテンツの優先順位を決める設問"),
    (74, 75, 77,
     "Q5. 他社サービスの利用・検討の経験はありますか？",
     "A) はじめて　B) 使ったことがある（満足）　C) 使ったことがある（不満）",
     "「不満あり」層には競合差別化コンテンツを優先配信"),
]
for q_id, a_id, aim_id, q, a, aim in QUESTIONS:
    set_lines(s, q_id, q)
    set_lines(s, a_id, a)
    set_lines(s, aim_id, aim)

# =====================================================================
# Slide 34: 具体施策（長期施策）／年間の企画カレンダー
#   転職市場の繁忙期 → どの業種にもある「年4回の山」へ
# =====================================================================
s = slides[33]
set_lines(s, 44, BADGE)
MOMENTS = [
    # (期間ID, ラベルID, 企画名ID, 本文ID, KPI ID, ラベル, 企画名, 本文3行, KPI)
    (8, 9, 10, 12,
     "年度末・新年度",
     "「新年度スタート応援キャンペーン」",
     ["1〜3月は「新しく始めたい」需要が高まる時期。",
      "年度替わり・新生活の準備で検討が加速する。",
      "施策：期間限定オファーをLINEで先行案内。「残り〇枠」を配信。"],
     "KPI目標：LINE経由の申込率 +30% 目標"),
    (17, 18, 19, 21,
     "GW前後",
     "「連休中にじっくり比較・検討」",
     ["大型連休は情報収集と比較に時間を割ける期間。",
      "検討ハードルが下がり、意思決定が進みやすい。",
      "施策：連休前に「失敗しない選び方ガイド」PDF特典を配信。"],
     "KPI目標：LINE登録数 ピーク月比 +20%"),
    (26, 27, 28, 30,
     "下半期スタート",
     "「下半期スタートダッシュ企画」",
     ["9月は下期の予算・計画が動き出す時期。",
      "先送りしていた検討が再始動しやすくなる。",
      "施策：「秋限定・〇〇の先行案内」をLINEで先出し。希少性で登録を促進。"],
     "KPI目標：先行案内経由の成約率 計測開始"),
    (35, 36, 37, 39,
     "年末",
     "「来年こそ、と考えている方へ」",
     ["年末は「来年こそ」という決意が高まるタイミング。",
      "競合の出稿が減り広告効率が良い穴場期間。",
      "施策：「来年の〇〇トレンド予測」を先行配布。年明けの商談を先に確保。"],
     "KPI目標：1月の商談・来店予約の先行獲得率"),
]
for label_id, title_id, body_id, kpi_id, label, title, body, kpi in MOMENTS:
    set_lines(s, label_id, label)
    set_lines(s, title_id, title)
    set_lines(s, body_id, body)
    set_lines(s, kpi_id, kpi)

# =====================================================================
# Slide 35: 具体施策（長期施策）／診断コンテンツ
#   転職の悩みを測る4問 → 商材を問わず「検討の温度」を測る4問へ
# =====================================================================
s = slides[34]
set_lines(s, 10, BADGE)
set_lines(s, 119,
          "目的：登録後すぐCVしないユーザーに"
          "“あなたの〇〇度”診断（ハードル：低）→現状可視化→セグメント追客")
DIAGNOSIS_Q = [
    (82, "Q1  「今のやり方で本当に良いか、迷うことはありますか？」"),
    (84, "Q2  「調べてはみたものの、比較しきれず止まっていませんか？」"),
    (86, "Q3  「1年後も今と同じ状態で良い、と言い切れますか？」"),
    (88, "Q4  「費用や手間がネックで、踏み出せずにいますか？」"),
]
for shape_id, text in DIAGNOSIS_Q:
    set_lines(s, shape_id, text)
DIAGNOSIS_R = [
    (92, 96, "今すぐ検討型",
     "個別相談を即座にオファー。限定枠・期限を明示したクロージングメッセージ"),
    (99, 103, "慎重比較型",
     "導入事例・他社比較データを週1配信。3〜4週間かけて温める"),
    (106, 110, "情報収集型",
     "業界動向・お役立ち情報を提供。半年スパンで関係構築"),
    (113, 117, "現状維持型",
     "無理にアプローチせず月1配信。「環境変化」で再燃を待つ長期戦略"),
]
for type_id, plan_id, type_name, plan in DIAGNOSIS_R:
    set_lines(s, type_id, type_name)
    set_lines(s, plan_id, plan)

# ---------- 全ページの中立化 → 不要ページの削除 ----------
n_neutralized = neutralize(prs)
drop_slides(prs, DROP)

prs.save(OUT)

# ---------- 検品 ----------
done = Presentation(OUT)
n_slides = len(list(done.slides))
assert n_slides == 36 - len(DROP), n_slides

# 業種固有・ダミーが残っている箇所を洗い出す。
# ページを削っていないので「残っていて当然」だが、
# **書き換えたはずの3ページに転職の語が残っていたら、それはバグ。**
NG_WORDS = ["転職", "求人", "面接", "年収", "入社", "退職", "キャリアアップ",
            "来場", "査定", "商圏", "定期購入",
            "あああ", "すうじ", "なまえ", "XXX"]
REWRITTEN = {33, 34, 35}  # 業種非依存に書き換えたページ（DROPが空のときの番号）

found = {}
for i, slide in enumerate(done.slides, 1):
    for sh in iter_shapes(slide.shapes):
        if not sh.has_text_frame:
            continue
        for word in NG_WORDS:
            if word in sh.text_frame.text:
                found.setdefault(i, set()).add(word)

print(f"saved: {OUT.name}")
print(f"slides: {n_slides}（FMT 36枚 − 削除 {len(DROP)}枚）")
print(f"中立化した語: {n_neutralized}箇所")

print("\n--- 差し替えチェックリスト（商談前に手当てするページ）---")
if DROP:
    print("  ※DROPが空でないためページ番号がずれています。FMT基準の番号で表示します。")
for page, note in sorted(NEEDS_WORK.items()):
    if page in DROP:
        continue
    words = "／".join(sorted(found.get(page, []))) or "—"
    print(f"  p{page:<2} {note}  [検出: {words}]")

# 書き換え済みページに業種固有の語が残っていないかだけは厳しく見る
bugs = {p: w for p, w in found.items() if p in REWRITTEN and not DROP}
if bugs:
    print("\n!! 書き換えたはずのページに業種固有の語が残っている:")
    for page, words in sorted(bugs.items()):
        print(f"   p{page} {'／'.join(sorted(words))}")
else:
    print("\n書き換えた3ページ（33-35）の業種固有語: なし")
